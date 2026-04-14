#!/usr/bin/env python3
"""
Deck Builder — Generates branded PPTX presentations from YAML definitions.

Uses a brand-specific PPTX template with precise placeholder mappings
derived from the template inspector (curated-layouts.json).

Usage:
    python3 build_deck.py <yaml_file> [--output <output.pptx>] [--upload]

The --upload flag will upload to Google Drive and convert to Google Slides.
"""

import argparse
import json
import math
import os
import subprocess
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Canvas scaling — supports templates with non-standard canvas sizes
# (e.g., Google Slides exports at 20"x11.25" instead of 10"x5.625")
# ---------------------------------------------------------------------------

_canvas_scale = 1.0  # Set by build_deck() from template dimensions


def SI(inches):
    """Scale-aware Inches — multiplies by canvas scale factor.

    Use this for all content placement coordinates. Templates with
    larger canvases (20"x11.25") get proportionally scaled content.
    """
    return Inches(inches * _canvas_scale)


def S(value):
    """Scale a raw float value by canvas scale factor.

    Use for dimension calculations that don't go through Inches().
    """
    return value * _canvas_scale


# ---------------------------------------------------------------------------
# Brand configuration
# ---------------------------------------------------------------------------

class BrandConfig:
    """Design tokens loaded from a brand YAML file.

    Provides colors, fonts, path resolution, and gradient interpolation
    used by all layout builders. Defaults to built-in brand if no brand is
    specified.
    """

    # Built-in defaults (used when no brand file is loaded)
    DEFAULTS = {
        "name": "Generic",
        "colors": {
            "primary": "#1A365D", "secondary": "#3182CE", "accent": "#63B3ED",
            "background_light": "#EBF4FF", "background_card": "#F7FAFC",
            "text_dark": "#2D3748", "text_gray": "#718096", "white": "#FFFFFF",
            "green": "#2F855A", "amber": "#C05621", "red": "#C53030",
            "divider": "#CBD5E0",
            "link": "#2B6CB0",
            "neutral_bg": "#F7FAFC",
            "gradient_start": "#1A365D",
            "gradient_end": "#3182CE",
            "gradient_light": "#A3C4EC",
            "gradient_ring_light": "#B3D4FC",
            "staircase_end": "#C0D0E0",
            "spoke_end": "#B0C0D0",
            "venn_tertiary": "#90CAF9",
        },
        "fonts": {"heading": "Arial", "body": "Arial"},
        "title_backgrounds": {
            "default": "title-assets/title-bg.jpg",
        },
        "agenda_backgrounds": {
            "default": "title-assets/agenda-left.jpg",
        },
    }

    def __init__(self, brand_data=None, brand_dir=None):
        data = dict(self.DEFAULTS)
        if brand_data:
            data.update(brand_data)
            # Merge nested dicts
            for nested_key in ("colors", "fonts", "title_backgrounds",
                               "agenda_backgrounds"):
                if nested_key in self.DEFAULTS and nested_key in brand_data:
                    merged = dict(self.DEFAULTS[nested_key])
                    merged.update(brand_data[nested_key])
                    data[nested_key] = merged
        self._data = data
        self._colors = data["colors"]
        self._fonts = data["fonts"]
        # brand_dir is the directory containing the brand.yaml (or the
        # deck-builder script dir for the default brand)
        self._brand_dir = brand_dir or os.path.dirname(os.path.abspath(__file__))
        # True when loaded from a file (vs. class defaults)
        self.is_explicit = brand_data is not None

    # --- Color accessors (return RGBColor) ---
    def _hex_to_rgb(self, hex_str):
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _hex_to_ints(self, hex_str):
        """Return (r, g, b) ints from a hex string."""
        h = hex_str.lstrip("#")
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

    @property
    def name(self): return self._data.get("name", "Unknown")

    # --- Core color properties ---
    @property
    def primary(self): return self._hex_to_rgb(self._colors["primary"])
    @property
    def secondary(self): return self._hex_to_rgb(self._colors["secondary"])
    @property
    def accent(self): return self._hex_to_rgb(self._colors["accent"])
    @property
    def bg_light(self): return self._hex_to_rgb(self._colors["background_light"])
    @property
    def bg_card(self): return self._hex_to_rgb(self._colors["background_card"])
    @property
    def text_dark(self): return self._hex_to_rgb(self._colors["text_dark"])
    @property
    def text_gray(self): return self._hex_to_rgb(self._colors["text_gray"])
    @property
    def white(self): return self._hex_to_rgb(self._colors["white"])
    @property
    def green(self): return self._hex_to_rgb(self._colors["green"])
    @property
    def amber(self): return self._hex_to_rgb(self._colors["amber"])
    @property
    def red(self): return self._hex_to_rgb(self._colors["red"])

    # --- Extended color properties (Phase 1) ---
    @property
    def divider(self): return self._hex_to_rgb(self._colors["divider"])
    @property
    def link(self): return self._hex_to_rgb(self._colors["link"])
    @property
    def neutral_bg(self): return self._hex_to_rgb(self._colors["neutral_bg"])

    # --- Hex accessors (for python-pptx table/XML APIs that need hex strings) ---
    @property
    def primary_hex(self): return self._colors["primary"].lstrip("#")
    @property
    def secondary_hex(self): return self._colors["secondary"].lstrip("#")
    @property
    def accent_hex(self): return self._colors["accent"].lstrip("#")
    @property
    def bg_light_hex(self): return self._colors["background_light"].lstrip("#")
    @property
    def bg_card_hex(self): return self._colors["background_card"].lstrip("#")
    @property
    def text_dark_hex(self): return self._colors["text_dark"].lstrip("#")
    @property
    def divider_hex(self): return self._colors["divider"].lstrip("#")

    def color_hex(self, name):
        """Get any color as a hex string (no '#' prefix) by name."""
        return self._colors.get(name, "#000000").lstrip("#")

    def all_color_hexes(self):
        """Return set of all brand color hex strings (no '#')."""
        return {v.lstrip("#") for v in self._colors.values()}

    # --- Gradient interpolation ---
    def interpolate(self, t, start="gradient_start", end="gradient_end"):
        """Interpolate between two brand colors. t=0 returns start, t=1 returns end.

        Args:
            t: Float 0.0-1.0 controlling interpolation position.
            start: Color key name (default: "gradient_start", which defaults to primary).
            end: Color key name (default: "gradient_end", which defaults to secondary).
        Returns:
            RGBColor interpolated between the two colors.
        """
        r1, g1, b1 = self._hex_to_ints(self._colors[start])
        r2, g2, b2 = self._hex_to_ints(self._colors[end])
        r = int(r1 + (r2 - r1) * t)
        g = int(g1 + (g2 - g1) * t)
        b = int(b1 + (b2 - b1) * t)
        return RGBColor(r, g, b)

    def interpolate_light(self, t):
        """Interpolate from a light version of primary to primary.

        Used by concentric_circles and pyramid where the gradient goes
        from light outer/bottom to dark inner/top.
        """
        # Light endpoint: blend primary toward white by ~75%
        r1, g1, b1 = self._hex_to_ints(self._colors["gradient_end"])
        r2, g2, b2 = self._hex_to_ints(self._colors["gradient_start"])
        # Lighten the end color further for a wider range
        lr = min(255, r1 + (255 - r1) // 3)
        lg = min(255, g1 + (255 - g1) // 3)
        lb = min(255, b1 + (255 - b1) // 3)
        r = int(lr + (r2 - lr) * t)
        g = int(lg + (g2 - lg) * t)
        b = int(lb + (b2 - lb) * t)
        return RGBColor(r, g, b)

    # --- Font accessors ---
    @property
    def heading_font(self): return self._fonts["heading"]
    @property
    def body_font(self): return self._fonts["body"]

    def all_font_names(self):
        """Return set of all brand font names."""
        return set(self._fonts.values())

    # --- Path resolution ---
    def resolve_path(self, relative_path):
        """Resolve a path relative to the brand package directory."""
        return os.path.join(self._brand_dir, relative_path)

    @property
    def template_path(self):
        """Absolute path to the brand's PPTX template."""
        tmpl = self._data.get("template")
        if tmpl:
            # Explicit template specified in brand.yaml
            if os.path.isabs(tmpl):
                return tmpl
            brand_rel = os.path.join(self._brand_dir, tmpl)
            if os.path.isfile(brand_rel):
                return brand_rel
            script_dir = os.path.dirname(os.path.abspath(__file__))
            return os.path.join(script_dir, "..", tmpl)
        # No explicit template — check for template.pptx in brand dir
        brand_default = os.path.join(self._brand_dir, "template.pptx")
        if os.path.isfile(brand_default):
            return brand_default
        # Fall back to generic brand template
        script_dir = os.path.dirname(os.path.abspath(__file__))
        return os.path.join(script_dir, "brands", "generic", "template.pptx")

    @property
    def icon_catalog_path(self):
        """Absolute path to the brand's icon-catalog.json."""
        return os.path.join(self._brand_dir, "icons", "icon-catalog.json")

    def get_title_background(self, key="default"):
        """Get absolute path to a title slide background image."""
        backgrounds = self._data.get("title_backgrounds", {})
        rel = backgrounds.get(key, backgrounds.get("default", ""))
        if not rel:
            return None
        return self.resolve_path(rel)

    def get_agenda_background(self, key="default"):
        """Get absolute path to an agenda slide background image."""
        backgrounds = self._data.get("agenda_backgrounds", {})
        rel = backgrounds.get(key, backgrounds.get("default", ""))
        if not rel:
            return None
        return self.resolve_path(rel)

    @classmethod
    def load(cls, brand_name_or_path=None):
        """Load a brand config from brands/ directory or a file path."""
        if not brand_name_or_path:
            return cls()  # Built-in defaults

        script_dir = os.path.dirname(os.path.abspath(__file__))

        # Try as a file path first
        if os.path.isfile(brand_name_or_path):
            path = brand_name_or_path
        else:
            # Try brands/ directory — first as brands/<name>/brand.yaml,
            # then as brands/<name>.yaml (legacy flat layout)
            pkg_path = os.path.join(script_dir, "brands", brand_name_or_path,
                                    "brand.yaml")
            flat_path = os.path.join(script_dir, "brands",
                                     f"{brand_name_or_path}.yaml")
            if os.path.isfile(pkg_path):
                path = pkg_path
            elif os.path.isfile(flat_path):
                path = flat_path
            else:
                print(f"    WARNING: Brand '{brand_name_or_path}' not found, using defaults")
                return cls()

        with open(path) as f:
            data = yaml.safe_load(f)
        brand_dir = os.path.dirname(os.path.abspath(path))
        return cls(data, brand_dir=brand_dir)


# ---------------------------------------------------------------------------
# Load curated layout catalog
# ---------------------------------------------------------------------------

def load_curated_layouts():
    """Load the curated layout catalog from JSON."""
    catalog_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "curated-layouts.json")
    with open(catalog_path) as f:
        catalog = json.load(f)
    return catalog


# ---------------------------------------------------------------------------
# YAML field → curated placeholder name aliases
# Allows existing YAMLs to work without changes.
# ---------------------------------------------------------------------------

FIELD_ALIASES = {
    # Three-section layouts: YAML uses section*_, catalog uses col*_
    "section1_title": "col1_title",
    "section1_body": "col1_body",
    "section2_title": "col2_title",
    "section2_body": "col2_body",
    "section3_title": "col3_title",
    "section3_body": "col3_body",
    # Big stat
    "description": "label",
}

# Layout name aliases for backward compatibility
LAYOUT_ALIASES = {
    "content_three_section": "three_column",
    "content_three_col_cards": "four_card",
    "content_four_cards": "four_card",
    "title_h1": "title_cover",
    "big_stat": "big_stat_manual",
}

# Separator characters that get dropped when splitting title text across lines
TITLE_SEPARATORS = [" — ", " – ", " - "]

# Split ratios: name → (diagram_fraction, text_fraction)
SPLIT_RATIOS = {
    "v-70/30": (0.70, 0.30),
    "v-60/40": (0.60, 0.40),
    "v-50/50": (0.50, 0.50),
    "v-40/60": (0.40, 0.60),
    "v-30/70": (0.30, 0.70),
    "h-60/40": (0.60, 0.40),
    "h-50/50": (0.50, 0.50),
    "h-40/60": (0.40, 0.60),
}

# Typography hierarchy — font sizes by level
TYPO_HIERARCHY = {
    "primary":   {"size": 12, "bold": True},
    "secondary": {"size": 10, "bold": False},
    "tertiary":  {"size": 8,  "bold": False},
}

# Adaptive sizing steps — try each in order if content doesn't fit
FONT_REDUCTION_STEPS = [
    {"secondary": -1},           # 10 → 9
    {"secondary": -1},           # 9 → 8
    {"primary": -1},             # 12 → 11
    {"primary": -1, "secondary": -1},  # 11 → 10, 8 → 7
]


# ---------------------------------------------------------------------------
# Placeholder text helpers
# ---------------------------------------------------------------------------

def set_placeholder_text(slide, ph_idx, text, preserve_formatting=True,
                         brand=None):
    """Set text on a placeholder by index.

    When brand is provided, uses brand heading font/color for headline placeholders
    (idx 0). Otherwise preserves the template's built-in formatting.
    Returns True if the placeholder was found and set.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame

            if preserve_formatting and tf.paragraphs:
                first_para = tf.paragraphs[0]

                # Save run-level formatting from template
                template_font = {}
                if first_para.runs:
                    run = first_para.runs[0]
                    if run.font.name:
                        template_font["name"] = run.font.name
                    if run.font.size:
                        template_font["size"] = run.font.size
                    if run.font.bold is not None:
                        template_font["bold"] = run.font.bold
                    if run.font.color and run.font.color.rgb:
                        template_font["color"] = run.font.color.rgb

                # Brand theme now handles headline font/color via the template.
                # No run-level override needed — theme inheritance works.

                template_alignment = first_para.alignment

                first_para.text = ""

                lines = str(text).split("\n")
                for li, line in enumerate(lines):
                    if li == 0:
                        para = first_para
                    else:
                        para = tf.add_paragraph()

                    run = para.add_run()
                    run.text = line

                    if template_font.get("name"):
                        run.font.name = template_font["name"]
                    if template_font.get("size"):
                        run.font.size = template_font["size"]
                    if template_font.get("bold") is not None:
                        run.font.bold = template_font["bold"]
                    if template_font.get("color"):
                        run.font.color.rgb = template_font["color"]
                    if template_alignment is not None:
                        para.alignment = template_alignment
            else:
                shape.text = str(text)

            return True
    return False


def set_placeholder_image(slide, ph_idx, image_path):
    """Insert an image into an image placeholder by index.

    Returns True if a placeholder was found and image inserted.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            try:
                shape.insert_picture(image_path)
                return True
            except (AttributeError, TypeError):
                # Not an image placeholder — fall through
                pass
    return False


def add_image_manual(slide, image_path, left, top, width, height=None):
    """Add an image as a shape at explicit coordinates (inches)."""
    left_emu = SI(left)
    top_emu = SI(top)
    width_emu = SI(width)
    if height:
        slide.shapes.add_picture(image_path, left_emu, top_emu, width_emu, SI(height))
    else:
        slide.shapes.add_picture(image_path, left_emu, top_emu, width_emu)


def estimate_text_height(text, width_inches, font_size_pt=10):
    """Estimate the rendered height of text in inches.

    Accounts for explicit newlines, blank-line spacing, and word wrapping
    based on approximate character width for proportional sans-serif fonts.
    """
    chars_per_inch = 13 * (10 / font_size_pt)
    chars_per_line = max(1, int(width_inches * chars_per_inch))
    line_height = font_size_pt * 1.4 / 72  # pt → inches with 1.4x leading

    total_lines = 0
    for paragraph in str(text).split("\n"):
        if not paragraph.strip():
            total_lines += 0.6  # blank lines are shorter
        else:
            total_lines += max(1, -(-len(paragraph) // chars_per_line))
    return total_lines * line_height


def fit_font_size(text, width_inches, max_pt=12, min_pt=7):
    """Find the largest font size where *text* fits on one line within *width_inches*.

    Returns the font size in points, clamped between *min_pt* and *max_pt*.
    """
    for pt in range(max_pt, min_pt - 1, -1):
        chars_per_inch = 13 * (10 / pt)
        chars_per_line = max(1, int(width_inches * chars_per_inch))
        if len(str(text)) <= chars_per_line:
            return pt
    return min_pt


# ---------------------------------------------------------------------------
# Table layout helpers
# ---------------------------------------------------------------------------

_TABLE_MIN_COL_WIDTH = 0.6      # inches — narrower columns are unusable
_TABLE_MAX_COL_FRACTION = 0.45  # no column > 45% of table width
_TABLE_MIN_ROW_H = 0.28         # inches
_TABLE_MAX_ROW_H = 0.60         # inches
_TABLE_CELL_PAD_H = 0.13        # marL(0.08) + marR(0.05)
_TABLE_CELL_PAD_V = 0.06        # marT + marB
_TABLE_FONT_CANDIDATES = [10, 9, 8, 7]
_TABLE_FONT_CANDIDATES_COMPACT = [9, 8, 7]
_TABLE_MIN_FONT = 7


def normalize_col_widths(col_widths, num_cols, table_width, columns=None, rows=None):
    """Return column widths in inches summing to *table_width*.

    If *col_widths* is provided (from YAML), values are treated as proportional
    weights and scaled to fill *table_width*.  ``[30, 70]`` and ``[3, 7]`` both
    produce a 30/70 split.  Existing inch-based specs like ``[2.0, 0.6, 0.9]``
    (sum ≈ table_width) are scaled proportionally — preserving ratios.

    If *col_widths* is ``None``, columns are auto-sized based on content length
    with square-root dampening so one long-text column doesn't dominate.
    """
    # --- Explicit weights → proportional scaling ---
    if col_widths:
        if len(col_widths) != num_cols:
            print(f"    WARNING: col_widths has {len(col_widths)} entries "
                  f"but table has {num_cols} columns — using auto-sizing")
        else:
            total = sum(col_widths)
            if total > 0:
                return [(w / total) * table_width for w in col_widths]

    # --- Content-aware auto-sizing ---
    if not columns and not rows:
        return [table_width / num_cols] * num_cols

    # Gather max character length per column (header + data)
    max_chars = [0] * num_cols
    for c in range(num_cols):
        if columns and c < len(columns):
            max_chars[c] = len(str(columns[c]))
        for row in (rows or []):
            if c < len(row):
                max_chars[c] = max(max_chars[c], len(str(row[c])))

    # Square-root dampening (sqrt(120)=10.95 vs sqrt(15)=3.87 → ~2.8x, not 8x)
    weights = [max(math.sqrt(max(mc, 1)), 1.0) for mc in max_chars]

    # First pass: proportional allocation
    total_weight = sum(weights)
    widths = [(w / total_weight) * table_width for w in weights]

    # Second pass: enforce min/max, redistribute surplus
    max_col_w = table_width * _TABLE_MAX_COL_FRACTION
    clamped = False
    for i in range(num_cols):
        if widths[i] < _TABLE_MIN_COL_WIDTH:
            widths[i] = _TABLE_MIN_COL_WIDTH
            clamped = True
        elif widths[i] > max_col_w:
            widths[i] = max_col_w
            clamped = True

    if clamped:
        current_total = sum(widths)
        if abs(current_total - table_width) > 0.01:
            unclamped = [i for i in range(num_cols)
                         if _TABLE_MIN_COL_WIDTH < widths[i] < max_col_w]
            if unclamped:
                diff = table_width - current_total
                unc_total = sum(widths[i] for i in unclamped)
                for i in unclamped:
                    widths[i] += diff * (widths[i] / unc_total)
            else:
                scale = table_width / current_total
                widths = [w * scale for w in widths]

    return widths


def compute_row_heights(columns, rows, col_widths_inches, font_size_pt,
                        avail_height, header_h_hint=0.35,
                        min_row_h=_TABLE_MIN_ROW_H, max_row_h=_TABLE_MAX_ROW_H,
                        compress_to_fit=False):
    """Compute per-row heights based on text wrapping in each cell.

    Returns ``(header_h, data_row_heights, total_height)``.
    """
    num_cols = len(columns)
    header_pt = min(font_size_pt + 1, 12)

    # Header height
    max_header_h = 0.0
    for c in range(num_cols):
        col_name = str(columns[c]) if c < len(columns) else ""
        usable_w = max(col_widths_inches[c] - _TABLE_CELL_PAD_H, 0.3)
        h = estimate_text_height(col_name, usable_w, header_pt)
        max_header_h = max(max_header_h, h)
    header_h = max(min(max_header_h + _TABLE_CELL_PAD_V, 0.45), min_row_h)

    # Data row heights
    data_row_heights = []
    for row in rows:
        max_cell_h = 0.0
        for c in range(num_cols):
            cell_text = str(row[c]) if c < len(row) else ""
            usable_w = max(col_widths_inches[c] - _TABLE_CELL_PAD_H, 0.3)
            h = estimate_text_height(cell_text, usable_w, font_size_pt)
            max_cell_h = max(max_cell_h, h)
        row_h = max(min(max_cell_h + _TABLE_CELL_PAD_V, max_row_h), min_row_h)
        data_row_heights.append(row_h)

    # Uniform row heights: when most rows are similar, use a single height
    # for all rows so the table looks clean.  Use the max height if it fits,
    # otherwise use the median.
    if data_row_heights:
        uniform_h = max(data_row_heights)
        uniform_total = header_h + uniform_h * len(data_row_heights)
        if uniform_total <= avail_height:
            data_row_heights = [uniform_h] * len(data_row_heights)

    total = header_h + sum(data_row_heights)

    # Proportional compression when content overflows
    if total > avail_height and compress_to_fit:
        scale = avail_height / total
        header_h = max(header_h * scale, 0.22)
        data_row_heights = [max(h * scale, 0.20) for h in data_row_heights]
        total = header_h + sum(data_row_heights)

    return (header_h, data_row_heights, total)


def select_font_size(columns, rows, col_widths_inches, avail_height,
                     header_h_hint=0.35, is_compact=False):
    """Step down through font sizes until the table fits *avail_height*.

    Returns ``(body_pt, header_pt, header_h, row_heights, total_height)``.
    """
    num_cols = len(columns)
    candidates = list(_TABLE_FONT_CANDIDATES_COMPACT if is_compact
                      else _TABLE_FONT_CANDIDATES)

    # Density check: heavy text → skip largest font
    all_lens = []
    for row in rows:
        for c in range(num_cols):
            all_lens.append(len(str(row[c])) if c < len(row) else 0)
    avg_chars = (sum(all_lens) / len(all_lens)) if all_lens else 0

    if avg_chars > 40 and not is_compact and candidates and candidates[0] == 10:
        candidates = candidates[1:]  # drop 10pt

    for body_pt in candidates:
        header_pt = min(body_pt + 1, 12)
        # Column-count pressure: many narrow columns need smaller fonts
        if num_cols >= 7:
            body_pt = max(body_pt - 1, _TABLE_MIN_FONT)
            header_pt = max(header_pt - 1, _TABLE_MIN_FONT)

        header_h, row_heights, total = compute_row_heights(
            columns, rows, col_widths_inches, body_pt,
            avail_height, header_h_hint)

        if total <= avail_height:
            return (body_pt, header_pt, header_h, row_heights, total)

    # Nothing fits — use minimum font with compression
    body_pt = _TABLE_MIN_FONT
    header_pt = _TABLE_MIN_FONT
    header_h, row_heights, total = compute_row_heights(
        columns, rows, col_widths_inches, body_pt,
        avail_height, header_h_hint, compress_to_fit=True)
    return (body_pt, header_pt, header_h, row_heights, total)


def split_table_rows(columns, rows, col_widths_inches, avail_height,
                     font_size_pt, header_h_hint=0.35):
    """Split *rows* into chunks that each fit within *avail_height*.

    Returns a list of row-lists: ``[chunk_1_rows, chunk_2_rows, ...]``.
    Each chunk's height (including a repeated header) fits *avail_height*.
    """
    if not rows:
        return [rows]

    header_pt = min(font_size_pt + 1, 12)
    # Compute header height once (same across all slides)
    header_h = 0.0
    for c in range(len(columns)):
        col_name = str(columns[c]) if c < len(columns) else ""
        usable_w = max(col_widths_inches[c] - _TABLE_CELL_PAD_H, 0.3)
        h = estimate_text_height(col_name, usable_w, header_pt)
        header_h = max(header_h, h)
    header_h = max(min(header_h + _TABLE_CELL_PAD_V, 0.45), _TABLE_MIN_ROW_H)

    # Greedily pack rows into slides
    chunks = []
    current_chunk = []
    current_h = header_h  # each slide starts with header

    for ri, row in enumerate(rows):
        # Compute this row's height
        max_cell_h = 0.0
        for c in range(len(columns)):
            cell_text = str(row[c]) if c < len(row) else ""
            usable_w = max(col_widths_inches[c] - _TABLE_CELL_PAD_H, 0.3)
            h = estimate_text_height(cell_text, usable_w, font_size_pt)
            max_cell_h = max(max_cell_h, h)
        row_h = max(min(max_cell_h + _TABLE_CELL_PAD_V, _TABLE_MAX_ROW_H),
                     _TABLE_MIN_ROW_H)

        if current_chunk and current_h + row_h > avail_height:
            # Start a new slide
            chunks.append(current_chunk)
            current_chunk = [row]
            current_h = header_h + row_h
        else:
            current_chunk.append(row)
            current_h += row_h

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def render_table_cell(cell, text, font_name, font_size_pt, color,
                      alignment=PP_ALIGN.LEFT):
    """Render text into a table *cell* with bullet and bold support.

    - Lines starting with ``- `` or ``• `` get a bullet character prefix.
    - ``**text**`` segments render as inline bold.
    - Alignment is applied to every paragraph.
    """
    import re
    BULLET_CHAR = "\u2022"
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True

    lines = str(text).split("\n")
    for li, line in enumerate(lines):
        para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        para.alignment = alignment

        # Bullet detection
        is_bullet = False
        if line.startswith("- ") or line.startswith("\u2022 "):
            is_bullet = True
            line = line[2:]

        if is_bullet:
            br = para.add_run()
            br.text = BULLET_CHAR + " "
            br.font.name = font_name
            br.font.size = Pt(font_size_pt)
            if color:
                br.font.color.rgb = color

        # Bold detection: split on **...**
        parts = re.split(r'(\*\*.*?\*\*)', line)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                r = para.add_run()
                r.text = part[2:-2]
                r.font.name = font_name
                r.font.size = Pt(font_size_pt)
                r.font.bold = True
                if color:
                    r.font.color.rgb = color
            elif part:
                r = para.add_run()
                r.text = part
                r.font.name = font_name
                r.font.size = Pt(font_size_pt)
                if color:
                    r.font.color.rgb = color


def _apply_table_borders(table, num_rows, num_cols, header_bg_hex="5F016F"):
    """Apply clean horizontal-rule borders to a table.

    Style: thin gray horizontal lines between rows, pink accent under header,
    no vertical lines.  This is the modern "Butterick" data-table look.
    """
    from pptx.oxml.ns import qn

    BORDER_COLOR = "D0D0D0"       # light gray for row separators
    HEADER_BORDER_COLOR = "FF80D4"  # pink accent under header
    BORDER_W = "6350"             # 0.5pt in EMU
    HEADER_BORDER_W = "19050"     # 1.5pt in EMU
    NO_LINE_W = "0"

    def _set_border(cell_tc, side, color, width):
        """Set a single border on a cell via OOXML."""
        tag = {"top": "a:lnT", "bottom": "a:lnB",
               "left": "a:lnL", "right": "a:lnR"}[side]
        tcPr = cell_tc.get_or_add_tcPr()
        # Remove existing border of this type
        for existing in tcPr.findall(qn(tag)):
            tcPr.remove(existing)
        if width == NO_LINE_W:
            ln = tcPr.makeelement(qn(tag), {"w": "0", "cmpd": "sng"})
            no_fill = ln.makeelement(qn("a:noFill"), {})
            ln.append(no_fill)
        else:
            ln = tcPr.makeelement(qn(tag), {"w": width, "cmpd": "sng"})
            sf = ln.makeelement(qn("a:solidFill"), {})
            clr = sf.makeelement(qn("a:srgbClr"), {"val": color})
            sf.append(clr)
            ln.append(sf)
        tcPr.append(ln)

    for ri in range(num_rows):
        for ci in range(num_cols):
            cell_tc = table.cell(ri, ci)._tc

            # Vertical borders — thin gray between all columns and outer edges
            _set_border(cell_tc, "left", BORDER_COLOR, BORDER_W)
            _set_border(cell_tc, "right", BORDER_COLOR, BORDER_W)

            # Horizontal borders
            if ri == 0:
                # Header row: no top border, pink accent bottom
                _set_border(cell_tc, "top", BORDER_COLOR, NO_LINE_W)
                _set_border(cell_tc, "bottom", HEADER_BORDER_COLOR,
                            HEADER_BORDER_W)
            elif ri == num_rows - 1:
                # Last row: gray top, slightly thicker gray bottom
                _set_border(cell_tc, "top", BORDER_COLOR, BORDER_W)
                _set_border(cell_tc, "bottom", BORDER_COLOR, "12700")  # 1pt
            else:
                # Middle rows: thin gray top and bottom
                _set_border(cell_tc, "top", BORDER_COLOR, BORDER_W)
                _set_border(cell_tc, "bottom", BORDER_COLOR, BORDER_W)


import re as _re

_URL_PATTERN = _re.compile(
    r'(https?://[^\s,)]+|(?:[a-zA-Z0-9-]+\.)+(?:com|net|org|io|dev|cloud|ai)(?:/[^\s,)]*)?)'
)

# Markdown link: [display text](url)
_MD_LINK_PATTERN = _re.compile(r'\[([^\]]+)\]\(([^)]+)\)')

# Inline bold: **text**
_BOLD_PATTERN = _re.compile(r'\*\*(.+?)\*\*')

BULLET_CHAR = "\u2022"  # •


def _set_hyperlink_on_run(run, url, paragraph):
    """Set a clickable hyperlink on a text run via OOXML.

    Uses paragraph.part to get the SlidePart, then registers
    the hyperlink as an external relationship on the run's rPr.
    """
    from pptx.oxml.ns import qn

    if not url.startswith("http"):
        url = "https://" + url

    # paragraph.part gives us the SlidePart directly via python-pptx's
    # proxy chain: _Paragraph -> TextFrame -> Shape -> SlidePart
    try:
        slide_part = paragraph.part
    except Exception:
        return  # Can't register relationship — hyperlink will be cosmetic only

    rPr = run._r.get_or_add_rPr()
    hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {})
    rel = slide_part.relate_to(
        url,
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
        is_external=True,
    )
    hlinkClick.set(qn('r:id'), rel)
    rPr.append(hlinkClick)


def _add_run_with_hyperlinks(paragraph, text, font_name, font_size_pt, color,
                             link_color=None):
    """Add text to a paragraph with inline markdown support.

    Parsing pipeline (applied in order):
        1. [Display Text](url)   — named hyperlink
        2. **bold text**         — inline bold
        3. https://example.com   — bare URL, auto-linked
        4. example.com/path      — bare domain, auto-linked

    All other text renders as plain text with the specified font/color.
    link_color defaults to brand.link (purple) via the BrandConfig default.
    """
    LINK_COLOR = link_color or BrandConfig().link

    def _make_run(txt, bold=False, link_url=None):
        """Create a styled run and optionally attach a hyperlink."""
        r = paragraph.add_run()
        r.text = txt
        r.font.name = font_name
        r.font.size = Pt(font_size_pt)
        if link_url:
            r.font.color.rgb = LINK_COLOR
            r.font.underline = True
            _set_hyperlink_on_run(r, link_url, paragraph)
        elif color:
            r.font.color.rgb = color
        if bold:
            r.font.bold = True

    def _emit_with_urls(txt, bold=False):
        """Split text on bare URLs and emit runs (plain or linked)."""
        parts = _URL_PATTERN.split(txt)
        for part in parts:
            if not part:
                continue
            if _URL_PATTERN.fullmatch(part):
                _make_run(part, bold=bold, link_url=part)
            else:
                _make_run(part, bold=bold)

    def _emit_with_bold(txt):
        """Split text on **bold** markers and emit runs."""
        last = 0
        for m in _BOLD_PATTERN.finditer(txt):
            if m.start() > last:
                _emit_with_urls(txt[last:m.start()], bold=False)
            _emit_with_urls(m.group(1), bold=True)
            last = m.end()
        if last < len(txt):
            _emit_with_urls(txt[last:], bold=False)

    # Step 1: Split on markdown links [text](url)
    last_end = 0
    for m in _MD_LINK_PATTERN.finditer(text):
        # Plain text before this link — process for bold + URLs
        if m.start() > last_end:
            _emit_with_bold(text[last_end:m.start()])
        # The markdown link itself
        _make_run(m.group(1), link_url=m.group(2))
        last_end = m.end()
    # Remaining text after last link
    if last_end < len(text):
        _emit_with_bold(text[last_end:])


def _render_body_text(text_frame, body_text, font_name, font_size_pt, color,
                      available_height=None):
    """Render multi-line body text into a text frame with markdown support.

    Supports:
        - Lines starting with '- ' or '• ' render as bulleted items
        - **bold text** renders as inline bold
        - [text](url) renders as named hyperlinks
        - Bare URLs are auto-linked
        - Blank lines are preserved as spacing

    If available_height (inches) is provided and the content doesn't fill it,
    space_after is distributed evenly across content paragraphs to fill the
    available vertical space — preventing text from bunching at the top.
    """
    lines = body_text.split("\n")
    for li, line in enumerate(lines):
        para = text_frame.paragraphs[0] if li == 0 else text_frame.add_paragraph()

        # Bullet detection: lines starting with "- " or "• "
        is_bullet = False
        if line.startswith("- ") or line.startswith("\u2022 "):
            is_bullet = True
            line = line[2:]  # strip the prefix

        if is_bullet:
            # Add bullet character as a separate run
            br = para.add_run()
            br.text = BULLET_CHAR + "  "
            br.font.name = font_name
            br.font.size = Pt(font_size_pt)
            if color:
                br.font.color.rgb = color

        _add_run_with_hyperlinks(para, line, font_name, font_size_pt, color)

    # --- Distribute vertical space ---
    # Typography-based paragraph spacing rules:
    #   Target: 75% of font size (Butterick / Duarte range for presentations)
    #   Maximum: 100% of font size (items feel disconnected beyond this)
    #   Minimum: 25% of font size (tighter than this loses paragraph separation)
    #
    # If available_height is provided, distribute space_after evenly up to these
    # limits. Any remaining space becomes bottom whitespace — we don't stretch
    # beyond the typographic maximum.
    if available_height is not None and available_height > 0:
        content_paras = [p for p in text_frame.paragraphs
                         if p.text.strip()]
        n_content = len(content_paras)
        if n_content > 1:
            line_height = font_size_pt * 1.3 / 72  # inches per line
            est_text_h = len(lines) * line_height
            spare = available_height - est_text_h

            # Typographic limits (in inches)
            target_gap = (font_size_pt * 0.75) / 72
            max_gap = font_size_pt / 72  # 100% of font size

            if spare > 0:
                even_gap = spare / (n_content - 1)
                # Use the lesser of even distribution and typographic target,
                # but never exceed the absolute maximum
                gap = min(even_gap, target_gap)
                gap = min(gap, max_gap)
                gap_emu = int(gap * 914400)  # inches to EMU
                for p in content_paras[:-1]:
                    p.space_after = Emu(gap_emu)


def add_text_box(slide, text, left, top, width, height, font_size=12,
                 font_name="Arial", bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a manual text box at explicit coordinates (inches)."""
    txBox = slide.shapes.add_textbox(SI(left), SI(top),
                                     SI(width), SI(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = str(text).split("\n")
    for li, line in enumerate(lines):
        if li == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        if font_name:
            run.font.name = font_name
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)
        para.alignment = alignment

    return txBox


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def resolve_image_path(image_path, base_dir):
    """Resolve a potentially relative image path."""
    if not os.path.isabs(image_path):
        image_path = os.path.join(base_dir, image_path)
    if not os.path.exists(image_path):
        print(f"    WARNING: Image not found: {image_path}")
        return None
    return image_path


def _generate_layout_note(layout_name, slide_def):
    """Generate a brief note explaining why this layout was chosen."""
    reasons = {
        "title_cover": None,  # structural, no explanation needed
        "closing": None,
        "agenda": None,
        "section_divider": None,
        "big_stat_manual": "Used big_stat: single hero metric commands attention.",
        "callout": "Used callout: single key message for decision-makers.",
        "quote": "Used quote: attributed statement carries authority.",
        "side_by_side": "Used side_by_side: exactly 2 topics to compare.",
        "three_column": "Used three_column: 3 parallel topics or pillars.",
        "before_after": "Used before_after: transformation narrative (current → target state).",
        "process_flow": lambda sd: f"Used process_flow: {len(sd.get('steps', []))} sequential steps.",
        "comparison_matrix": lambda sd: f"Used comparison_matrix: {len(sd.get('columns', []))} options evaluated across {len(sd.get('rows', []))} criteria.",
        "content_table": lambda sd: f"Used content_table: {len(sd.get('columns', []))} columns x {len(sd.get('rows', []))} rows — tabular data needs a table.",
        "kpi_dashboard": lambda sd: f"Used kpi_dashboard: {len(sd.get('metrics', []))} metrics in a scannable grid.",
        "status_board": lambda sd: f"Used status_board: {len(sd.get('items', []))} workstreams with RAG status.",
        "four_card": lambda sd: f"Used four_card: {sum(1 for i in range(1,9) if sd.get(f'card{i}_title'))} discrete items, each self-contained.",
        "numbered_list": lambda sd: f"Used numbered_list: {len(sd.get('items', []))} ordered steps where sequence matters.",
        "pros_cons": "Used pros_cons: strengths vs risks evaluation for decision support.",
        "quadrant": "Used quadrant: 2x2 matrix for strategic positioning.",
        "team_profiles": lambda sd: f"Used team_profiles: {len(sd.get('profiles', []))} people to introduce.",
        "roadmap": "Used roadmap: time-based plan with durations and milestones.",
        "waterfall": "Used waterfall: showing how contributions build from start to end value.",
        "donut_rings": lambda sd: f"Used donut_rings: {len(sd.get('rings', []))} progress percentages.",
        "gauge_dashboard": "Used gauge_dashboard: performance against threshold ranges.",
        "risk_heat_map": "Used risk_heat_map: likelihood x impact grid with risk register.",
        "radar_chart": "Used radar_chart: multi-dimensional capability comparison.",
        "combo_chart": "Used combo_chart: two metrics (bars + line) on shared x-axis.",
        "bubble_chart": "Used bubble_chart: three-variable scatter (x, y, size).",
        "tornado_chart": "Used tornado_chart: sensitivity analysis showing variable impact range.",
        "staircase": "Used staircase: maturity/capability progression levels.",
        "cycle_diagram": "Used cycle_diagram: repeating process with no start/end.",
        "hub_spoke": "Used hub_spoke: central concept with radiating components.",
        "pyramid": "Used pyramid: hierarchical layers from foundation to peak.",
        "venn": "Used venn: overlapping concepts showing intersection.",
        "concentric_circles": "Used concentric_circles: nested scope (broad → narrow).",
        "pricing_table": "Used pricing_table: tier comparison with feature checklist.",
        "bold_bullet": "Used bold_bullet: executive summary — assertions backed by evidence.",
        "bento_grid": "Used bento_grid: mixed-size overview with hero + supporting tiles.",
        "dashboard_panel": "Used dashboard_panel: combined KPIs + chart + summary.",
        "funnel": "Used funnel: pipeline stages showing progressive filtering.",
        "left_nav_sidebar": "Used left_nav_sidebar: step-by-step walkthrough with navigation.",
        "image_text_hero": "Used image_text_hero: visual storytelling with narrative.",
    }
    note = reasons.get(layout_name)
    if note is None:
        return ""
    if callable(note):
        return f"[Layout choice] {note(slide_def)}"
    return f"[Layout choice] {note}"


def build_slide(prs, slide_def, deck_meta, catalog):
    """Build a single slide from a YAML slide definition."""
    brand = deck_meta.get("brand", BrandConfig())
    layout_name = slide_def.get("layout", "content_generic")

    # Resolve layout aliases
    resolved_name = LAYOUT_ALIASES.get(layout_name, layout_name)
    layouts = catalog["layouts"]

    if resolved_name not in layouts:
        print(f"    WARNING: Unknown layout '{layout_name}', using content_generic")
        resolved_name = "content_generic"

    layout_info = layouts[resolved_name]
    layout_idx = layout_info["layout_idx"]
    placeholders = layout_info["placeholders"]
    footer_phs = layout_info.get("footer_placeholders", {})

    # Add slide with the specified layout (fall back to blank canvas if index doesn't exist)
    blank_idx = deck_meta.get("layout_mapping", {}).get("blank_canvas_idx", 2)
    if layout_idx < len(prs.slide_layouts):
        slide_layout = prs.slide_layouts[layout_idx]
    else:
        slide_layout = prs.slide_layouts[blank_idx]
    slide = prs.slides.add_slide(slide_layout)

    # --- Set footer placeholders (skip for title_cover — footers overlap the title) ---
    no_footer_layouts = {"title_cover"}
    if resolved_name not in no_footer_layouts:
        _set_footers(slide, footer_phs, deck_meta)

    _variants_handled_notes = False

    # --- Handle special layouts ---
    if resolved_name == "title_cover":
        _build_title_cover(slide, slide_def, layout_info, deck_meta)
        # Clear any placeholder text that conflicts with the manual title overlay
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx != 2:  # Don't clear the background image placeholder
                try:
                    if hasattr(shape, 'text_frame'):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ""
                except Exception:
                    pass
    elif resolved_name == "agenda":
        agenda_title = slide_def.get("headline", "Agenda")
        set_placeholder_text(slide, 0, agenda_title, brand=brand)
        _build_agenda(slide, slide_def, layout_info, deck_meta)
    elif resolved_name == "content_diagram_text":
        # Diagram + text — swap to content_generic canvas, generate via diagram engine
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

        generic_footers = catalog["layouts"]["content_generic"].get("footer_placeholders", {})
        existing_notes = slide_def.get("notes", "")
        visual_def = slide_def.get("visual", {})
        diagram_type = visual_def.get("type", "org-hierarchy")
        diagram_data = visual_def.get("data", {})
        diagram_approaches = visual_def.get("approaches", None)
        legacy_diagram = slide_def.get("diagram", "")

        # Deck-level defaults with per-slide overrides
        deck_style = deck_meta.get("style", "corporate")
        slide_style = visual_def.get("style", deck_style)
        max_variants = deck_meta.get("diagram_variants", 3)

        # Layout options (one per direction)
        layout_configs = [
            ("v-50/50", "Vertical — diagram top, text below"),
            ("h-50/50", "Horizontal — diagram left, text right"),
            ("h-40/60", "Horizontal — smaller diagram, more text"),
        ]

        # Build the variant plan based on max_variants:
        #   1 variant  → best layout, primary style, native approach
        #   3 variants → one per layout, primary style, native approach
        #   6 variants → one per layout × (primary style + tech-gradient), native
        #   9 variants → + blueprint style
        #   12+ variants → + draw.io/AI approaches
        alt_styles = ["tech-gradient", "blueprint"]
        variant_plan = []

        if max_variants <= 1:
            variant_plan.append(("v-50/50", slide_style, ["native"]))
        elif max_variants <= 3:
            for split, label in layout_configs:
                variant_plan.append((split, slide_style, ["native"]))
        elif max_variants <= 6:
            for split, label in layout_configs:
                variant_plan.append((split, slide_style, ["native"]))
            for split, label in layout_configs:
                variant_plan.append((split, alt_styles[0], ["native"]))
        elif max_variants <= 9:
            for split, label in layout_configs:
                variant_plan.append((split, slide_style, ["native"]))
            for alt in alt_styles:
                for split, label in layout_configs:
                    variant_plan.append((split, alt, ["native"]))
        else:
            styles = [slide_style] + alt_styles
            approaches = diagram_approaches or ["native", "drawio", "ai"]
            for split, label in layout_configs:
                for sty in styles:
                    variant_plan.append((split, sty, approaches))

        # Cap to max_variants
        variant_plan = variant_plan[:max_variants]

        from diagrams import DiagramEngine
        engine = DiagramEngine(brand=brand)

        variant_num = 0
        for split_ratio, style_name, approaches_list in variant_plan:
            direction = "v" if split_ratio.startswith("v-") else "h"
            top_start = 0.85
            total_avail = 5.0 - top_start
            slide_w = 9.3
            gap = 0.2

            if direction == "v":
                frac = SPLIT_RATIOS[split_ratio][0]
                slot_w = slide_w
                slot_h = (total_avail - gap) * frac
            else:
                frac = SPLIT_RATIOS[split_ratio][0]
                slot_w = (slide_w - gap) * frac
                slot_h = total_avail

            # Get diagram from engine
            diagram_results = []
            if diagram_data:
                diagram_results = engine.generate(
                    diagram_type=diagram_type,
                    data=diagram_data,
                    target_width_in=slot_w,
                    target_height_in=slot_h,
                    approaches=approaches_list,
                    style=style_name,
                )

            if not diagram_results and legacy_diagram:
                diagram_results = [{"approach": "legacy", "type": "png_path",
                                    "label": "Existing diagram file"}]

            if not diagram_results:
                diagram_results = [{"approach": "none", "type": "none", "label": "No diagram"}]

            # One slide per diagram result
            for dr in diagram_results:
                variant_num += 1
                if variant_num > max_variants:
                    break

                split_label = dict(layout_configs).get(split_ratio, split_ratio)
                variant_slide = prs.slides.add_slide(prs.slide_layouts[deck_meta["layout_mapping"]["blank_canvas_idx"]])
                _set_footers(variant_slide, generic_footers, deck_meta)
                set_placeholder_text(variant_slide, 0, slide_def.get("headline", ""), brand=brand)

                variant_def = dict(slide_def)
                variant_def["split"] = split_ratio
                variant_def["_diagram_result"] = dr

                _build_content_diagram_text(variant_slide, variant_def, deck_meta)

                approach_label = dr.get("label", dr.get("approach", ""))
                style_tag = f" [{style_name}]" if style_name != slide_style else ""
                note = f"[Variant {variant_num}/{max_variants}: {split_label} | {approach_label}{style_tag}] {existing_notes}".strip()
                variant_slide.notes_slide.notes_text_frame.text = note

            if variant_num >= max_variants:
                break

        slide = variant_slide
        _variants_handled_notes = True

    elif resolved_name in ("side_by_side", "three_column", "four_card",
                           "big_stat_manual", "closing", "content_table",
                           "content_table_bullets", "section_divider", "quote",
                           "kpi_dashboard", "roadmap", "before_after",
                           "numbered_list", "status_board", "image_showcase",
                           "matrix", "funnel", "callout",
                           "process_flow", "comparison_matrix", "quadrant",
                           "team_profiles", "pros_cons",
                           "staircase", "donut_rings", "pyramid", "venn",
                           "waterfall", "pricing_table", "concentric_circles",
                           "bold_bullet", "cycle_diagram", "hub_spoke",
                           "gauge_dashboard", "risk_heat_map", "tornado_chart",
                           "radar_chart", "combo_chart", "bubble_chart",
                           "bento_grid", "dashboard_panel", "left_nav_sidebar",
                           "image_text_hero"):
        # Manual-render layouts — swap to content_generic canvas
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

        generic_layout = prs.slide_layouts[deck_meta["layout_mapping"]["blank_canvas_idx"]]  # content_generic / CUSTOM_15
        slide = prs.slides.add_slide(generic_layout)
        generic_footers = catalog["layouts"]["content_generic"].get("footer_placeholders", {})
        _set_footers(slide, generic_footers, deck_meta)

        headline = slide_def.get("headline", "")
        # Layouts that handle their own headline rendering
        no_headline_layouts = ("closing", "section_divider", "callout")
        if resolved_name not in no_headline_layouts:
            set_placeholder_text(slide, 0, headline, brand=brand)
        else:
            # Remove the empty headline placeholder to avoid a visible blank box
            for shape in list(slide.placeholders):
                if shape.placeholder_format.idx == 0:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    break

        if resolved_name == "side_by_side":
            _build_side_by_side(slide, slide_def, deck_meta)
        elif resolved_name == "three_column":
            _build_three_column(slide, slide_def, deck_meta)
        elif resolved_name == "four_card":
            _build_four_card(slide, slide_def, deck_meta)
        elif resolved_name == "big_stat_manual":
            _build_big_stat(slide, slide_def, deck_meta)
        elif resolved_name == "closing":
            _build_closing(slide, slide_def, deck_meta)
        elif resolved_name == "content_table":
            _build_content_table(slide, slide_def, deck_meta, prs=prs, catalog=catalog)
        elif resolved_name == "content_table_bullets":
            _build_content_table_bullets(slide, slide_def, deck_meta)
        elif resolved_name == "section_divider":
            _build_section_divider(slide, slide_def, deck_meta)
        elif resolved_name == "quote":
            _build_quote(slide, slide_def, deck_meta)
        elif resolved_name == "kpi_dashboard":
            _build_kpi_dashboard(slide, slide_def, deck_meta)
        elif resolved_name == "roadmap":
            _build_roadmap(slide, slide_def, deck_meta)
        elif resolved_name == "before_after":
            _build_before_after(slide, slide_def, deck_meta)
        elif resolved_name == "numbered_list":
            _build_numbered_list(slide, slide_def, deck_meta)
        elif resolved_name == "status_board":
            _build_status_board(slide, slide_def, deck_meta)
        elif resolved_name == "image_showcase":
            _build_image_showcase(slide, slide_def, deck_meta)
        elif resolved_name == "matrix":
            _build_matrix(slide, slide_def, deck_meta)
        elif resolved_name == "funnel":
            _build_funnel(slide, slide_def, deck_meta)
        elif resolved_name == "callout":
            _build_callout(slide, slide_def, deck_meta)
        elif resolved_name == "process_flow":
            _build_process_flow(slide, slide_def, deck_meta)
        elif resolved_name == "comparison_matrix":
            _build_comparison_matrix(slide, slide_def, deck_meta)
        elif resolved_name == "quadrant":
            _build_quadrant(slide, slide_def, deck_meta)
        elif resolved_name == "team_profiles":
            _build_team_profiles(slide, slide_def, deck_meta)
        elif resolved_name == "pros_cons":
            _build_pros_cons(slide, slide_def, deck_meta)
        elif resolved_name == "staircase":
            _build_staircase(slide, slide_def, deck_meta)
        elif resolved_name == "donut_rings":
            _build_donut_rings(slide, slide_def, deck_meta)
        elif resolved_name == "pyramid":
            _build_pyramid(slide, slide_def, deck_meta)
        elif resolved_name == "venn":
            _build_venn(slide, slide_def, deck_meta)
        elif resolved_name == "waterfall":
            _build_waterfall(slide, slide_def, deck_meta)
        elif resolved_name == "pricing_table":
            _build_pricing_table(slide, slide_def, deck_meta)
        elif resolved_name == "concentric_circles":
            _build_concentric_circles(slide, slide_def, deck_meta)
        elif resolved_name == "bold_bullet":
            _build_bold_bullet(slide, slide_def, deck_meta)
        elif resolved_name == "cycle_diagram":
            _build_cycle_diagram(slide, slide_def, deck_meta)
        elif resolved_name == "hub_spoke":
            _build_hub_spoke(slide, slide_def, deck_meta)
        elif resolved_name == "gauge_dashboard":
            _build_gauge_dashboard(slide, slide_def, deck_meta)
        elif resolved_name == "risk_heat_map":
            _build_risk_heat_map(slide, slide_def, deck_meta)
        elif resolved_name == "tornado_chart":
            _build_tornado_chart(slide, slide_def, deck_meta)
        elif resolved_name == "radar_chart":
            _build_radar_chart(slide, slide_def, deck_meta)
        elif resolved_name == "combo_chart":
            _build_combo_chart(slide, slide_def, deck_meta)
        elif resolved_name == "bubble_chart":
            _build_bubble_chart(slide, slide_def, deck_meta)
        elif resolved_name == "bento_grid":
            _build_bento_grid(slide, slide_def, deck_meta)
        elif resolved_name == "dashboard_panel":
            _build_dashboard_panel(slide, slide_def, deck_meta)
        elif resolved_name == "left_nav_sidebar":
            _build_left_nav_sidebar(slide, slide_def, deck_meta)
        elif resolved_name == "image_text_hero":
            _build_image_text_hero(slide, slide_def, deck_meta)

    elif resolved_name in ("content_two_col", "content_h3_two_col") and slide_def.get("left_stats"):
        # Multi-variant stacked layout — generate graphic style variants as consecutive slides.
        # Per-slide `variants` field overrides: 1 = first style only, default = 3.
        # Remove the placeholder-based slide we just added; we'll create fresh ones.
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

        generic_footers = catalog["layouts"]["content_generic"].get("footer_placeholders", {})
        graphic_styles = ["stats_icons", "stats_purple", "stats_pink"]
        style_labels = ["Stats with icons", "Stats on purple", "Stats on light pink"]

        num_variants = slide_def.get("variants", 3)
        for si, (style, label) in enumerate(zip(graphic_styles[:num_variants], style_labels[:num_variants])):
            variant_slide = prs.slides.add_slide(prs.slide_layouts[deck_meta["layout_mapping"]["blank_canvas_idx"]])
            _set_footers(variant_slide, generic_footers, deck_meta)
            headline = slide_def.get("headline", "")
            set_placeholder_text(variant_slide, 0, headline, brand=brand)
            _build_content_stacked(variant_slide, slide_def, deck_meta, graphic_style=style)
            # Add note identifying the variant
            existing_notes = slide_def.get("notes", "")
            if num_variants > 1:
                variant_note = f"[Variant {si+1}/{num_variants}: {label}] {existing_notes}".strip()
            else:
                variant_note = existing_notes
            variant_slide.notes_slide.notes_text_frame.text = variant_note

        # Return the last slide created (notes handled above, skip default notes below)
        slide = variant_slide
        _variants_handled_notes = True
    else:
        # --- Set content placeholders from YAML fields ---
        for field_name, ph_info in placeholders.items():
            # Skip image and bg_image placeholders (handled separately)
            if field_name in ("image", "bg_image", "image_area"):
                continue

            ph_idx = ph_info["idx"]

            # Try the field name directly, then check aliases
            value = slide_def.get(field_name)
            if value is None:
                # Check if a YAML alias maps to this catalog name
                for yaml_name, catalog_name in FIELD_ALIASES.items():
                    if catalog_name == field_name:
                        value = slide_def.get(yaml_name)
                        if value is not None:
                            break

            if value is None:
                continue

            set_placeholder_text(slide, ph_idx, str(value).strip())

        # --- Handle images ---
        image_path = slide_def.get("image")
        if image_path:
            image_path = resolve_image_path(image_path, deck_meta.get("base_dir", "."))
            if image_path:
                _insert_image(slide, image_path, placeholders, layout_info)

    # --- Handle speaker notes (skip if variants already set their own notes) ---
    if not _variants_handled_notes:
        notes = slide_def.get("notes", "")
        # Add layout recommendation note
        layout_note = _generate_layout_note(resolved_name, slide_def)
        full_notes = notes
        if layout_note:
            full_notes = (notes + "\n\n" + layout_note).strip() if notes else layout_note
        if full_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = full_notes

    return slide


def _set_footers(slide, footer_phs, deck_meta):
    """Set footer placeholders using layout mapping or legacy catalog format.

    Supports two formats:
      1. Legacy (curated-layouts.json): footer_phs dict with footer_left/footer_right keys
      2. Layout mapping: deck_meta["layout_mapping"]["footers"] list with type/default fields

    Footer elements are filled with deck-level values:
      {title} → deck title, {date} → deck date, auto → slide number (left as-is)
    """
    brand = deck_meta.get("brand", BrandConfig())

    def _style_footer(shape, text):
        shape.text = text
        if hasattr(shape, "text_frame"):
            shape.text_frame.word_wrap = False
            for para in shape.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER

    def _resolve_template(template_str):
        """Replace {title} and {date} with deck values."""
        if template_str == 'auto':
            return None  # Leave slide number as-is
        return (template_str
                .replace('{title}', deck_meta.get('title', ''))
                .replace('{date}', deck_meta.get('date', '')))

    # Use layout mapping footers only if NO legacy catalog footers are provided.
    # Legacy catalog has precise repositioning; layout mapping is for
    # templates that don't have a curated catalog entry.
    layout_mapping = deck_meta.get("layout_mapping", {})
    mapping_footers = layout_mapping.get("footers", [])

    if mapping_footers and not footer_phs:
        # New format: list of footer elements with placeholder_idx and default template
        footer_idx_map = {f['placeholder_idx']: f for f in mapping_footers}
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx in footer_idx_map:
                footer = footer_idx_map[idx]
                text = _resolve_template(footer.get('default', ''))
                if text is not None:
                    _style_footer(shape, text)
        return

    # Legacy format: curated-layouts.json with footer_left/footer_right
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx

        if "footer_left" in footer_phs and idx == footer_phs["footer_left"]["idx"]:
            footer_top = Inches(footer_phs["footer_left"]["top"])
            footer_height = Inches(footer_phs["footer_left"]["height"])
            shape.left = Inches(0.35)
            shape.top = footer_top
            shape.width = Inches(4.3)
            shape.height = footer_height
            _style_footer(shape, deck_meta.get("title", ""))

        elif "footer_right" in footer_phs and idx == footer_phs["footer_right"]["idx"]:
            footer_top = Inches(footer_phs["footer_right"]["top"])
            footer_height = Inches(footer_phs["footer_right"]["height"])
            shape.left = Inches(5.35)
            shape.top = footer_top
            shape.width = Inches(4.3)
            shape.height = footer_height
            _style_footer(shape, deck_meta.get("date", ""))


def _split_title(headline):
    """Split a title into lines using separator rules.

    Rules:
    - If headline contains a separator (em dash, en dash, hyphen), split there
      and DROP the separator.
    - Max two lines.
    - Returns (lines, used_split) where used_split indicates whether we split.
    """
    for sep in TITLE_SEPARATORS:
        if sep in headline:
            parts = headline.split(sep, 1)
            return [p.strip() for p in parts], True
    return [headline], False


def _build_title_cover(slide, slide_def, layout_info, deck_meta):
    """Build a title/cover slide with manual text overlay on background image."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    brand = deck_meta.get("brand", BrandConfig())
    rendering = layout_info["title_rendering"]

    # --- Background image (prefer brand config, fall back to gradient) ---
    bg_key = slide_def.get("background", "default")

    bg_path = brand.get_title_background(bg_key)
    if not bg_path or not os.path.exists(bg_path):
        bg_path = brand.get_title_background("default")
    if not bg_path or not os.path.exists(bg_path):
        # Try raw path (user-provided absolute or relative path)
        if os.path.exists(bg_key):
            bg_path = bg_key
        else:
            bg_path = None

    if bg_path and os.path.exists(bg_path):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 2:
                shape.insert_picture(bg_path)
                break
    else:
        # No background image — draw a gradient fill
        from pptx.enum.shapes import MSO_SHAPE
        bg_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SI(0), SI(0), SI(10.0), SI(5.625))
        bg_rect.line.fill.background()
        fill = bg_rect.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = brand.primary
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = brand.interpolate(0.6)
        fill.gradient_stops[1].position = 1.0

    # --- Title text ---
    headline = slide_def.get("headline", "")
    subheader = slide_def.get("subheader", "")

    lines, did_split = _split_title(headline)

    if did_split and len(lines) == 2:
        style = rendering["two_line"]
    else:
        style = rendering["single_line"]

    pos = style["position"]
    font_size = Pt(style["font_size_pt"])
    # Brand fonts/colors always win — catalog values are legacy fallbacks
    font_name = brand.heading_font
    line1_color = brand.white
    line2_color = brand.secondary
    sub_color = brand.white

    txBox = slide.shapes.add_textbox(
        SI(pos["left"]), SI(pos["top"]),
        SI(pos["width"]), SI(pos["height"])
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for li, line_text in enumerate(lines):
        para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line_text
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = line1_color if li == 0 else line2_color
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(0)
        para.space_after = Pt(4)

    # Subtitle
    if subheader:
        sub_para = tf.add_paragraph()
        sub_run = sub_para.add_run()
        sub_run.text = subheader
        sub_run.font.name = brand.body_font
        sub_run.font.size = Pt(style["subtitle_size_pt"])
        sub_run.font.color.rgb = sub_color
        sub_para.alignment = PP_ALIGN.LEFT
        sub_para.space_before = Pt(16)


def _build_agenda(slide, slide_def, layout_info, deck_meta):
    """Build the agenda slide with manual text boxes, numbers, and divider lines.

    Uses content_generic as a clean canvas and draws everything as shapes
    for full control over spacing and line wrapping.
    """
    script_dir = os.path.dirname(os.path.abspath(__file__))
    items = slide_def.get("items", [])

    # Content starts below the headline (bottom edge at ~0.71")
    top_start = 0.85       # first element y, below headline
    bottom_limit = 4.80    # don't go below this

    # --- Left side image (from brand config) ---
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    agenda_image = slide_def.get("agenda_image", "default")

    img_path = brand.get_agenda_background(agenda_image)
    if not img_path or not os.path.exists(img_path):
        # Fall back to hardcoded paths (legacy)
        agenda_images = {
            "default": "title-assets/agenda-left-p22.jpg",
            "p21": "title-assets/agenda-left-p21.jpg",
            "p22": "title-assets/agenda-left-p22.jpg",
            "p23": "title-assets/agenda-left-p23.png",
        }
        img_rel = agenda_images.get(agenda_image, agenda_image)
        img_path = os.path.join(script_dir, img_rel) if not os.path.isabs(img_rel) else img_rel

    img_height = bottom_limit - top_start
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            SI(0.34), SI(top_start),
            SI(3.94), SI(img_height)
        )

    # --- Right side: numbers + items + divider lines ---
    PURPLE = brand.primary

    right_left = 5.0       # where lines start
    num_left = 5.1         # number x position
    text_left = 5.75       # item text x position
    right_edge = 9.62      # where lines end
    text_width = right_edge - text_left

    item_count = len(items)
    if item_count == 0:
        return

    # Font sizes — scale down slightly for many items
    if item_count <= 5:
        num_size, text_size = 14, 14
    elif item_count <= 7:
        num_size, text_size = 12, 12
    else:
        num_size, text_size = 11, 11

    # Estimate line count per item (based on approximate chars per line at this font/width)
    chars_per_line = int(text_width / (text_size * 0.007))  # rough: 1pt ≈ 0.007" width
    line_counts = []
    for item_text in items:
        text = item_text if isinstance(item_text, str) else str(item_text)
        lines = max(1, -(-len(text) // chars_per_line))  # ceiling division
        line_counts.append(lines)

    # Allocate space proportionally: each "line unit" gets equal space, plus padding
    total_units = sum(line_counts)
    available = bottom_limit - top_start
    padding_per_item = 0.18  # minimum gap between items (includes divider line clearance)
    total_padding = padding_per_item * item_count
    space_for_text = available - total_padding
    unit_height = space_for_text / max(total_units, 1)

    # Compute row tops and heights
    row_tops = []
    row_heights = []
    y = top_start
    for lc in line_counts:
        row_tops.append(y)
        h = (lc * unit_height) + padding_per_item
        row_heights.append(h)
        y += h

    from pptx.enum.shapes import MSO_SHAPE

    for i, item_text in enumerate(items):
        row_top = row_tops[i]
        row_height = row_heights[i]

        # Divider line above each row
        line = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            SI(right_left), SI(row_top),
            SI(right_edge), SI(row_top)
        )
        line.line.color.rgb = PURPLE
        line.line.width = Pt(0.75)

        # Number
        text_y = row_top + 0.12  # padding below divider line to prevent overlap
        num_box = slide.shapes.add_textbox(
            SI(num_left), SI(text_y),
            SI(0.55), SI(row_height - 0.08)
        )
        ntf = num_box.text_frame
        ntf.word_wrap = False
        np = ntf.paragraphs[0]
        nr = np.add_run()
        nr.text = f"{i + 1:02d}"
        nr.font.name = HEADING
        nr.font.size = Pt(num_size)
        nr.font.color.rgb = PURPLE
        np.alignment = PP_ALIGN.LEFT

        # Item text
        txt_box = slide.shapes.add_textbox(
            SI(text_left), SI(text_y),
            SI(text_width), SI(row_height - 0.08)
        )
        ttf = txt_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = item_text
        tr.font.name = BODY
        tr.font.size = Pt(text_size)
        tr.font.color.rgb = PURPLE
        tp.alignment = PP_ALIGN.LEFT

    # Final divider line below last item
    final_y = top_start + (item_count * row_height)
    if final_y < bottom_limit:
        line = slide.shapes.add_connector(
            1, SI(right_left), SI(final_y),
            SI(right_edge), SI(final_y)
        )
        line.line.color.rgb = PURPLE
        line.line.width = Pt(0.75)


def _render_text_hierarchical(slide, text_def, x, y, w, max_h, base_sizes=None,
                              deck_meta=None):
    """Render text with typography hierarchy (primary/secondary/tertiary).

    text_def can be:
      - A plain string (rendered as secondary)
      - A dict with 'primary', 'secondary', 'tertiary' keys

    Returns the actual height used.
    """
    sizes = dict(TYPO_HIERARCHY)
    if base_sizes:
        for level, overrides in base_sizes.items():
            if level in sizes:
                sizes[level].update(overrides)

    brand = (deck_meta or {}).get("brand", BrandConfig())
    font_name = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark
    GRAY = brand.text_gray

    level_colors = {
        "primary": PURPLE,
        "secondary": DARK,
        "tertiary": GRAY,
    }

    # Normalize to dict
    if isinstance(text_def, str):
        text_def = {"secondary": text_def}

    # Calculate total height needed
    total_h = 0
    sections = []
    for level in ["primary", "secondary", "tertiary"]:
        text = text_def.get(level, "").strip()
        if not text:
            continue
        sz = sizes[level]["size"]
        h = estimate_text_height(text, w, sz)
        sections.append((level, text, sz, h))
        total_h += h + 0.05  # gap between levels

    if not sections:
        return 0

    # Render
    cur_y = y
    for level, text, sz, h in sections:
        is_bold = sizes[level]["bold"]
        color = level_colors[level]

        txb = slide.shapes.add_textbox(SI(x), SI(cur_y), SI(w), SI(h))
        tf = txb.text_frame
        tf.word_wrap = True

        for li, line in enumerate(text.split("\n")):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            r = para.add_run()
            r.text = line
            r.font.name = font_name
            r.font.size = Pt(sz)
            r.font.bold = is_bold
            r.font.color.rgb = color

        cur_y += h + 0.05

    return cur_y - y


def _calc_split(split_name, total_avail, min_text_h):
    """Calculate diagram and text heights from a split ratio.

    If text doesn't fit at the requested ratio, bumps toward more text.
    Returns (diagram_h, text_h, actual_split_name, was_bumped).
    """
    # Parse direction
    direction = "v" if split_name.startswith("v-") else "h"

    # Try the requested ratio first, then progressively favor text
    if direction == "v":
        candidates = ["v-70/30", "v-60/40", "v-50/50", "v-40/60", "v-30/70"]
    else:
        candidates = ["h-60/40", "h-50/50", "h-40/60"]

    # Start from the requested ratio
    try:
        start_idx = candidates.index(split_name)
    except ValueError:
        start_idx = len(candidates) // 2  # default to middle

    gap = 0.2  # gap between diagram and text zones

    for i in range(start_idx, len(candidates)):
        ratio_name = candidates[i]
        diag_frac, text_frac = SPLIT_RATIOS[ratio_name]
        diag_h = (total_avail - gap) * diag_frac
        text_h = (total_avail - gap) * text_frac

        if text_h >= min_text_h:
            was_bumped = (i != start_idx)
            return diag_h, text_h, ratio_name, was_bumped

    # Last resort: give text whatever it needs, diagram gets the rest
    text_h = min_text_h
    diag_h = total_avail - gap - text_h
    return max(diag_h, 0.5), text_h, "auto", True


def _build_content_diagram_text(slide, slide_def, deck_meta):
    """Build a slide with diagram + two text sections using split ratios.

    Supports vertical splits (diagram on top) and horizontal splits (diagram left/right).
    Uses typography hierarchy for text content.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark

    icon_size = 0.4
    title_h = 0.30
    pad = 0.12

    split_name = slide_def.get("split", "v-50/50")
    direction = "v" if split_name.startswith("v-") else "h"

    # Resolve diagram source — engine result, legacy path, or none
    diagram_result = slide_def.get("_diagram_result", None)
    native_shapes_fn = None
    diagram_path = None

    if diagram_result:
        if diagram_result.get("type") == "shapes":
            native_shapes_fn = diagram_result["shapes_fn"]
        elif diagram_result.get("type") == "png":
            diagram_path = diagram_result["path"]
        elif diagram_result.get("type") == "png_path":
            # Legacy fallback
            dp = slide_def.get("diagram", "")
            if dp and not os.path.isabs(dp):
                dp = os.path.join(deck_meta.get("base_dir", "."), dp)
            diagram_path = dp if os.path.exists(dp) else None
    else:
        dp = slide_def.get("diagram", "")
        if dp:
            if not os.path.isabs(dp):
                dp = os.path.join(deck_meta.get("base_dir", "."), dp)
            diagram_path = dp if os.path.exists(dp) else None

    left_title = slide_def.get("left_title", "")
    left_body = slide_def.get("left_body", "")
    right_title = slide_def.get("right_title", "")
    right_body = slide_def.get("right_body", "")
    left_icon_name = slide_def.get("left_icon", "")
    right_icon_name = slide_def.get("right_icon", "")

    left_icon = _resolve_icon(left_icon_name) if left_icon_name else None
    right_icon = _resolve_icon(right_icon_name) if right_icon_name else None

    # Get diagram aspect ratio (for PNG-based diagrams)
    diag_aspect = 16 / 9  # default
    if diagram_path and os.path.exists(diagram_path):
        try:
            with Image.open(diagram_path) as img:
                diag_aspect = img.width / img.height
        except Exception:
            pass

    top_start = 0.85
    bottom_limit = 5.0
    total_avail = bottom_limit - top_start
    slide_content_w = 9.3  # usable width

    if direction == "v":
        # --- VERTICAL: diagram on top, text sections below ---
        col_w = 4.0
        text_indent = 0.55  # after icon

        # Calculate minimum text height
        left_h = estimate_text_height(
            left_body if isinstance(left_body, str) else left_body.get("secondary", ""),
            col_w, 10)
        right_h = estimate_text_height(
            right_body if isinstance(right_body, str) else right_body.get("secondary", ""),
            col_w, 10)
        max_body_h = max(left_h, right_h)
        min_text_h = icon_size + pad + max_body_h

        diag_h, text_h, actual_split, bumped = _calc_split(
            split_name, total_avail, min_text_h)

        if bumped:
            print(f"    NOTE: Split bumped from {split_name} to {actual_split} to fit text")

        # Size diagram to fit the allocated height while maintaining aspect ratio
        diag_render_h = diag_h
        diag_render_w = diag_render_h * diag_aspect
        if diag_render_w > slide_content_w:
            diag_render_w = slide_content_w
            diag_render_h = diag_render_w / diag_aspect

        # Center diagram horizontally
        diag_left = 0.35 + (slide_content_w - diag_render_w) / 2
        diag_top = top_start + (diag_h - diag_render_h) / 2  # center in zone

        if native_shapes_fn:
            # Native shapes render within the full diagram zone — centered internally
            native_shapes_fn(slide, 0.35, top_start)
        elif diagram_path and os.path.exists(diagram_path):
            slide.shapes.add_picture(
                diagram_path,
                SI(diag_left), SI(diag_top),
                SI(diag_render_w), SI(diag_render_h))

        # Text sections below
        gap = 0.2
        sec_top = top_start + diag_h + gap

        # Left section
        if left_icon and os.path.exists(left_icon):
            slide.shapes.add_picture(left_icon,
                                     SI(0.35), SI(sec_top),
                                     SI(icon_size), SI(icon_size))
        txb = slide.shapes.add_textbox(
            SI(0.9), SI(sec_top + 0.05),
            SI(col_w), SI(title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = left_title; r.font.name = HEADING
        r.font.size = Pt(14); r.font.color.rgb = PURPLE

        body_top = sec_top + icon_size + pad
        if isinstance(left_body, dict):
            _render_text_hierarchical(slide, left_body, 0.9, body_top, col_w, max_body_h, deck_meta=deck_meta)
        else:
            txb = slide.shapes.add_textbox(
                SI(0.9), SI(body_top),
                SI(col_w), SI(max_body_h))
            tf = txb.text_frame; tf.word_wrap = True
            _render_body_text(tf, left_body.strip(), BODY, 10, DARK)

        # Right section
        if right_icon and os.path.exists(right_icon):
            slide.shapes.add_picture(right_icon,
                                     SI(5.1), SI(sec_top),
                                     SI(icon_size), SI(icon_size))
        txb = slide.shapes.add_textbox(
            SI(5.65), SI(sec_top + 0.05),
            SI(col_w), SI(title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = right_title; r.font.name = HEADING
        r.font.size = Pt(14); r.font.color.rgb = PURPLE

        if isinstance(right_body, dict):
            _render_text_hierarchical(slide, right_body, 5.65, body_top, col_w, max_body_h, deck_meta=deck_meta)
        else:
            txb = slide.shapes.add_textbox(
                SI(5.65), SI(body_top),
                SI(col_w), SI(max_body_h))
            tf = txb.text_frame; tf.word_wrap = True
            _render_body_text(tf, right_body.strip(), BODY, 10, DARK)

    else:
        # --- HORIZONTAL: diagram left, text right (stacked) ---
        diag_frac, text_frac = SPLIT_RATIOS.get(split_name, (0.5, 0.5))
        gap = 0.2

        diag_zone_w = (slide_content_w - gap) * diag_frac
        text_zone_w = (slide_content_w - gap) * text_frac
        text_w = text_zone_w - 0.55  # after icon indent

        right_x = 0.35 + diag_zone_w + gap

        # Diagram: fit within left zone
        diag_render_w = diag_zone_w
        diag_render_h = diag_render_w / diag_aspect
        if diag_render_h > total_avail:
            diag_render_h = total_avail
            diag_render_w = diag_render_h * diag_aspect

        diag_left = 0.35 + (diag_zone_w - diag_render_w) / 2
        diag_top = top_start + (total_avail - diag_render_h) / 2

        if native_shapes_fn:
            native_shapes_fn(slide, 0.35, top_start)
        elif diagram_path and os.path.exists(diagram_path):
            slide.shapes.add_picture(
                diagram_path,
                SI(diag_left), SI(diag_top),
                SI(diag_render_w), SI(diag_render_h))

        # Text sections stacked on right
        sec1_body_h = estimate_text_height(
            left_body if isinstance(left_body, str) else left_body.get("secondary", ""),
            text_w, 10)
        sec2_body_h = estimate_text_height(
            right_body if isinstance(right_body, str) else right_body.get("secondary", ""),
            text_w, 10)

        right_total = title_h + 0.05 + sec1_body_h + pad*2 + title_h + 0.05 + sec2_body_h
        offset = top_start + (total_avail - right_total) / 2
        offset = max(offset, top_start)

        sec1_top = offset
        sec1_tb = sec1_top + title_h + 0.05
        div_y = sec1_tb + sec1_body_h + pad
        sec2_top = div_y + pad

        if left_icon and os.path.exists(left_icon):
            slide.shapes.add_picture(left_icon, SI(right_x), SI(sec1_top),
                                     SI(icon_size), SI(icon_size))
        txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec1_top + 0.05),
                                       SI(text_w), SI(title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = left_title; r.font.name = HEADING
        r.font.size = Pt(14); r.font.color.rgb = PURPLE

        if isinstance(left_body, dict):
            _render_text_hierarchical(slide, left_body, right_x + 0.55, sec1_tb, text_w, sec1_body_h, deck_meta=deck_meta)
        else:
            txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec1_tb),
                                           SI(text_w), SI(sec1_body_h))
            tf = txb.text_frame; tf.word_wrap = True
            _render_body_text(tf, left_body.strip(), BODY, 10, DARK)

        line = slide.shapes.add_connector(1, SI(right_x), SI(div_y),
                                          SI(right_x + text_zone_w), SI(div_y))
        line.line.color.rgb = PURPLE; line.line.width = Pt(0.75)

        if right_icon and os.path.exists(right_icon):
            slide.shapes.add_picture(right_icon, SI(right_x), SI(sec2_top),
                                     SI(icon_size), SI(icon_size))
        txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec2_top + 0.05),
                                       SI(text_w), SI(title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = right_title; r.font.name = HEADING
        r.font.size = Pt(14); r.font.color.rgb = PURPLE

        sec2_body_top = sec2_top + title_h + 0.05
        if isinstance(right_body, dict):
            _render_text_hierarchical(slide, right_body, right_x + 0.55, sec2_body_top, text_w, sec2_body_h, deck_meta=deck_meta)
        else:
            txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec2_body_top),
                                           SI(text_w), SI(sec2_body_h))
            tf = txb.text_frame; tf.word_wrap = True
            _render_body_text(tf, right_body.strip(), BODY, 10, DARK)


# Module-level icon directory — set by build_deck() from brand config
_icon_dir = None


def _resolve_icon(icon_name):
    """Resolve an icon name to a file path using the icon catalog."""
    global _icon_dir
    # Use brand icon dir if set, otherwise fall back to script_dir/icons
    if _icon_dir:
        icons_dir = _icon_dir
    else:
        icons_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "icons")
    catalog_path = os.path.join(icons_dir, "icon-catalog.json")
    try:
        with open(catalog_path) as f:
            icon_catalog = json.load(f)
        info = icon_catalog["icons"].get(icon_name, {})
        if info.get("file"):
            path = os.path.join(icons_dir, info["file"] + ".png")
            if os.path.exists(path):
                return path
            print(f"    WARNING: Icon file missing for '{icon_name}': {path}")
            return None
        print(f"    WARNING: Icon '{icon_name}' not found in icon-catalog.json")
    except Exception as e:
        print(f"    WARNING: Icon catalog error for '{icon_name}': {e}")
    return None


def _build_content_stacked(slide, slide_def, deck_meta, graphic_style="stats_icons"):
    """Build a two-section content slide with stacked layout.

    Left side: stats graphic (style varies by graphic_style parameter).
    Right side: two sections stacked vertically with icons and divider.
    Vertically centered in available space.

    graphic_style options:
        "stats_icons"  — icons + purple numbers on light purple bg (preferred)
        "stats_purple" — white numbers on brand primary bg
        "stats_pink"   — pink numbers on light purple bg
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    DARK = brand.text_dark
    WHITE = WHITE
    PINK = brand.secondary
    LIGHT_PINK = brand.accent
    LIGHT_BG = brand.bg_light

    icon_size = 0.4
    title_size_pt = 14
    body_size_pt = 10
    title_h = 0.30
    pad = 0.15

    left_title = slide_def.get("left_title", "")
    left_body = slide_def.get("left_body", "").strip()
    right_title = slide_def.get("right_title", "")
    right_body = slide_def.get("right_body", "").strip()
    left_stats = slide_def.get("left_stats", [])
    left_icon_name = slide_def.get("left_icon", "")
    right_icon_name = slide_def.get("right_icon", "")

    # --- Calculate heights for BOTH sides before positioning anything ---
    right_x = 4.5
    section_w = 5.1
    text_w = section_w - 0.55  # after icon indent

    sec1_body_h = estimate_text_height(left_body, text_w, body_size_pt)
    sec2_body_h = estimate_text_height(right_body, text_w, body_size_pt)

    # Right side total height
    right_total_h = (
        title_h + 0.05 + sec1_body_h  # section 1
        + pad + pad                     # divider gaps
        + title_h + 0.05 + sec2_body_h  # section 2
    )

    # Left side total height (stats)
    show_icons = (graphic_style == "stats_icons")
    stat_num_font = 30 if show_icons else 36
    stat_num_h = stat_num_font * 1.4 / 72 + 0.05
    stat_label_gap = 0.04
    stat_label_h = 0.25
    stat_divider_gap = 0.16  # 0.08 above + 0.08 below
    stat_count = len(left_stats) if left_stats else 0

    if stat_count > 0:
        single_stat_h = stat_num_h + stat_label_gap + stat_label_h
        left_total_h = (
            stat_count * single_stat_h
            + (stat_count - 1) * stat_divider_gap
            + 0.30  # top/bottom padding inside background
        )
    else:
        left_total_h = 0

    # Use the taller side to size the content block
    content_block_h = max(left_total_h, right_total_h)

    # Vertically center in available space (below headline at 0.71, above footer at ~5.0)
    avail = 4.8 - 0.85
    offset = 0.85 + (avail - content_block_h) / 2
    offset = max(offset, 0.85)

    # --- Right side positions ---
    sec1_top = offset
    sec1_title_bottom = sec1_top + title_h + 0.05
    sec1_body_bottom = sec1_title_bottom + sec1_body_h
    div_y = sec1_body_bottom + pad
    sec2_top = div_y + pad

    # Resolve icons
    left_icon = _resolve_icon(left_icon_name) if left_icon_name else None
    right_icon = _resolve_icon(right_icon_name) if right_icon_name else None

    # Section 1
    if left_icon and os.path.exists(left_icon):
        slide.shapes.add_picture(left_icon, SI(right_x), SI(sec1_top),
                                 SI(icon_size), SI(icon_size))
    txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec1_top + 0.05),
                                   SI(text_w), SI(title_h))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = left_title; r.font.name = HEADING
    r.font.size = Pt(title_size_pt); r.font.color.rgb = PURPLE

    txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec1_title_bottom),
                                   SI(text_w), SI(sec1_body_h))
    tf = txb.text_frame; tf.word_wrap = True
    _render_body_text(tf, left_body, BODY, body_size_pt, DARK)

    # Divider
    line = slide.shapes.add_connector(
        1, SI(right_x), SI(div_y),
        SI(right_x + section_w), SI(div_y))
    line.line.color.rgb = PURPLE; line.line.width = Pt(0.75)

    # Section 2
    if right_icon and os.path.exists(right_icon):
        slide.shapes.add_picture(right_icon, SI(right_x), SI(sec2_top),
                                 SI(icon_size), SI(icon_size))
    txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec2_top + 0.05),
                                   SI(text_w), SI(title_h))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = right_title; r.font.name = HEADING
    r.font.size = Pt(title_size_pt); r.font.color.rgb = PURPLE

    sec2_body_top = sec2_top + title_h + 0.05
    txb = slide.shapes.add_textbox(SI(right_x + 0.55), SI(sec2_body_top),
                                   SI(text_w), SI(sec2_body_h))
    tf = txb.text_frame; tf.word_wrap = True
    _render_body_text(tf, right_body, BODY, body_size_pt, DARK)

    content_bottom = sec2_body_top + sec2_body_h

    # --- Left side: stats graphic ---
    graphic_x = 0.35
    graphic_w = 3.8
    graphic_top = offset - 0.15
    graphic_bottom = offset + content_block_h + 0.15

    # Background color depends on style
    if graphic_style == "stats_purple":
        bg_color = PURPLE
    else:
        bg_color = LIGHT_BG

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SI(graphic_x), SI(graphic_top),
        SI(graphic_w), SI(graphic_bottom - graphic_top))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    if left_stats:
        # Center stats vertically within graphic area
        single_stat_h_actual = stat_num_h + stat_label_gap + stat_label_h
        total_stats_h = stat_count * single_stat_h_actual + (stat_count - 1) * stat_divider_gap
        graphic_area_h = graphic_bottom - graphic_top
        stat_start = graphic_top + (graphic_area_h - total_stats_h) / 2
        stat_spacing = single_stat_h_actual + stat_divider_gap

        # Style-specific colors
        if graphic_style == "stats_purple":
            num_color = WHITE
            label_color = LIGHT_PINK
            divider_color = brand.interpolate(0.3)
            show_icons = False
        elif graphic_style == "stats_pink":
            num_color = PINK
            label_color = PURPLE
            divider_color = brand.divider
            show_icons = False
        else:  # stats_icons (default/preferred)
            num_color = PURPLE
            label_color = DARK
            divider_color = brand.divider
            show_icons = True

        for i, stat in enumerate(left_stats):
            sy = stat_start + i * stat_spacing

            if show_icons:
                # Icon + left-aligned number/label
                stat_icon_name = stat.get("icon", "")
                stat_icon = _resolve_icon(stat_icon_name) if stat_icon_name else None
                if stat_icon and os.path.exists(stat_icon):
                    slide.shapes.add_picture(stat_icon,
                                             SI(graphic_x + 0.3), SI(sy + 0.1),
                                             SI(0.5), SI(0.5))
                num_x = graphic_x + 1.0
                num_w = 2.8
                align = PP_ALIGN.LEFT
            else:
                # Centered number/label, no icon
                num_x = graphic_x
                num_w = graphic_w
                align = PP_ALIGN.CENTER

            # Size depends on style
            num_font_pt = 30 if show_icons else 36
            num_h = num_font_pt * 1.4 / 72 + 0.05  # line height + padding
            label_gap = 0.04
            label_font_pt = 10 if show_icons else 11
            label_h = 0.25

            # Number
            txb = slide.shapes.add_textbox(
                SI(num_x), SI(sy), SI(num_w), SI(num_h))
            tf = txb.text_frame; tf.word_wrap = False
            r = tf.paragraphs[0].add_run()
            r.text = stat.get("number", "")
            r.font.name = HEADING
            r.font.size = Pt(num_font_pt)
            r.font.color.rgb = num_color
            tf.paragraphs[0].alignment = align

            # Label — positioned below number with proper gap
            label_y = sy + num_h + label_gap
            txb = slide.shapes.add_textbox(
                SI(num_x), SI(label_y), SI(num_w), SI(label_h))
            tf = txb.text_frame; tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = stat.get("label", "")
            r.font.name = BODY
            r.font.size = Pt(label_font_pt)
            r.font.color.rgb = label_color
            tf.paragraphs[0].alignment = align

            # Divider between stats
            if i < stat_count - 1:
                dy = label_y + stat_label_h + 0.08
                if show_icons:
                    dl_left = SI(graphic_x + 0.3)
                    dl_right = SI(graphic_x + graphic_w - 0.3)
                else:
                    dl_left = SI(graphic_x + 0.5)
                    dl_right = SI(graphic_x + graphic_w - 0.5)
                dl = slide.shapes.add_connector(1, dl_left, SI(dy), dl_right, SI(dy))
                dl.line.color.rgb = divider_color
                dl.line.width = Pt(0.5)


def _build_side_by_side(slide, slide_def, deck_meta):
    """Build a two-column side-by-side slide with icons, vertically centered.

    Fields: headline, left_icon, left_title, left_body,
            right_icon, right_title, right_body.
    Uses content_generic (layout 2) as canvas.
    """
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark

    icon_size = 0.4
    title_h = 0.30
    pad = 0.12
    title_size_pt = 14
    body_size_pt = 10

    col_w = 4.0
    left_x = 0.35
    right_x = 5.1
    text_indent = 0.55  # icon width + gap

    left_title = slide_def.get("left_title", "")
    left_body = slide_def.get("left_body", "").strip()
    right_title = slide_def.get("right_title", "")
    right_body = slide_def.get("right_body", "").strip()
    left_icon_name = slide_def.get("left_icon", "")
    right_icon_name = slide_def.get("right_icon", "")

    left_icon = _resolve_icon(left_icon_name) if left_icon_name else None
    right_icon = _resolve_icon(right_icon_name) if right_icon_name else None

    # Calculate text heights
    left_body_h = estimate_text_height(left_body, col_w - text_indent, body_size_pt)
    right_body_h = estimate_text_height(right_body, col_w - text_indent, body_size_pt)
    max_body_h = max(left_body_h, right_body_h)

    # Total content height: icon row (title beside icon) + gap + body
    content_h = icon_size + pad + max_body_h

    # Vertically center between headline bottom (0.71") and footer top (~5.0")
    avail_top = 0.85
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - content_h) / 2
    offset = max(offset, avail_top)

    # Align icon center with title text center vertically
    title_top = offset + (icon_size - title_h) / 2
    body_top = offset + icon_size + pad

    # --- Left column ---
    if left_icon and os.path.exists(left_icon):
        slide.shapes.add_picture(left_icon,
                                 SI(left_x), SI(offset),
                                 SI(icon_size), SI(icon_size))
    txb = slide.shapes.add_textbox(
        SI(left_x + text_indent), SI(title_top),
        SI(col_w - text_indent), SI(title_h))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = left_title; r.font.name = HEADING
    r.font.size = Pt(title_size_pt); r.font.color.rgb = PURPLE

    # Use full available height for body text boxes — allows even spacing
    body_avail_h = avail_bottom - body_top

    txb = slide.shapes.add_textbox(
        SI(left_x + text_indent), SI(body_top),
        SI(col_w - text_indent), SI(body_avail_h))
    tf = txb.text_frame; tf.word_wrap = True
    _render_body_text(tf, left_body, BODY, body_size_pt, DARK,
                      available_height=body_avail_h)

    # --- Right column ---
    if right_icon and os.path.exists(right_icon):
        slide.shapes.add_picture(right_icon,
                                 SI(right_x), SI(offset),
                                 SI(icon_size), SI(icon_size))
    txb = slide.shapes.add_textbox(
        SI(right_x + text_indent), SI(title_top),
        SI(col_w - text_indent), SI(title_h))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = right_title; r.font.name = HEADING
    r.font.size = Pt(title_size_pt); r.font.color.rgb = PURPLE

    txb = slide.shapes.add_textbox(
        SI(right_x + text_indent), SI(body_top),
        SI(col_w - text_indent), SI(body_avail_h))
    tf = txb.text_frame; tf.word_wrap = True
    _render_body_text(tf, right_body, BODY, body_size_pt, DARK,
                      available_height=body_avail_h)


def _build_three_column(slide, slide_def, deck_meta):
    """Build a three-column slide with icons, content-driven vertical positioning.

    Fields: headline, col1_icon, col1_title, col1_body,
            col2_icon, col2_title, col2_body, col3_icon, col3_title, col3_body.
    Also supports YAML aliases: section1_title->col1_title, etc.
    Uses content_generic (layout 2) as canvas.
    """
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark

    icon_size = 0.4
    title_h = 0.30
    pad = 0.12
    title_size_pt = 12
    body_size_pt = 9

    col_positions = [0.35, 3.55, 6.75]
    col_w = 2.9
    text_indent = 0.55

    # Resolve fields with alias support
    columns = []
    for i in range(1, 4):
        icon_name = slide_def.get(f"col{i}_icon", slide_def.get(f"section{i}_icon", ""))
        title = slide_def.get(f"col{i}_title", slide_def.get(f"section{i}_title", ""))
        body = slide_def.get(f"col{i}_body", slide_def.get(f"section{i}_body", "")).strip()
        icon = _resolve_icon(icon_name) if icon_name else None
        columns.append({"icon": icon, "title": title, "body": body})

    # Calculate text heights
    body_heights = []
    for col in columns:
        h = estimate_text_height(col["body"], col_w - text_indent, body_size_pt)
        body_heights.append(h)
    max_body_h = max(body_heights) if body_heights else 0.5

    # Total content height: icon row + gap + body
    content_h = icon_size + pad + max_body_h

    # Vertically center between headline bottom (0.71") and footer top (~5.0")
    avail_top = 0.85
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - content_h) / 2
    offset = max(offset, avail_top)

    # Align icon center with title text center vertically
    title_top = offset + (icon_size - title_h) / 2
    body_top = offset + icon_size + pad

    for ci, col in enumerate(columns):
        cx = col_positions[ci]

        # Icon
        if col["icon"] and os.path.exists(col["icon"]):
            slide.shapes.add_picture(col["icon"],
                                     SI(cx), SI(offset),
                                     SI(icon_size), SI(icon_size))

        # Title (beside icon, vertically centered with icon)
        txb = slide.shapes.add_textbox(
            SI(cx + text_indent), SI(title_top),
            SI(col_w - text_indent), SI(title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = col["title"]; r.font.name = HEADING
        r.font.size = Pt(title_size_pt); r.font.color.rgb = PURPLE

        # Body
        col_body_h = estimate_text_height(col["body"], col_w - text_indent, body_size_pt)
        txb = slide.shapes.add_textbox(
            SI(cx + text_indent), SI(body_top),
            SI(col_w - text_indent), SI(col_body_h))
        tf = txb.text_frame; tf.word_wrap = True
        _render_body_text(tf, col["body"], BODY, body_size_pt, DARK)


def _build_four_card(slide, slide_def, deck_meta):
    """Build a dynamic card slide — renders only cards that have content.

    Scans card1_title through card8_title, renders only non-empty cards.
    Card width adjusts automatically based on count.
    Font scales down for 6+ cards.

    Fields: headline, card1_title, card1_body, card2_title, card2_body, ...
    Uses content_generic (layout 2) as canvas.
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark
    LIGHT_BG = brand.bg_light

    gap = 0.1
    card_pad_x = 0.12
    card_pad_top = 0.12
    card_pad_bottom = 0.12
    title_h = 0.28
    usable_w = 9.3  # total usable width
    start_x = 0.35

    # Gather cards — support both array format and flat key format
    cards = []
    if "cards" in slide_def and isinstance(slide_def["cards"], list):
        # Array format: cards: [{title: ..., body: ...}, ...]
        for card in slide_def["cards"]:
            title = card.get("title", "")
            body = card.get("body", "").strip() if card.get("body") else ""
            if title:
                cards.append({"title": title, "body": body})
    else:
        # Flat key format: card1_title, card1_body, card2_title, ...
        for i in range(1, 9):  # support up to 8 cards
            title = slide_def.get(f"card{i}_title", "")
            body = slide_def.get(f"card{i}_body", "").strip()
            if title:
                cards.append({"title": title, "body": body})

    n = len(cards)
    if n == 0:
        return

    # Dynamic card width based on count
    card_w = (usable_w - (n - 1) * gap) / n
    card_w = min(card_w, 3.0)  # cap at 3" for 1-2 cards

    # Adaptive font sizing
    if n <= 4:
        title_size_pt = 11
        body_size_pt = 9
    elif n <= 6:
        title_size_pt = 10
        body_size_pt = 8
    else:
        title_size_pt = 9
        body_size_pt = 7

    # Recenter if cards don't fill full width
    total_w = n * card_w + (n - 1) * gap
    start_x = 0.35 + (usable_w - total_w) / 2

    # Calculate body heights for content-driven card sizing
    text_w = card_w - 2 * card_pad_x
    body_heights = []
    for card in cards:
        h = estimate_text_height(card["body"], text_w, body_size_pt)
        body_heights.append(h)
    max_body_h = max(body_heights) if body_heights else 0.5

    # Card height
    card_h = card_pad_top + title_h + 0.08 + max_body_h + card_pad_bottom

    # Vertically center
    avail_top = 0.85
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - card_h) / 2
    offset = max(offset, avail_top)

    for ci, card in enumerate(cards):
        cx = start_x + ci * (card_w + gap)

        # Background rectangle
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SI(cx), SI(offset),
            SI(card_w), SI(card_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()

        # Title
        title_top = offset + card_pad_top
        txb = slide.shapes.add_textbox(
            SI(cx + card_pad_x), SI(title_top),
            SI(text_w), SI(title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = card["title"]; r.font.name = HEADING
        r.font.size = Pt(title_size_pt); r.font.color.rgb = PURPLE

        # Body
        body_top = title_top + title_h + 0.08
        body_h = estimate_text_height(card["body"], text_w, body_size_pt)
        txb = slide.shapes.add_textbox(
            SI(cx + card_pad_x), SI(body_top),
            SI(text_w), SI(body_h))
        tf = txb.text_frame; tf.word_wrap = True
        _render_body_text(tf, card["body"], BODY, body_size_pt, DARK)


def _build_big_stat(slide, slide_def, deck_meta):
    """Build a hero number/stat slide, centered on the canvas.

    Fields: headline, number, label, icon.
    Uses content_generic (layout 2) as canvas.
    """
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark

    icon_size = 0.6
    number_size_pt = 72
    label_size_pt = 16

    number_text = slide_def.get("number", "")
    label_text = slide_def.get("label", slide_def.get("description", ""))
    icon_name = slide_def.get("icon", "")

    icon_path = _resolve_icon(icon_name) if icon_name else None

    # Calculate heights
    number_h = estimate_text_height(str(number_text), 9.0, number_size_pt)
    label_h = estimate_text_height(str(label_text), 6.0, label_size_pt) if label_text else 0
    icon_h = icon_size + 0.15 if (icon_path and os.path.exists(icon_path)) else 0

    total_h = icon_h + number_h + 0.1 + label_h

    # Vertically center between headline bottom (0.71") and footer top (~5.0")
    avail_top = 0.85
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - total_h) / 2
    offset = max(offset, avail_top)

    cur_y = offset

    # Optional icon — centered
    if icon_path and os.path.exists(icon_path):
        icon_left = (10.0 - icon_size) / 2
        slide.shapes.add_picture(icon_path,
                                 SI(icon_left), SI(cur_y),
                                 SI(icon_size), SI(icon_size))
        cur_y += icon_size + 0.15

    # Number — large, centered
    txb = slide.shapes.add_textbox(
        SI(0.5), SI(cur_y),
        SI(9.0), SI(number_h))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = str(number_text); r.font.name = HEADING
    r.font.size = Pt(number_size_pt); r.font.color.rgb = PURPLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    cur_y += number_h + 0.1

    # Label — below number, centered
    if label_text:
        txb = slide.shapes.add_textbox(
            SI(2.0), SI(cur_y),
            SI(6.0), SI(label_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = str(label_text); r.font.name = BODY
        r.font.size = Pt(label_size_pt); r.font.color.rgb = DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _build_content_table(slide, slide_def, deck_meta, prs=None, catalog=None):
    """Build a slide with a styled table.

    Fields: headline, columns (list of header strings),
            rows (list of lists — each inner list is one row of cell values),
            col_widths (optional proportional weights, auto-calculated if omitted),
            col_align (optional — single string or per-column list: left/center/right),
            header_bg (optional hex color, default brand primary),
            header_style (optional — "bold" (default), "allcaps"),
            stripe (optional bool, default True — alternate row shading).

    Cell text supports ``**bold**`` and ``- `` bullet prefixes.
    Tables that overflow are automatically split across continuation slides.
    Uses content_generic (layout 2) as canvas.
    """
    from pptx.oxml.ns import qn

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    WHITE = brand.white
    DARK = brand.text_dark

    columns = slide_def.get("columns", [])
    rows = slide_def.get("rows", [])
    if not columns:
        return

    num_cols = len(columns)

    # Table position
    left_margin = 0.35
    avail_top = 0.85
    avail_bottom = 5.05
    table_width = 10.0 - left_margin - 0.35
    avail_h = avail_bottom - avail_top

    # Column widths — proportional normalization or content-aware auto-sizing
    raw_col_widths = slide_def.get("col_widths", None)
    col_widths = normalize_col_widths(raw_col_widths, num_cols, table_width,
                                      columns, rows)

    # Column alignment
    col_align_raw = slide_def.get("col_align", "left")
    if isinstance(col_align_raw, str):
        col_align_raw = [col_align_raw] * num_cols
    _align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                  "right": PP_ALIGN.RIGHT}
    col_align = [_align_map.get(a, PP_ALIGN.LEFT) for a in col_align_raw]

    # Style options
    header_bg_hex = slide_def.get("header_bg", brand.primary_hex)
    stripe = slide_def.get("stripe", True)
    header_style = slide_def.get("header_style", "bold")

    # Font selection + content-aware row heights
    body_font_size_pt, header_font_size_pt, header_h, row_heights, total_h = \
        select_font_size(columns, rows, col_widths, avail_h)

    # --- Multi-slide splitting ---
    if total_h > avail_h and prs is not None and len(rows) > 1:
        chunks = split_table_rows(columns, rows, col_widths, avail_h,
                                  body_font_size_pt)
        # Render first chunk on this slide
        rows = chunks[0]
        # Build continuation slides for remaining chunks
        for chunk_rows in chunks[1:]:
            headline = slide_def.get("headline", "")
            if not headline.endswith("(continued)"):
                headline = headline + " (continued)"
            cont_def = {**slide_def, "rows": chunk_rows, "headline": headline}
            build_slide(prs, cont_def, deck_meta, catalog)
        # Recompute heights for the (now smaller) first chunk
        body_font_size_pt, header_font_size_pt, header_h, row_heights, total_h = \
            select_font_size(columns, rows, col_widths, avail_h)
    elif total_h > avail_h:
        print(f"    WARNING: Table overflow on "
              f"'{slide_def.get('headline', '')}' — "
              f"{len(rows)+1} rows at {body_font_size_pt}pt = "
              f"{total_h:.2f}\" exceeds {avail_h:.2f}\" available")

    num_rows = len(rows) + 1  # +1 for header
    table_height = total_h
    header_font_size = Pt(header_font_size_pt)
    body_font_size = Pt(body_font_size_pt)

    # Vertically center the table in available space
    table_top = avail_top + (avail_h - table_height) / 2
    table_top = max(table_top, avail_top)

    # Add the table shape
    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        SI(left_margin), SI(table_top),
        SI(table_width), SI(table_height))
    table = tbl_shape.table

    # Set column widths (already in inches from normalize_col_widths)
    for i, w in enumerate(col_widths):
        table.columns[i].width = SI(w)

    # Set per-row heights
    table.rows[0].height = SI(header_h)
    for ri, rh in enumerate(row_heights):
        table.rows[ri + 1].height = SI(rh)

    # --- Header row ---
    for ci, col_name in enumerate(columns):
        cell = table.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        header_text = str(col_name)
        if header_style == "allcaps":
            header_text = header_text.upper()
        r.text = header_text
        r.font.name = HEADING
        r.font.size = header_font_size
        r.font.color.rgb = WHITE
        r.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        # Background fill
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': header_bg_hex})
        solidFill.append(srgb)
        tcPr.append(solidFill)
        # Cell margins
        tcPr.set('marL', str(Emu(SI(0.08))))
        tcPr.set('marR', str(Emu(SI(0.05))))
        tcPr.set('marT', str(Emu(SI(0.04))))
        tcPr.set('marB', str(Emu(SI(0.04))))

    # --- Data rows ---
    for ri, row_data in enumerate(rows):
        for ci in range(num_cols):
            cell_val = str(row_data[ci]) if ci < len(row_data) else ""
            cell = table.cell(ri + 1, ci)
            render_table_cell(cell, cell_val, BODY, body_font_size_pt,
                              DARK, alignment=col_align[ci] if ci < len(col_align)
                              else PP_ALIGN.LEFT)
            # Stripe background
            if stripe and ri % 2 == 1:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFF0F8'})
                solidFill.append(srgb)
                tcPr.append(solidFill)
            # Cell margins
            tcPr = cell._tc.get_or_add_tcPr()
            tcPr.set('marL', str(Emu(SI(0.08))))
            tcPr.set('marR', str(Emu(SI(0.05))))
            tcPr.set('marT', str(Emu(SI(0.03))))
            tcPr.set('marB', str(Emu(SI(0.03))))

    # Bold first column (row labels)
    for ri in range(1, num_rows):
        cell = table.cell(ri, 0)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = PURPLE

    # Remove default table style borders by setting tblPr
    tbl_xml = table._tbl
    tblPr = tbl_xml.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('bandRow', '0')
        tblPr.set('bandCol', '0')
        tblPr.set('firstRow', '0')
        tblPr.set('lastRow', '0')
        tblPr.set('firstCol', '0')
        tblPr.set('lastCol', '0')

    # Apply clean horizontal-rule borders
    _apply_table_borders(table, num_rows, num_cols, header_bg_hex)


def _build_content_table_bullets(slide, slide_def, deck_meta):
    """Build a slide with a table in the top portion and bullet points below.

    Fields: headline, columns, rows, col_widths (proportional weights),
            col_align (single string or per-column list: left/center/right),
            header_bg, header_style ("bold"/"allcaps"), stripe,
            bullets_title (optional heading above bullets),
            bullets (list of strings — supports "bold part: rest" syntax),
            cell_colors (optional dict mapping cell values to hex colors for
            color-coded table cells, e.g. {"A": "4EC98B", "G": "FFD766", "N": "E85D5D"}).
    Uses content_generic (layout 2) as canvas.
    """
    from pptx.oxml.ns import qn

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    WHITE = brand.white
    DARK = brand.text_dark
    GRAY = brand.text_gray

    columns = slide_def.get("columns", [])
    rows = slide_def.get("rows", [])
    bullets = slide_def.get("bullets", [])
    bullets_title = slide_def.get("bullets_title", "")
    cell_colors = slide_def.get("cell_colors", {})
    if not columns:
        return

    num_cols = len(columns)
    num_rows = len(rows) + 1  # +1 for header

    # Table position — compressed to leave room for bullets
    left_margin = 0.35
    table_top = 0.85
    table_width = 10.0 - left_margin - 0.35
    max_table_height = 3.4  # leave room for legend + bullets

    # Column widths — proportional normalization or content-aware auto-sizing
    raw_col_widths = slide_def.get("col_widths", None)
    col_widths = normalize_col_widths(raw_col_widths, num_cols, table_width,
                                      columns, rows)

    # Column alignment
    col_align_raw = slide_def.get("col_align", "left")
    if isinstance(col_align_raw, str):
        col_align_raw = [col_align_raw] * num_cols
    _align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                  "right": PP_ALIGN.RIGHT}
    col_align = [_align_map.get(a, PP_ALIGN.LEFT) for a in col_align_raw]

    # Style options
    header_bg_hex = slide_def.get("header_bg", brand.primary_hex)
    stripe = slide_def.get("stripe", True)
    header_style = slide_def.get("header_style", "bold")

    # Font selection + content-aware row heights (compact mode)
    body_font_size_pt, header_font_size_pt, header_h, row_heights, table_height = \
        select_font_size(columns, rows, col_widths, max_table_height,
                         header_h_hint=0.30, is_compact=True)

    if table_height > max_table_height:
        print(f"    WARNING: Table on '{slide_def.get('headline', '')}' "
              f"leaves limited room for bullets "
              f"({5.05 - table_top - table_height:.2f}\" remaining)")

    header_font_size = Pt(header_font_size_pt)
    body_font_size = Pt(body_font_size_pt)

    # Add the table shape
    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        SI(left_margin), SI(table_top),
        SI(table_width), SI(table_height))
    table = tbl_shape.table

    # Set column widths (already in inches from normalize_col_widths)
    for i, w in enumerate(col_widths):
        table.columns[i].width = SI(w)

    # Set per-row heights
    table.rows[0].height = SI(header_h)
    for ri, rh in enumerate(row_heights):
        table.rows[ri + 1].height = SI(rh)

    # --- Header row ---
    for ci, col_name in enumerate(columns):
        cell = table.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        header_text = str(col_name)
        if header_style == "allcaps":
            header_text = header_text.upper()
        r.text = header_text
        r.font.name = HEADING
        r.font.size = header_font_size
        r.font.color.rgb = WHITE
        r.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': header_bg_hex})
        solidFill.append(srgb)
        tcPr.append(solidFill)
        tcPr.set('marL', str(Emu(SI(0.06))))
        tcPr.set('marR', str(Emu(SI(0.04))))
        tcPr.set('marT', str(Emu(SI(0.02))))
        tcPr.set('marB', str(Emu(SI(0.02))))

    # --- Data rows ---
    # Color map for A|G|N pipe-delimited cells
    agn_colors = {
        "A": RGBColor.from_string(cell_colors.get("A", "2E7D32")),
        "G": RGBColor.from_string(cell_colors.get("G", "F57F17")),
        "N": RGBColor.from_string(cell_colors.get("N", "C62828")),
    }
    agn_labels = ["A", "G", "N"]

    for ri, row_data in enumerate(rows):
        for ci in range(num_cols):
            cell_val = str(row_data[ci]) if ci < len(row_data) else ""
            cell = table.cell(ri + 1, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            cell_alignment = col_align[ci] if ci < len(col_align) else PP_ALIGN.LEFT
            p.alignment = cell_alignment

            # Detect A|G|N pipe format (e.g. "10|0|2")
            if "|" in cell_val:
                parts = [v.strip() for v in cell_val.split("|")]
                # Check for 100% A or 100% N
                if len(parts) == 3 and parts[1] == "0" and parts[2] == "0" and parts[0] != "0":
                    r = p.add_run()
                    r.text = "\u2713"
                    r.font.name = BODY
                    r.font.size = Pt(12)
                    r.font.bold = True
                    r.font.color.rgb = agn_colors["A"]
                    p.alignment = PP_ALIGN.CENTER
                elif len(parts) == 3 and parts[0] == "0" and parts[1] == "0" and parts[2] != "0":
                    r = p.add_run()
                    r.text = "\u2717"
                    r.font.name = BODY
                    r.font.size = Pt(12)
                    r.font.bold = True
                    r.font.color.rgb = agn_colors["N"]
                    p.alignment = PP_ALIGN.CENTER
                else:
                    for pi, part in enumerate(parts):
                        val = part.strip()
                        label = agn_labels[pi] if pi < len(agn_labels) else ""
                        r = p.add_run()
                        r.text = val
                        r.font.name = BODY
                        r.font.size = body_font_size
                        r.font.bold = True
                        r.font.color.rgb = agn_colors.get(label, DARK)
                        if pi < len(parts) - 1:
                            sep = p.add_run()
                            sep.text = " | "
                            sep.font.name = BODY
                            sep.font.size = body_font_size
                            sep.font.bold = False
                            sep.font.color.rgb = GRAY
            elif cell_val in cell_colors:
                r = p.add_run()
                r.text = cell_val
                r.font.name = BODY
                r.font.size = body_font_size
                r.font.color.rgb = RGBColor.from_string(cell_colors[cell_val])
                r.font.bold = True
            elif ci == 0:
                # First column — bold purple (use render_table_cell for bullet/bold support)
                render_table_cell(cell, cell_val, BODY, body_font_size_pt,
                                  PURPLE, alignment=cell_alignment)
                # Re-bold all runs in first column
                for pp in cell.text_frame.paragraphs:
                    for rr in pp.runs:
                        rr.font.bold = True
            else:
                render_table_cell(cell, cell_val, BODY, body_font_size_pt,
                                  DARK, alignment=cell_alignment)

            # Stripe background
            if stripe and ri % 2 == 1:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFF0F8'})
                solidFill.append(srgb)
                tcPr.append(solidFill)

            tcPr = cell._tc.get_or_add_tcPr()
            tcPr.set('marL', str(Emu(SI(0.06))))
            tcPr.set('marR', str(Emu(SI(0.04))))
            tcPr.set('marT', str(Emu(SI(0.02))))
            tcPr.set('marB', str(Emu(SI(0.02))))

    # Remove default table style borders
    tbl_xml = table._tbl
    tblPr = tbl_xml.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('bandRow', '0')
        tblPr.set('bandCol', '0')
        tblPr.set('firstRow', '0')
        tblPr.set('lastRow', '0')
        tblPr.set('firstCol', '0')
        tblPr.set('lastCol', '0')

    # Apply clean horizontal-rule borders
    _apply_table_borders(table, num_rows, num_cols, header_bg_hex)

    # --- Legend ---
    legend = slide_def.get("legend", "")
    legend_top = table_top + table_height + 0.08

    if legend:
        legend_box = slide.shapes.add_textbox(
            SI(left_margin), SI(legend_top),
            SI(table_width), SI(0.25))
        ltf = legend_box.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.RIGHT

        legend_parts = legend.split("|")
        for lpi, part in enumerate(legend_parts):
            part = part.strip()
            for token in agn_labels:
                prefix = token + " "
                if part.startswith(prefix):
                    lr = lp.add_run()
                    lr.text = token
                    lr.font.name = BODY
                    lr.font.size = Pt(7)
                    lr.font.bold = True
                    lr.font.color.rgb = agn_colors.get(token, DARK)
                    lr2 = lp.add_run()
                    lr2.text = part[len(token):]
                    lr2.font.name = BODY
                    lr2.font.size = Pt(7)
                    lr2.font.bold = False
                    lr2.font.color.rgb = GRAY
                    break
            else:
                lr = lp.add_run()
                lr.text = part
                lr.font.name = BODY
                lr.font.size = Pt(7)
                lr.font.color.rgb = GRAY

            if lpi < len(legend_parts) - 1:
                sep = lp.add_run()
                sep.text = "   |   "
                sep.font.name = BODY
                sep.font.size = Pt(7)
                sep.font.color.rgb = GRAY

        legend_top += 0.25

    # --- Bullets section below table ---
    bullet_top = legend_top + 0.05

    if bullets_title:
        add_text_box(slide, bullets_title,
                     left_margin, bullet_top, table_width, 0.3,
                     font_size=11, font_name=HEADING,
                     bold=True, color=(0x5F, 0x01, 0x6F))
        bullet_top += 0.3

    if bullets:
        txBox = slide.shapes.add_textbox(
            SI(left_margin), SI(bullet_top),
            SI(table_width), SI(5.05 - bullet_top))
        tf = txBox.text_frame
        tf.word_wrap = True

        for bi, item in enumerate(bullets):
            para = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            para.space_after = Pt(3)

            if ":" in item:
                bold_part, rest = item.split(":", 1)
                r1 = para.add_run()
                r1.text = bold_part + ":"
                r1.font.name = BODY
                r1.font.size = Pt(9)
                r1.font.bold = True
                r1.font.color.rgb = PURPLE
                r2 = para.add_run()
                r2.text = rest
                r2.font.name = BODY
                r2.font.size = Pt(9)
                r2.font.bold = False
                r2.font.color.rgb = DARK
            else:
                r = para.add_run()
                r.text = item
                r.font.name = BODY
                r.font.size = Pt(9)
                r.font.color.rgb = DARK


def _build_closing(slide, slide_def, deck_meta):
    """Build a closing/CTA slide with large centered headline.

    Fields: headline (default "Questions?"), subheader, contact_info.
    Uses content_generic (layout 2) as canvas. Attempts to use the
    title_cover background (p12) if available.
    """
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    WHITE = brand.white
    DARK = brand.text_dark

    headline = slide_def.get("headline", "Questions?")
    subheader = slide_def.get("subheader", "")
    contact_info = slide_def.get("contact_info", "")

    headline_size_pt = 48
    sub_size_pt = 16
    contact_size_pt = 12

    # Try to add background image (from brand config)
    bg_key = slide_def.get("background", "default")
    bg_path = brand.get_title_background(bg_key)
    if not bg_path or not os.path.exists(bg_path):
        bg_path = brand.get_title_background("default")

    has_bg = False
    if bg_path and os.path.exists(bg_path):
        slide.shapes.add_picture(bg_path,
                                 SI(0), SI(0),
                                 SI(10.0), SI(5.62))
        has_bg = True
    else:
        # No background image — draw a gradient fill
        from pptx.enum.shapes import MSO_SHAPE
        bg_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SI(0), SI(0), SI(10.0), SI(5.625))
        bg_rect.line.fill.background()
        fill = bg_rect.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = brand.primary
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = brand.interpolate(0.6)
        fill.gradient_stops[1].position = 1.0
        has_bg = True

    # Hide the headline placeholder (we draw our own)
    set_placeholder_text(slide, 0, "")

    # Choose text color based on background
    headline_color = WHITE if has_bg else PURPLE
    sub_color = WHITE if has_bg else DARK

    # Calculate heights
    headline_h = estimate_text_height(headline, 8.0, headline_size_pt)
    sub_h = estimate_text_height(subheader, 7.0, sub_size_pt) if subheader else 0
    contact_h = estimate_text_height(contact_info, 6.0, contact_size_pt) if contact_info else 0

    total_h = headline_h
    if subheader:
        total_h += 0.2 + sub_h
    if contact_info:
        total_h += 0.3 + contact_h

    # Vertically center
    avail_top = 0.5
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - total_h) / 2
    offset = max(offset, avail_top)

    cur_y = offset

    # Headline — large centered
    txb = slide.shapes.add_textbox(
        SI(1.0), SI(cur_y),
        SI(8.0), SI(headline_h))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = headline; r.font.name = HEADING
    r.font.size = Pt(headline_size_pt); r.font.color.rgb = headline_color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    cur_y += headline_h

    # Subheader
    if subheader:
        cur_y += 0.2
        sub_h_est = estimate_text_height(subheader, 7.0, sub_size_pt)
        txb = slide.shapes.add_textbox(
            SI(1.5), SI(cur_y),
            SI(7.0), SI(sub_h_est))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = subheader; r.font.name = BODY
        r.font.size = Pt(sub_size_pt); r.font.color.rgb = sub_color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cur_y += sub_h_est

    # Contact info
    if contact_info:
        cur_y += 0.3
        contact_h_est = estimate_text_height(contact_info, 6.0, contact_size_pt)
        txb = slide.shapes.add_textbox(
            SI(2.0), SI(cur_y),
            SI(6.0), SI(contact_h_est))
        tf = txb.text_frame; tf.word_wrap = True
        for li, line in enumerate(str(contact_info).split("\n")):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            _add_run_with_hyperlinks(para, line, BODY, contact_size_pt, sub_color)
            para.alignment = PP_ALIGN.CENTER


def _build_section_divider(slide, slide_def, deck_meta):
    """Build a section divider slide — visual break between deck sections.

    Fields: section_number (optional), headline, subheader (optional),
            background ("purple"|"light"|"image"), icon (optional).
    No footers rendered.
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    DARK = brand.text_dark
    PURPLE = brand.primary
    PINK = brand.secondary
    WHITE = WHITE
    LIGHT_BG = brand.bg_light
    DARK = DARK

    bg_style = slide_def.get("background", "purple")
    section_number = slide_def.get("section_number", "")
    headline = slide_def.get("headline", "")
    subheader = slide_def.get("subheader", "")
    icon_name = slide_def.get("icon", "")

    # Hide default headline placeholder
    set_placeholder_text(slide, 0, "")

    # Background
    if bg_style == "purple":
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SI(0), SI(0), SI(10.0), SI(5.625))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PURPLE
        bg.line.fill.background()
        num_color = PINK
        head_color = WHITE
        sub_color = brand.accent  # light pink
    elif bg_style == "light":
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SI(0), SI(0), SI(10.0), SI(5.625))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()
        num_color = PURPLE
        head_color = PURPLE
        sub_color = DARK
    else:
        # image background — try brand config, fall back to hardcoded paths
        script_dir = os.path.dirname(os.path.abspath(__file__))
        bg_path = brand.get_title_background(bg_style)
        if not bg_path or not os.path.exists(bg_path):
            bg_map = {
                "p12": "title-assets/title-bg-p12.jpg",
                "p13": "title-assets/title-bg-p13.jpg",
                "p17": "title-assets/title-bg-p17.png",
            }
            bg_rel = bg_map.get(bg_style, bg_map.get("p12"))
            bg_path = os.path.join(script_dir, bg_rel)
        if os.path.exists(bg_path):
            slide.shapes.add_picture(bg_path, SI(0), SI(0),
                                     SI(10.0), SI(5.625))
        num_color = PINK
        head_color = WHITE
        sub_color = brand.accent

    # Calculate content heights
    icon_path = _resolve_icon(icon_name) if icon_name else None
    icon_h = 0.7 if (icon_path and os.path.exists(icon_path)) else 0
    num_h = estimate_text_height(str(section_number), 8.0, 72) if section_number else 0
    head_h = estimate_text_height(headline, 8.0, 36)
    sub_h = estimate_text_height(subheader, 7.0, 16) if subheader else 0

    total_h = icon_h + num_h + head_h
    if subheader:
        total_h += 0.15 + sub_h

    offset = (5.625 - total_h) / 2
    offset = max(offset, 0.5)
    cur_y = offset

    # Icon
    if icon_path and os.path.exists(icon_path):
        icon_left = (10.0 - 0.6) / 2
        slide.shapes.add_picture(icon_path, SI(icon_left), SI(cur_y),
                                 SI(0.6), SI(0.6))
        cur_y += 0.7

    # Section number
    if section_number:
        txb = slide.shapes.add_textbox(SI(1.0), SI(cur_y), SI(8.0), SI(num_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = str(section_number); r.font.name = HEADING
        r.font.size = Pt(72); r.font.color.rgb = num_color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cur_y += num_h

    # Headline
    txb = slide.shapes.add_textbox(SI(1.0), SI(cur_y), SI(8.0), SI(head_h))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = headline; r.font.name = HEADING
    r.font.size = Pt(36); r.font.color.rgb = head_color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    cur_y += head_h

    # Subheader
    if subheader:
        cur_y += 0.15
        txb = slide.shapes.add_textbox(SI(1.5), SI(cur_y), SI(7.0), SI(sub_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = subheader; r.font.name = BODY
        r.font.size = Pt(16); r.font.color.rgb = sub_color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _build_quote(slide, slide_def, deck_meta):
    """Build a quote/testimonial slide with decorative quotation mark.

    Fields: headline (optional), quote_text, attribution, attribution_title (optional),
            icon (optional), style ("centered"|"left-aligned").
    """
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    DARK = brand.text_dark
    GRAY = brand.text_gray

    quote_text = slide_def.get("quote_text", "")
    attribution = slide_def.get("attribution", "")
    attribution_title = slide_def.get("attribution_title", "")
    style = slide_def.get("style", "centered")
    icon_name = slide_def.get("icon", "")

    align = PP_ALIGN.CENTER if style == "centered" else PP_ALIGN.LEFT
    margin_x = 1.5 if style == "centered" else 0.8
    text_w = 10.0 - 2 * margin_x

    # Calculate heights
    icon_path = _resolve_icon(icon_name) if icon_name else None
    icon_h = 0.6 if (icon_path and os.path.exists(icon_path)) else 0
    mark_h = 0.8  # decorative quotation mark
    quote_h = estimate_text_height(quote_text, text_w, 18)
    rule_h = 0.3  # thin line + spacing
    attr_h = estimate_text_height(attribution, text_w, 12) if attribution else 0
    attr_title_h = estimate_text_height(attribution_title, text_w, 10) if attribution_title else 0

    total_h = icon_h + mark_h + quote_h + rule_h + attr_h + attr_title_h

    avail_top = 0.85
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - total_h) / 2
    offset = max(offset, avail_top)
    cur_y = offset

    # Icon
    if icon_path and os.path.exists(icon_path):
        icon_left = (10.0 - 0.5) / 2 if style == "centered" else margin_x
        slide.shapes.add_picture(icon_path, SI(icon_left), SI(cur_y),
                                 SI(0.5), SI(0.5))
        cur_y += 0.6

    # Decorative quotation mark
    txb = slide.shapes.add_textbox(SI(margin_x), SI(cur_y),
                                   SI(text_w), SI(mark_h))
    tf = txb.text_frame; tf.word_wrap = False
    r = tf.paragraphs[0].add_run()
    r.text = "\u201C"  # left double quotation mark
    r.font.name = HEADING; r.font.size = Pt(96)
    r.font.color.rgb = PINK
    tf.paragraphs[0].alignment = align
    cur_y += mark_h + 0.05  # position quote text below the mark

    # Quote text
    txb = slide.shapes.add_textbox(SI(margin_x), SI(cur_y),
                                   SI(text_w), SI(quote_h))
    tf = txb.text_frame; tf.word_wrap = True
    for li, line in enumerate(quote_text.split("\n")):
        para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        r = para.add_run(); r.text = line
        r.font.name = BODY; r.font.size = Pt(18)
        r.font.color.rgb = PURPLE
        para.alignment = align
    cur_y += quote_h

    # Pink rule
    cur_y += 0.1
    rule_w = 2.0
    rule_left = (10.0 - rule_w) / 2 if style == "centered" else margin_x
    line = slide.shapes.add_connector(
        1, SI(rule_left), SI(cur_y),
        SI(rule_left + rule_w), SI(cur_y))
    line.line.color.rgb = PINK; line.line.width = Pt(1.5)
    cur_y += 0.2

    # Attribution
    if attribution:
        txb = slide.shapes.add_textbox(SI(margin_x), SI(cur_y),
                                       SI(text_w), SI(attr_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = attribution; r.font.name = HEADING
        r.font.size = Pt(12); r.font.color.rgb = DARK
        tf.paragraphs[0].alignment = align
        cur_y += attr_h

    # Attribution title (second line)
    if attribution_title:
        txb = slide.shapes.add_textbox(SI(margin_x), SI(cur_y),
                                       SI(text_w), SI(attr_title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = attribution_title; r.font.name = BODY
        r.font.size = Pt(10); r.font.color.rgb = GRAY
        tf.paragraphs[0].alignment = align


def _build_kpi_dashboard(slide, slide_def, deck_meta):
    """Build a KPI dashboard with a grid of metric cards.

    Fields: headline, metrics (list of {number, label, trend, target, icon, highlight}),
            columns (optional, auto-detected from count).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    DARK = brand.text_dark
    WHITE = WHITE
    LIGHT_BG = brand.bg_light
    GREEN = RGBColor(0x4E, 0xC9, 0x8B)
    RED = RGBColor(0xE8, 0x5D, 0x5D)
    GRAY = brand.text_gray

    metrics = slide_def.get("metrics", []) or slide_def.get("kpis", [])
    n = len(metrics)
    if n == 0:
        return

    # Determine grid layout
    cols = slide_def.get("columns", None)
    if not cols:
        if n <= 4:
            cols = n
        elif n <= 6:
            cols = 3
        else:
            cols = 4
    rows = -(-n // cols)  # ceiling division

    # Card sizing
    gap = 0.12
    usable_w = 9.3
    usable_h = 3.8  # available content height
    card_w = (usable_w - (cols - 1) * gap) / cols
    card_h = (usable_h - (rows - 1) * gap) / rows
    card_h = min(card_h, 1.8)  # cap height

    # Center the grid
    total_w = cols * card_w + (cols - 1) * gap
    total_h = rows * card_h + (rows - 1) * gap
    start_x = 0.35 + (usable_w - total_w) / 2
    avail_top = 0.85
    avail_bottom = 5.0
    start_y = avail_top + ((avail_bottom - avail_top) - total_h) / 2
    start_y = max(start_y, avail_top)

    # Base font sizing by card count (adjusted per-metric below for long text)
    if n <= 4:
        base_num_size = 32
        label_size = 10
        trend_size = 14
    elif n <= 6:
        base_num_size = 28
        label_size = 9
        trend_size = 12
    else:
        base_num_size = 24
        label_size = 8
        trend_size = 10

    # Pre-compute consistent label Y position across all cards.
    # The label is anchored at a fixed distance from the card bottom so it
    # aligns across cards regardless of number font auto-sizing.
    pad = 0.12
    max_label_h = 0
    for metric in metrics:
        inner_w = card_w - 2 * pad
        label_text = metric.get("label", "")
        lh = estimate_text_height(label_text, inner_w, label_size)
        max_label_h = max(max_label_h, lh)
    # Label top = card bottom - padding - label height
    label_offset_from_bottom = pad + max_label_h

    for i, metric in enumerate(metrics):
        row = i // cols
        col = i % cols
        cx = start_x + col * (card_w + gap)
        cy = start_y + row * (card_h + gap)

        # Card background
        is_highlight = metric.get("highlight", False)
        bg_color = PURPLE if is_highlight else LIGHT_BG
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(cx), SI(cy), SI(card_w), SI(card_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()

        num_color = WHITE if is_highlight else PURPLE
        label_color = brand.accent if is_highlight else DARK
        trend_color_map = {"up": GREEN, "down": RED, "flat": GRAY}

        inner_y = cy + pad
        inner_w = card_w - 2 * pad

        # Optional icon (0.45" minimum — smaller icons get dropped by Google Slides import)
        icon_name = metric.get("icon", "")
        icon_path = _resolve_icon(icon_name) if icon_name else None
        if icon_path and os.path.exists(icon_path):
            icon_sz = 0.45
            icon_left = cx + (card_w - icon_sz) / 2
            slide.shapes.add_picture(icon_path, SI(icon_left), SI(inner_y),
                                     SI(icon_sz), SI(icon_sz))
            inner_y += icon_sz + 0.04
        elif icon_name:
            print(f"    WARNING: KPI icon '{icon_name}' not found in catalog")

        # Number (with optional trend arrow)
        number_text = metric.get("number", "")
        trend = metric.get("trend", "")
        trend_symbol = {"up": " \u2191", "down": " \u2193", "flat": " \u2192"}.get(trend, "")
        target = metric.get("target", "")

        display_text = str(number_text)
        if target:
            display_text += f" / {target}"

        # Auto-scale number font to fit card width.
        # Estimate: heading font at N pt ≈ N * 0.6 / 72 inches per char.
        # Scale down if text would exceed the inner card width.
        num_size = base_num_size
        chars_per_inch_at_size = lambda sz: 72 / (sz * 0.6)
        text_width_est = len(display_text) / chars_per_inch_at_size(num_size)
        while text_width_est > inner_w and num_size > 14:
            num_size -= 2
            text_width_est = len(display_text) / chars_per_inch_at_size(num_size)

        # Number text fills the space between icon bottom and label top
        label_y = cy + card_h - label_offset_from_bottom
        num_avail_h = label_y - inner_y - 0.04
        num_h = min(estimate_text_height(display_text, inner_w, num_size), num_avail_h)

        txb = slide.shapes.add_textbox(SI(cx + pad), SI(inner_y),
                                       SI(inner_w), SI(num_h))
        tf = txb.text_frame; tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = display_text; r.font.name = HEADING
        r.font.size = Pt(num_size); r.font.color.rgb = num_color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Trend arrow (as separate run in same box or separate box)
        if trend_symbol:
            r2 = tf.paragraphs[0].add_run()
            r2.text = trend_symbol; r2.font.name = BODY
            r2.font.size = Pt(trend_size)
            r2.font.color.rgb = trend_color_map.get(trend, GRAY)

        # Label — anchored to card bottom for consistent alignment across cards
        label_text = metric.get("label", "")
        label_h = estimate_text_height(label_text, inner_w, label_size)
        txb = slide.shapes.add_textbox(SI(cx + pad), SI(label_y),
                                       SI(inner_w), SI(label_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = label_text; r.font.name = BODY
        r.font.size = Pt(label_size); r.font.color.rgb = label_color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _build_roadmap(slide, slide_def, deck_meta):
    """Build a roadmap slide with time axis, swimlanes, and milestones.

    Fields: headline, time_axis (list of period labels),
            swimlanes (list of {name, items: [{label, start, end, status}]}),
            milestones (optional list of {date, label}).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    DARK = brand.text_dark
    WHITE = WHITE
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    GREEN = RGBColor(0x4E, 0xC9, 0x8B)
    GRAY = brand.text_gray

    time_axis = slide_def.get("time_axis", [])
    swimlanes = slide_def.get("swimlanes", []) or slide_def.get("lanes", [])
    milestones = slide_def.get("milestones", [])

    # Auto-generate time_axis from lane item positions if not provided
    if not time_axis and swimlanes:
        max_pos = 0
        for lane in swimlanes:
            for item in lane.get("items", []):
                max_pos = max(max_pos, item.get("end", 0), item.get("start", 0))
        if max_pos > 0:
            time_axis = [f"Phase {i+1}" for i in range(max_pos)]

    if not time_axis:
        return

    from datetime import datetime

    n_periods = len(time_axis)
    n_lanes = len(swimlanes)

    # Layout constants
    left_margin = 0.35
    lane_label_w = 1.2
    chart_left = left_margin + lane_label_w + 0.1
    chart_right = 9.65
    chart_w = chart_right - chart_left

    top_start = 0.85
    axis_h = 0.35  # time axis row
    bar_area_top = top_start + axis_h + 0.1
    bottom_limit = 4.9
    bar_area_h = bottom_limit - bar_area_top

    if n_lanes > 0:
        lane_h = min(bar_area_h / n_lanes, 1.2)
    else:
        lane_h = bar_area_h

    # Build period index for start/end mapping (string label -> int index)
    period_index = {label: i for i, label in enumerate(time_axis)}

    def _resolve_idx(value):
        """Resolve a start/end value to an integer index. Handles int or string."""
        if isinstance(value, int):
            return max(0, min(value, n_periods - 1))
        if isinstance(value, str) and value in period_index:
            return period_index[value]
        # Try numeric string
        try:
            return max(0, min(int(value), n_periods - 1))
        except (ValueError, TypeError):
            return 0

    # --- Proportional date-based spacing ---
    # Try to parse time_axis labels as dates for proportional column widths.
    # Falls back to equal-width columns if parsing fails.
    def _try_parse_dates(labels):
        formats = ["%b %d", "%B %d", "%Y-%m-%d", "%b %d, %Y", "%d %b",
                   "%b %d-%d", "%b %d-%b %d"]  # "Apr 14-15" style
        dates = []
        for label in labels:
            # Strip ranges — use the start date (e.g., "Apr 14-15" -> "Apr 14")
            clean = label.split("-")[0].strip() if "-" in label and not label.startswith("2") else label.strip()
            parsed = None
            for fmt in formats:
                try:
                    d = datetime.strptime(clean, fmt)
                    if d.year == 1900:
                        d = d.replace(year=2026)
                    parsed = d
                    break
                except Exception:
                    continue
            if parsed is None:
                return None
            dates.append(parsed)
        return dates

    parsed_dates = _try_parse_dates(time_axis)

    if parsed_dates and len(parsed_dates) >= 2:
        min_d = parsed_dates[0]
        max_d = parsed_dates[-1]
        total_days = (max_d - min_d).days
        if total_days > 0:
            fractions = [(d - min_d).days / total_days for d in parsed_dates]
        else:
            fractions = [i / max(1, n_periods - 1) for i in range(n_periods)]
        # Column left edges at proportional positions
        col_lefts = [chart_left + f * chart_w for f in fractions]
    else:
        # Equal-width fallback
        period_w = chart_w / n_periods
        col_lefts = [chart_left + i * period_w for i in range(n_periods)]

    # Column right edges (each column extends to the next column's left, or chart_right)
    col_rights = col_lefts[1:] + [chart_right]
    col_widths = [col_rights[i] - col_lefts[i] for i in range(n_periods)]

    # --- Time axis labels ---
    for pi, label in enumerate(time_axis):
        px = col_lefts[pi]
        pw = col_widths[pi]
        txb = slide.shapes.add_textbox(SI(px), SI(top_start),
                                       SI(pw), SI(axis_h))
        tf = txb.text_frame; tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = label; r.font.name = HEADING
        r.font.size = Pt(10); r.font.color.rgb = PURPLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Axis line
    axis_line_y = top_start + axis_h
    line = slide.shapes.add_connector(
        1, SI(chart_left), SI(axis_line_y),
        SI(chart_right), SI(axis_line_y))
    line.line.color.rgb = PURPLE; line.line.width = Pt(1.0)

    # Vertical period dividers (light)
    for pi in range(1, n_periods):
        px = col_lefts[pi]
        line = slide.shapes.add_connector(
            1, SI(px), SI(axis_line_y),
            SI(px), SI(bottom_limit))
        line.line.color.rgb = brand.divider
        line.line.width = Pt(0.5)

    # Status colors
    status_colors = {
        "complete": GREEN,
        "active": PURPLE,
        "planned": LIGHT_BG,
    }
    status_text_colors = {
        "complete": WHITE,
        "active": WHITE,
        "planned": PURPLE,
    }

    # --- Swimlanes ---
    for li, lane in enumerate(swimlanes):
        lane_top = bar_area_top + li * lane_h
        lane_name = lane.get("name", "")

        # Lane label
        txb = slide.shapes.add_textbox(SI(left_margin), SI(lane_top),
                                       SI(lane_label_w), SI(lane_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = lane_name; r.font.name = HEADING
        r.font.size = Pt(10); r.font.color.rgb = PURPLE

        # Lane divider
        if li > 0:
            line = slide.shapes.add_connector(
                1, SI(left_margin), SI(lane_top),
                SI(chart_right), SI(lane_top))
            line.line.color.rgb = brand.divider
            line.line.width = Pt(0.5)

        # Bars
        items = lane.get("items", [])
        bar_font_pt = 8
        bar_gap = 0.06

        # Determine bar height: check if any label needs wrapping
        # A label needs wrapping if it's wider than the bar at the given font size
        def _label_needs_wrap(label_text, available_w):
            chars_per_inch = 13 * (10 / bar_font_pt)
            text_w_needed = len(label_text) / chars_per_inch
            return text_w_needed > available_w

        # Pre-scan to determine if any bar in this lane needs 2-line height
        any_wraps = False
        for item in items:
            label = item.get("label", "")
            si_idx = _resolve_idx(item.get("start", 0))
            ei_idx = _resolve_idx(item.get("end", si_idx))
            bw = (col_rights[ei_idx] - col_lefts[si_idx]) - 0.08 - 0.12
            if _label_needs_wrap(label, bw):
                any_wraps = True
                break

        bar_h = 0.40 if any_wraps else 0.28

        for ii, item in enumerate(items):
            label = item.get("label", "")
            start = item.get("start", 0)
            end = item.get("end", start)
            status = item.get("status", "planned")

            start_idx = _resolve_idx(start)
            end_idx = _resolve_idx(end)

            bar_left = col_lefts[start_idx] + 0.04
            bar_right = col_rights[end_idx] - 0.04
            bar_w = max(bar_right - bar_left, 0.3)  # minimum bar width
            bar_top = lane_top + 0.08 + ii * (bar_h + bar_gap)

            bg_color = status_colors.get(status, LIGHT_BG)
            text_color = status_text_colors.get(status, PURPLE)

            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                SI(bar_left), SI(bar_top),
                SI(bar_w), SI(bar_h))
            bar.fill.solid()
            bar.fill.fore_color.rgb = bg_color
            if status == "planned":
                bar.line.color.rgb = PURPLE
                bar.line.width = Pt(0.75)
            else:
                bar.line.fill.background()

            # Bar label — word wrap enabled for long labels
            text_w = bar_w - 0.12
            txb = slide.shapes.add_textbox(SI(bar_left + 0.06), SI(bar_top),
                                           SI(text_w), SI(bar_h))
            tf = txb.text_frame; tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = label; r.font.name = BODY
            r.font.size = Pt(bar_font_pt); r.font.color.rgb = text_color
            r.font.bold = True

    # --- Milestones ---
    for ms in milestones:
        ms_date = ms.get("date", "")
        ms_label = ms.get("label", "")
        # Resolve milestone position — try string label first, then index
        ms_idx = period_index.get(ms_date, None)
        if ms_idx is None:
            if isinstance(ms_date, int) and 0 <= ms_date < n_periods:
                ms_idx = ms_date
            else:
                continue

        mx = (col_lefts[ms_idx] + col_rights[ms_idx]) / 2  # center of column
        my = bottom_limit - 0.15

        # Diamond marker
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            SI(mx - 0.08), SI(my - 0.08),
            SI(0.16), SI(0.16))
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = PINK
        diamond.line.fill.background()

        # Label below
        txb = slide.shapes.add_textbox(SI(mx - 0.5), SI(my + 0.1),
                                       SI(1.0), SI(0.2))
        tf = txb.text_frame; tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = ms_label; r.font.name = BODY
        r.font.size = Pt(7); r.font.color.rgb = DARK; r.font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _build_before_after(slide, slide_def, deck_meta):
    """Build a before/after transformation slide with two columns and center arrow.

    Fields: headline, before ({label, icon, items}), after ({label, icon, items}),
            arrow_label (optional).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    DARK = brand.text_dark
    WHITE = WHITE
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    NEUTRAL_BG = brand.neutral_bg

    before = slide_def.get("before", {})
    after = slide_def.get("after", {})
    arrow_label = slide_def.get("arrow_label", "")

    col_w = 3.8
    arrow_zone_w = 1.4
    left_x = 0.35
    right_x = left_x + col_w + arrow_zone_w
    arrow_cx = left_x + col_w + arrow_zone_w / 2

    avail_top = 0.85
    avail_bottom = 5.0

    # Calculate content height from items
    before_items = before.get("items", [])
    after_items = after.get("items", [])
    before_text = "\n".join(f"\u2022 {item}" for item in before_items)
    after_text = "\n".join(f"\u2022 {item}" for item in after_items)

    before_h = estimate_text_height(before_text, col_w - 0.4, 10)
    after_h = estimate_text_height(after_text, col_w - 0.4, 10)
    max_items_h = max(before_h, after_h)

    label_h = 0.35
    icon_h = 0.5
    content_h = icon_h + label_h + 0.1 + max_items_h + 0.4  # + padding

    avail = avail_bottom - avail_top
    offset = avail_top + (avail - content_h) / 2
    offset = max(offset, avail_top)

    # --- Background rectangles ---
    bg_left = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SI(left_x), SI(offset), SI(col_w), SI(content_h))
    bg_left.fill.solid(); bg_left.fill.fore_color.rgb = NEUTRAL_BG
    bg_left.line.fill.background()

    bg_right = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SI(right_x), SI(offset), SI(col_w), SI(content_h))
    bg_right.fill.solid(); bg_right.fill.fore_color.rgb = LIGHT_BG
    bg_right.line.fill.background()

    # --- Center arrow ---
    arrow_y = offset + content_h / 2
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        SI(arrow_cx - 0.35), SI(arrow_y - 0.2),
        SI(0.7), SI(0.4))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = PINK
    arrow.line.fill.background()

    if arrow_label:
        txb = slide.shapes.add_textbox(SI(arrow_cx - 0.6), SI(arrow_y + 0.25),
                                       SI(1.2), SI(0.2))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = arrow_label; r.font.name = BODY
        r.font.size = Pt(8); r.font.color.rgb = DARK; r.font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Render each column ---
    for side, data, sx in [("before", before, left_x), ("after", after, right_x)]:
        cur_y = offset + 0.15
        side_label = data.get("label", side.title())
        icon_name = data.get("icon", "")
        items = data.get("items", [])

        # Icon
        icon_path = _resolve_icon(icon_name) if icon_name else None
        if icon_path and os.path.exists(icon_path):
            icon_left = sx + (col_w - 0.4) / 2
            slide.shapes.add_picture(icon_path, SI(icon_left), SI(cur_y),
                                     SI(0.4), SI(0.4))
            cur_y += 0.5

        # Label
        txb = slide.shapes.add_textbox(SI(sx + 0.2), SI(cur_y),
                                       SI(col_w - 0.4), SI(label_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = side_label; r.font.name = HEADING
        r.font.size = Pt(14); r.font.color.rgb = PURPLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cur_y += label_h + 0.1

        # Items
        items_text = "\n".join(f"\u2022 {item}" for item in items)
        items_h = estimate_text_height(items_text, col_w - 0.4, 10)
        txb = slide.shapes.add_textbox(SI(sx + 0.2), SI(cur_y),
                                       SI(col_w - 0.4), SI(items_h))
        tf = txb.text_frame; tf.word_wrap = True
        for li, item in enumerate(items):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            r = para.add_run(); r.text = f"\u2022 {item}"
            r.font.name = BODY; r.font.size = Pt(10)
            r.font.color.rgb = DARK


def _build_numbered_list(slide, slide_def, deck_meta):
    """Build a numbered list slide with large numbers and title/body per item.

    Fields: headline, items (list of {title, body, icon}).
    """
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark

    items = slide_def.get("items", [])
    n = len(items)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    usable_w = 9.3
    num_w = 0.6
    text_left = 0.35 + num_w + 0.15
    text_w = usable_w - num_w - 0.15

    # Calculate per-item height
    item_heights = []
    for item in items:
        title = item.get("title", "") if isinstance(item, dict) else str(item)
        body = item.get("body", "") if isinstance(item, dict) else ""
        title_h = estimate_text_height(title, text_w, 12)
        body_h = estimate_text_height(body, text_w, 10) if body else 0
        item_heights.append(title_h + body_h + 0.05)

    total_h = sum(item_heights) + (n - 1) * 0.15  # divider gaps
    offset = avail_top + (avail_h - total_h) / 2
    offset = max(offset, avail_top)

    cur_y = offset
    for i, item in enumerate(items):
        if isinstance(item, dict):
            title = item.get("title", "")
            body = item.get("body", "")
        else:
            title = str(item)
            body = ""

        # Large number
        num_h = item_heights[i]
        txb = slide.shapes.add_textbox(SI(0.35), SI(cur_y),
                                       SI(num_w), SI(0.5))
        tf = txb.text_frame; tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = str(i + 1); r.font.name = HEADING
        r.font.size = Pt(28); r.font.color.rgb = PURPLE
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Title
        title_h = estimate_text_height(title, text_w, 12)
        txb = slide.shapes.add_textbox(SI(text_left), SI(cur_y),
                                       SI(text_w), SI(title_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = title; r.font.name = HEADING
        r.font.size = Pt(12); r.font.color.rgb = PURPLE
        cur_y += title_h

        # Body
        if body:
            body_h = estimate_text_height(body, text_w, 10)
            txb = slide.shapes.add_textbox(SI(text_left), SI(cur_y + 0.02),
                                           SI(text_w), SI(body_h))
            tf = txb.text_frame; tf.word_wrap = True
            _render_body_text(tf, body, BODY, 10, DARK)
            cur_y += body_h + 0.05

        cur_y += 0.03  # small gap before divider

        # Divider between items (not after last)
        if i < n - 1:
            line = slide.shapes.add_connector(
                1, SI(0.35), SI(cur_y + 0.06),
                SI(9.65), SI(cur_y + 0.06))
            line.line.color.rgb = brand.divider
            line.line.width = Pt(0.5)
            cur_y += 0.15


def _build_status_board(slide, slide_def, deck_meta):
    """Build a RAG status board with colored status circles per item.

    Fields: headline, items (list of {name, status, summary}),
            as_of (optional date string).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark
    GRAY = brand.text_gray
    GREEN = RGBColor(0x4E, 0xC9, 0x8B)
    AMBER = RGBColor(0xFF, 0xD7, 0x66)
    RED = RGBColor(0xE8, 0x5D, 0x5D)

    status_colors = {"green": GREEN, "amber": AMBER, "red": RED}

    items = slide_def.get("items", [])
    as_of = slide_def.get("as_of", "")
    n = len(items)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    circle_size = 0.3
    name_w = 2.5
    summary_w = 5.8
    row_left = 0.55

    # Calculate row heights
    row_heights = []
    for item in items:
        summary = item.get("summary", "")
        sh = estimate_text_height(summary, summary_w, 10)
        row_heights.append(max(sh, circle_size + 0.1) + 0.15)

    total_h = sum(row_heights) + (0.25 if as_of else 0)
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - total_h) / 2
    offset = max(offset, avail_top)
    cur_y = offset

    for i, item in enumerate(items):
        name = item.get("name", "")
        status = item.get("status", "green")
        summary = item.get("summary", "")
        rh = row_heights[i]

        # Status circle
        color = status_colors.get(status, GRAY)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SI(row_left), SI(cur_y + 0.04),
            SI(circle_size), SI(circle_size))
        circle.fill.solid(); circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Name
        txb = slide.shapes.add_textbox(
            SI(row_left + circle_size + 0.15), SI(cur_y),
            SI(name_w), SI(0.35))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = name; r.font.name = HEADING
        r.font.size = Pt(12); r.font.color.rgb = PURPLE

        # Summary
        summary_h = estimate_text_height(summary, summary_w, 10)
        txb = slide.shapes.add_textbox(
            SI(row_left + circle_size + 0.15 + name_w + 0.15), SI(cur_y),
            SI(summary_w), SI(summary_h))
        tf = txb.text_frame; tf.word_wrap = True
        _render_body_text(tf, summary, BODY, 10, DARK)

        cur_y += rh

        # Divider
        if i < n - 1:
            line = slide.shapes.add_connector(
                1, SI(row_left), SI(cur_y - 0.05),
                SI(9.65), SI(cur_y - 0.05))
            line.line.color.rgb = brand.divider
            line.line.width = Pt(0.5)

    # As-of date stamp
    if as_of:
        txb = slide.shapes.add_textbox(SI(7.0), SI(cur_y + 0.05),
                                       SI(2.65), SI(0.2))
        tf = txb.text_frame; tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = f"As of {as_of}"; r.font.name = BODY
        r.font.size = Pt(8); r.font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _build_image_showcase(slide, slide_def, deck_meta):
    """Build an image showcase slide — full-slide image with optional caption.

    Fields: headline, image (path), caption (optional), border (optional bool).
    """
    from PIL import Image as PILImage

    DARK = DARK
    GRAY = brand.text_gray

    image_path = slide_def.get("image", "")
    caption = slide_def.get("caption", "")
    show_border = slide_def.get("border", False)

    if not image_path:
        return

    # Resolve path
    if not os.path.isabs(image_path):
        base_dir = deck_meta.get("base_dir", ".")
        image_path = os.path.join(base_dir, image_path)

    if not os.path.exists(image_path):
        print(f"    WARNING: Image not found: {image_path}")
        return

    # Available area
    left_margin = 0.35
    right_margin = 0.35
    avail_top = 0.85
    avail_bottom = 4.7 if caption else 5.0
    max_w = 10.0 - left_margin - right_margin
    max_h = avail_bottom - avail_top

    # Read image aspect ratio
    try:
        with PILImage.open(image_path) as img:
            aspect = img.width / img.height
    except Exception:
        aspect = 16 / 9

    # Fit within bounds
    if max_w / max_h > aspect:
        img_h = max_h
        img_w = img_h * aspect
    else:
        img_w = max_w
        img_h = img_w / aspect

    img_left = left_margin + (max_w - img_w) / 2
    img_top = avail_top + (max_h - img_h) / 2

    pic = slide.shapes.add_picture(
        image_path,
        SI(img_left), SI(img_top),
        SI(img_w), SI(img_h))

    # Optional border
    if show_border:
        pic.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        pic.line.width = Pt(0.5)

    # Caption
    if caption:
        cap_h = estimate_text_height(caption, max_w, 8)
        txb = slide.shapes.add_textbox(
            SI(left_margin), SI(img_top + img_h + 0.08),
            SI(max_w), SI(cap_h))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = caption; r.font.name = BODY
        r.font.size = Pt(8); r.font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _build_matrix(slide, slide_def, deck_meta):
    """Build a color-coded matrix slide (RACI, capability grid, etc.).

    Fields: headline, row_header, col_header, rows (list of {label, values}),
            columns (list of column names), cell_colors (dict mapping value → hex),
            legend (optional string).
    """
    from pptx.oxml.ns import qn

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    WHITE = brand.white
    DARK = brand.text_dark
    GRAY = brand.text_gray

    columns = slide_def.get("columns", [])
    rows = slide_def.get("rows", [])
    cell_colors = slide_def.get("cell_colors", {})
    legend = slide_def.get("legend", "")

    if not columns or not rows:
        return

    num_cols = len(columns) + 1  # +1 for row labels
    num_rows = len(rows) + 1  # +1 for header

    left_margin = 0.35
    table_top = 0.85
    table_width = 9.3
    row_label_w = 2.0
    data_col_w = (table_width - row_label_w) / len(columns)

    header_h = 0.35
    data_row_h = min(0.4, (4.0 - header_h) / max(len(rows), 1))
    table_height = header_h + data_row_h * len(rows)

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        SI(left_margin), SI(table_top),
        SI(table_width), SI(table_height))
    table = tbl_shape.table

    # Column widths
    table.columns[0].width = SI(row_label_w)
    for ci in range(1, num_cols):
        table.columns[ci].width = SI(data_col_w)

    # Header row
    # First cell: row header label
    cell = table.cell(0, 0)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = slide_def.get("row_header", "")
    r.font.name = HEADING; r.font.size = Pt(9)
    r.font.color.rgb = WHITE; r.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '5F016F'})
    solidFill.append(srgb); tcPr.append(solidFill)

    for ci, col_name in enumerate(columns):
        cell = table.cell(0, ci + 1)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(col_name)
        r.font.name = HEADING; r.font.size = Pt(9)
        r.font.color.rgb = WHITE; r.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '5F016F'})
        solidFill.append(srgb); tcPr.append(solidFill)

    # Data rows
    for ri, row in enumerate(rows):
        label = row.get("label", "")
        values = row.get("values", [])

        # Row label
        cell = table.cell(ri + 1, 0)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label; r.font.name = BODY
        r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.LEFT

        # Values
        for ci, val in enumerate(values):
            cell = table.cell(ri + 1, ci + 1)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val); r.font.name = HEADING
            r.font.size = Pt(11); r.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Color the cell background if value matches cell_colors
            if str(val) in cell_colors:
                hex_color = cell_colors[str(val)]
                r.font.color.rgb = WHITE
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': hex_color.lstrip('#')})
                solidFill.append(srgb); tcPr.append(solidFill)
            else:
                r.font.color.rgb = DARK

    # Remove default table style
    tbl_xml = table._tbl
    tblPr = tbl_xml.find(qn('a:tblPr'))
    if tblPr is not None:
        for attr in ('bandRow', 'bandCol', 'firstRow', 'lastRow', 'firstCol', 'lastCol'):
            tblPr.set(attr, '0')

    # Legend
    if legend:
        legend_top = table_top + table_height + 0.12
        txb = slide.shapes.add_textbox(
            SI(left_margin), SI(legend_top),
            SI(table_width), SI(0.25))
        tf = txb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = legend; r.font.name = BODY
        r.font.size = Pt(7); r.font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _build_funnel(slide, slide_def, deck_meta):
    """Build a funnel/pipeline slide with progressively narrowing bars.

    Fields: headline, stages (list of {label, value, width (0-100)}).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    PINK = brand.secondary
    WHITE = WHITE

    stages = slide_def.get("stages", [])
    n = len(stages)
    if n == 0:
        return

    max_w = 8.0
    bar_h = 0.55
    gap = 0.08
    total_h = n * bar_h + (n - 1) * gap

    avail_top = 0.85
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - total_h) / 2
    offset = max(offset, avail_top)

    # Color gradient from dark purple to pink
    for i, stage in enumerate(stages):
        label = stage.get("label", "")
        value = stage.get("value", "")
        width_pct = stage.get("width", 100)
        bar_w = max_w * (width_pct / 100)
        bar_left = (10.0 - bar_w) / 2  # centered
        bar_top = offset + i * (bar_h + gap)

        # Interpolate color from primary to secondary
        t = i / max(n - 1, 1)
        bar_color = brand.interpolate(t)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(bar_left), SI(bar_top),
            SI(bar_w), SI(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        # Single text box inside bar with label left + value right
        text_w = bar_w - 0.16
        full_text = label + ("  " + str(value) if value else "")
        font_pt = fit_font_size(full_text, text_w, max_pt=12, min_pt=6)
        txb = slide.shapes.add_textbox(
            SI(bar_left + 0.08), SI(bar_top + 0.05),
            SI(text_w), SI(bar_h - 0.1))
        tf = txb.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label; r.font.name = HEADING
        r.font.size = Pt(font_pt); r.font.color.rgb = WHITE
        if value:
            r2 = p.add_run()
            r2.text = "  " + str(value); r2.font.name = BODY
            r2.font.size = Pt(max(font_pt - 1, 6)); r2.font.color.rgb = WHITE


def _build_callout(slide, slide_def, deck_meta):
    """Build a callout/highlight slide for key takeaways.

    Fields: headline, callout_text, supporting_text (optional),
            icon (optional), style ("boxed"|"open").
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    WHITE = brand.white
    DARK = brand.text_dark
    LIGHT_PINK = brand.accent

    callout_text = slide_def.get("callout_text", "")
    supporting_text = slide_def.get("supporting_text", "")
    icon_name = slide_def.get("icon", "")
    box_style = slide_def.get("style", "boxed")

    is_boxed = (box_style == "boxed")
    text_w = 7.0
    margin_x = (10.0 - text_w) / 2

    # Calculate heights
    icon_path = _resolve_icon(icon_name) if icon_name else None
    icon_h = 0.6 if (icon_path and os.path.exists(icon_path)) else 0
    callout_h = estimate_text_height(callout_text, text_w - 0.6, 20)
    support_h = estimate_text_height(supporting_text, text_w - 0.6, 12) if supporting_text else 0
    rule_h = 0.3 if (supporting_text and not is_boxed) else 0

    inner_h = icon_h + callout_h + 0.15 + rule_h + support_h
    total_h = inner_h + (0.8 if is_boxed else 0)  # padding for box (0.4 top + 0.4 bottom)

    avail_top = 0.85
    avail_bottom = 5.0
    avail = avail_bottom - avail_top
    offset = avail_top + (avail - total_h) / 2
    offset = max(offset, avail_top)

    if is_boxed:
        # Purple background card
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(margin_x), SI(offset),
            SI(text_w), SI(total_h))
        bg.fill.solid(); bg.fill.fore_color.rgb = PURPLE
        bg.line.fill.background()
        callout_color = WHITE
        support_color = LIGHT_PINK
    else:
        callout_color = PURPLE
        support_color = DARK

    pad = 0.4 if is_boxed else 0
    cur_y = offset + pad

    # Icon
    if icon_path and os.path.exists(icon_path):
        icon_left = (10.0 - 0.5) / 2
        slide.shapes.add_picture(icon_path, SI(icon_left), SI(cur_y),
                                 SI(0.5), SI(0.5))
        cur_y += 0.6

    # Callout text
    txb = slide.shapes.add_textbox(SI(margin_x + 0.3), SI(cur_y),
                                   SI(text_w - 0.6), SI(callout_h))
    tf = txb.text_frame; tf.word_wrap = True
    for li, line in enumerate(callout_text.split("\n")):
        para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        _add_run_with_hyperlinks(para, line, HEADING, 20, callout_color)
        para.alignment = PP_ALIGN.CENTER
    cur_y += callout_h + 0.1

    # Rule (open style only)
    if supporting_text and not is_boxed:
        rule_w = 2.0
        rule_left = (10.0 - rule_w) / 2
        line = slide.shapes.add_connector(
            1, SI(rule_left), SI(cur_y),
            SI(rule_left + rule_w), SI(cur_y))
        line.line.color.rgb = PINK; line.line.width = Pt(1.5)
        cur_y += 0.2

    # Supporting text
    if supporting_text:
        txb = slide.shapes.add_textbox(SI(margin_x + 0.3), SI(cur_y),
                                       SI(text_w - 0.6), SI(support_h))
        tf = txb.text_frame; tf.word_wrap = True
        for li, line in enumerate(supporting_text.split("\n")):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            _add_run_with_hyperlinks(para, line, BODY, 12, support_color)
            para.alignment = PP_ALIGN.CENTER


def _build_process_flow(slide, slide_def, deck_meta):
    """Build a horizontal process-flow slide with connected chevron or circle steps.

    Fields: headline, steps (list of {label, body (optional), status (optional)}),
            style ("chevron"|"circles", default "chevron").
    Status values: complete (filled dark), active (filled pink), default (outlined).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark
    GRAY = RGBColor(0xBB, 0xBB, 0xBB)

    steps = slide_def.get("steps", [])
    style = slide_def.get("style", "chevron")
    n = len(steps)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.5
    right_margin = 0.5
    total_w = 10.0 - left_margin - right_margin  # 9.0"

    # Step dimensions
    gap = 0.15
    step_w = (total_w - (n - 1) * gap) / n
    step_h = 0.65
    body_h = 0.60

    # Vertical centering
    block_h = step_h + 0.10 + body_h  # step + gap + body
    avail_h = avail_bottom - avail_top
    start_y = avail_top + (avail_h - block_h) / 2
    start_y = max(start_y, avail_top)

    for i, step in enumerate(steps):
        label = step.get("label", "")
        body = step.get("body", "")
        status = step.get("status", "")
        x = left_margin + i * (step_w + gap)

        # Determine colors based on status
        if status == "complete":
            fill_color, text_color = PURPLE, WHITE
        elif status == "active":
            fill_color, text_color = PINK, WHITE
        else:
            fill_color, text_color = LIGHT_BG, PURPLE

        if style == "circles":
            # Circle with number
            circle_size = min(step_w * 0.65, step_h)
            cx = x + (step_w - circle_size) / 2
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, SI(cx), SI(start_y),
                SI(circle_size), SI(circle_size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.fill.background()
            # Number inside circle
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(i + 1)
            r.font.name = HEADING
            r.font.size = Pt(16)
            r.font.color.rgb = text_color
        else:
            # Chevron/rounded rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, SI(x), SI(start_y),
                SI(step_w), SI(step_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.fill.background()
            # Label inside shape
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            step_font = fit_font_size(label, step_w - 0.20, max_pt=11 if n <= 5 else 9, min_pt=6)
            r = p.add_run()
            r.text = label
            r.font.name = HEADING
            r.font.size = Pt(step_font)
            r.font.color.rgb = text_color

        # Connector arrow between steps
        if i < n - 1:
            arrow_x = x + step_w + 0.02
            arrow_y = start_y + step_h / 2 - 0.06
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, SI(arrow_x), SI(arrow_y),
                SI(gap - 0.04), SI(0.12))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

        # Body text below
        desc_text = label if style == "circles" else ""
        if body:
            desc_text = (desc_text + "\n" + body).strip() if desc_text else body
        if desc_text:
            txb = slide.shapes.add_textbox(
                SI(x), SI(start_y + step_h + 0.10),
                SI(step_w), SI(body_h))
            tf = txb.text_frame
            tf.word_wrap = True
            lines = desc_text.split("\n")
            for li, line in enumerate(lines):
                para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                para.alignment = PP_ALIGN.CENTER
                r = para.add_run()
                r.text = line
                if li == 0 and style == "circles":
                    r.font.name = HEADING
                    r.font.size = Pt(10 if n <= 5 else 8)
                    r.font.color.rgb = PURPLE
                else:
                    r.font.name = BODY
                    r.font.size = Pt(9 if n <= 5 else 7)
                    r.font.color.rgb = DARK


def _build_comparison_matrix(slide, slide_def, deck_meta):
    """Build a feature comparison matrix with Harvey balls or checkmarks.

    Fields: headline, row_header (label for first column),
            columns (list of option names),
            highlight_column (optional — column name to accent),
            rows (list of {label, values: [level per column]}),
            legend (optional dict mapping level names to descriptions).

    Value levels: "full", "three-quarter", "half", "quarter", "none",
                  "check", "cross", or short text.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark
    GREEN = brand.green
    RED = brand.red

    columns = slide_def.get("columns", [])
    rows = slide_def.get("rows", [])
    row_header = slide_def.get("row_header", "Criteria")
    highlight_col = slide_def.get("highlight_column", "")
    if not columns or not rows:
        return

    num_cols = len(columns) + 1  # +1 for row header column
    num_rows = len(rows) + 1     # +1 for header row

    left_margin = 0.35
    avail_top = 0.85
    avail_bottom = 5.05
    table_width = 10.0 - left_margin - 0.35

    # Column widths: row header gets 30%, rest split equally
    header_w = table_width * 0.30
    data_w = (table_width - header_w) / len(columns)
    col_widths = [header_w] + [data_w] * len(columns)

    # Row heights
    row_h = min(0.42, (avail_bottom - avail_top) / num_rows)
    header_h = 0.35
    table_height = header_h + row_h * len(rows)

    # Center vertically
    avail_h = avail_bottom - avail_top
    table_top = avail_top + (avail_h - table_height) / 2
    table_top = max(table_top, avail_top)

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        SI(left_margin), SI(table_top),
        SI(table_width), SI(table_height))
    table = tbl_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = SI(w)
    table.rows[0].height = SI(header_h)
    for ri in range(1, num_rows):
        table.rows[ri].height = SI(row_h)

    # Header row
    all_headers = [row_header] + columns
    header_bg_hex = "5F016F"
    for ci, hdr in enumerate(all_headers):
        cell = table.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(hdr)
        r.font.name = HEADING
        r.font.size = Pt(9)
        r.font.color.rgb = WHITE
        r.font.bold = True
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': header_bg_hex})
        solidFill.append(srgb)
        tcPr.append(solidFill)

    # Harvey ball symbols
    HARVEY = {
        "full": "\u25CF",           # ●
        "three-quarter": "\u25D5",  # ◕ (not available everywhere, use ●)
        "three_quarter": "\u25D5",
        "half": "\u25D1",           # ◑
        "quarter": "\u25D4",        # ◔
        "none": "\u25CB",           # ○
        "check": "\u2713",          # ✓
        "cross": "\u2717",          # ✗
    }
    HARVEY_COLOR = {
        "full": PURPLE, "three-quarter": PURPLE, "three_quarter": PURPLE,
        "half": PURPLE, "quarter": PURPLE, "none": DARK,
        "check": GREEN, "cross": RED,
    }

    # Data rows
    highlight_ci = None
    if highlight_col:
        for ci, c in enumerate(columns):
            if c == highlight_col:
                highlight_ci = ci + 1  # +1 for row header offset

    for ri, row in enumerate(rows):
        label = row.get("label", "")
        values = row.get("values", [])

        # Row header cell
        cell = table.cell(ri + 1, 0)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(label)
        r.font.name = BODY
        r.font.size = Pt(9)
        r.font.color.rgb = PURPLE
        r.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Value cells
        for ci, val in enumerate(values):
            cell = table.cell(ri + 1, ci + 1)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            val_str = str(val).lower().replace("-", "_")

            if val_str in HARVEY:
                r = p.add_run()
                r.text = HARVEY[val_str]
                r.font.name = BODY
                r.font.size = Pt(16)
                r.font.color.rgb = HARVEY_COLOR.get(val_str, DARK)
            else:
                r = p.add_run()
                r.text = str(val)
                r.font.name = BODY
                r.font.size = Pt(9)
                r.font.color.rgb = DARK

            # Highlight column background
            if highlight_ci and ci + 1 == highlight_ci:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'F0E8F5'})
                solidFill.append(srgb)
                tcPr.append(solidFill)
            elif ri % 2 == 1:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFF0F8'})
                solidFill.append(srgb)
                tcPr.append(solidFill)

        # Stripe the row header cell too
        if highlight_ci is None and ri % 2 == 1:
            tcPr = table.cell(ri + 1, 0)._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn('a:solidFill'), {})
            srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFF0F8'})
            solidFill.append(srgb)
            tcPr.append(solidFill)

    # Remove default style + apply borders
    tbl_xml = table._tbl
    tblPr = tbl_xml.find(qn('a:tblPr'))
    if tblPr is not None:
        for attr in ('bandRow', 'bandCol', 'firstRow', 'lastRow', 'firstCol', 'lastCol'):
            tblPr.set(attr, '0')
    _apply_table_borders(table, num_rows, num_cols, header_bg_hex)


def _build_quadrant(slide, slide_def, deck_meta):
    """Build a 2x2 quadrant matrix (SWOT, priority, positioning).

    Fields: headline, x_axis (label), y_axis (label),
            quadrants (list of 4 dicts: {position, title, color (optional), items}).
    Position values: top-left, top-right, bottom-left, bottom-right.
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    PINK = brand.secondary
    DARK = brand.text_dark
    WHITE = WHITE

    quadrants = slide_def.get("quadrants", [])
    x_axis = slide_def.get("x_axis", "")
    y_axis = slide_def.get("y_axis", "")
    if len(quadrants) < 4:
        return

    # Layout constants
    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.90  # room for y-axis label (wider to prevent wrap)
    right_margin = 0.35
    grid_w = 10.0 - left_margin - right_margin
    grid_h = avail_bottom - avail_top - 0.30  # room for x-axis label
    quad_w = (grid_w - 0.08) / 2   # 0.08 gap
    quad_h = (grid_h - 0.08) / 2
    grid_top = avail_top

    # Default colors per position
    DEFAULT_COLORS = {
        "top-left": "4CAF50",     # green
        "top-right": "2196F3",    # blue
        "bottom-left": "FFC107",  # amber
        "bottom-right": "F44336", # red
    }
    POS_COORDS = {
        "top-left": (left_margin, grid_top),
        "top-right": (left_margin + quad_w + 0.08, grid_top),
        "bottom-left": (left_margin, grid_top + quad_h + 0.08),
        "bottom-right": (left_margin + quad_w + 0.08, grid_top + quad_h + 0.08),
    }

    for q in quadrants:
        pos = q.get("position", "top-left")
        title = q.get("title", "")
        items = q.get("items", [])
        color_hex = q.get("color", DEFAULT_COLORS.get(pos, brand.primary_hex))
        if color_hex.startswith("#"):
            color_hex = color_hex[1:]
        fill_color = RGBColor.from_string(color_hex)

        x, y = POS_COORDS.get(pos, (left_margin, grid_top))

        # Quadrant background
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, SI(x), SI(y),
            SI(quad_w), SI(quad_h))
        rect.fill.solid()
        # Use a light tint of the color (mix with white)
        r_val = min(255, fill_color[0] + (255 - fill_color[0]) * 80 // 100)
        g_val = min(255, fill_color[1] + (255 - fill_color[1]) * 80 // 100)
        b_val = min(255, fill_color[2] + (255 - fill_color[2]) * 80 // 100)
        rect.fill.fore_color.rgb = RGBColor(r_val, g_val, b_val)
        rect.line.color.rgb = RGBColor.from_string("E0E0E0")
        rect.line.width = Pt(0.5)

        # Title bar at top of quadrant
        title_h = 0.30
        title_box = slide.shapes.add_textbox(
            SI(x + 0.10), SI(y + 0.08),
            SI(quad_w - 0.20), SI(title_h))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = HEADING
        r.font.size = Pt(10)
        r.font.color.rgb = fill_color

        # Items
        if items:
            items_box = slide.shapes.add_textbox(
                SI(x + 0.10), SI(y + title_h + 0.10),
                SI(quad_w - 0.20), SI(quad_h - title_h - 0.20))
            tf = items_box.text_frame
            tf.word_wrap = True
            for ii, item in enumerate(items):
                para = tf.paragraphs[0] if ii == 0 else tf.add_paragraph()
                r = para.add_run()
                r.text = "\u2022 " + str(item)
                r.font.name = BODY
                r.font.size = Pt(8)
                r.font.color.rgb = DARK

    # X-axis label (bottom center)
    if x_axis:
        add_text_box(slide, x_axis,
                     left_margin, grid_top + grid_h + 0.05, grid_w, 0.25,
                     font_size=9, font_name=BODY, bold=True,
                     color=DARK, alignment=PP_ALIGN.CENTER)
        # Arrow indicators
        arrow_y = grid_top + grid_h + 0.02
        add_text_box(slide, "\u2190 Low", left_margin, arrow_y, 1.0, 0.20,
                     font_size=7, font_name=BODY, color=DARK)
        add_text_box(slide, "High \u2192", left_margin + grid_w - 1.0, arrow_y, 1.0, 0.20,
                     font_size=7, font_name=BODY, color=DARK,
                     alignment=PP_ALIGN.RIGHT)

    # Y-axis label (left side, vertical — rendered as horizontal short text)
    if y_axis:
        add_text_box(slide, "\u2191 High", 0.05, avail_top, 0.80, 0.20,
                     font_size=7, font_name=BODY, color=DARK)
        add_text_box(slide, y_axis, 0.05, avail_top + grid_h / 2 - 0.10, 0.80, 0.25,
                     font_size=8, font_name=BODY, bold=True, color=DARK)
        add_text_box(slide, "\u2193 Low", 0.05, avail_top + grid_h - 0.20, 0.80, 0.20,
                     font_size=7, font_name=BODY, color=DARK)


def _build_team_profiles(slide, slide_def, deck_meta):
    """Build a team/people profile cards slide.

    Fields: headline, profiles (list of {name, role, icon (optional), context (optional)}).
    Supports 2-8 profiles arranged in rows.
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    DARK = brand.text_dark
    WHITE = brand.white

    profiles = slide_def.get("profiles", [])
    n = len(profiles)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.35
    total_w = 10.0 - left_margin * 2
    avail_h = avail_bottom - avail_top

    # Grid arrangement
    if n <= 4:
        cols = n
        rows_count = 1
    elif n <= 6:
        cols = 3
        rows_count = 2
    else:
        cols = 4
        rows_count = 2

    gap_x = 0.15
    gap_y = 0.15
    card_w = (total_w - (cols - 1) * gap_x) / cols
    card_h = (avail_h - (rows_count - 1) * gap_y) / rows_count
    card_h = min(card_h, 2.0)  # cap height

    # Center vertically
    total_h = rows_count * card_h + (rows_count - 1) * gap_y
    start_y = avail_top + (avail_h - total_h) / 2
    start_y = max(start_y, avail_top)

    for idx, profile in enumerate(profiles):
        row = idx // cols
        col = idx % cols
        x = left_margin + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        name = profile.get("name", "")
        role = profile.get("role", "")
        icon_name = profile.get("icon", "")
        context = profile.get("context", "")

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, SI(x), SI(y),
            SI(card_w), SI(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.fill.background()

        # Icon circle at top center
        icon_size = 0.50
        icon_x = x + (card_w - icon_size) / 2
        icon_y_pos = y + 0.15
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, SI(icon_x), SI(icon_y_pos),
            SI(icon_size), SI(icon_size))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PURPLE
        circle.line.fill.background()

        # Icon or initials inside circle
        if icon_name:
            icon_path = _resolve_icon(icon_name)
            if icon_path:
                icon_img_size = icon_size * 0.55
                icon_img_x = icon_x + (icon_size - icon_img_size) / 2
                icon_img_y = icon_y_pos + (icon_size - icon_img_size) / 2
                try:
                    slide.shapes.add_picture(
                        icon_path,
                        SI(icon_img_x), SI(icon_img_y),
                        SI(icon_img_size), SI(icon_img_size))
                except Exception:
                    pass
        if not icon_name or not _resolve_icon(icon_name):
            # Fallback: initials
            initials = "".join(w[0] for w in name.split()[:2]).upper()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = initials
            r.font.name = HEADING
            r.font.size = Pt(14)
            r.font.color.rgb = WHITE

        # Name
        text_top = icon_y_pos + icon_size + 0.10
        add_text_box(slide, name, x + 0.08, text_top, card_w - 0.16, 0.25,
                     font_size=11, font_name=HEADING, bold=True,
                     color=PURPLE, alignment=PP_ALIGN.CENTER)

        # Role
        add_text_box(slide, role, x + 0.08, text_top + 0.25, card_w - 0.16, 0.20,
                     font_size=9, font_name=BODY, color=PINK,
                     alignment=PP_ALIGN.CENTER)

        # Context
        if context:
            add_text_box(slide, context, x + 0.08, text_top + 0.50,
                         card_w - 0.16, card_h - (text_top - y) - 0.60,
                         font_size=8, font_name=BODY, color=DARK,
                         alignment=PP_ALIGN.CENTER)


def _build_pros_cons(slide, slide_def, deck_meta):
    """Build a pros/cons evaluation slide with green/red columns.

    Fields: headline, pros_label (default "Strengths"), cons_label (default "Risks"),
            pros (list of strings), cons (list of strings),
            recommendation (optional string at bottom).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    DARK = brand.text_dark
    PURPLE = brand.primary
    GREEN = brand.green
    RED = brand.red
    GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
    RED_LIGHT = RGBColor(0xFD, 0xE8, 0xE8)
    WHITE = WHITE
    DARK = DARK
    PINK = brand.secondary

    pros = slide_def.get("pros", [])
    cons = slide_def.get("cons", [])
    pros_label = slide_def.get("pros_label", "Strengths")
    cons_label = slide_def.get("cons_label", "Risks")
    recommendation = slide_def.get("recommendation", "")

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.35
    total_w = 10.0 - left_margin * 2
    col_w = (total_w - 0.20) / 2  # 0.20 gap between columns
    reco_h = 0.50 if recommendation else 0
    avail_h = avail_bottom - avail_top - reco_h

    # Column header height
    header_h = 0.40
    item_h = 0.30
    max_items = max(len(pros), len(cons))
    content_h = header_h + max_items * item_h
    content_h = min(content_h, avail_h)

    start_y = avail_top + (avail_h - content_h) / 2
    start_y = max(start_y, avail_top)

    # --- Pros column (left) ---
    pros_x = left_margin

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, SI(pros_x), SI(start_y),
        SI(col_w), SI(header_h))
    header.fill.solid()
    header.fill.fore_color.rgb = GREEN
    header.line.fill.background()
    tf = header.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\u2713  " + pros_label
    r.font.name = HEADING
    r.font.size = Pt(12)
    r.font.color.rgb = WHITE

    # Pros background
    pros_bg_h = content_h - header_h
    if pros_bg_h > 0:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, SI(pros_x), SI(start_y + header_h),
            SI(col_w), SI(pros_bg_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = GREEN_LIGHT
        bg.line.fill.background()

    # Pros items
    for i, item in enumerate(pros):
        item_y = start_y + header_h + 0.06 + i * item_h
        txb = slide.shapes.add_textbox(
            SI(pros_x + 0.12), SI(item_y),
            SI(col_w - 0.24), SI(item_h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # Checkmark
        r = p.add_run()
        r.text = "\u2713  "
        r.font.name = BODY
        r.font.size = Pt(9)
        r.font.color.rgb = GREEN
        r.font.bold = True
        # Text
        r2 = p.add_run()
        r2.text = str(item)
        r2.font.name = BODY
        r2.font.size = Pt(9)
        r2.font.color.rgb = DARK

    # --- Cons column (right) ---
    cons_x = left_margin + col_w + 0.20

    header2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, SI(cons_x), SI(start_y),
        SI(col_w), SI(header_h))
    header2.fill.solid()
    header2.fill.fore_color.rgb = RED
    header2.line.fill.background()
    tf = header2.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\u2717  " + cons_label
    r.font.name = HEADING
    r.font.size = Pt(12)
    r.font.color.rgb = WHITE

    cons_bg_h = content_h - header_h
    if cons_bg_h > 0:
        bg2 = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, SI(cons_x), SI(start_y + header_h),
            SI(col_w), SI(cons_bg_h))
        bg2.fill.solid()
        bg2.fill.fore_color.rgb = RED_LIGHT
        bg2.line.fill.background()

    for i, item in enumerate(cons):
        item_y = start_y + header_h + 0.06 + i * item_h
        txb = slide.shapes.add_textbox(
            SI(cons_x + 0.12), SI(item_y),
            SI(col_w - 0.24), SI(item_h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "\u2717  "
        r.font.name = BODY
        r.font.size = Pt(9)
        r.font.color.rgb = RED
        r.font.bold = True
        r2 = p.add_run()
        r2.text = str(item)
        r2.font.name = BODY
        r2.font.size = Pt(9)
        r2.font.color.rgb = DARK

    # --- Recommendation bar ---
    if recommendation:
        reco_y = start_y + content_h + 0.15
        reco_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(left_margin), SI(reco_y),
            SI(total_w), SI(0.35))
        reco_box.fill.solid()
        reco_box.fill.fore_color.rgb = PURPLE
        reco_box.line.fill.background()
        tf = reco_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = recommendation
        r.font.name = BODY
        r.font.size = Pt(10)
        r.font.color.rgb = WHITE
        r.font.bold = True


def _build_staircase(slide, slide_def, deck_meta):
    """Build a staircase/maturity model with ascending levels.

    Fields: headline, levels (list of {label, body}), current_level (1-based index).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark

    levels = slide_def.get("levels", [])
    current = slide_def.get("current_level", 0)
    n = len(levels)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.50
    total_w = 10.0 - left_margin * 2
    avail_h = avail_bottom - avail_top

    step_w = total_w / n
    step_base_h = 0.50
    step_increment = (avail_h - step_base_h) / max(n - 1, 1)

    for i, level in enumerate(levels):
        label = level.get("label", "")
        body = level.get("body", "")
        step_h = step_base_h + i * step_increment
        x = left_margin + i * step_w
        y = avail_bottom - step_h

        is_current = (i + 1) == current
        fill = PINK if is_current else LIGHT_BG
        text_color = WHITE if is_current else PURPLE
        # Gradient from light to dark
        if not is_current:
            t = i / max(n - 1, 1)
            fill = brand.interpolate(t, "background_light", "staircase_end")
            text_color = PURPLE

        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SI(x), SI(y),
            SI(step_w - 0.04), SI(step_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.fill.background()

        # Level number + label at top of step
        txb = slide.shapes.add_textbox(
            SI(x + 0.08), SI(y + 0.06),
            SI(step_w - 0.20), SI(0.30))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{i+1}. {label}"
        r.font.name = HEADING
        r.font.size = Pt(9 if n <= 5 else 7)
        r.font.color.rgb = text_color
        r.font.bold = True

        if body:
            txb2 = slide.shapes.add_textbox(
                SI(x + 0.08), SI(y + 0.36),
                SI(step_w - 0.20), SI(step_h - 0.44))
            tf2 = txb2.text_frame; tf2.word_wrap = True
            p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = body
            r2.font.name = BODY
            r2.font.size = Pt(7 if n >= 5 else 8)
            r2.font.color.rgb = DARK if not is_current else WHITE


def _build_donut_rings(slide, slide_def, deck_meta):
    """Build a progress ring dashboard with donut arcs.

    Fields: headline, rings (list of {value (0-100), label, color (green/amber/red/purple)}).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark
    GRAY_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)

    COLOR_MAP = {
        "green": brand.green,
        "amber": brand.amber,
        "red": brand.red,
        "purple": PURPLE,
        "pink": brand.secondary,
    }

    rings = slide_def.get("rings", [])
    n = len(rings)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.50
    total_w = 10.0 - left_margin * 2
    avail_h = avail_bottom - avail_top

    # Layout: rings in a single row
    ring_size = min(total_w / n - 0.15, avail_h - 0.60, 1.6)
    gap = (total_w - n * ring_size) / max(n + 1, 1)
    start_y = avail_top + (avail_h - ring_size - 0.50) / 2

    for i, ring in enumerate(rings):
        value = ring.get("value", 0)
        label = ring.get("label", "")
        color_name = ring.get("color", "purple")

        # Auto-color by threshold if not specified
        if color_name == "auto":
            color_name = "green" if value >= 80 else "amber" if value >= 50 else "red"
        ring_color = COLOR_MAP.get(color_name, PURPLE)

        cx = left_margin + gap + i * (ring_size + gap)
        cy = start_y

        # Background ring (full donut, gray)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.DONUT, SI(cx), SI(cy),
            SI(ring_size), SI(ring_size))
        bg.fill.solid()
        bg.fill.fore_color.rgb = GRAY_LIGHT
        bg.line.fill.background()
        bg.adjustments[0] = 0.25  # ring thickness

        # Foreground progress arc — use BLOCK_ARC for a thick colored arc
        if value > 0 and value <= 100:
            fg = slide.shapes.add_shape(
                MSO_SHAPE.BLOCK_ARC, SI(cx), SI(cy),
                SI(ring_size), SI(ring_size))
            fg.fill.solid()
            fg.fill.fore_color.rgb = ring_color
            fg.line.fill.background()
            # BLOCK_ARC adjustments:
            # adj1 = start angle (fraction of 360, 0=top/right in OOXML, but
            #         for visual "start at top" we use 270/360 = 0.75)
            # adj2 = ring thickness (same as donut)
            sweep_frac = value / 100.0
            # Start at top (270 degrees OOXML = 12 o'clock position)
            start_frac = 0.75  # 270/360
            end_frac = start_frac + sweep_frac
            if end_frac > 1.0:
                end_frac -= 1.0
            fg.adjustments[0] = start_frac
            fg.adjustments[1] = 0.25  # ring thickness matching background

        # Value text in center of the donut hole
        # Inner hole starts at 25% ring thickness from each edge
        inner_offset = ring_size * 0.25
        inner_size = ring_size - 2 * inner_offset
        val_text = f"{value}%"
        val_font = fit_font_size(val_text, inner_size,
                                  max_pt=min(18, int(ring_size * 11)), min_pt=8)
        # Text box spans full ring width for centering (text is smaller than hole)
        val_box = slide.shapes.add_textbox(
            SI(cx + inner_offset * 0.5), SI(cy + ring_size * 0.33),
            SI(ring_size - inner_offset), SI(inner_size * 0.50))
        tf = val_box.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{value}%"
        r.font.name = HEADING
        r.font.size = Pt(val_font)
        r.font.color.rgb = ring_color

        # Label below
        lbl_box = slide.shapes.add_textbox(
            SI(cx - 0.10), SI(cy + ring_size + 0.08),
            SI(ring_size + 0.20), SI(0.35))
        tf2 = lbl_box.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = BODY
        r2.font.size = Pt(9)
        r2.font.color.rgb = DARK


def _build_pyramid(slide, slide_def, deck_meta):
    """Build a pyramid/hierarchy with 3-5 horizontal tiers.

    Fields: headline, tiers (list of {label, body} from top to bottom).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    WHITE = brand.white
    DARK = brand.text_dark

    tiers = slide_def.get("tiers", [])
    n = len(tiers)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    center_x = 5.0  # slide center
    max_w = 8.0
    min_w = 2.0
    tier_h = min(avail_h / n - 0.04, 0.75)
    total_h = n * tier_h + (n - 1) * 0.04
    start_y = avail_top + (avail_h - total_h) / 2

    for i, tier in enumerate(tiers):
        label = tier.get("label", "")
        body = tier.get("body", "")
        # Width narrows toward top (tier 0 = narrowest = top)
        t = i / max(n - 1, 1)
        w = min_w + (max_w - min_w) * t
        x = center_x - w / 2
        y = start_y + i * (tier_h + 0.04)

        # Color gradient: darkest at top, lightest at bottom
        t_color = i / max(n - 1, 1)
        fill_color = brand.interpolate(t_color, "gradient_start", "gradient_light")
        text_color = WHITE if t_color < 0.5 else DARK

        # Trapezoid shape (use TRAPEZOID or RECTANGLE with adjusted corners)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SI(x), SI(y),
            SI(w), SI(tier_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill_color
        rect.line.fill.background()

        # Label (bold) + body text
        txb = slide.shapes.add_textbox(
            SI(x + 0.15), SI(y + 0.04),
            SI(w - 0.30), SI(tier_h - 0.08))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = HEADING
        r.font.size = Pt(11 if n <= 4 else 9)
        r.font.color.rgb = text_color
        r.font.bold = True
        if body:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = body
            r2.font.name = BODY
            r2.font.size = Pt(8 if n <= 4 else 7)
            r2.font.color.rgb = text_color


def _build_venn(slide, slide_def, deck_meta):
    """Build a Venn diagram with 2-3 overlapping circles.

    Fields: headline, circles (list of {label, color}),
            intersections (list of {regions (index list), label}).
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    DARK = brand.text_dark
    WHITE = WHITE

    COLOR_MAP = {
        "purple": brand.primary_hex, "pink": brand.secondary_hex,
        "light_purple": brand.color_hex("venn_tertiary"),
        "blue": "2196F3", "green": "4CAF50", "amber": "FFC107",
        "red": "F44336", "teal": "009688",
    }

    circles_def = slide_def.get("circles", [])
    intersections = slide_def.get("intersections", [])
    n = len(circles_def)
    if n < 2:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    center_x = 5.0
    center_y = avail_top + avail_h / 2
    circle_r = min(avail_h * 0.42, 1.8)
    overlap = circle_r * 0.55  # how much circles overlap

    # Circle center positions
    if n == 2:
        positions = [
            (center_x - overlap / 2, center_y),
            (center_x + overlap / 2, center_y),
        ]
    else:  # 3 circles
        import math as _math
        positions = [
            (center_x, center_y - overlap * 0.55),
            (center_x - overlap * 0.50, center_y + overlap * 0.35),
            (center_x + overlap * 0.50, center_y + overlap * 0.35),
        ]

    for i, cdef in enumerate(circles_def):
        label = cdef.get("label", "")
        color_name = cdef.get("color", "purple")
        color_hex = COLOR_MAP.get(color_name, color_name.replace("#", ""))

        cx, cy = positions[i]
        x = cx - circle_r
        y = cy - circle_r

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, SI(x), SI(y),
            SI(circle_r * 2), SI(circle_r * 2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        # Set transparency via OOXML (60% transparent)
        try:
            sp_elem = oval._element
            spPr = sp_elem.find(qn('a:solidFill')) or sp_elem.find('.//' + qn('a:solidFill'))
            if spPr is None:
                for sf in sp_elem.iter(qn('a:solidFill')):
                    spPr = sf
                    break
            if spPr is not None:
                srgb = spPr.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha = srgb.makeelement(qn('a:alpha'), {'val': '40000'})
                    srgb.append(alpha)
        except Exception:
            pass  # transparency is cosmetic, don't fail the build
        oval.line.fill.background()

        # Label in non-overlapping region
        if n == 2:
            lx = cx + (-circle_r * 0.55 if i == 0 else circle_r * 0.55) - 0.60
        elif i == 0:
            lx = cx - 0.60
        elif i == 1:
            lx = cx - circle_r * 0.7 - 0.30
        else:
            lx = cx + circle_r * 0.3

        ly = cy - 0.15 if (n == 3 and i == 0) else cy - 0.12
        if n == 3 and i == 0:
            ly = cy - circle_r * 0.65

        add_text_box(slide, label, lx, ly, 1.20, 0.30,
                     font_size=9, font_name=HEADING,
                     bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

    # Intersection labels
    for inter in intersections:
        regions = inter.get("regions", [])
        label = inter.get("label", "")
        if not label:
            continue
        # Calculate center of intersection
        ix = sum(positions[r][0] for r in regions if r < n) / len(regions)
        iy = sum(positions[r][1] for r in regions if r < n) / len(regions)
        add_text_box(slide, label, ix - 0.65, iy - 0.12, 1.30, 0.30,
                     font_size=8, font_name=BODY, bold=True,
                     color=DARK, alignment=PP_ALIGN.CENTER)


def _build_waterfall(slide, slide_def, deck_meta):
    """Build a waterfall/bridge chart showing contributions from start to end.

    Fields: headline, start ({label, value}), items (list of {label, value, type}),
            end ({label, value}), source (optional).
    Type: "positive" (green) or "negative" (red). Value should include sign.
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    GREEN = brand.green
    RED = brand.red
    GRAY = RGBColor(0x99, 0x99, 0x99)
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    DARK = brand.text_dark
    WHITE = brand.white

    start = slide_def.get("start", {})
    items = slide_def.get("items", [])
    end = slide_def.get("end", {})
    source = slide_def.get("source", "")

    all_bars = [start] + items + [end]
    n = len(all_bars)
    if n < 3:
        return

    avail_top = 1.10
    avail_bottom = 4.60
    left_margin = 0.60
    total_w = 10.0 - left_margin * 2
    bar_w = (total_w - (n - 1) * 0.08) / n
    avail_h = avail_bottom - avail_top

    # Find value range for scaling
    start_val = abs(start.get("value", 0))
    end_val = abs(end.get("value", 0))
    max_val = max(start_val, end_val)
    for item in items:
        max_val = max(max_val, abs(item.get("value", 0)))
    if max_val == 0:
        max_val = 1

    # Running total for waterfall positioning
    running = start.get("value", 0)
    baseline = avail_bottom  # bottom of chart area

    # Scale: pixels per unit value
    scale = avail_h * 0.75 / max(start_val, end_val, 1)

    for i, bar_def in enumerate(all_bars):
        label = bar_def.get("label", "")
        value = bar_def.get("value", 0)
        bar_type = bar_def.get("type", "")
        x = left_margin + i * (bar_w + 0.08)

        if i == 0:  # Start bar
            bar_h = abs(value) * scale
            bar_h = max(bar_h, 0.15)
            bar_top = baseline - bar_h
            fill = PURPLE
            running = value
        elif i == n - 1:  # End bar
            bar_h = abs(value) * scale
            bar_h = max(bar_h, 0.15)
            bar_top = baseline - bar_h
            fill = PURPLE
        else:  # Delta bars
            delta = value
            bar_h = abs(delta) * scale
            bar_h = max(bar_h, 0.10)
            if delta >= 0:
                bar_top = baseline - running * scale - bar_h
                fill = GREEN
            else:
                bar_top = baseline - running * scale
                fill = RED
            running += delta

        bar_top = max(bar_top, avail_top)
        bar_h = min(bar_h, avail_bottom - bar_top)

        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SI(x), SI(bar_top),
            SI(bar_w), SI(bar_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.fill.background()

        # Value label above/below bar
        val_text = f"{value:+,}" if i > 0 and i < n - 1 else f"{value:,}"
        val_y = bar_top - 0.22 if bar_top > avail_top + 0.25 else bar_top + bar_h + 0.02
        add_text_box(slide, str(val_text), x, val_y, bar_w, 0.20,
                     font_size=8, font_name=BODY, bold=True,
                     color=DARK, alignment=PP_ALIGN.CENTER)

        # Label below chart
        add_text_box(slide, label, x - 0.05, avail_bottom + 0.05, bar_w + 0.10, 0.30,
                     font_size=7, font_name=BODY, color=DARK,
                     alignment=PP_ALIGN.CENTER)

        # Connector line to next bar
        if i < n - 1 and i > 0:
            conn_y = bar_top if value >= 0 else bar_top + bar_h
            conn_x = x + bar_w
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, SI(conn_x), SI(conn_y),
                SI(0.08), SI(0.01))
            conn.fill.solid()
            conn.fill.fore_color.rgb = GRAY
            conn.line.fill.background()

    # Source citation
    if source:
        add_text_box(slide, source, left_margin, avail_bottom + 0.35,
                     total_w, 0.20, font_size=7, font_name=BODY,
                     color=GRAY)


def _build_pricing_table(slide, slide_def, deck_meta):
    """Build a pricing/tier comparison with 3-4 vertical cards.

    Fields: headline, tiers (list of {name, price, features (list), highlight (bool)}).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark
    GREEN = brand.green

    tiers = slide_def.get("tiers", [])
    n = len(tiers)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.50
    total_w = 10.0 - left_margin * 2
    gap = 0.15
    card_w = (total_w - (n - 1) * gap) / n
    card_h = avail_bottom - avail_top

    for i, tier in enumerate(tiers):
        name = tier.get("name", "")
        price = tier.get("price", "")
        features = tier.get("features", [])
        highlight = tier.get("highlight", False)

        x = left_margin + i * (card_w + gap)
        y = avail_top

        # Card background
        fill = PURPLE if highlight else LIGHT_BG
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, SI(x), SI(y),
            SI(card_w), SI(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.fill.background()

        text_color = WHITE if highlight else DARK
        accent = PINK if highlight else PURPLE

        # Tier name
        add_text_box(slide, name, x + 0.10, y + 0.15, card_w - 0.20, 0.30,
                     font_size=12, font_name=HEADING,
                     bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

        # Price
        add_text_box(slide, price, x + 0.10, y + 0.50, card_w - 0.20, 0.35,
                     font_size=18, font_name=HEADING,
                     bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Divider line
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SI(x + 0.15), SI(y + 0.90),
            SI(card_w - 0.30), SI(0.01))
        div.fill.solid()
        div.fill.fore_color.rgb = accent
        div.line.fill.background()

        # Features list
        feat_box = slide.shapes.add_textbox(
            SI(x + 0.15), SI(y + 1.00),
            SI(card_w - 0.30), SI(card_h - 1.15))
        tf = feat_box.text_frame; tf.word_wrap = True
        for fi, feat in enumerate(features):
            p = tf.paragraphs[0] if fi == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = "\u2713  " + str(feat)
            r.font.name = BODY
            r.font.size = Pt(9)
            r.font.color.rgb = text_color
            p.space_after = Pt(4)


def _build_concentric_circles(slide, slide_def, deck_meta):
    """Build concentric circles (TAM/SAM/SOM or layered strategy).

    Fields: headline, rings (list of {label, value} from outermost to innermost).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    WHITE = brand.white
    DARK = brand.text_dark

    rings = slide_def.get("rings", [])
    n = len(rings)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    center_x = 5.0
    center_y = avail_top + avail_h / 2
    max_r = min(avail_h * 0.45, 2.0)

    for i, ring in enumerate(rings):
        label = ring.get("label", "")
        value = ring.get("value", "")
        # Outermost ring is index 0, innermost is last
        r = max_r * (1.0 - i * 0.28)
        r = max(r, 0.40)

        # Color gradient: lightest outer, darkest inner
        t = i / max(n - 1, 1)
        fill = brand.interpolate(t, "gradient_ring_light", "gradient_start")
        text_color = WHITE if t > 0.5 else DARK

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SI(center_x - r), SI(center_y - r),
            SI(r * 2), SI(r * 2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = fill
        oval.line.fill.background()

    # Labels (rendered after shapes so they're on top)
    for i, ring in enumerate(rings):
        label = ring.get("label", "")
        value = ring.get("value", "")
        r = max_r * (1.0 - i * 0.28)
        r = max(r, 0.40)
        t = i / max(n - 1, 1)
        text_color = WHITE if t > 0.5 else DARK

        # Position label at top of each ring
        ly = center_y - r + 0.08
        text = f"{label}\n{value}" if value else label
        add_text_box(slide, text, center_x - 0.80, ly, 1.60, 0.40,
                     font_size=9 if i < 2 else 8, font_name=HEADING,
                     bold=True, color=text_color, alignment=PP_ALIGN.CENTER)


def _build_bold_bullet(slide, slide_def, deck_meta):
    """Build an executive summary with bold assertion sentences and indented evidence.

    Fields: headline, points (list of {assertion, evidence (list of strings)}).
    """
    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    DARK = brand.text_dark
    PINK = brand.secondary

    points = slide_def.get("points", [])
    if not points:
        return

    avail_top = 0.85
    left_margin = 0.50
    total_w = 10.0 - left_margin * 2
    y = avail_top + 0.10

    for point in points:
        assertion = point.get("assertion", "")
        evidence = point.get("evidence", [])

        # Bold assertion
        txb = slide.shapes.add_textbox(
            SI(left_margin), SI(y),
            SI(total_w), SI(0.30))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = assertion
        r.font.name = BODY
        r.font.size = Pt(11)
        r.font.color.rgb = PURPLE
        r.font.bold = True
        y += estimate_text_height(assertion, total_w, 11) + 0.06

        # Indented evidence bullets
        if evidence:
            ev_box = slide.shapes.add_textbox(
                SI(left_margin + 0.30), SI(y),
                SI(total_w - 0.30), SI(0.20 * len(evidence) + 0.10))
            tf2 = ev_box.text_frame; tf2.word_wrap = True
            for ei, ev in enumerate(evidence):
                p2 = tf2.paragraphs[0] if ei == 0 else tf2.add_paragraph()
                r2 = p2.add_run()
                r2.text = "\u2022  " + str(ev)
                r2.font.name = BODY
                r2.font.size = Pt(9)
                r2.font.color.rgb = DARK
            y += 0.20 * len(evidence) + 0.12

        # Subtle divider
        if point != points[-1]:
            div = slide.shapes.add_textbox(
                SI(left_margin), SI(y),
                SI(total_w), SI(0.01))
            y += 0.12


def _build_cycle_diagram(slide, slide_def, deck_meta):
    """Build a circular cycle diagram with 3-6 nodes.

    Fields: headline, nodes (list of {label}), center_label (optional).
    """
    from pptx.enum.shapes import MSO_SHAPE
    import math as _math

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark

    nodes = slide_def.get("nodes", [])
    center_label = slide_def.get("center_label", "")
    n = len(nodes)
    if n < 3:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    cx = 5.0
    cy = avail_top + avail_h / 2
    orbit_r = min(avail_h * 0.38, 1.70)
    node_r = min(0.55, orbit_r * 0.50)

    # Center circle
    if center_label:
        c_size = max(node_r * 1.5, 0.85)
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SI(cx - c_size / 2), SI(cy - c_size / 2),
            SI(c_size), SI(c_size))
        center.fill.solid()
        center.fill.fore_color.rgb = PURPLE
        center.line.fill.background()
        center_font = fit_font_size(center_label, c_size - 0.15, max_pt=10, min_pt=5)
        ctf = center.text_frame; ctf.word_wrap = False
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = center_label
        cr.font.name = HEADING
        cr.font.size = Pt(center_font)
        cr.font.bold = True
        cr.font.color.rgb = WHITE

    # Nodes around the circle
    for i, node in enumerate(nodes):
        label = node.get("label", "") if isinstance(node, dict) else str(node)
        angle = -90 + (360 / n) * i  # start at top
        rad = _math.radians(angle)
        nx = cx + orbit_r * _math.cos(rad) - node_r
        ny = cy + orbit_r * _math.sin(rad) - node_r

        # Color gradient around the circle
        t = i / n

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SI(nx), SI(ny),
            SI(node_r * 2), SI(node_r * 2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = brand.interpolate(t)
        oval.line.fill.background()

        # Label inside
        node_font = fit_font_size(label, node_r * 1.4, max_pt=10, min_pt=6)
        tf = oval.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = HEADING
        r.font.size = Pt(node_font)
        r.font.color.rgb = WHITE

        # Arrow to next node
        next_i = (i + 1) % n
        next_angle = -90 + (360 / n) * next_i
        next_rad = _math.radians(next_angle)
        # Arrow midpoint between nodes
        mid_angle = (angle + next_angle) / 2
        if next_angle < angle:
            mid_angle = (angle + next_angle + 360) / 2
        mid_rad = _math.radians(mid_angle)
        ax = cx + (orbit_r + node_r * 0.3) * _math.cos(mid_rad)
        ay = cy + (orbit_r + node_r * 0.3) * _math.sin(mid_rad)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            SI(ax - 0.10), SI(ay - 0.06),
            SI(0.20), SI(0.12))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        arrow.line.fill.background()
        # Rotate arrow to point in direction of travel
        arrow.rotation = mid_angle + 90


def _build_hub_spoke(slide, slide_def, deck_meta):
    """Build a hub-and-spoke diagram with a central circle and radiating nodes.

    Fields: headline, hub (string — center label), spokes (list of {label, body}).
    """
    from pptx.enum.shapes import MSO_SHAPE
    import math as _math

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark
    GRAY = RGBColor(0xCC, 0xCC, 0xCC)

    hub_label = slide_def.get("hub", "")
    spokes = slide_def.get("spokes", [])
    n = len(spokes)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    cx = 5.0
    cy = avail_top + avail_h / 2
    # Scale hub and spoke sizes based on count
    hub_r = 0.65 if n <= 6 else 0.55
    orbit_r = min(avail_h * 0.40, 1.80)
    spoke_r = 0.48 if n <= 6 else 0.38

    # Hub
    hub = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        SI(cx - hub_r), SI(cy - hub_r),
        SI(hub_r * 2), SI(hub_r * 2))
    hub.fill.solid()
    hub.fill.fore_color.rgb = PURPLE
    hub.line.fill.background()
    hub_font = fit_font_size(hub_label, hub_r * 1.4, max_pt=11, min_pt=7)
    tf = hub.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = hub_label
    r.font.name = HEADING
    r.font.size = Pt(hub_font)
    r.font.color.rgb = WHITE

    # Spokes
    for i, spoke in enumerate(spokes):
        label = spoke.get("label", "") if isinstance(spoke, dict) else str(spoke)
        body = spoke.get("body", "") if isinstance(spoke, dict) else ""
        angle = -90 + (360 / n) * i
        rad = _math.radians(angle)
        sx = cx + orbit_r * _math.cos(rad)
        sy = cy + orbit_r * _math.sin(rad)

        # Connector line from hub edge to spoke (freeform line)
        line_start_x = cx + hub_r * _math.cos(rad)
        line_start_y = cy + hub_r * _math.sin(rad)
        line_end_x = sx - spoke_r * _math.cos(rad)
        line_end_y = sy - spoke_r * _math.sin(rad)
        freeform = slide.shapes.build_freeform(
            SI(line_start_x), SI(line_start_y))
        freeform.add_line_segments([
            (SI(line_end_x), SI(line_end_y)),
        ])
        connector = freeform.convert_to_shape()
        connector.line.color.rgb = GRAY
        connector.line.width = Pt(1.5)
        connector.fill.background()

        # Spoke circle
        t = i / max(n - 1, 1)
        spoke_fill = brand.interpolate(t, "background_light", "spoke_end")

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SI(sx - spoke_r), SI(sy - spoke_r),
            SI(spoke_r * 2), SI(spoke_r * 2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = spoke_fill
        oval.line.fill.background()

        spoke_font = fit_font_size(label, spoke_r * 1.4, max_pt=9, min_pt=5)
        tf = oval.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = HEADING
        r.font.size = Pt(spoke_font)
        r.font.color.rgb = PURPLE


def _build_gauge_dashboard(slide, slide_def, deck_meta):
    """Build a gauge/speedometer dashboard with 2-4 dials.

    Fields: headline, gauges (list of {value (0-100), label, ranges (optional)}).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    DARK = brand.text_dark
    PURPLE = brand.primary
    GREEN = brand.green
    AMBER = brand.amber
    RED = brand.red
    GRAY_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
    DARK = DARK

    gauges = slide_def.get("gauges", [])
    n = len(gauges)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.50
    total_w = 10.0 - left_margin * 2
    avail_h = avail_bottom - avail_top

    gauge_size = min((total_w - (n - 1) * 0.20) / n, avail_h - 0.60, 2.0)
    gap = (total_w - n * gauge_size) / max(n + 1, 1)
    start_y = avail_top + (avail_h - gauge_size - 0.40) / 3

    for i, gauge in enumerate(gauges):
        value = gauge.get("value", 0)
        label = gauge.get("label", "")
        gx = left_margin + gap + i * (gauge_size + gap)
        gy = start_y

        # Determine color by value
        color = GREEN if value >= 75 else AMBER if value >= 40 else RED

        # Background semicircle (gray block arc for top half)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.BLOCK_ARC, SI(gx), SI(gy),
            SI(gauge_size), SI(gauge_size))
        bg.fill.solid()
        bg.fill.fore_color.rgb = GRAY_LIGHT
        bg.line.fill.background()
        # BLOCK_ARC: adj1=start angle (fraction), adj2=thickness
        # 0.5 = 180 degrees (left/9 o'clock), 0.0 = 0 degrees (right/3 o'clock)
        bg.adjustments[0] = 0.5   # start at 180 degrees (left)
        bg.adjustments[1] = 0.25  # ring thickness

        # Foreground colored arc
        if value > 0:
            fg = slide.shapes.add_shape(
                MSO_SHAPE.BLOCK_ARC, SI(gx), SI(gy),
                SI(gauge_size), SI(gauge_size))
            fg.fill.solid()
            fg.fill.fore_color.rgb = color
            fg.line.fill.background()
            # Sweep from left (180°) clockwise by value% of 180°
            sweep_frac = value / 100.0 * 0.5  # 0.5 = 180 degrees
            fg.adjustments[0] = 0.5  # start at left
            fg.adjustments[1] = 0.25  # thickness

        # Value text — centered in the gauge
        val_font = fit_font_size(f"{value}%", gauge_size * 0.50,
                                  max_pt=min(20, int(gauge_size * 14)), min_pt=10)
        add_text_box(slide, f"{value}%",
                     gx + gauge_size * 0.15, gy + gauge_size * 0.30,
                     gauge_size * 0.70, gauge_size * 0.30,
                     font_size=val_font,
                     font_name=HEADING, bold=True,
                     color=color, alignment=PP_ALIGN.CENTER)

        # Label — below the gauge (gauge is top-half semicircle)
        add_text_box(slide, label,
                     gx - 0.10, gy + gauge_size + 0.08,
                     gauge_size + 0.20, 0.30,
                     font_size=9, font_name=BODY, color=DARK,
                     alignment=PP_ALIGN.CENTER)


def _build_risk_heat_map(slide, slide_def, deck_meta):
    """Build a 5x5 risk heat map with risk detail panel.

    Fields: headline, x_axis (default "Impact"), y_axis (default "Likelihood"),
            items (list of {label, x (1-5), y (1-5), description (optional)}).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    DARK = brand.text_dark
    WHITE = WHITE
    LIGHT_BG = brand.bg_light

    HEAT_COLORS = [
        ["4CAF50", "8BC34A", "FFC107", "FF9800", "F44336"],
        ["8BC34A", "FFC107", "FF9800", "F44336", "D32F2F"],
        ["CDDC39", "FFC107", "FF9800", "F44336", "D32F2F"],
        ["4CAF50", "8BC34A", "FFC107", "FF9800", "F44336"],
        ["4CAF50", "4CAF50", "8BC34A", "FFC107", "FF9800"],
    ]

    items = slide_def.get("items", [])
    x_axis = slide_def.get("x_axis", "Impact")
    y_axis = slide_def.get("y_axis", "Likelihood")

    avail_top = 0.85
    avail_bottom = 4.90
    avail_h = avail_bottom - avail_top

    # Layout: grid on left ~55%, detail panel on right ~40%
    y_label_w = 0.55
    grid_left = y_label_w + 0.10
    grid_size = min(avail_h - 0.30, 3.60)  # keep grid square
    cell_size = grid_size / 5
    panel_left = grid_left + grid_size + 0.50
    panel_w = 10.0 - panel_left - 0.35

    grid_top = avail_top + (avail_h - grid_size - 0.30) / 2

    # Draw 5x5 grid
    for row in range(5):
        for col in range(5):
            x = grid_left + col * cell_size
            y = grid_top + row * cell_size
            color_hex = HEAT_COLORS[row][col]
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, SI(x), SI(y),
                SI(cell_size), SI(cell_size))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            cell.line.color.rgb = WHITE
            cell.line.width = Pt(1)

    # Plot items on grid
    for item in items:
        label = item.get("label", "")
        ix = item.get("x", 1)
        iy = item.get("y", 1)
        # Place a small dot marker
        dot_r = 0.12
        dot_x = grid_left + (ix - 0.5) * cell_size - dot_r
        dot_y = grid_top + (5 - iy + 0.5) * cell_size - dot_r
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, SI(dot_x), SI(dot_y),
            SI(dot_r * 2), SI(dot_r * 2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = DARK
        dot.line.width = Pt(1.5)
        # Label next to dot
        add_text_box(slide, label,
                     dot_x - 0.10, dot_y - 0.18,
                     cell_size * 0.9, 0.16,
                     font_size=6, font_name=BODY, bold=True,
                     color=WHITE, alignment=PP_ALIGN.CENTER)

    # Axis labels
    grid_bottom = grid_top + grid_size
    add_text_box(slide, x_axis, grid_left, grid_bottom + 0.05, grid_size, 0.20,
                 font_size=9, font_name=BODY, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)
    for i in range(5):
        add_text_box(slide, str(5 - i), grid_left - 0.25,
                     grid_top + i * cell_size + cell_size / 2 - 0.08, 0.20, 0.16,
                     font_size=7, font_name=BODY, color=DARK,
                     alignment=PP_ALIGN.CENTER)
    for i in range(5):
        add_text_box(slide, str(i + 1),
                     grid_left + i * cell_size + cell_size / 2 - 0.08,
                     grid_bottom + 0.01, 0.16, 0.16,
                     font_size=7, font_name=BODY, color=DARK,
                     alignment=PP_ALIGN.CENTER)
    add_text_box(slide, y_axis, 0.05, grid_top + grid_size / 2 - 0.10, y_label_w, 0.20,
                 font_size=9, font_name=BODY, bold=True, color=DARK)

    # --- Risk detail panel on the right ---
    if items and panel_w > 1.5:
        # Panel title
        add_text_box(slide, "Risk Register", panel_left, grid_top, panel_w, 0.30,
                     font_size=11, font_name=HEADING,
                     bold=True, color=PURPLE)

        # Risk items list
        item_y = grid_top + 0.35
        item_h = 0.32
        for item in items:
            label = item.get("label", "")
            ix = item.get("x", 1)
            iy = item.get("y", 1)
            desc = item.get("description", "")
            risk_score = ix * iy

            # Risk level color
            if risk_score >= 15:
                level_color = RGBColor(0xD3, 0x2F, 0x2F)  # red
                level = "Critical"
            elif risk_score >= 8:
                level_color = RGBColor(0xFF, 0x98, 0x00)  # orange
                level = "High"
            elif risk_score >= 4:
                level_color = RGBColor(0xFF, 0xC1, 0x07)  # yellow
                level = "Medium"
            else:
                level_color = RGBColor(0x4C, 0xAF, 0x50)  # green
                level = "Low"

            # Color dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                SI(panel_left), SI(item_y + 0.04),
                SI(0.10), SI(0.10))
            dot.fill.solid()
            dot.fill.fore_color.rgb = level_color
            dot.line.fill.background()

            # Label + level
            text = f"{label} ({level})"
            if desc:
                text += f"\n{desc}"
            txb = slide.shapes.add_textbox(
                SI(panel_left + 0.18), SI(item_y),
                SI(panel_w - 0.20), SI(item_h))
            tf = txb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = label
            r.font.name = BODY
            r.font.size = Pt(9)
            r.font.color.rgb = DARK
            r.font.bold = True
            r2 = p.add_run()
            r2.text = f"  {level}"
            r2.font.name = BODY
            r2.font.size = Pt(8)
            r2.font.color.rgb = level_color
            r2.font.bold = True
            if desc:
                p2 = tf.add_paragraph()
                r3 = p2.add_run()
                r3.text = desc
                r3.font.name = BODY
                r3.font.size = Pt(7)
                r3.font.color.rgb = DARK

            item_y += item_h + 0.05


def _build_tornado_chart(slide, slide_def, deck_meta):
    """Build a tornado/sensitivity chart with bars extending left and right.

    Fields: headline, items (list of {label, low, high}), center_label (optional).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    PINK = brand.secondary
    DARK = brand.text_dark
    WHITE = WHITE

    items = slide_def.get("items", [])
    center_label = slide_def.get("center_label", "")
    n = len(items)
    if n == 0:
        return

    avail_top = 0.85
    avail_bottom = 4.80
    label_w = 1.80
    left_margin = 0.35 + label_w
    right_margin = 0.35
    bar_area_w = (10.0 - left_margin - right_margin) / 2  # each side
    center_x = left_margin + bar_area_w
    bar_h = min((avail_bottom - avail_top) / n - 0.06, 0.40)
    gap = 0.06

    # Find max magnitude for scaling
    max_mag = max(max(abs(it.get("low", 0)), abs(it.get("high", 0))) for it in items)
    if max_mag == 0:
        max_mag = 1

    # Sort by magnitude (largest swing first)
    sorted_items = sorted(items, key=lambda x: abs(x.get("high", 0)) + abs(x.get("low", 0)), reverse=True)

    total_h = n * (bar_h + gap) - gap
    start_y = avail_top + ((avail_bottom - avail_top) - total_h) / 2

    # Center axis line
    axis = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SI(center_x - 0.005), SI(start_y - 0.05),
        SI(0.01), SI(total_h + 0.10))
    axis.fill.solid()
    axis.fill.fore_color.rgb = DARK
    axis.line.fill.background()

    for i, item in enumerate(sorted_items):
        label = item.get("label", "")
        low = item.get("low", 0)
        high = item.get("high", 0)
        y = start_y + i * (bar_h + gap)

        # Low bar (extends left)
        low_w = abs(low) / max_mag * bar_area_w
        if low_w > 0.02:
            bar_l = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                SI(center_x - low_w), SI(y),
                SI(low_w), SI(bar_h))
            bar_l.fill.solid()
            bar_l.fill.fore_color.rgb = PURPLE
            bar_l.line.fill.background()

        # High bar (extends right)
        high_w = abs(high) / max_mag * bar_area_w
        if high_w > 0.02:
            bar_r = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                SI(center_x), SI(y),
                SI(high_w), SI(bar_h))
            bar_r.fill.solid()
            bar_r.fill.fore_color.rgb = PINK
            bar_r.line.fill.background()

        # Label to the left
        add_text_box(slide, label, 0.35, y + bar_h / 2 - 0.10, label_w - 0.10, 0.20,
                     font_size=8, font_name=BODY, bold=True,
                     color=DARK, alignment=PP_ALIGN.RIGHT)

    if center_label:
        add_text_box(slide, center_label, center_x - 0.50, avail_bottom + 0.05, 1.0, 0.20,
                     font_size=8, font_name=BODY, color=DARK,
                     alignment=PP_ALIGN.CENTER)


def _build_radar_chart(slide, slide_def, deck_meta):
    """Build a radar/spider chart for multi-dimensional comparison.

    Fields: headline, axes (list of label strings),
            series (list of {name, values (list of 0-100), color}).
    """
    from pptx.enum.shapes import MSO_SHAPE
    import math as _math

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    DARK = brand.text_dark
    GRAY = RGBColor(0xDD, 0xDD, 0xDD)

    axes = slide_def.get("axes", [])
    series = slide_def.get("series", [])
    n = len(axes)
    if n < 3:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    cx = 5.0
    cy = avail_top + avail_h / 2
    max_r = min(avail_h * 0.40, 1.80)

    SERIES_COLORS = [PURPLE, PINK, brand.green, brand.amber]

    # Draw grid rings as freeform polygons + axis spoke lines
    for ring_pct in [0.33, 0.66, 1.0]:
        ring_r = max_r * ring_pct
        # Build polygon ring
        pts = []
        for j in range(n):
            angle = -90 + (360 / n) * j
            rad = _math.radians(angle)
            pts.append((SI(cx + ring_r * _math.cos(rad)),
                         SI(cy + ring_r * _math.sin(rad))))
        if pts:
            ff = slide.shapes.build_freeform(pts[0][0], pts[0][1])
            ff.add_line_segments(pts[1:] + [pts[0]])  # close the polygon
            ring_shape = ff.convert_to_shape()
            ring_shape.line.color.rgb = GRAY
            ring_shape.line.width = Pt(0.5)
            ring_shape.fill.background()

    # Axis spoke lines from center to each vertex
    for j in range(n):
        angle = -90 + (360 / n) * j
        rad = _math.radians(angle)
        ex = cx + max_r * _math.cos(rad)
        ey = cy + max_r * _math.sin(rad)
        ff = slide.shapes.build_freeform(SI(cx), SI(cy))
        ff.add_line_segments([(SI(ex), SI(ey))])
        spoke = ff.convert_to_shape()
        spoke.line.color.rgb = GRAY
        spoke.line.width = Pt(0.5)
        spoke.fill.background()

    # Axis labels
    for j in range(n):
        angle = -90 + (360 / n) * j
        rad = _math.radians(angle)
        lx = cx + (max_r + 0.25) * _math.cos(rad) - 0.50
        ly = cy + (max_r + 0.25) * _math.sin(rad) - 0.10
        add_text_box(slide, axes[j], lx, ly, 1.0, 0.22,
                     font_size=8, font_name=BODY, bold=True,
                     color=DARK, alignment=PP_ALIGN.CENTER)

    # Data series — plot as dots at axis positions
    for si, s in enumerate(series):
        values = s.get("values", [])
        name = s.get("name", "")
        color = SERIES_COLORS[si % len(SERIES_COLORS)]

        for j, val in enumerate(values):
            if j >= n:
                break
            angle = -90 + (360 / n) * j
            rad = _math.radians(angle)
            pr = max_r * (val / 100.0)
            px = cx + pr * _math.cos(rad) - 0.06
            py = cy + pr * _math.sin(rad) - 0.06
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, SI(px), SI(py),
                SI(0.12), SI(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

    # Legend
    if series:
        for si, s in enumerate(series):
            color = SERIES_COLORS[si % len(SERIES_COLORS)]
            ly = avail_bottom - 0.20 * (len(series) - si)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, SI(8.5), SI(ly),
                SI(0.10), SI(0.10))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            add_text_box(slide, s.get("name", ""), 8.65, ly - 0.02, 1.0, 0.15,
                         font_size=8, font_name=BODY, color=DARK)


def _build_combo_chart(slide, slide_def, deck_meta):
    """Build a combo chart with bars and a line overlay.

    Fields: headline, categories (list of x-axis labels),
            bars ({name, values, color}), line ({name, values, color}),
            y_axis_left (label), y_axis_right (label).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    DARK = brand.text_dark
    GRAY = RGBColor(0xDD, 0xDD, 0xDD)

    categories = slide_def.get("categories", [])
    bars_def = slide_def.get("bars", {})
    line_def = slide_def.get("line", {})
    n = len(categories)
    if n == 0:
        return

    avail_top = 1.00
    avail_bottom = 4.50
    left_margin = 1.00
    right_margin = 1.00
    chart_w = 10.0 - left_margin - right_margin
    chart_h = avail_bottom - avail_top

    bar_values = bars_def.get("values", [])
    line_values = line_def.get("values", [])
    bar_color_hex = bars_def.get("color", brand.primary_hex)
    line_color_hex = line_def.get("color", brand.secondary_hex)

    bar_color = RGBColor.from_string(bar_color_hex.replace("#", ""))
    line_color = RGBColor.from_string(line_color_hex.replace("#", ""))

    max_bar = max(bar_values) if bar_values else 1
    max_line = max(line_values) if line_values else 1
    if max_bar == 0: max_bar = 1
    if max_line == 0: max_line = 1

    bar_w = chart_w / n * 0.6
    bar_gap = chart_w / n

    # Bars
    for i, val in enumerate(bar_values):
        if i >= n:
            break
        h = (val / max_bar) * chart_h * 0.85
        x = left_margin + i * bar_gap + (bar_gap - bar_w) / 2
        y = avail_bottom - h
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SI(x), SI(y),
            SI(bar_w), SI(h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bar_color
        rect.line.fill.background()

    # Line connecting dots
    line_pts = []
    for i, val in enumerate(line_values):
        if i >= n:
            break
        lx = left_margin + i * bar_gap + bar_gap / 2
        ly = avail_bottom - (val / max_line) * chart_h * 0.85
        line_pts.append((SI(lx), SI(ly)))

    if len(line_pts) >= 2:
        ff = slide.shapes.build_freeform(line_pts[0][0], line_pts[0][1])
        ff.add_line_segments(line_pts[1:])
        line_shape = ff.convert_to_shape()
        line_shape.line.color.rgb = line_color
        line_shape.line.width = Pt(2.0)
        line_shape.fill.background()

    # Line dots
    for i, val in enumerate(line_values):
        if i >= n:
            break
        lx = left_margin + i * bar_gap + bar_gap / 2
        ly = avail_bottom - (val / max_line) * chart_h * 0.85
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SI(lx - 0.06), SI(ly - 0.06),
            SI(0.12), SI(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = line_color
        dot.line.fill.background()

    # Category labels
    for i, cat in enumerate(categories):
        x = left_margin + i * bar_gap
        add_text_box(slide, cat, x, avail_bottom + 0.05, bar_gap, 0.25,
                     font_size=7, font_name=BODY, color=DARK,
                     alignment=PP_ALIGN.CENTER)

    # Axis labels
    y_left = slide_def.get("y_axis_left", bars_def.get("name", ""))
    y_right = slide_def.get("y_axis_right", line_def.get("name", ""))
    if y_left:
        add_text_box(slide, y_left, 0.10, avail_top + chart_h / 2 - 0.10, 0.85, 0.25,
                     font_size=8, font_name=BODY, bold=True, color=bar_color)
    if y_right:
        add_text_box(slide, y_right, 9.05, avail_top + chart_h / 2 - 0.10, 0.85, 0.25,
                     font_size=8, font_name=BODY, bold=True, color=line_color)


def _build_bubble_chart(slide, slide_def, deck_meta):
    """Build a bubble chart (scatter with variable-size circles).

    Fields: headline, x_axis (label), y_axis (label),
            bubbles (list of {label, x, y, size (relative 1-100)}).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    WHITE = brand.white
    PURPLE = brand.primary
    PINK = brand.secondary
    DARK = brand.text_dark
    WHITE = WHITE

    bubbles = slide_def.get("bubbles", [])
    x_label = slide_def.get("x_axis", "")
    y_label = slide_def.get("y_axis", "")
    if not bubbles:
        return

    avail_top = 0.85
    avail_bottom = 4.70
    left_margin = 1.10  # room for y-axis label
    right_margin = 0.50
    chart_w = 10.0 - left_margin - right_margin
    chart_h = avail_bottom - avail_top

    max_size = max(b.get("size", 50) for b in bubbles)
    max_x = max(b.get("x", 0) for b in bubbles)
    max_y = max(b.get("y", 0) for b in bubbles)
    min_x = min(b.get("x", 0) for b in bubbles)
    min_y = min(b.get("y", 0) for b in bubbles)
    x_range = max(max_x - min_x, 1)
    y_range = max(max_y - min_y, 1)

    for i, b in enumerate(bubbles):
        bx = b.get("x", 0)
        by = b.get("y", 0)
        bsize = b.get("size", 50)
        label = b.get("label", "")

        # Position on chart
        px = left_margin + ((bx - min_x) / x_range) * chart_w * 0.85 + chart_w * 0.05
        py = avail_bottom - ((by - min_y) / y_range) * chart_h * 0.85 - chart_h * 0.05
        # Size — larger bubbles for readability
        r = 0.25 + (bsize / max_size) * 0.55

        # Color cycle
        t = i / max(len(bubbles) - 1, 1)

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SI(px - r), SI(py - r),
            SI(r * 2), SI(r * 2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = brand.interpolate(t)
        oval.line.fill.background()
        # Set 60% opacity via OOXML
        try:
            from pptx.oxml.ns import qn as _qn
            for sf in oval._element.iter(_qn('a:solidFill')):
                clr = sf.find(_qn('a:srgbClr'))
                if clr is not None:
                    alpha = clr.makeelement(_qn('a:alpha'), {'val': '60000'})
                    clr.append(alpha)
                    break
        except Exception:
            pass

        if label:
            # Label inside the bubble — text box overlay for proper word wrapping
            inner_w = r * 1.4  # usable text width inside circle
            lbl_font = fit_font_size(label, inner_w, max_pt=10, min_pt=6)
            lbl_box = slide.shapes.add_textbox(
                SI(px - inner_w / 2), SI(py - r * 0.35),
                SI(inner_w), SI(r * 0.80))
            ltf = lbl_box.text_frame; ltf.word_wrap = True
            lp = ltf.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
            lr = lp.add_run()
            lr.text = label
            lr.font.name = BODY
            lr.font.size = Pt(lbl_font)
            lr.font.bold = True
            lr.font.color.rgb = WHITE

    # Axis labels
    if x_label:
        add_text_box(slide, x_label, left_margin, avail_bottom + 0.08, chart_w, 0.25,
                     font_size=9, font_name=BODY, bold=True,
                     color=DARK, alignment=PP_ALIGN.CENTER)
    if y_label:
        yl_font = fit_font_size(y_label, 1.0, max_pt=9, min_pt=6)
        add_text_box(slide, y_label, 0.05, avail_top + chart_h / 2 - 0.10, 1.00, 0.40,
                     font_size=yl_font, font_name=BODY, bold=True, color=DARK)


def _build_bento_grid(slide, slide_def, deck_meta):
    """Build a bento grid with mixed-size tiles.

    Fields: headline, tiles (list of {title, body, size ("large"|"medium"|"small"), icon}).
    First "large" tile gets hero treatment. Others fill remaining space.
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark

    tiles = slide_def.get("tiles", [])
    if not tiles:
        return

    avail_top = 0.85
    avail_bottom = 5.0
    left_margin = 0.35
    total_w = 10.0 - left_margin * 2
    avail_h = avail_bottom - avail_top
    gap = 0.10

    # Separate hero tile from rest
    hero = None
    smalls = []
    for t in tiles:
        if t.get("size", "small") == "large" and hero is None:
            hero = t
        else:
            smalls.append(t)

    if hero:
        # Hero takes left 50%, smalls stack on right
        hero_w = total_w * 0.48
        small_w = total_w - hero_w - gap
        small_x = left_margin + hero_w + gap

        # Hero tile
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(left_margin), SI(avail_top),
            SI(hero_w), SI(avail_h))
        card.fill.solid()
        card.fill.fore_color.rgb = PURPLE
        card.line.fill.background()
        add_text_box(slide, hero.get("title", ""),
                     left_margin + 0.20, avail_top + 0.25,
                     hero_w - 0.40, 0.35,
                     font_size=16, font_name=HEADING,
                     bold=True, color=WHITE)
        add_text_box(slide, hero.get("body", ""),
                     left_margin + 0.20, avail_top + 0.65,
                     hero_w - 0.40, avail_h - 0.90,
                     font_size=10, font_name=BODY, color=WHITE)

        # Small tiles stacked on right
        n_small = len(smalls)
        if n_small > 0:
            small_h = (avail_h - (n_small - 1) * gap) / n_small
            for i, st in enumerate(smalls):
                sy = avail_top + i * (small_h + gap)
                sc = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    SI(small_x), SI(sy),
                    SI(small_w), SI(small_h))
                sc.fill.solid()
                sc.fill.fore_color.rgb = LIGHT_BG
                sc.line.fill.background()
                add_text_box(slide, st.get("title", ""),
                             small_x + 0.12, sy + 0.08,
                             small_w - 0.24, 0.25,
                             font_size=10, font_name=HEADING,
                             bold=True, color=PURPLE)
                add_text_box(slide, st.get("body", ""),
                             small_x + 0.12, sy + 0.33,
                             small_w - 0.24, small_h - 0.45,
                             font_size=8, font_name=BODY, color=DARK)
    else:
        # All tiles in a grid
        cols = min(3, len(smalls))
        rows_count = -(-len(smalls) // cols)
        tile_w = (total_w - (cols - 1) * gap) / cols
        tile_h = (avail_h - (rows_count - 1) * gap) / rows_count
        for i, st in enumerate(smalls):
            col = i % cols
            row = i // cols
            tx = left_margin + col * (tile_w + gap)
            ty = avail_top + row * (tile_h + gap)
            sc = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                SI(tx), SI(ty), SI(tile_w), SI(tile_h))
            sc.fill.solid()
            sc.fill.fore_color.rgb = LIGHT_BG
            sc.line.fill.background()
            add_text_box(slide, st.get("title", ""),
                         tx + 0.10, ty + 0.08, tile_w - 0.20, 0.25,
                         font_size=10, font_name=HEADING,
                         bold=True, color=PURPLE)
            add_text_box(slide, st.get("body", ""),
                         tx + 0.10, ty + 0.35, tile_w - 0.20, tile_h - 0.45,
                         font_size=8, font_name=BODY, color=DARK)


def _build_dashboard_panel(slide, slide_def, deck_meta):
    """Build a dashboard panel with KPI tiles + chart area + summary.

    Fields: headline, kpis (list of {number, label}),
            chart_title, chart_data (list of {label, value}),
            summary (text).
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark

    kpis = slide_def.get("kpis", [])
    chart_title = slide_def.get("chart_title", "")
    chart_data = slide_def.get("chart_data", [])
    summary = slide_def.get("summary", "")

    avail_top = 0.85
    left_margin = 0.35
    total_w = 10.0 - left_margin * 2

    # KPI tiles across top
    kpi_h = 0.65
    n_kpi = len(kpis)
    if n_kpi > 0:
        kpi_w = (total_w - (n_kpi - 1) * 0.10) / n_kpi
        for i, kpi in enumerate(kpis):
            kx = left_margin + i * (kpi_w + 0.10)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                SI(kx), SI(avail_top),
                SI(kpi_w), SI(kpi_h))
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_BG
            card.line.fill.background()
            add_text_box(slide, str(kpi.get("number", "")),
                         kx + 0.08, avail_top + 0.05, kpi_w - 0.16, 0.35,
                         font_size=18, font_name=HEADING,
                         bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, kpi.get("label", ""),
                         kx + 0.08, avail_top + 0.40, kpi_w - 0.16, 0.20,
                         font_size=8, font_name=BODY, color=DARK,
                         alignment=PP_ALIGN.CENTER)

    # Chart area (left 65%) + Summary (right 35%)
    chart_top = avail_top + kpi_h + 0.15
    chart_bottom = 5.0
    chart_area_h = chart_bottom - chart_top
    chart_area_w = total_w * 0.62
    summary_w = total_w - chart_area_w - 0.10

    # Simple bar chart
    if chart_data:
        # Chart background
        cb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(left_margin), SI(chart_top),
            SI(chart_area_w), SI(chart_area_h))
        cb.fill.solid()
        cb.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        cb.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        cb.line.width = Pt(0.5)

        if chart_title:
            add_text_box(slide, chart_title,
                         left_margin + 0.10, chart_top + 0.05,
                         chart_area_w - 0.20, 0.25,
                         font_size=9, font_name=HEADING,
                         bold=True, color=PURPLE)

        n_bars = len(chart_data)
        max_val = max(d.get("value", 0) for d in chart_data) or 1
        bar_area_top = chart_top + 0.35
        bar_area_h = chart_area_h - 0.60
        bar_w = (chart_area_w - 0.40) / n_bars * 0.6
        bar_gap = (chart_area_w - 0.40) / n_bars

        for j, d in enumerate(chart_data):
            val = d.get("value", 0)
            h = (val / max_val) * bar_area_h * 0.85
            bx = left_margin + 0.20 + j * bar_gap + (bar_gap - bar_w) / 2
            by = bar_area_top + bar_area_h - h
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, SI(bx), SI(by),
                SI(bar_w), SI(h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = PURPLE
            rect.line.fill.background()
            # Label
            add_text_box(slide, d.get("label", ""),
                         bx - 0.10, bar_area_top + bar_area_h + 0.02,
                         bar_w + 0.20, 0.20,
                         font_size=7, font_name=BODY, color=DARK,
                         alignment=PP_ALIGN.CENTER)

    # Summary panel
    if summary:
        sx = left_margin + chart_area_w + 0.10
        sb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(sx), SI(chart_top),
            SI(summary_w), SI(chart_area_h))
        sb.fill.solid()
        sb.fill.fore_color.rgb = LIGHT_BG
        sb.line.fill.background()
        add_text_box(slide, "Summary",
                     sx + 0.10, chart_top + 0.08, summary_w - 0.20, 0.25,
                     font_size=10, font_name=HEADING,
                     bold=True, color=PURPLE)
        add_text_box(slide, summary,
                     sx + 0.10, chart_top + 0.38, summary_w - 0.20, chart_area_h - 0.50,
                     font_size=9, font_name=BODY, color=DARK)


def _build_left_nav_sidebar(slide, slide_def, deck_meta):
    """Build a slide with a left navigation sidebar and main content area.

    Fields: headline, nav_items (list of {label, active (bool)}),
            content_title, content_body.
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    LIGHT_BG = brand.bg_light
    WHITE = brand.white
    DARK = brand.text_dark

    nav_items = slide_def.get("nav_items", [])
    content_title = slide_def.get("content_title", "")
    content_body = slide_def.get("content_body", "")

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    sidebar_w = 2.20
    main_x = sidebar_w + 0.20
    main_w = 10.0 - main_x - 0.35

    # Sidebar background
    sb = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SI(0), SI(avail_top),
        SI(sidebar_w), SI(avail_h))
    sb.fill.solid()
    sb.fill.fore_color.rgb = PURPLE
    sb.line.fill.background()

    # Nav items
    item_h = 0.40
    for i, nav in enumerate(nav_items):
        label = nav.get("label", "") if isinstance(nav, dict) else str(nav)
        active = nav.get("active", False) if isinstance(nav, dict) else False
        ny = avail_top + 0.15 + i * item_h

        if active:
            # Active indicator
            indicator = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                SI(0), SI(ny),
                SI(0.05), SI(item_h - 0.05))
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = PINK
            indicator.line.fill.background()

        add_text_box(slide, label, 0.15, ny + 0.05, sidebar_w - 0.30, 0.25,
                     font_size=10 if active else 9,
                     font_name=HEADING if active else BODY,
                     bold=active, color=WHITE)

    # Main content area
    if content_title:
        add_text_box(slide, content_title, main_x, avail_top + 0.10,
                     main_w, 0.35, font_size=14,
                     font_name=HEADING, bold=True, color=PURPLE)
    if content_body:
        body_top = avail_top + 0.55 if content_title else avail_top + 0.10
        txb = slide.shapes.add_textbox(
            SI(main_x), SI(body_top),
            SI(main_w), SI(avail_bottom - body_top))
        tf = txb.text_frame; tf.word_wrap = True
        _render_body_text(tf, content_body, BODY, 10, DARK,
                          available_height=avail_bottom - body_top)


def _build_image_text_hero(slide, slide_def, deck_meta):
    """Build a hero image slide with text overlay strip.

    Fields: headline, image (path), body, text_position ("left"|"right"|"bottom").
    """
    from pptx.enum.shapes import MSO_SHAPE

    brand = deck_meta.get("brand", BrandConfig())
    HEADING = brand.heading_font
    BODY = brand.body_font
    PURPLE = brand.primary
    PINK = brand.secondary
    WHITE = brand.white
    DARK = brand.text_dark

    image = slide_def.get("image", "")
    body = slide_def.get("body", "")
    text_pos = slide_def.get("text_position", "right")

    avail_top = 0.85
    avail_bottom = 5.0
    avail_h = avail_bottom - avail_top
    left_margin = 0.35

    if text_pos == "left":
        text_x = left_margin
        text_w = 3.50
        img_x = text_x + text_w + 0.15
        img_w = 10.0 - img_x - 0.35
    elif text_pos == "bottom":
        img_x = left_margin
        img_w = 10.0 - left_margin * 2
        img_h = avail_h * 0.60
        text_x = left_margin
        text_w = img_w
    else:  # right (default)
        img_x = left_margin
        img_w = 5.50
        text_x = img_x + img_w + 0.15
        text_w = 10.0 - text_x - 0.35

    # Image placeholder (purple rectangle if no image)
    if text_pos == "bottom":
        img_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(img_x), SI(avail_top),
            SI(img_w), SI(img_h))
    else:
        img_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SI(img_x), SI(avail_top),
            SI(img_w), SI(avail_h))

    # Try to load actual image
    if image:
        base_dir = deck_meta.get("base_dir", "")
        img_path = os.path.join(base_dir, image) if not os.path.isabs(image) else image
        if os.path.exists(img_path):
            try:
                if text_pos == "bottom":
                    slide.shapes.add_picture(img_path,
                        SI(img_x), SI(avail_top),
                        SI(img_w), SI(img_h))
                else:
                    slide.shapes.add_picture(img_path,
                        SI(img_x), SI(avail_top),
                        SI(img_w), SI(avail_h))
            except Exception:
                pass
        else:
            img_shape.fill.solid()
            img_shape.fill.fore_color.rgb = brand.bg_light
    else:
        img_shape.fill.solid()
        img_shape.fill.fore_color.rgb = brand.bg_light
    img_shape.line.fill.background()

    # Text area
    if text_pos == "bottom":
        text_y = avail_top + img_h + 0.10
        text_h = avail_bottom - text_y
    else:
        text_y = avail_top + 0.20
        text_h = avail_h - 0.40

    if body:
        txb = slide.shapes.add_textbox(
            SI(text_x), SI(text_y),
            SI(text_w), SI(text_h))
        tf = txb.text_frame; tf.word_wrap = True
        _render_body_text(tf, body, BODY, 11, DARK, available_height=text_h)


def _insert_image(slide, image_path, placeholders, layout_info):
    """Insert an image into the appropriate location on a slide."""
    # Try to use the 'image' placeholder if the layout has one
    if "image" in placeholders:
        ph_idx = placeholders["image"]["idx"]
        if set_placeholder_image(slide, ph_idx, image_path):
            return

        # If placeholder insertion failed, place manually at the placeholder's coords
        pi = placeholders["image"]
        add_image_manual(slide, image_path,
                         pi["left"], pi["top"], pi["width"], pi["height"])
        return

    # For content_generic and other layouts without an image placeholder,
    # use the manual content area or center below the headline
    manual_area = layout_info.get("manual_content_area", {}).get("safe_area")
    if manual_area:
        # Place image in the safe content area, centered
        area_w = manual_area["width"]
        area_h = manual_area["height"]

        # Read image dimensions to maintain aspect ratio
        from PIL import Image
        try:
            with Image.open(image_path) as img:
                img_w, img_h = img.size
            aspect = img_w / img_h
        except Exception:
            # Fallback: assume 16:9
            aspect = 16 / 9

        # Fit image within the safe area
        if area_w / area_h > aspect:
            # Height-constrained
            img_height = area_h
            img_width = img_height * aspect
        else:
            # Width-constrained
            img_width = area_w
            img_height = img_width / aspect

        img_left = manual_area["left"] + (area_w - img_width) / 2
        img_top = manual_area["top"] + (area_h - img_height) / 2

        add_image_manual(slide, image_path, img_left, img_top, img_width, img_height)
    else:
        # Last resort: center on slide
        sw = 10.0
        sh = 5.62
        img_width = sw * 0.80
        img_left = (sw - img_width) / 2
        img_top = sh * 0.18
        add_image_manual(slide, image_path, img_left, img_top, img_width)


# ---------------------------------------------------------------------------
# Main build pipeline
# ---------------------------------------------------------------------------

def build_deck(yaml_path, output_path=None, upload=False):
    """Build a complete deck from a YAML definition file."""
    script_dir = os.path.dirname(os.path.abspath(__file__))

    catalog = load_curated_layouts()

    # Load YAML
    with open(yaml_path) as f:
        deck_def = yaml.safe_load(f)

    title = deck_def.get("title", "Untitled Presentation")
    date = deck_def.get("date", "")
    slides_def = deck_def.get("slides", [])

    if not output_path:
        output_path = os.path.splitext(yaml_path)[0] + ".pptx"

    # Load brand configuration (before template, so we can use brand.template_path)
    brand_name = deck_def.get("brand", None)
    brand = BrandConfig.load(brand_name)
    if brand_name:
        print(f"Brand: {brand.name}")

    # Set module-level icon directory from brand config
    global _icon_dir
    icon_catalog = brand.icon_catalog_path
    if os.path.isfile(icon_catalog):
        _icon_dir = os.path.dirname(icon_catalog)
    else:
        _icon_dir = os.path.join(script_dir, "icons")

    template_path = brand.template_path
    if not os.path.exists(template_path):
        print(f"ERROR: Template not found at {template_path}")
        sys.exit(1)

    print(f"Building: {title}")
    print(f"Template: {template_path}")
    print(f"Output: {output_path}")
    print(f"Slides: {len(slides_def)}")
    print()

    # Open template and set canvas scale
    prs = Presentation(template_path)

    # Detect canvas scale — Google Slides exports use 20"x11.25" (2x standard)
    global _canvas_scale
    canvas_w_in = prs.slide_width / 914400
    _canvas_scale = canvas_w_in / 10.0
    if _canvas_scale != 1.0:
        print(f"Canvas: {canvas_w_in:.1f}\" (scale: {_canvas_scale:.1f}x)")

    # Remove all existing slides

    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Load or extract layout mapping
    from extract_layout_mapping import extract_layout_mapping
    layout_mapping_path = os.path.join(brand._brand_dir, "layout_mapping.json")
    if os.path.isfile(layout_mapping_path):
        with open(layout_mapping_path) as f:
            layout_mapping = json.load(f)
    else:
        layout_mapping = extract_layout_mapping(template_path)

    deck_meta = {
        "title": title,
        "date": date,
        "base_dir": os.path.dirname(os.path.abspath(yaml_path)),
        "style": deck_def.get("style", "corporate"),
        "diagram_variants": deck_def.get("diagram_variants", 3),
        "brand": brand,
        "audience": deck_def.get("audience", ""),  # c-suite, manager, ic
        "layout_mapping": layout_mapping,
    }

    # Build each slide
    for i, slide_def in enumerate(slides_def):
        layout_name = slide_def.get("layout", "content_generic")
        slide_title = slide_def.get("headline", slide_def.get("presentation_title", ""))
        print(f"  Slide {i+1}: {layout_name} — {slide_title[:50]}")
        build_slide(prs, slide_def, deck_meta, catalog)

    prs.save(output_path)
    print(f"\nSaved: {output_path}")

    if upload:
        upload_to_drive(output_path, title)

    return output_path


def upload_to_drive(file_path, title):
    """Upload PPTX to Google Drive and convert to Google Slides.

    Requires the `gog` CLI tool and GOG_ACCOUNT environment variable.
    Set GOG_ACCOUNT to your Google Workspace email address.
    """
    account = os.environ.get("GOG_ACCOUNT", "")
    if not account:
        print("Upload failed: Set GOG_ACCOUNT environment variable to your Google Workspace email")
        return
    print("\nUploading to Google Drive...")
    result = subprocess.run(
        [
            "gog", "-a", account,
            "drive", "upload", file_path,
            "--name", title,
            "--convert-to", "slides",
        ],
        capture_output=True, text=True
    )
    if result.returncode == 0:
        print(result.stdout)
    else:
        print(f"Upload failed: {result.stderr}")


def main():
    parser = argparse.ArgumentParser(description="Deck Builder — YAML to branded PPTX")
    parser.add_argument("yaml_file", help="YAML deck definition file")
    parser.add_argument("--output", "-o", help="Output PPTX file path")
    parser.add_argument("--upload", action="store_true",
                        help="Upload to Google Drive and convert to Google Slides")
    args = parser.parse_args()

    build_deck(args.yaml_file, args.output, args.upload)


if __name__ == "__main__":
    main()
