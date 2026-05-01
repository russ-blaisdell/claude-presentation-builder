#!/usr/bin/env python3
"""
Numbered List Layout Tests

Tests the structural correctness of the numbered_list layout:
- Slot heights are uniform across items (parallel structure)
- Number markers don't overlap each other vertically
- Item text never wraps beyond 2 lines (authoring rule)

Regression: claude-costs-deck slide 25 had 5 numbered items with mixed line
counts. Item 1 wrapped to 2 lines; items 2-5 fit on 1 line. The renderer
allocated per-item slot heights based on estimate_text_height(), which
over-counted lines for items 3-5 (~100 chars at 8.4" 12pt → estimator says
2 lines, DM Sans actually fits 1). The result was visibly inconsistent
inter-item pitch — items 2→3 looked tight while 3→4 and 4→5 looked loose.

Fix: use uniform slot height (max actual height across all items) so the
list reads as a parallel sequence regardless of per-item text length.

Run: python3 -m pytest test_numbered_list_layout.py -v
  or: python3 test_numbered_list_layout.py
"""

import os
import sys
import tempfile
import subprocess
from pathlib import Path

try:
    import pytest
except ImportError:
    pytest = None

try:
    from pptx import Presentation
    from pptx.util import Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


SCRIPT_DIR = Path(__file__).parent.resolve()
PYTHON = sys.executable
EMU_PER_INCH = 914400


# Fixture YAML: 5 items engineered to span the wrap boundary.
# Lengths chosen to match the slide-25 regression (one 2-line item + four
# items of length ~100-110 chars that the buggy estimator over-counts as
# 2-line but actually render as 1 line).
FIXTURE_YAML = """\
title: "Numbered List Layout Test"
brand: generic
slides:
- layout: numbered_list
  headline: "Sequencing — Adoption First"
  items:
  - "Months 1-6 (Mar-Aug 2026) — Drive adoption; capture baseline productivity data; EM calibration on objective metrics"
  - "Q3 2026 — Performance-managed reduction of bottom 10% (140 heads)"
  - "Q3 2026 onward — Activate retention program: comp refresh, promotion velocity, equity for top decile"
  - "Months 9-18 — Selective backfill at high bar (refill ~30% of departures); reinvest savings into retention"
  - "Year 2 (2027) — Steady state. Top-skewed org of ~900 producing 1,485 effective output (vs today's 1,400)"
"""


def _build_fixture(tmp_dir):
    """Build the fixture deck in tmp_dir and return path to the .pptx."""
    yaml_path = tmp_dir / "fixture.yaml"
    yaml_path.write_text(FIXTURE_YAML)

    pptx_path = tmp_dir / "fixture.pptx"

    # Build via build_deck.py main()
    sys.path.insert(0, str(SCRIPT_DIR))
    try:
        import build_deck
        build_deck.build_deck(str(yaml_path), str(pptx_path))
    finally:
        sys.path.pop(0)

    return pptx_path


def _find_number_markers(slide):
    """Find numbered marker shapes (single-digit textboxes, narrow width).

    Returns list sorted by y-position (top to bottom).
    """
    markers = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # A numbered_list marker is a single digit (1-9) in a narrow box (<1")
        if text in {str(n) for n in range(1, 10)}:
            width_in = shape.width / EMU_PER_INCH
            if width_in < 1.0:
                markers.append(shape)
    return sorted(markers, key=lambda s: s.top)


def _find_title_textboxes(slide, markers):
    """Find the title textbox associated with each number marker.

    Returns list of (marker, title_textbox) pairs in marker order.
    """
    pairs = []
    other_shapes = [s for s in slide.shapes
                    if s.has_text_frame and s not in markers]

    for marker in markers:
        # The title is the textbox to the right of the marker, at the same y
        marker_y = marker.top
        candidates = [s for s in other_shapes
                      if abs(s.top - marker_y) < EMU_PER_INCH * 0.1
                      and s.left > marker.left]
        if candidates:
            # Closest to marker.top
            title = min(candidates, key=lambda s: abs(s.top - marker_y))
            pairs.append((marker, title))
    return pairs


# ─── Tests ──────────────────────────────────────────────────


def test_numbered_list_uniform_pitch():
    """Number markers must have uniform vertical pitch (max/min ≤ 1.05).

    A numbered_list represents a parallel sequence (step 1, step 2, ...) and
    must read with consistent rhythm regardless of per-item text length.
    """
    if not HAS_PPTX:
        if pytest:
            pytest.skip("python-pptx not available")
        return

    with tempfile.TemporaryDirectory() as td:
        pptx = _build_fixture(Path(td))
        prs = Presentation(pptx)
        markers = _find_number_markers(prs.slides[0])

    assert len(markers) == 5, (
        f"Expected 5 number markers, found {len(markers)}"
    )

    gaps_emu = [markers[i+1].top - markers[i].top
                for i in range(len(markers) - 1)]
    gaps_in = [g / EMU_PER_INCH for g in gaps_emu]

    max_gap = max(gaps_emu)
    min_gap = min(gaps_emu)
    ratio = max_gap / min_gap

    assert ratio <= 1.05, (
        f"\nNumbered list pitch is not uniform.\n"
        f"  Inter-marker gaps (inches): "
        f"{[f'{g:.3f}' for g in gaps_in]}\n"
        f"  max/min ratio: {ratio:.2f} (max allowed: 1.05)\n"
        f"\n"
        f"  Expected: every numbered item occupies the same vertical slot,\n"
        f"  so all gaps are equal within rounding.\n"
        f"\n"
        f"  Likely cause: per-item slots sized to estimated text height,\n"
        f"  causing items with different (or mis-estimated) line counts to\n"
        f"  produce visibly different inter-item pitches.\n"
        f"\n"
        f"  Fix: use uniform slot height = max(item_height) across all items."
    )


def test_numbered_list_markers_dont_overlap():
    """Number marker bounding boxes must not overlap vertically.

    A 0.5-inch tall '1' textbox followed by a slot only 0.4 inches tall
    causes consecutive markers to overlap, even if text content fits.
    """
    if not HAS_PPTX:
        if pytest:
            pytest.skip("python-pptx not available")
        return

    with tempfile.TemporaryDirectory() as td:
        pptx = _build_fixture(Path(td))
        prs = Presentation(pptx)
        markers = _find_number_markers(prs.slides[0])

    for i in range(len(markers) - 1):
        bottom = markers[i].top + markers[i].height
        top_next = markers[i+1].top
        overlap_emu = bottom - top_next
        overlap_in = overlap_emu / EMU_PER_INCH

        assert overlap_emu <= 0, (
            f"\nNumber markers {i+1} and {i+2} overlap by "
            f"{overlap_in:.3f} inches.\n"
            f"\n"
            f"  Cause: number textbox height exceeds the per-item slot\n"
            f"  height for short single-line items.\n"
            f"\n"
            f"  Fix: size number textbox to match slot height, or make\n"
            f"  slot height = max(title_h, marker_h)."
        )


def test_numbered_list_max_two_lines_authoring_rule():
    """Authoring rule: numbered_list items must not exceed 2 lines.

    Long items break the parallel-step reading model and force inconsistent
    vertical layout. Authors should either shorten the text or split into
    multiple items.
    """
    if not HAS_PPTX:
        if pytest:
            pytest.skip("python-pptx not available")
        return

    sys.path.insert(0, str(SCRIPT_DIR))
    try:
        from build_deck import estimate_text_height
    finally:
        sys.path.pop(0)

    # Mirror the renderer's text width and font size for numbered_list
    text_w = 8.4   # see _build_numbered_list: 9.3 - 0.6 - 0.15 = 8.55, conservative
    font_pt = 12
    line_h = font_pt * 1.4 / 72  # 0.233"
    max_lines = 2
    max_h = max_lines * line_h * 1.10  # +10% slack for estimator/renderer drift

    import yaml
    fixture = yaml.safe_load(FIXTURE_YAML)
    items = fixture["slides"][0]["items"]

    violations = []
    for i, item in enumerate(items):
        text = item if isinstance(item, str) else item.get("title", "")
        h = estimate_text_height(text, text_w, font_pt)
        est_lines = round(h / line_h)
        if h > max_h:
            violations.append((i + 1, len(text), est_lines, text[:60]))

    assert not violations, (
        f"\n{len(violations)} numbered_list item(s) exceed the 2-line limit:\n"
        + "\n".join(
            f"  Item {n}: {chars} chars, ~{lines} lines — \"{preview}...\""
            for n, chars, lines, preview in violations
        )
        + f"\n\n  Authoring rule: numbered_list items must fit in ≤ 2 lines.\n"
        f"  Either shorten the text or split into multiple items."
    )


# ─── Number marker rendering ───────────────────────────────


# Fixture: all items fit on a single line. This stresses the case where
# slot_h would otherwise be smaller than the 28pt number marker's natural
# height — which used to clip the markers to invisibility.
SHORT_ITEMS_YAML = """\
title: "Short Items Numbered List"
brand: generic
slides:
- layout: numbered_list
  headline: "All Single-Line Items"
  items:
  - "Sponsor AWS EDP and Google CUD layered commits"
  - "Approve smaller Anthropic-direct commit"
  - "Approve Q3 2026 performance reduction framework"
  - "Sponsor active top-decile retention investment"
  - "Endorse 18-month adoption-first sequencing"
"""


def test_number_markers_render_with_short_items():
    """Number markers must remain visible when all items fit on 1 line.

    Regression: when slot_h was sized strictly to max(item_heights), short
    single-line items produced a slot ~0.28" tall, but the 28pt number
    marker needs ~0.54". The marker textbox shrunk below its content,
    causing the digit to clip to invisibility in PowerPoint and LibreOffice
    rendering. Fix: enforce min_slot_h that accommodates the marker.
    """
    if not HAS_PPTX:
        if pytest:
            pytest.skip("python-pptx not available")
        return

    with tempfile.TemporaryDirectory() as td:
        yaml_path = Path(td) / "short.yaml"
        yaml_path.write_text(SHORT_ITEMS_YAML)
        pptx_path = Path(td) / "short.pptx"

        sys.path.insert(0, str(SCRIPT_DIR))
        try:
            import build_deck
            build_deck.build_deck(str(yaml_path), str(pptx_path))
        finally:
            sys.path.pop(0)

        prs = Presentation(pptx_path)
        markers = _find_number_markers(prs.slides[0])

    assert len(markers) == 5, (
        f"Expected 5 number markers, found {len(markers)}.\n"
        f"This typically means short single-line items produced slots\n"
        f"too small for the 28pt marker, causing it to render but become\n"
        f"invisible. Enforce a minimum slot height ≥ 0.6\"."
    )

    # Each marker should have a height that fits its 28pt content
    # (28pt × 1.4 / 72 = 0.544"). Allow some renderer slack but flag
    # anything below 0.5" as a likely clipping risk.
    too_small = []
    for i, m in enumerate(markers):
        h_in = m.height / EMU_PER_INCH
        if h_in < 0.5:
            too_small.append((i + 1, h_in))

    assert not too_small, (
        f"\n{len(too_small)} number marker(s) sized too small for 28pt content:\n"
        + "\n".join(
            f"  Marker {n}: height {h:.3f}\" (need ≥ 0.5\" for 28pt)"
            for n, h in too_small
        )
        + "\n\n  Fix: enforce min_slot_h ≥ 0.6\" so the marker fits."
    )


# ─── Vertical centering ────────────────────────────────────


def test_title_textframes_are_vertically_centered():
    """Title textboxes (when no body) must be vertically centered in their slot.

    With uniform slot heights, single-line items would otherwise sit at the
    top of their (taller) slot with whitespace below — visibly off vs the
    multi-line items that fill the slot. MSO_ANCHOR.MIDDLE on the title text
    frame keeps the rhythm consistent regardless of per-item line count.
    """
    if not HAS_PPTX:
        if pytest:
            pytest.skip("python-pptx not available")
        return

    from pptx.enum.text import MSO_ANCHOR

    with tempfile.TemporaryDirectory() as td:
        pptx = _build_fixture(Path(td))
        prs = Presentation(pptx)
        slide = prs.slides[0]
        markers = _find_number_markers(slide)
        pairs = _find_title_textboxes(slide, markers)

    assert len(pairs) == 5, (
        f"Expected 5 (marker, title) pairs, found {len(pairs)}"
    )

    not_centered = []
    for i, (_marker, title) in enumerate(pairs):
        anchor = title.text_frame.vertical_anchor
        if anchor != MSO_ANCHOR.MIDDLE:
            not_centered.append((i + 1, anchor))

    assert not not_centered, (
        f"\n{len(not_centered)} title textbox(es) not vertically centered:\n"
        + "\n".join(
            f"  Item {n}: vertical_anchor = {a} (expected MIDDLE)"
            for n, a in not_centered
        )
        + "\n\n  Fix: set tf.vertical_anchor = MSO_ANCHOR.MIDDLE on the\n"
        "  title textbox in _build_numbered_list when there's no body."
    )


# ─── Build-time validation ─────────────────────────────────


def test_build_rejects_overlong_items():
    """Build must raise ValueError when any item exceeds 2 lines.

    Authoring rule enforced at build time so authors fail fast rather than
    discovering it visually after rendering.
    """
    bad_yaml = """\
title: "Bad Numbered List"
brand: generic
slides:
- layout: numbered_list
  headline: "Items Too Long"
  items:
  - "Months 1-6 (Mar-Aug 2026) — Drive adoption hard. Training, AI Academy, fluency targets, workflow integrations across the entire engineering organization with full coverage and reporting"
  - "Q3 2026 — Performance-managed reduction"
"""
    sys.path.insert(0, str(SCRIPT_DIR))
    try:
        import build_deck
        with tempfile.TemporaryDirectory() as td:
            yaml_path = Path(td) / "bad.yaml"
            pptx_path = Path(td) / "bad.pptx"
            yaml_path.write_text(bad_yaml)

            raised = False
            try:
                build_deck.build_deck(str(yaml_path), str(pptx_path))
            except ValueError as e:
                raised = True
                msg = str(e)
                assert "2-line" in msg or "2 line" in msg, (
                    f"ValueError raised but message doesn't mention the rule: {msg}"
                )
                assert "Item 1" in msg, (
                    f"ValueError should identify which item violated: {msg}"
                )

            assert raised, (
                "Build should have raised ValueError for an item exceeding "
                "the 2-line limit, but built successfully."
            )
    finally:
        sys.path.pop(0)


# ─── Pre-fix sanity check ──────────────────────────────────


def test_fixture_reproduces_slide25_bug():
    """Sentinel: this test should FAIL on the buggy renderer and PASS after
    the fix. It's the regression contract for the slide-25 issue.

    Marked xfail before fix lands; flip to expecting pass once renderer is
    updated to use uniform slot heights.
    """
    # This is just an alias for test_numbered_list_uniform_pitch — it
    # documents that the fixture is specifically designed to reproduce
    # the slide-25 regression and that a passing build means the bug is
    # fixed.
    test_numbered_list_uniform_pitch()


# ─── CLI runner ─────────────────────────────────────────────


if __name__ == "__main__":
    tests = [
        test_numbered_list_uniform_pitch,
        test_numbered_list_markers_dont_overlap,
        test_numbered_list_max_two_lines_authoring_rule,
        test_number_markers_render_with_short_items,
        test_title_textframes_are_vertically_centered,
        test_build_rejects_overlong_items,
    ]
    failures = 0
    for t in tests:
        name = t.__name__
        try:
            t()
            print(f"  PASS  {name}")
        except AssertionError as e:
            failures += 1
            print(f"  FAIL  {name}")
            print(str(e))
        except Exception as e:
            failures += 1
            print(f"  ERROR {name}: {e}")
    sys.exit(1 if failures else 0)
