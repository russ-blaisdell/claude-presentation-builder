#!/usr/bin/env python3
"""PPTX Theme Manager — Extract, derive, and inject OOXML themes.

Three core operations:
  1. extract_theme(pptx_path) — read the 12 color slots + 2 font slots from a PPTX
  2. derive_theme(colors, fonts) — compute a full theme from brand tokens
  3. inject_theme(pptx_path, theme) — write a complete theme into a PPTX

Usage:
    # Extract theme from a template
    python3 patch_template_theme.py extract --template brands/generic/template.pptx

    # Inject a stored theme into a template
    python3 patch_template_theme.py inject --brand startup

    # Derive + inject from brand.yaml (no source template)
    python3 patch_template_theme.py derive --brand startup
"""

import argparse
import json
import os
import shutil
import sys
import tempfile
import zipfile

import yaml
from lxml import etree

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# The 12 OOXML color scheme slots (fixed by spec, every PPTX has exactly these)
COLOR_SLOTS = [
    'dk1', 'dk2', 'lt1', 'lt2',
    'accent1', 'accent2', 'accent3', 'accent4',
    'accent5', 'accent6', 'hlink', 'folHlink',
]


# ---------------------------------------------------------------------------
# Extract — read theme from a PPTX
# ---------------------------------------------------------------------------

def extract_theme(pptx_path):
    """Extract the complete theme from a PPTX file.

    Returns dict with 'colors' (12 slots) and 'fonts' (heading, body).
    Reads theme XML for colors, then scans actual slide runs for the
    real fonts in use (template themes often declare generic fonts while
    placeholders override with the actual brand fonts).
    """
    theme = {'colors': {}, 'fonts': {'heading': 'Arial', 'body': 'Arial'}}

    with zipfile.ZipFile(pptx_path, 'r') as z:
        # Find theme1.xml (primary theme)
        theme_files = [f for f in z.namelist()
                       if 'theme/theme1.xml' in f.lower()]
        if not theme_files:
            return theme

        data = z.read(theme_files[0])
        root = etree.fromstring(data)

        # Extract colors
        for clr_scheme in root.iter(f'{{{A_NS}}}clrScheme'):
            for slot in COLOR_SLOTS:
                el = clr_scheme.find(f'{{{A_NS}}}{slot}')
                if el is not None:
                    srgb = el.find(f'{{{A_NS}}}srgbClr')
                    if srgb is not None:
                        theme['colors'][slot] = f"#{srgb.get('val')}"
                    else:
                        # sysClr fallback (e.g., dk1 sometimes uses system color)
                        sys_clr = el.find(f'{{{A_NS}}}sysClr')
                        if sys_clr is not None:
                            last = sys_clr.get('lastClr', '000000')
                            theme['colors'][slot] = f"#{last}"
            break  # Only first clrScheme

        # Extract fonts
        for major in root.iter(f'{{{A_NS}}}majorFont'):
            latin = major.find(f'{{{A_NS}}}latin')
            if latin is not None and latin.get('typeface'):
                theme['fonts']['heading'] = latin.get('typeface')
            break
        for minor in root.iter(f'{{{A_NS}}}minorFont'):
            latin = minor.find(f'{{{A_NS}}}latin')
            if latin is not None and latin.get('typeface'):
                theme['fonts']['body'] = latin.get('typeface')
            break

    # Second pass: scan actual slide runs for the real fonts in use.
    # Template themes often declare generic fonts (Arial) while the actual
    # placeholders use brand-specific fonts that differ from the theme declaration.
    try:
        from pptx import Presentation as _Prs
        from collections import defaultdict as _defaultdict
        prs = _Prs(pptx_path)
        font_stats = _defaultdict(lambda: {'count': 0, 'total_size': 0, 'sized_count': 0})

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        fname = run.font.name
                        if not fname:
                            continue
                        font_stats[fname]['count'] += 1
                        if run.font.size:
                            font_stats[fname]['total_size'] += run.font.size / 12700  # EMU to pt
                            font_stats[fname]['sized_count'] += 1

        # Heading = font with highest average size (among fonts with 5+ runs)
        # Body = most-used font that isn't the heading font
        candidates = {f: s for f, s in font_stats.items() if s['count'] >= 5}
        if candidates:
            # Find heading: highest avg size
            heading_font = max(candidates, key=lambda f:
                candidates[f]['total_size'] / max(candidates[f]['sized_count'], 1))
            theme['fonts']['heading'] = heading_font

            # Find body: most runs, excluding heading font
            body_candidates = {f: s for f, s in candidates.items() if f != heading_font}
            if body_candidates:
                body_font = max(body_candidates, key=lambda f: body_candidates[f]['count'])
                theme['fonts']['body'] = body_font
    except Exception:
        pass  # Fall back to theme XML fonts if scan fails

    return theme


# ---------------------------------------------------------------------------
# Derive — compute a theme from brand tokens
# ---------------------------------------------------------------------------

def _hex_to_rgb(h):
    h = h.lstrip('#')
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _rgb_to_hex(r, g, b):
    return f"#{r:02X}{g:02X}{b:02X}"


def _lighten(hex_color, factor):
    """Lighten a color toward white. factor=0 unchanged, factor=1 pure white."""
    r, g, b = _hex_to_rgb(hex_color)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return _rgb_to_hex(r, g, b)


def _darken(hex_color, factor):
    """Darken a color toward black. factor=0 unchanged, factor=1 pure black."""
    r, g, b = _hex_to_rgb(hex_color)
    return _rgb_to_hex(int(r * (1 - factor)), int(g * (1 - factor)),
                       int(b * (1 - factor)))


def derive_theme(colors, fonts):
    """Compute a full 12+2 theme from brand color/font tokens.

    Used when no source template is available — computes all slots using
    the functional role derivation rules.
    """
    primary = colors.get('primary', '#1A365D')
    secondary = colors.get('secondary', '#3182CE')
    accent = colors.get('accent', '#63B3ED')
    bg_light = colors.get('background_light', '#F5F5F5')
    link = colors.get('link', primary)

    return {
        'colors': {
            'dk1': primary,
            'dk2': secondary,
            'lt1': '#FFFFFF',
            'lt2': bg_light,
            'accent1': primary,
            'accent2': _lighten(primary, 0.80),      # Light decorative bg
            'accent3': accent,
            'accent4': '#FFFFFF',                      # Neutral fills (white)
            'accent5': _darken(primary, 0.40),         # Dark accent
            'accent6': bg_light,                       # Muted background
            'hlink': link,
            'folHlink': _darken(link, 0.30),
        },
        'fonts': {
            'heading': fonts.get('heading', 'Arial'),
            'body': fonts.get('body', 'Arial'),
        },
    }


# ---------------------------------------------------------------------------
# Inject — write a theme into a PPTX
# ---------------------------------------------------------------------------

def _inject_theme_xml(theme_xml_bytes, theme):
    """Inject a complete theme into theme XML. Returns patched bytes + count."""
    root = etree.fromstring(theme_xml_bytes)
    patched = 0

    colors = theme.get('colors', {})
    fonts = theme.get('fonts', {})

    # Inject fonts
    heading = fonts.get('heading', 'Arial')
    body = fonts.get('body', 'Arial')

    for major in root.iter(f'{{{A_NS}}}majorFont'):
        latin = major.find(f'{{{A_NS}}}latin')
        if latin is not None:
            latin.set('typeface', heading)
            patched += 1

    for minor in root.iter(f'{{{A_NS}}}minorFont'):
        latin = minor.find(f'{{{A_NS}}}latin')
        if latin is not None:
            latin.set('typeface', body)
            patched += 1

    # Inject all 12 color slots
    for clr_scheme in root.iter(f'{{{A_NS}}}clrScheme'):
        for slot in COLOR_SLOTS:
            if slot in colors:
                hex_clean = colors[slot].lstrip('#')
                el = clr_scheme.find(f'{{{A_NS}}}{slot}')
                if el is not None:
                    for child in list(el):
                        el.remove(child)
                    srgb = etree.SubElement(el, f'{{{A_NS}}}srgbClr')
                    srgb.set('val', hex_clean)
                    patched += 1
        break  # Only first clrScheme

    return etree.tostring(root, xml_declaration=True, encoding='UTF-8',
                          standalone=True), patched


def inject_theme(pptx_path, theme):
    """Inject a complete theme into all theme XML files in a PPTX."""
    total_patched = 0
    tmp_fd, tmp_path = tempfile.mkstemp(suffix='.pptx')
    os.close(tmp_fd)

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if 'theme' in item.filename.lower() and item.filename.endswith('.xml'):
                    patched_data, count = _inject_theme_xml(data, theme)
                    zout.writestr(item, patched_data)
                    total_patched += count
                else:
                    zout.writestr(item, data)

    shutil.move(tmp_path, pptx_path)
    return total_patched


# ---------------------------------------------------------------------------
# File I/O helpers
# ---------------------------------------------------------------------------

def inject_logo(pptx_path, logo_blob, logo_ext='png', position=None,
                pic_xml_template=None):
    """Inject a logo image into every slide layout in a PPTX template.

    If pic_xml_template is provided (from extract_logo), uses the exact
    original XML structure for pixel-perfect reproduction. Otherwise builds
    a new pic element from position dict.

    Args:
        pptx_path: Path to template PPTX
        logo_blob: Logo image bytes
        logo_ext: Image extension ('png' or 'jpg')
        position: Dict with left_emu, top_emu, width_emu, height_emu
        pic_xml_template: Original lxml Element from extract_logo for exact reproduction
    """
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _In
    from pptx.oxml.ns import qn as _qn
    import io as _io
    import copy as _copy

    prs = _Prs(pptx_path)
    injected = 0

    # Register the image in the package via a temp slide
    temp_layout = prs.slide_layouts[0]
    temp_slide = prs.slides.add_slide(temp_layout)
    pic_shape = temp_slide.shapes.add_picture(
        _io.BytesIO(logo_blob), _In(0), _In(0), _In(1), _In(1))

    blip = pic_shape._element.findall('.//' + _qn('a:blip'))[0]
    embed_rId = blip.get(_qn('r:embed'))
    image_part = temp_slide.part.rels[embed_rId].target_part

    for layout in prs.slide_layouts:
        # Skip layouts that already have a logo
        has_logo = False
        for shape in layout.shapes:
            if hasattr(shape, 'image'):
                w = shape.width / 914400
                h = shape.height / 914400
                if w < 1.5 and h < 0.5:
                    has_logo = True
                    break
        if has_logo:
            continue

        # Add image relationship to this layout
        rel = layout.part.relate_to(
            image_part,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image',
        )

        if pic_xml_template is not None:
            # Use the exact original XML structure
            pic_el = _copy.deepcopy(pic_xml_template)
            # Update the relationship ID to point to our image
            for blip_el in pic_el.findall('.//' + _qn('a:blip')):
                blip_el.set(_qn('r:embed'), rel)
        else:
            # Build from position
            pos = position or {'left_emu': 301752, 'top_emu': 4681728,
                               'width_emu': 640080, 'height_emu': 182880}
            a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            pic_xml = (
                f'<p:pic xmlns:p="{p_ns}" xmlns:a="{a_ns}" xmlns:r="{r_ns}">'
                f'<p:nvPicPr><p:cNvPr id="0" name="Logo"/>'
                f'<p:cNvPicPr preferRelativeResize="0"/><p:nvPr/></p:nvPicPr>'
                f'<p:blipFill rotWithShape="1"><a:blip r:embed="{rel}">'
                f'<a:alphaModFix/></a:blip><a:srcRect b="0" l="0" r="0" t="0"/>'
                f'<a:stretch/></p:blipFill>'
                f'<p:spPr><a:xfrm>'
                f'<a:off x="{pos["left_emu"]}" y="{pos["top_emu"]}"/>'
                f'<a:ext cx="{pos["width_emu"]}" cy="{pos["height_emu"]}"/>'
                f'</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f'<a:noFill/><a:ln><a:noFill/></a:ln></p:spPr></p:pic>'
            )
            pic_el = etree.fromstring(pic_xml)

        sp_tree = layout.element.find(_qn('p:cSld')).find(_qn('p:spTree'))
        sp_tree.append(pic_el)
        injected += 1

    # Remove temp slide
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

    prs.save(pptx_path)
    return injected


def extract_logo(pptx_path):
    """Extract the logo from a PPTX template.

    Finds the most common small image across slide layouts.
    Returns (blob, ext, position_dict, pic_xml_template) or (None, None, None, None).
    Position uses exact EMU values to preserve pixel-perfect placement.
    pic_xml_template is the original XML element for faithful reproduction.
    """
    from pptx import Presentation as _Prs
    from collections import Counter as _Counter
    import hashlib as _hashlib
    import copy as _copy

    prs = _Prs(pptx_path)
    logo_candidates = _Counter()  # hash -> count
    logo_data = {}  # hash -> data dict

    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if not hasattr(shape, 'image'):
                continue
            w = shape.width / 914400
            h = shape.height / 914400
            if w < 1.5 and h < 0.5 and w > 0.2:
                blob = shape.image.blob
                h_key = _hashlib.md5(blob).hexdigest()[:12]
                logo_candidates[h_key] += 1
                if h_key not in logo_data:
                    ext = 'png' if 'png' in shape.image.content_type else 'jpg'
                    logo_data[h_key] = {
                        'blob': blob,
                        'ext': ext,
                        # Exact EMU values for pixel-perfect placement
                        'left_emu': shape.left,
                        'top_emu': shape.top,
                        'width_emu': shape.width,
                        'height_emu': shape.height,
                        # Also store the full pic XML for faithful reproduction
                        'pic_xml': _copy.deepcopy(shape._element),
                    }

    if not logo_candidates:
        return None, None, None, None

    best_hash = logo_candidates.most_common(1)[0][0]
    data = logo_data[best_hash]
    return data['blob'], data['ext'], {
        'left_emu': data['left_emu'],
        'top_emu': data['top_emu'],
        'width_emu': data['width_emu'],
        'height_emu': data['height_emu'],
    }, data['pic_xml']


def save_theme(theme, output_path):
    """Save a theme dict to a JSON file."""
    with open(output_path, 'w') as f:
        json.dump(theme, f, indent=2)


def load_theme(theme_path):
    """Load a theme dict from a JSON file."""
    with open(theme_path) as f:
        return json.load(f)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="PPTX Theme Manager")
    sub = parser.add_subparsers(dest='command')

    # Extract
    ext = sub.add_parser('extract', help='Extract theme from a PPTX template')
    ext.add_argument('--template', required=True, help='Path to PPTX template')
    ext.add_argument('--output', help='Output theme.json path')

    # Derive
    der = sub.add_parser('derive', help='Derive theme from brand.yaml')
    der.add_argument('--brand', required=True, help='Brand name (under brands/)')
    der.add_argument('--output', help='Output theme.json path')

    # Inject
    inj = sub.add_parser('inject', help='Inject theme.json into brand template')
    inj.add_argument('--brand', required=True, help='Brand name (under brands/)')
    inj.add_argument('--theme', help='theme.json path (default: brands/<name>/theme.json)')

    args = parser.parse_args()
    script_dir = os.path.dirname(os.path.abspath(__file__))

    if args.command == 'extract':
        theme = extract_theme(args.template)
        output = args.output or os.path.splitext(args.template)[0] + '-theme.json'
        save_theme(theme, output)
        print(f"Extracted theme to {output}")
        print(f"  dk1: {theme['colors'].get('dk1', '?')}")
        print(f"  Heading: {theme['fonts']['heading']}")
        print(f"  Body: {theme['fonts']['body']}")

    elif args.command == 'derive':
        brand_dir = os.path.join(script_dir, 'brands', args.brand)
        yaml_path = os.path.join(brand_dir, 'brand.yaml')
        with open(yaml_path) as f:
            data = yaml.safe_load(f)
        theme = derive_theme(data.get('colors', {}), data.get('fonts', {}))
        output = args.output or os.path.join(brand_dir, 'theme.json')
        save_theme(theme, output)
        print(f"Derived theme for {data.get('name', args.brand)}")
        print(f"  dk1: {theme['colors']['dk1']}")
        print(f"  Saved to {output}")

    elif args.command == 'inject':
        brand_dir = os.path.join(script_dir, 'brands', args.brand)
        theme_path = args.theme or os.path.join(brand_dir, 'theme.json')
        template_path = os.path.join(brand_dir, 'template.pptx')
        theme = load_theme(theme_path)
        count = inject_theme(template_path, theme)
        print(f"Injected {count} theme values into {template_path}")

    else:
        parser.print_help()


if __name__ == '__main__':
    main()
