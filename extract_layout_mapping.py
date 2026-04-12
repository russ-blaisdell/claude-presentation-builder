#!/usr/bin/env python3
"""Extract layout mapping from a PPTX template.

Analyzes a template to determine:
  1. Which layout index to use as the blank canvas
  2. What placeholders exist on that canvas (headline, footers, etc.)
  3. Footer element configuration

Usage:
    python3 extract_layout_mapping.py --template path/to/template.pptx

Output: layout_mapping dict suitable for brand.yaml
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


# OOXML placeholder type IDs
PH_TYPES = {
    0: 'title',          # PP_PLACEHOLDER.TITLE
    1: 'body',           # PP_PLACEHOLDER.BODY
    2: 'center_title',   # PP_PLACEHOLDER.CENTER_TITLE
    3: 'subtitle',       # PP_PLACEHOLDER.SUBTITLE
    4: 'subtitle',       # Alt subtitle
    10: 'date',
    11: 'footer',        # Not in enum, but common
    12: 'slide_number',  # Not always in enum
    13: 'slide_number',  # PP_PLACEHOLDER.SLIDE_NUMBER
    14: 'header',
    15: 'footer',        # PP_PLACEHOLDER.FOOTER
    16: 'date',          # PP_PLACEHOLDER.DATE
}


def _classify_placeholder(ph):
    """Classify a placeholder by its type, index, position, and size."""
    idx = ph.placeholder_format.idx
    ph_type = ph.placeholder_format.type
    w_in = ph.width / 914400
    h_in = ph.height / 914400
    t_in = ph.top / 914400
    l_in = ph.left / 914400
    canvas_h = 5.625  # Will be adjusted by caller

    # Determine role based on type enum + position heuristics
    type_val = ph_type.real if ph_type else None

    if type_val in (1, 15) or idx == 0:
        if w_in > 5.0 and h_in > 0.3 and t_in < 1.0:
            role = 'headline'
        elif t_in > canvas_h * 0.85:
            role = 'footer'
        else:
            role = 'headline' if idx == 0 else 'body'
    elif type_val == 13 or idx == 12:
        role = 'slide_number'
    elif type_val == 16 or idx == 10:
        role = 'date'
    elif type_val == 15 or idx == 11:
        role = 'footer'
    elif type_val in (3, 4):
        # Subtitle — could be footer if at bottom
        if t_in > canvas_h * 0.85:
            role = 'footer'
        else:
            role = 'subtitle'
    else:
        role = 'unknown'

    return {
        'idx': idx,
        'role': role,
        'type': str(ph_type) if ph_type else None,
        'type_val': type_val,
        'left_in': round(l_in, 2),
        'top_in': round(t_in, 2),
        'width_in': round(w_in, 2),
        'height_in': round(h_in, 2),
    }


def find_blank_canvas(prs):
    """Find the best blank canvas layout for content rendering.

    Strategy:
      1. Check which layout the actual slides use most (if slides exist)
      2. Look for a layout named 'BLANK'
      3. Fall back to the layout with fewest placeholders
    """
    # Check actual slide usage
    layout_usage = {}
    for slide in prs.slides:
        name = slide.slide_layout.name
        layout_usage[name] = layout_usage.get(name, 0) + 1

    if layout_usage:
        most_used_name = max(layout_usage, key=layout_usage.get)
        for i, layout in enumerate(prs.slide_layouts):
            if layout.name == most_used_name:
                return i, layout.name, 'most_used'

    # Look for BLANK
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name.upper() == 'BLANK':
            return i, layout.name, 'name_match'

    # Fall back to fewest placeholders
    best_idx = 0
    best_count = 999
    for i, layout in enumerate(prs.slide_layouts):
        n = len(list(layout.placeholders))
        if n < best_count:
            best_count = n
            best_idx = i
    return best_idx, prs.slide_layouts[best_idx].name, 'fewest_placeholders'


def extract_layout_mapping(pptx_path):
    """Extract the complete layout mapping from a PPTX template.

    Returns dict with:
      - blank_canvas_idx: layout index for content slides
      - canvas: width/height in inches
      - headline: placeholder info for headline (or None)
      - footers: list of footer element configs
      - all_placeholders: raw placeholder data for the blank canvas
    """
    prs = Presentation(pptx_path)
    canvas_w = prs.slide_width / 914400
    canvas_h = prs.slide_height / 914400

    # Find blank canvas
    canvas_idx, canvas_name, method = find_blank_canvas(prs)
    layout = prs.slide_layouts[canvas_idx]

    # Classify all placeholders on the blank canvas
    placeholders = []
    for ph in layout.placeholders:
        info = _classify_placeholder(ph)
        # Adjust for actual canvas height
        if info['top_in'] > canvas_h * 0.85 and info['role'] in ('body', 'subtitle', 'unknown'):
            info['role'] = 'footer'
        placeholders.append(info)

    # Extract headline
    headline = None
    for ph in placeholders:
        if ph['role'] == 'headline':
            headline = {
                'placeholder_idx': ph['idx'],
                'width_in': ph['width_in'],
                'height_in': ph['height_in'],
                'top_in': ph['top_in'],
            }
            break

    # Extract footer elements
    footers = []
    for ph in placeholders:
        if ph['role'] in ('footer', 'date', 'slide_number'):
            element = {
                'placeholder_idx': ph['idx'],
                'type': ph['role'],
                'left_in': ph['left_in'],
                'top_in': ph['top_in'],
                'width_in': ph['width_in'],
                'height_in': ph['height_in'],
            }

            # Assign default content based on type
            if ph['role'] == 'date':
                element['default'] = '{date}'
            elif ph['role'] == 'footer':
                # Could be deck title or custom text
                # Heuristic: wider footer = title, narrower = custom
                if ph['width_in'] > 3.0:
                    element['default'] = '{title}'
                else:
                    element['default'] = '{title}'
            elif ph['role'] == 'slide_number':
                element['default'] = 'auto'

            footers.append(element)

    # Also check for subtitle placeholders used as footers in some templates
    for ph in placeholders:
        if ph['role'] == 'subtitle' and ph['top_in'] > canvas_h * 0.8:
            footers.append({
                'placeholder_idx': ph['idx'],
                'type': 'custom',
                'left_in': ph['left_in'],
                'top_in': ph['top_in'],
                'width_in': ph['width_in'],
                'height_in': ph['height_in'],
                'default': '{title}' if ph['left_in'] < canvas_w / 2 else '{date}',
            })

    # Sort footers by left position
    footers.sort(key=lambda f: f['left_in'])

    # Also scan other layouts for potential title/section divider use
    title_layout_idx = None
    section_layout_idx = None
    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name.upper()
        if name in ('TITLE', 'TITLE SLIDE', 'TITLE_SLIDE') or \
           (any(ph.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE
                for ph in layout.placeholders if ph.placeholder_format.type)):
            title_layout_idx = i
        if name in ('SECTION_HEADER', 'SECTION HEADER', 'START', 'DIVIDER'):
            section_layout_idx = i

    return {
        'blank_canvas_idx': canvas_idx,
        'blank_canvas_name': canvas_name,
        'detection_method': method,
        'canvas': {
            'width_in': round(canvas_w, 3),
            'height_in': round(canvas_h, 3),
        },
        'headline': headline,
        'footers': footers,
        'title_layout_idx': title_layout_idx,
        'section_layout_idx': section_layout_idx,
        'all_placeholders': placeholders,
    }


def print_mapping(mapping):
    """Pretty-print a layout mapping."""
    print(f"  Blank canvas: [{mapping['blank_canvas_idx']}] {mapping['blank_canvas_name']}"
          f" (detected via: {mapping['detection_method']})")
    print(f"  Canvas: {mapping['canvas']['width_in']}\" x {mapping['canvas']['height_in']}\"")

    if mapping['headline']:
        h = mapping['headline']
        print(f"  Headline: ph[{h['placeholder_idx']}] {h['width_in']}\"x{h['height_in']}\" at top={h['top_in']}\"")
    else:
        print(f"  Headline: none (builder will add textbox)")

    if mapping['footers']:
        print(f"  Footers ({len(mapping['footers'])}):")
        for f in mapping['footers']:
            print(f"    ph[{f['placeholder_idx']}] {f['type']:12s} -> {f['default']:10s}"
                  f"  at ({f['left_in']}\",{f['top_in']}\") {f['width_in']}\"x{f['height_in']}\"")
    else:
        print(f"  Footers: none")

    if mapping['title_layout_idx'] is not None:
        print(f"  Title layout: [{mapping['title_layout_idx']}]")
    if mapping['section_layout_idx'] is not None:
        print(f"  Section layout: [{mapping['section_layout_idx']}]")


def main():
    parser = argparse.ArgumentParser(description="Extract layout mapping from PPTX template")
    parser.add_argument("--template", required=True, help="Path to PPTX template")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    args = parser.parse_args()

    mapping = extract_layout_mapping(args.template)

    if args.json:
        # Remove all_placeholders for cleaner output
        clean = {k: v for k, v in mapping.items() if k != 'all_placeholders'}
        print(json.dumps(clean, indent=2))
    else:
        print_mapping(mapping)


if __name__ == "__main__":
    main()
