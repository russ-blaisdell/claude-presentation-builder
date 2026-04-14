#!/usr/bin/env python3
"""Brand Onboarding Pipeline — Create a complete brand package.

Four onboarding paths:
  1. Template + corpus — extract theme from template, enrich from corpus
  2. Template only — extract theme from template
  3. Corpus only — synthesize theme from slide analysis, build template
  4. Brand YAML only — derive theme from color tokens, build template

Each path produces:
  brands/<name>/brand.yaml      — brand tokens
  brands/<name>/theme.json      — stored 12-slot OOXML theme
  brands/<name>/template.pptx   — template with theme baked in
  brands/<name>/title-assets/   — background images + logo

Interactive image selection:
  When a template or corpus contains images, the user is prompted to
  assign them to roles: title backgrounds, agenda panels, section dividers,
  closing backgrounds, and logo.

Usage:
    python3 onboard-brand.py --name "acme" --template path/to/template.pptx
    python3 onboard-brand.py --name "acme" --corpus path/to/decks/
    python3 onboard-brand.py --name "acme" --template t.pptx --corpus path/to/decks/
    python3 onboard-brand.py --name "acme"  # uses existing brand.yaml
    python3 onboard-brand.py --name "acme" --auto  # non-interactive (auto-select)
"""

import argparse
import json
import os
import shutil
import subprocess
import sys
from collections import Counter
from pathlib import Path

import yaml

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx")
    sys.exit(1)

from patch_template_theme import (
    extract_theme, derive_theme, inject_theme, save_theme
)
from extract_brand_images import extract_all_images


# ---------------------------------------------------------------------------
# Corpus analysis
# ---------------------------------------------------------------------------

def extract_corpus_colors(corpus_dir):
    """Extract dominant colors from a directory of PPTX files."""
    color_counter = Counter()
    if not corpus_dir or not os.path.isdir(corpus_dir):
        return {}

    for pptx_file in Path(corpus_dir).glob("*.pptx"):
        try:
            prs = Presentation(str(pptx_file))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'fill') and shape.fill.type is not None:
                        try:
                            rgb = shape.fill.fore_color.rgb
                            hex_str = f"#{rgb}"
                            if hex_str not in ("#FFFFFF", "#000000"):
                                color_counter[hex_str] += 1
                        except Exception:
                            pass
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color and run.font.color.rgb:
                                        hex_str = f"#{run.font.color.rgb}"
                                        if hex_str not in ("#FFFFFF", "#000000"):
                                            color_counter[hex_str] += 1
                                except Exception:
                                    pass
        except Exception:
            continue

    return dict(color_counter.most_common(10))


def extract_corpus_images(corpus_dir):
    """Extract images from all PPTX files in a corpus directory."""
    from extract_brand_images import extract_all_images as _extract
    all_images = {'images': [], 'stats': {}}

    for pptx_file in sorted(Path(corpus_dir).glob("*.pptx")):
        try:
            result = _extract(str(pptx_file))
            # Tag each image with its source file
            for img in result['images']:
                img['source'] = pptx_file.name
            all_images['images'].extend(result['images'])
        except Exception:
            continue

    # Deduplicate by hash
    seen = {}
    deduped = []
    for img in all_images['images']:
        if img['hash'] not in seen:
            seen[img['hash']] = img
            deduped.append(img)
        else:
            # Merge slide counts
            seen[img['hash']]['num_slides'] += img['num_slides']
    all_images['images'] = deduped

    stats = Counter(r['category'] for r in deduped)
    all_images['stats'] = dict(stats)
    return all_images


def extract_canvas_size(prs):
    """Extract slide dimensions in inches."""
    return {
        "width_inches": round(prs.slide_width / 914400, 3),
        "height_inches": round(prs.slide_height / 914400, 3),
    }


# ---------------------------------------------------------------------------
# Interactive image selection
# ---------------------------------------------------------------------------

def _prompt_selection(candidates, role, allow_multiple=True, allow_none=True):
    """Prompt user to select images for a role.

    Returns list of selected image dicts, or empty list.
    """
    if not candidates:
        print(f"  No candidates found for {role}.")
        return []

    print(f"\n  {role.upper()} — {len(candidates)} candidates:")
    for i, img in enumerate(candidates):
        slides = ', '.join(str(s) for s in img.get('slides', [])[:3])
        source = f" from {img['source']}" if 'source' in img else ''
        print(f"    [{i+1}] {img['width']:.0f}\"x{img['height']:.0f}\" "
              f"{img['content_type'].split('/')[-1]} {img['size_kb']}KB "
              f"— slides: {slides}{source}")

    prompt_parts = []
    if allow_multiple:
        prompt_parts.append("pick numbers (e.g. 1,3)")
    else:
        prompt_parts.append("pick one number")
    if allow_none:
        prompt_parts.append("'none' to skip")
    prompt_parts.append("'all' for all")

    while True:
        choice = input(f"  Select for {role} ({', '.join(prompt_parts)}): ").strip().lower()

        if choice == 'none' and allow_none:
            return []
        if choice == 'all':
            return candidates
        if choice == '':
            # Default: first item
            return [candidates[0]] if candidates else []

        try:
            indices = [int(x.strip()) - 1 for x in choice.split(',')]
            selected = [candidates[i] for i in indices if 0 <= i < len(candidates)]
            if selected:
                if not allow_multiple:
                    return [selected[0]]
                return selected
        except (ValueError, IndexError):
            pass
        print(f"    Invalid input. Try again.")


def _auto_select(candidates, role, max_items=3):
    """Auto-select images for a role (non-interactive mode).

    Picks the largest/most-used candidates.
    """
    if not candidates:
        return []
    # Sort by size (largest first) then by usage frequency
    ranked = sorted(candidates, key=lambda x: (-x['size_kb'], -x['num_slides']))
    return ranked[:max_items]


def select_brand_images(extracted, interactive=True):
    """Present extracted images to user and get role assignments.

    Returns dict with:
        'title': list of image dicts for title backgrounds
        'agenda': list for agenda panels
        'section': list for section divider backgrounds
        'closing': list for closing backgrounds
        'logo': single image dict or None
    """
    full_bleed = [img for img in extracted['images'] if img['category'] == 'full-bleed']
    panels = [img for img in extracted['images'] if img['category'] == 'panel']
    logos = [img for img in extracted['images'] if img['category'] == 'logo']

    print(f"\n  Found: {len(full_bleed)} backgrounds, {len(panels)} panels, {len(logos)} logos")

    if interactive:
        title = _prompt_selection(full_bleed, "TITLE slide backgrounds")
        agenda = _prompt_selection(panels, "AGENDA side panels")
        section = _prompt_selection(full_bleed, "SECTION DIVIDER backgrounds",
                                    allow_none=True)
        closing_prompt = "same as title"
        if title:
            closing_choice = input(f"\n  CLOSING backgrounds — same as title? (yes/no): ").strip().lower()
            if closing_choice in ('yes', 'y', ''):
                closing = title
            else:
                closing = _prompt_selection(full_bleed, "CLOSING backgrounds")
        else:
            closing = _prompt_selection(full_bleed, "CLOSING backgrounds")

        logo = None
        if logos:
            logo_sel = _prompt_selection(logos, "LOGO", allow_multiple=False)
            logo = logo_sel[0] if logo_sel else None
    else:
        # Auto mode
        title = _auto_select(full_bleed, "title", 3)
        agenda = _auto_select(panels, "agenda", 3)
        section = title[:1] if title else []  # reuse first title bg
        closing = title[:1] if title else []
        logo = logos[0] if logos else None

        print(f"  Auto-selected: {len(title)} title, {len(agenda)} agenda, "
              f"{len(section)} section, {len(closing)} closing"
              f"{', 1 logo' if logo else ''}")

    return {
        'title': title,
        'agenda': agenda,
        'section': section,
        'closing': closing,
        'logo': logo,
    }


def save_brand_images(selections, output_dir):
    """Save selected images to brand package title-assets/ directory.

    Returns updated paths dict for brand.yaml.
    """
    assets_dir = os.path.join(output_dir, "title-assets")
    os.makedirs(assets_dir, exist_ok=True)

    paths = {
        'title_backgrounds': {},
        'agenda_backgrounds': {},
        'section_backgrounds': {},
        'closing_backgrounds': {},
        'logo': None,
    }

    for role, key in [('title', 'title_backgrounds'), ('agenda', 'agenda_backgrounds'),
                      ('section', 'section_backgrounds'), ('closing', 'closing_backgrounds')]:
        images = selections.get(role, [])
        for i, img in enumerate(images):
            suffix = f"-{i+1:02d}" if len(images) > 1 else ""
            filename = f"{role}-bg{suffix}.{img['ext']}"
            filepath = os.path.join(assets_dir, filename)
            with open(filepath, 'wb') as f:
                f.write(img['blob'])
            label = f"alt{i}" if i > 0 else "default"
            paths[key][label] = f"title-assets/{filename}"

    if selections.get('logo'):
        logo = selections['logo']
        filename = f"logo.{logo['ext']}"
        filepath = os.path.join(assets_dir, filename)
        with open(filepath, 'wb') as f:
            f.write(logo['blob'])
        paths['logo'] = f"title-assets/{filename}"

    return paths


# ---------------------------------------------------------------------------
# Brand.yaml generation
# ---------------------------------------------------------------------------

def theme_to_brand_yaml(name, theme, canvas=None, template_rel="template.pptx",
                        image_paths=None):
    """Generate brand.yaml content from an extracted/derived theme."""
    colors = theme['colors']
    fonts = theme['fonts']
    canvas = canvas or {"width_inches": 10.0, "height_inches": 5.625}

    data = {
        "name": name,
        "colors": {
            "primary": colors.get('dk1', '#1A365D'),
            "secondary": colors.get('dk2', '#3182CE'),
            "accent": colors.get('accent3', '#63B3ED'),
            "background_light": colors.get('lt2', '#F5F5F5'),
            "background_card": colors.get('accent2', '#FAFAFA'),
            "text_dark": "#333333",
            "text_gray": "#888888",
            "white": "#FFFFFF",
            "divider": colors.get('accent6', '#D0D0D0'),
            "link": colors.get('hlink', colors.get('dk1', '#1A365D')),
            "neutral_bg": "#F5F5F5",
            "gradient_start": colors.get('dk1', '#1A365D'),
            "gradient_end": colors.get('dk2', '#3182CE'),
            "gradient_light": colors.get('accent2', '#E0E0E0'),
            "gradient_ring_light": colors.get('accent2', '#E0E0E0'),
            "staircase_end": colors.get('accent6', '#D0D0D0'),
            "spoke_end": colors.get('accent6', '#D0D0D0'),
            "venn_tertiary": colors.get('accent3', '#B0B0B0'),
            "green": "#2E7D32",
            "amber": "#F57F17",
            "red": "#C62828",
        },
        "fonts": {
            "heading": fonts.get('heading', 'Arial'),
            "body": fonts.get('body', 'Arial'),
            "mono": "Courier New",
        },
        "canvas": canvas,
        "template": template_rel,
    }

    # Image paths
    if image_paths:
        if image_paths.get('title_backgrounds'):
            data['title_backgrounds'] = image_paths['title_backgrounds']
        else:
            data['title_backgrounds'] = {"default": "title-assets/title-bg.jpg"}
        if image_paths.get('agenda_backgrounds'):
            data['agenda_backgrounds'] = image_paths['agenda_backgrounds']
        else:
            data['agenda_backgrounds'] = {"default": "title-assets/agenda-left.jpg"}
        if image_paths.get('section_backgrounds'):
            data['section_backgrounds'] = image_paths['section_backgrounds']
        if image_paths.get('closing_backgrounds'):
            data['closing_backgrounds'] = image_paths['closing_backgrounds']
        if image_paths.get('logo'):
            data['logo'] = image_paths['logo']
    else:
        data['title_backgrounds'] = {"default": "title-assets/title-bg.jpg"}
        data['agenda_backgrounds'] = {"default": "title-assets/agenda-left.jpg"}

    return data


# ---------------------------------------------------------------------------
# Background generation (fallback when no images extracted)
# ---------------------------------------------------------------------------

def generate_backgrounds(brand_dir, primary_hex, secondary_hex):
    """Generate gradient background images for title and agenda slides."""
    try:
        from PIL import Image, ImageDraw, ImageFilter
    except ImportError:
        print("    Skipping background generation (PIL not installed)")
        return

    def hex_to_rgb(h):
        h = h.lstrip('#')
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

    def darken(rgb, factor=0.5):
        return tuple(int(c * (1 - factor)) for c in rgb)

    assets_dir = os.path.join(brand_dir, "title-assets")
    os.makedirs(assets_dir, exist_ok=True)

    primary = hex_to_rgb(primary_hex)
    secondary = hex_to_rgb(secondary_hex)
    dark = darken(primary)

    # Title background
    w, h = 1920, 1080
    img = Image.new("RGB", (w, h))
    draw = ImageDraw.Draw(img)
    for y in range(h):
        t = y / h
        if t < 0.5:
            t2 = t * 2
            r = int(dark[0] + (primary[0] - dark[0]) * t2)
            g = int(dark[1] + (primary[1] - dark[1]) * t2)
            b = int(dark[2] + (primary[2] - dark[2]) * t2)
        else:
            t2 = (t - 0.5) * 2
            r = int(primary[0] + (secondary[0] - primary[0]) * t2 * 0.3)
            g = int(primary[1] + (secondary[1] - primary[1]) * t2 * 0.3)
            b = int(primary[2] + (secondary[2] - primary[2]) * t2 * 0.3)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    img = img.filter(ImageFilter.GaussianBlur(radius=2))
    img.save(os.path.join(assets_dir, "title-bg.jpg"), quality=90)

    # Agenda left panel
    aw, ah = 760, 760
    img2 = Image.new("RGB", (aw, ah))
    draw2 = ImageDraw.Draw(img2)
    light = tuple(min(255, c + (255 - c) * 2 // 5) for c in primary)
    for y in range(ah):
        t = y / ah
        r = int(light[0] + (primary[0] - light[0]) * t)
        g = int(light[1] + (primary[1] - light[1]) * t)
        b = int(light[2] + (primary[2] - light[2]) * t)
        draw2.line([(0, y), (aw, y)], fill=(r, g, b))
    img2 = img2.filter(ImageFilter.GaussianBlur(radius=1))
    img2.save(os.path.join(assets_dir, "agenda-left.jpg"), quality=90)


# ---------------------------------------------------------------------------
# Main onboarding pipeline
# ---------------------------------------------------------------------------

def onboard_brand(name, template_path=None, corpus_dir=None, output_dir=None,
                  brand_yaml_path=None, interactive=True):
    """Run the full brand onboarding pipeline."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = output_dir or os.path.join(script_dir, "brands", name)
    os.makedirs(output_dir, exist_ok=True)

    print(f"\n{'='*60}")
    print(f"  Brand Onboarding: {name}")
    print(f"{'='*60}")

    theme = None
    canvas = None
    image_paths = None

    # =================================================================
    # STEP 1: Extract or derive theme
    # =================================================================
    print(f"\nSTEP 1: Theme")

    if template_path and os.path.isfile(template_path):
        print(f"  Source: template ({os.path.basename(template_path)})")
        theme = extract_theme(template_path)
        prs = Presentation(template_path)
        canvas = extract_canvas_size(prs)
        print(f"  dk1={theme['colors'].get('dk1')} heading={theme['fonts']['heading']}")

        if corpus_dir:
            corpus_colors = extract_corpus_colors(corpus_dir)
            if corpus_colors:
                print(f"  Corpus: {len(corpus_colors)} dominant colors enrichment")

        dest_template = os.path.join(output_dir, "template.pptx")
        shutil.copy2(template_path, dest_template)

    elif corpus_dir and os.path.isdir(corpus_dir):
        print(f"  Source: corpus ({corpus_dir})")
        corpus_colors = extract_corpus_colors(corpus_dir)
        top = list(corpus_colors.keys())
        colors = {
            "primary": top[0] if len(top) >= 1 else "#1A365D",
            "secondary": top[1] if len(top) >= 2 else "#3182CE",
            "accent": top[2] if len(top) >= 3 else "#63B3ED",
        }
        theme = derive_theme(colors, {"heading": "Arial", "body": "Arial"})
        print(f"  Synthesized: dk1={theme['colors']['dk1']}")

        generic_template = os.path.join(script_dir, "brands", "generic", "template.pptx")
        dest_template = os.path.join(output_dir, "template.pptx")
        shutil.copy2(generic_template, dest_template)
        inject_theme(dest_template, theme)

    elif brand_yaml_path and os.path.isfile(brand_yaml_path):
        print(f"  Source: brand.yaml")
        with open(brand_yaml_path) as f:
            data = yaml.safe_load(f)
        theme = derive_theme(data.get('colors', {}), data.get('fonts', {}))
        print(f"  Derived: dk1={theme['colors']['dk1']} heading={theme['fonts']['heading']}")

        generic_template = os.path.join(script_dir, "brands", "generic", "template.pptx")
        dest_template = os.path.join(output_dir, "template.pptx")
        shutil.copy2(generic_template, dest_template)
        inject_theme(dest_template, theme)

    else:
        existing_yaml = os.path.join(output_dir, "brand.yaml")
        if os.path.isfile(existing_yaml):
            print(f"  Source: existing brand.yaml")
            with open(existing_yaml) as f:
                data = yaml.safe_load(f)
            theme = derive_theme(data.get('colors', {}), data.get('fonts', {}))
            print(f"  Derived: dk1={theme['colors']['dk1']}")

            generic_template = os.path.join(script_dir, "brands", "generic", "template.pptx")
            dest_template = os.path.join(output_dir, "template.pptx")
            shutil.copy2(generic_template, dest_template)
            inject_theme(dest_template, theme)
        else:
            print("  ERROR: No template, corpus, or brand.yaml provided")
            return None

    if not theme:
        return None

    # Store theme
    save_theme(theme, os.path.join(output_dir, "theme.json"))
    print(f"  Saved theme.json")

    # =================================================================
    # STEP 2: Extract and select images
    # =================================================================
    print(f"\nSTEP 2: Images")

    has_images = False
    if template_path and os.path.isfile(template_path):
        print(f"  Scanning template for images...")
        extracted = extract_all_images(template_path)
        non_icon = [img for img in extracted['images']
                    if img['category'] not in ('icon', 'small')]
        if non_icon:
            has_images = True
            selections = select_brand_images(extracted, interactive=interactive)
            image_paths = save_brand_images(selections, output_dir)
            n_saved = sum(len(v) if isinstance(v, dict) else (1 if v else 0)
                         for v in image_paths.values())
            print(f"  Saved {n_saved} images to title-assets/")

    if corpus_dir and os.path.isdir(corpus_dir) and not has_images:
        print(f"  Scanning corpus for images...")
        extracted = extract_corpus_images(corpus_dir)
        non_icon = [img for img in extracted['images']
                    if img['category'] not in ('icon', 'small')]
        if non_icon:
            has_images = True
            selections = select_brand_images(extracted, interactive=interactive)
            image_paths = save_brand_images(selections, output_dir)
            print(f"  Saved images from corpus")

    if not has_images:
        print(f"  No images found — generating gradient backgrounds")
        primary = theme['colors'].get('dk1', '#1A365D')
        secondary = theme['colors'].get('dk2', '#3182CE')
        generate_backgrounds(output_dir, primary, secondary)
        image_paths = {
            'title_backgrounds': {"default": "title-assets/title-bg.jpg"},
            'agenda_backgrounds': {"default": "title-assets/agenda-left.jpg"},
        }

    # =================================================================
    # STEP 3: Generate brand.yaml
    # =================================================================
    print(f"\nSTEP 3: brand.yaml")

    brand_yaml_out = os.path.join(output_dir, "brand.yaml")
    if not os.path.isfile(brand_yaml_out):
        brand_data = theme_to_brand_yaml(name, theme, canvas,
                                         image_paths=image_paths)
        with open(brand_yaml_out, 'w') as f:
            yaml.dump(brand_data, f, default_flow_style=False, sort_keys=False)
        print(f"  Generated brand.yaml")
    else:
        # Update image paths in existing brand.yaml
        with open(brand_yaml_out) as f:
            existing = yaml.safe_load(f)
        if image_paths:
            for key in ('title_backgrounds', 'agenda_backgrounds',
                        'section_backgrounds', 'closing_backgrounds', 'logo'):
                if image_paths.get(key):
                    existing[key] = image_paths[key]
            with open(brand_yaml_out, 'w') as f:
                yaml.dump(existing, f, default_flow_style=False, sort_keys=False)
        print(f"  Updated brand.yaml with image paths")

    # =================================================================
    # STEP 4: Validation build
    # =================================================================
    print(f"\nSTEP 4: Validation")

    test_yaml = os.path.join(script_dir, "examples", "showcase-generic.yaml")
    if os.path.isfile(test_yaml):
        result = subprocess.run(
            [sys.executable, os.path.join(script_dir, "build_deck.py"),
             test_yaml, "--output", os.path.join(output_dir, "validation.pptx")],
            capture_output=True, text=True,
        )
        if result.returncode == 0:
            print(f"  Validation build: SUCCESS")
        else:
            print(f"  Validation build: FAILED")
            if result.stderr:
                print(f"    {result.stderr[:300]}")

    print(f"\n{'='*60}")
    print(f"  Brand package: {output_dir}/")
    print(f"  Contents: brand.yaml, theme.json, template.pptx, title-assets/")
    print(f"  Use: add 'brand: {name}' to your YAML deck definition")
    print(f"{'='*60}\n")
    return output_dir


def main():
    parser = argparse.ArgumentParser(
        description="Onboard a new brand — extract theme and images from template/corpus")
    parser.add_argument("--name", required=True,
                        help="Brand name (used as directory name)")
    parser.add_argument("--template", default=None,
                        help="Path to PPTX template file")
    parser.add_argument("--corpus", default=None,
                        help="Directory of example PPTX files")
    parser.add_argument("--output", default=None,
                        help="Output directory (default: brands/<name>/)")
    parser.add_argument("--auto", action="store_true",
                        help="Non-interactive mode (auto-select images)")
    args = parser.parse_args()

    onboard_brand(args.name, args.template, args.corpus, args.output,
                  interactive=not args.auto)


if __name__ == "__main__":
    main()
