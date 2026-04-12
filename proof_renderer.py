#!/usr/bin/env python3
"""
Visual Proof Renderer — Renders each slide of a PPTX as a PNG for visual QA.

Draws all shapes at their PPTX coordinates using PIL, overlays red outlines
around elements with layout issues (overlaps, containment violations, margin
violations), and produces a summary JSON of all issues found.

Usage (standalone):
    python3 proof_renderer.py <input.pptx> [--output-dir <dir>] [--dpi 150]

Usage (as module):
    from proof_renderer import render_proof_images
    issues = render_proof_images("deck.pptx", "proof_dir/")
"""

import io
import json
import os
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

DEFAULT_DPI = 150
SLIDE_WIDTH_INCHES = 10.0
SLIDE_HEIGHT_INCHES = 5.625  # 16:9
IMG_W = int(SLIDE_WIDTH_INCHES * DEFAULT_DPI)   # 1500
IMG_H = int(SLIDE_HEIGHT_INCHES * DEFAULT_DPI)  # ~843

MARGIN_THRESHOLD_IN = 0.5   # elements closer than this to slide edge = violation
OVERLAP_THRESHOLD_IN = 0.02  # vertical overlap below this is ignored (sub-pixel)

RED = (220, 38, 38)
GREEN = (34, 197, 94)
ISSUE_OUTLINE_COLOR = (220, 38, 38, 200)  # semi-transparent red
ISSUE_OUTLINE_WIDTH = 2

# Background color for proof rendering
SLIDE_BG = (255, 255, 255)


# ---------------------------------------------------------------------------
# Font helpers — best-effort system font loading
# ---------------------------------------------------------------------------

_font_cache = {}


def _load_font(size, bold=False):
    """Load a TrueType font at the given size, falling back gracefully."""
    key = (size, bold)
    if key in _font_cache:
        return _font_cache[key]

    candidates = []
    if bold:
        candidates = [
            "/System/Library/Fonts/Helvetica.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
            "/System/Library/Fonts/HelveticaNeue.ttc",
        ]
    else:
        candidates = [
            "/System/Library/Fonts/Helvetica.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
            "/System/Library/Fonts/HelveticaNeue.ttc",
        ]

    for path in candidates:
        try:
            font = ImageFont.truetype(path, size)
            _font_cache[key] = font
            return font
        except (OSError, IOError):
            continue

    # Ultimate fallback
    font = ImageFont.load_default()
    _font_cache[key] = font
    return font


# ---------------------------------------------------------------------------
# Coordinate conversion helpers
# ---------------------------------------------------------------------------

def emu_to_px(emu, dpi=DEFAULT_DPI):
    """Convert EMU (English Metric Units) to pixels at given DPI."""
    inches = emu / 914400.0
    return inches * dpi


def emu_to_in(emu):
    """Convert EMU to inches."""
    return emu / 914400.0


def in_to_px(inches, dpi=DEFAULT_DPI):
    """Convert inches to pixels."""
    return inches * dpi


# ---------------------------------------------------------------------------
# Color extraction from python-pptx objects
# ---------------------------------------------------------------------------

def _rgb_tuple(rgb_color):
    """Extract (r, g, b) tuple from a pptx RGBColor or return a default."""
    if rgb_color is None:
        return None
    try:
        return (rgb_color[0], rgb_color[1], rgb_color[2])
    except (TypeError, IndexError):
        pass
    try:
        s = str(rgb_color)
        if len(s) == 6:
            return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except (ValueError, TypeError):
        pass
    return None


def _shape_fill_color(shape):
    """Try to extract the solid fill color from a shape."""
    try:
        fill = shape.fill
        if fill is not None and fill.type is not None:
            fc = fill.fore_color
            if fc is not None and fc.rgb is not None:
                return _rgb_tuple(fc.rgb)
    except (AttributeError, TypeError, KeyError):
        pass
    return None


def _shape_line_color(shape):
    """Try to extract the line/border color from a shape."""
    try:
        ln = shape.line
        if ln is not None and ln.color is not None and ln.color.rgb is not None:
            return _rgb_tuple(ln.color.rgb)
    except (AttributeError, TypeError, KeyError):
        pass
    return None


def _shape_line_width(shape):
    """Get line width in pixels (approximate)."""
    try:
        ln = shape.line
        if ln is not None and ln.width is not None:
            # line.width is in EMU
            return max(1, int(emu_to_px(ln.width)))
    except (AttributeError, TypeError):
        pass
    return 1


def _get_text_color(shape):
    """Get color of the first run's font in a shape."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    return _rgb_tuple(run.font.color.rgb)
    except (AttributeError, TypeError):
        pass
    return (51, 51, 51)  # default dark gray


def _get_font_size_pt(shape, slide=None):
    """Get font size in points, resolving inherited sizes from layout/master.

    Resolution order:
    1. Explicit run.font.size on the text run
    2. Paragraph-level defRPr sz
    3. Layout placeholder defRPr sz (for placeholder shapes)
    4. Master placeholder defRPr sz
    5. Fallback: 10pt
    """
    # Step 1: Check explicit run font size
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return run.font.size.pt
    except (AttributeError, TypeError):
        pass

    # Step 2: Check paragraph-level default
    try:
        for para in shape.text_frame.paragraphs:
            pPr = para._p.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
            if pPr is not None:
                sz = pPr.get('sz')
                if sz:
                    return int(sz) / 100  # hundredths of a point to points
    except (AttributeError, TypeError):
        pass

    # Step 3: Check layout placeholder default (for placeholder shapes)
    if slide is not None:
        try:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                layout = slide.slide_layout
                for layout_ph in layout.placeholders:
                    if layout_ph.placeholder_format.idx == ph_idx:
                        # Check layout placeholder's defRPr
                        from lxml import etree
                        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                        for defRPr in layout_ph._element.iter(f'{ns}defRPr'):
                            sz = defRPr.get('sz')
                            if sz:
                                return int(sz) / 100
                        break
        except (AttributeError, TypeError, ValueError):
            pass

    # Step 4: Check master placeholder default
    if slide is not None:
        try:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                master = slide.slide_layout.slide_master
                for master_ph in master.placeholders:
                    if master_ph.placeholder_format.idx == ph_idx:
                        from lxml import etree
                        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                        for defRPr in master_ph._element.iter(f'{ns}defRPr'):
                            sz = defRPr.get('sz')
                            if sz:
                                return int(sz) / 100
                        break
        except (AttributeError, TypeError, ValueError):
            pass

    return 10  # final fallback


# ---------------------------------------------------------------------------
# Shape rendering
# ---------------------------------------------------------------------------

def _get_geom_type(shape):
    """Get the preset geometry type string (e.g., 'ellipse', 'arc', 'donut')."""
    try:
        from pptx.oxml.ns import qn
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is not None:
            prstGeom = spPr.find(qn('a:prstGeom'))
            if prstGeom is not None:
                return prstGeom.get('prst', 'rect')
    except Exception:
        pass
    # Fallback: check auto_shape_type
    try:
        from pptx.enum.shapes import MSO_SHAPE
        ast = shape.auto_shape_type
        if ast == MSO_SHAPE.OVAL:
            return "ellipse"
        elif ast == MSO_SHAPE.ROUNDED_RECTANGLE:
            return "roundRect"
        elif ast == MSO_SHAPE.RIGHT_ARROW:
            return "rightArrow"
        elif ast == MSO_SHAPE.CHEVRON:
            return "chevron"
    except Exception:
        pass
    return "rect"


def _render_rectangle(draw, shape, dpi):
    """Draw a rectangle/rounded rectangle shape."""
    x = emu_to_px(shape.left, dpi)
    y = emu_to_px(shape.top, dpi)
    w = emu_to_px(shape.width, dpi)
    h = emu_to_px(shape.height, dpi)

    fill = _shape_fill_color(shape)
    line_color = _shape_line_color(shape)
    line_w = _shape_line_width(shape)

    # Determine if this is a rounded rectangle
    is_rounded = False
    try:
        shape_name = shape.auto_shape_type
        if shape_name is not None:
            from pptx.enum.shapes import MSO_SHAPE
            if shape_name == MSO_SHAPE.ROUNDED_RECTANGLE:
                is_rounded = True
    except (AttributeError, ValueError):
        pass

    if is_rounded and min(w, h) > 10:
        radius = min(10, int(min(w, h) * 0.15))
        _draw_rounded_rect(draw, x, y, x + w, y + h, radius, fill, line_color, line_w)
    else:
        if fill:
            draw.rectangle([x, y, x + w, y + h], fill=fill)
        if line_color:
            draw.rectangle([x, y, x + w, y + h], outline=line_color, width=line_w)

    # Draw text if present
    if hasattr(shape, "text") and shape.text.strip():
        _render_text_in_box(draw, shape, x, y, w, h, dpi)


def _draw_rounded_rect(draw, x0, y0, x1, y1, radius, fill, outline, line_w):
    """Draw a rounded rectangle using arcs and rectangles."""
    r = radius
    if fill:
        # Fill the main body
        draw.rectangle([x0 + r, y0, x1 - r, y1], fill=fill)
        draw.rectangle([x0, y0 + r, x1, y1 - r], fill=fill)
        # Fill corners
        draw.pieslice([x0, y0, x0 + 2*r, y0 + 2*r], 180, 270, fill=fill)
        draw.pieslice([x1 - 2*r, y0, x1, y0 + 2*r], 270, 360, fill=fill)
        draw.pieslice([x0, y1 - 2*r, x0 + 2*r, y1], 90, 180, fill=fill)
        draw.pieslice([x1 - 2*r, y1 - 2*r, x1, y1], 0, 90, fill=fill)

    if outline:
        draw.arc([x0, y0, x0 + 2*r, y0 + 2*r], 180, 270, fill=outline, width=line_w)
        draw.arc([x1 - 2*r, y0, x1, y0 + 2*r], 270, 360, fill=outline, width=line_w)
        draw.arc([x0, y1 - 2*r, x0 + 2*r, y1], 90, 180, fill=outline, width=line_w)
        draw.arc([x1 - 2*r, y1 - 2*r, x1, y1], 0, 90, fill=outline, width=line_w)
        draw.line([x0 + r, y0, x1 - r, y0], fill=outline, width=line_w)
        draw.line([x0 + r, y1, x1 - r, y1], fill=outline, width=line_w)
        draw.line([x0, y0 + r, x0, y1 - r], fill=outline, width=line_w)
        draw.line([x1, y0 + r, x1, y1 - r], fill=outline, width=line_w)


def _get_text_alignment(shape):
    """Get text alignment from the first paragraph (0=left, 1=center, 2=right)."""
    try:
        from pptx.enum.text import PP_ALIGN
        for para in shape.text_frame.paragraphs:
            if para.alignment == PP_ALIGN.CENTER:
                return 1
            elif para.alignment == PP_ALIGN.RIGHT:
                return 2
            elif para.alignment is not None:
                return 0
    except (AttributeError, TypeError):
        pass
    return 0  # default left


def _render_text_in_box(draw, shape, box_x, box_y, box_w, box_h, dpi, slide=None):
    """Render text content within a bounding box, respecting alignment."""
    text_color = _get_text_color(shape)
    font_size_pt = _get_font_size_pt(shape, slide=slide)
    alignment = _get_text_alignment(shape)  # 0=left, 1=center, 2=right

    # Convert pt to pixels: 1pt = 1/72 inch
    font_size_px = int(font_size_pt * dpi / 72)
    font_size_px = max(8, min(font_size_px, 150))  # clamp (150px = 72pt at 150dpi)

    font = _load_font(font_size_px)

    # Padding inside the box — keep tight to avoid eating vertical space
    pad = max(3, min(8, int(box_w * 0.02)))

    # Get all paragraph texts
    lines = []
    try:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)
    except (AttributeError, TypeError):
        text = shape.text.strip()
        if text:
            lines = text.split("\n")

    if not lines:
        return

    # Word-wrap lines to fit the box width
    max_text_w = box_w - 2 * pad
    wrapped = []
    for line in lines:
        wrapped.extend(_word_wrap(draw, line, font, max_text_w))

    # Draw each line
    line_height = font_size_px * 1.3
    # Reduce padding for tight boxes to avoid clipping single-line text
    if len(wrapped) <= 2 and box_h < line_height * 2 + pad * 2:
        pad = max(1, min(pad, int((box_h - line_height) / 2)))
    y = box_y + pad
    for wline in wrapped:
        if y + line_height > box_y + box_h + line_height * 0.3:
            break  # don't overflow the box visually (allow slight overshoot)

        # Calculate x position based on alignment
        if alignment == 1:  # center
            try:
                text_w = draw.textlength(wline, font=font)
            except (AttributeError, TypeError):
                text_w = len(wline) * font_size_px * 0.6
            x = box_x + (box_w - text_w) / 2
        elif alignment == 2:  # right
            try:
                text_w = draw.textlength(wline, font=font)
            except (AttributeError, TypeError):
                text_w = len(wline) * font_size_px * 0.6
            x = box_x + box_w - pad - text_w
        else:  # left
            x = box_x + pad

        draw.text((x, y), wline, fill=text_color, font=font)
        y += line_height


def _word_wrap(draw, text, font, max_width):
    """Simple word-wrap that fits text within max_width pixels."""
    if max_width <= 0:
        return [text]

    words = text.split()
    if not words:
        return [""]

    lines = []
    current = words[0]

    for word in words[1:]:
        test = current + " " + word
        try:
            bbox = draw.textbbox((0, 0), test, font=font)
            tw = bbox[2] - bbox[0]
        except (AttributeError, TypeError):
            tw = len(test) * 7  # rough fallback
        if tw <= max_width:
            current = test
        else:
            lines.append(current)
            current = word

    lines.append(current)
    return lines


def _render_line(draw, shape, dpi):
    """Draw a line/connector shape."""
    x1 = emu_to_px(shape.left, dpi)
    y1 = emu_to_px(shape.top, dpi)
    x2 = x1 + emu_to_px(shape.width, dpi)
    y2 = y1 + emu_to_px(shape.height, dpi)

    color = _shape_line_color(shape) or (100, 100, 100)
    width = _shape_line_width(shape)

    draw.line([x1, y1, x2, y2], fill=color, width=width)


def _render_image(draw, img_canvas, shape, prs, dpi):
    """Paste an embedded image onto the canvas at the shape's position."""
    x = int(emu_to_px(shape.left, dpi))
    y = int(emu_to_px(shape.top, dpi))
    w = int(emu_to_px(shape.width, dpi))
    h = int(emu_to_px(shape.height, dpi))

    try:
        blob = shape.image.blob
        img = Image.open(io.BytesIO(blob))
        img = img.convert("RGBA")
        img = img.resize((max(1, w), max(1, h)), Image.LANCZOS)
        img_canvas.paste(img, (x, y), img)
        return True
    except Exception:
        # Draw a placeholder box for images we can't extract
        draw.rectangle([x, y, x + w, y + h], outline=(180, 180, 180), width=1)
        draw.line([x, y, x + w, y + h], fill=(180, 180, 180), width=1)
        draw.line([x + w, y, x, y + h], fill=(180, 180, 180), width=1)
        return False


def _render_textbox(draw, shape, dpi, slide=None):
    """Draw a text box shape (no background fill unless explicitly set)."""
    x = emu_to_px(shape.left, dpi)
    y = emu_to_px(shape.top, dpi)
    w = emu_to_px(shape.width, dpi)
    h = emu_to_px(shape.height, dpi)

    fill = _shape_fill_color(shape)
    if fill:
        draw.rectangle([x, y, x + w, y + h], fill=fill)

    line_color = _shape_line_color(shape)
    if line_color:
        draw.rectangle([x, y, x + w, y + h], outline=line_color, width=1)

    if hasattr(shape, "text") and shape.text.strip():
        _render_text_in_box(draw, shape, x, y, w, h, dpi, slide=slide)


def _render_group(draw, img_canvas, group_shape, prs, dpi):
    """Recursively render shapes within a group."""
    try:
        for shape in group_shape.shapes:
            _render_shape(draw, img_canvas, shape, prs, dpi)
    except (AttributeError, TypeError):
        pass


def _render_table(draw, shape, dpi):
    """Render a TABLE shape with header row, cell text, and grid lines."""
    try:
        table = shape.table
    except Exception:
        return

    x0 = emu_to_px(shape.left, dpi)
    y0 = emu_to_px(shape.top, dpi)

    # Calculate column positions
    col_positions = [x0]
    for ci in range(len(table.columns)):
        col_positions.append(col_positions[-1] + emu_to_px(table.columns[ci].width, dpi))

    # Calculate row positions
    row_positions = [y0]
    for ri in range(len(table.rows)):
        row_positions.append(row_positions[-1] + emu_to_px(table.rows[ri].height, dpi))

    # Header colors
    HEADER_BG = (95, 1, 111)      # Purple
    HEADER_TEXT = (255, 255, 255)  # White
    BODY_TEXT = (51, 51, 51)       # Dark
    STRIPE_BG = (255, 240, 248)   # Light pink stripe
    GRID_COLOR = (208, 192, 216)  # Light purple grid

    # Default fonts (used when cell has no explicit size)
    header_font_px = int(10 * dpi / 72)  # 10pt default for headers
    body_font_px = int(9 * dpi / 72)     # 9pt default for body cells

    for ri in range(len(table.rows)):
        for ci in range(len(table.columns)):
            cx = col_positions[ci]
            cy = row_positions[ri]
            cx2 = col_positions[ci + 1]
            cy2 = row_positions[ri + 1]

            # Cell background
            if ri == 0:
                draw.rectangle([cx, cy, cx2, cy2], fill=HEADER_BG)
            elif ri % 2 == 0:
                draw.rectangle([cx, cy, cx2, cy2], fill=STRIPE_BG)

            # Cell border
            draw.rectangle([cx, cy, cx2, cy2], outline=GRID_COLOR, width=1)

            # Cell text
            try:
                cell = table.cell(ri, ci)
                text = cell.text.strip()
                if text:
                    color = HEADER_TEXT if ri == 0 else BODY_TEXT
                    # Read actual font size from cell runs
                    cell_font_px = header_font_px if ri == 0 else body_font_px
                    cell_bold = (ri == 0)
                    try:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    cell_font_px = int(run.font.size.pt * dpi / 72)
                                if run.font.bold:
                                    cell_bold = True
                                # Check font color
                                try:
                                    if run.font.color.type is not None and run.font.color.rgb:
                                        color = _rgb_tuple(run.font.color.rgb)
                                except (AttributeError, TypeError):
                                    pass
                                break
                            break
                    except (AttributeError, TypeError):
                        pass
                    cell_font_px = max(8, min(cell_font_px, 48))
                    font = _load_font(cell_font_px, cell_bold)
                    pad = 3
                    draw.text((cx + pad, cy + pad), text[:40], fill=color, font=font)
            except Exception:
                pass


def _render_shape(draw, img_canvas, shape, prs, dpi, slide=None):
    """Dispatch rendering for a single shape based on its type."""
    shape_type = shape.shape_type

    if shape_type == MSO_SHAPE_TYPE.PICTURE or shape_type == 13:
        _render_image(draw, img_canvas, shape, prs, dpi)
    elif shape_type == MSO_SHAPE_TYPE.LINE or shape_type == 9:
        _render_line(draw, shape, dpi)
    elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape_type == 17:
        _render_textbox(draw, shape, dpi, slide=slide)
    elif shape_type == MSO_SHAPE_TYPE.GROUP or shape_type == 6:
        _render_group(draw, img_canvas, shape, prs, dpi)
    elif shape_type == MSO_SHAPE_TYPE.FREEFORM or shape_type == 5:
        _render_rectangle(draw, shape, dpi)
    elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape_type == 1:
        # Determine actual geometry type from OOXML
        geom_type = _get_geom_type(shape)
        x = int(emu_to_px(shape.left, dpi))
        y = int(emu_to_px(shape.top, dpi))
        w = int(emu_to_px(shape.width, dpi))
        h = int(emu_to_px(shape.height, dpi))
        fill = _shape_fill_color(shape)
        line_color = _shape_line_color(shape)
        line_w = _shape_line_width(shape)

        if geom_type in ("ellipse", "oval"):
            # Draw circle/ellipse
            if fill:
                draw.ellipse([x, y, x + w, y + h], fill=fill)
            if line_color:
                draw.ellipse([x, y, x + w, y + h], outline=line_color, width=line_w)
            if hasattr(shape, "text") and shape.text.strip():
                _render_text_in_box(draw, shape, x, y, w, h, dpi)

        elif geom_type == "donut":
            # Draw donut as thick circle outline
            ring_w = max(int(min(w, h) * 0.15), 3)
            color = fill or line_color or (200, 200, 200)
            draw.ellipse([x, y, x + w, y + h], outline=color, width=ring_w)

        elif geom_type == "arc":
            # Draw arc as partial ellipse
            color = fill or line_color or (100, 100, 100)
            arc_w = max(int(min(w, h) * 0.12), 3)
            # Try to get arc angles from adjustments
            try:
                adj1 = shape.adjustments[0]  # start fraction
                adj2 = shape.adjustments[1]  # end fraction
                start_deg = int(adj1 * 360) - 90
                end_deg = int(adj2 * 360) - 90
            except Exception:
                start_deg = 0
                end_deg = 270
            draw.arc([x, y, x + w, y + h], start_deg, end_deg,
                     fill=color, width=arc_w)

        elif geom_type in ("chevron", "homePlate"):
            # Draw chevron (pentagon arrow shape)
            if fill:
                notch = w // 5
                pts = [(x, y), (x + w - notch, y), (x + w, y + h // 2),
                       (x + w - notch, y + h), (x, y + h), (x + notch, y + h // 2)]
                draw.polygon(pts, fill=fill)
            if hasattr(shape, "text") and shape.text.strip():
                _render_text_in_box(draw, shape, x, y, w, h, dpi)

        elif geom_type == "rightArrow":
            # Draw arrow shape
            color = fill or (187, 187, 187)
            shaft_h = max(h // 3, 2)
            head_w = min(w // 2, h // 2)
            shaft_y = y + (h - shaft_h) // 2
            # Shaft
            draw.rectangle([x, shaft_y, x + w - head_w, shaft_y + shaft_h], fill=color)
            # Arrowhead
            pts = [(x + w - head_w, y), (x + w, y + h // 2), (x + w - head_w, y + h)]
            draw.polygon(pts, fill=color)

        else:
            # Default: render as rectangle
            _render_rectangle(draw, shape, dpi)
    elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER or shape_type == 14:
        # Check if this placeholder contains an image (picture placeholder)
        has_image = False
        try:
            if hasattr(shape, "image") and shape.image and shape.image.blob:
                has_image = True
        except Exception:
            pass
        if has_image:
            _render_image(draw, img_canvas, shape, prs, dpi)
        else:
            _render_textbox(draw, shape, dpi, slide=slide)
    elif shape_type == 19 or (hasattr(shape, "has_table") and shape.has_table):
        # TABLE shape
        _render_table(draw, shape, dpi)
    else:
        # Fallback: try to render as a rectangle/text box
        try:
            if hasattr(shape, "has_table") and shape.has_table:
                _render_table(draw, shape, dpi)
            elif hasattr(shape, "text") and shape.text.strip():
                _render_textbox(draw, shape, dpi)
            elif hasattr(shape, "width") and shape.width and shape.height:
                x = emu_to_px(shape.left, dpi)
                y = emu_to_px(shape.top, dpi)
                w = emu_to_px(shape.width, dpi)
                h = emu_to_px(shape.height, dpi)
                fill = _shape_fill_color(shape)
                if fill:
                    draw.rectangle([x, y, x + w, y + h], fill=fill)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Slide background extraction
# ---------------------------------------------------------------------------

def _get_slide_background(slide, prs, dpi):
    """Try to get the slide background color."""
    # Try slide background
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc is not None and fc.rgb is not None:
                return _rgb_tuple(fc.rgb)
    except (AttributeError, TypeError, KeyError):
        pass

    # Try slide layout background
    try:
        bg = slide.slide_layout.background
        fill = bg.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc is not None and fc.rgb is not None:
                return _rgb_tuple(fc.rgb)
    except (AttributeError, TypeError, KeyError):
        pass

    return SLIDE_BG


# ---------------------------------------------------------------------------
# Issue detection
# ---------------------------------------------------------------------------

def _collect_text_elements(slide, dpi):
    """Collect all text-bearing elements with their bounding boxes (in inches)."""
    elements = []
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        if shape.shape_type == 1:
            # Check if it has meaningful text (skip background rects with no text)
            if not shape.text.strip():
                continue

        top_in = emu_to_in(shape.top)
        left_in = emu_to_in(shape.left)
        bottom_in = top_in + emu_to_in(shape.height)
        right_in = left_in + emu_to_in(shape.width)

        elements.append({
            "shape": shape,
            "top": top_in,
            "left": left_in,
            "bottom": bottom_in,
            "right": right_in,
            "text": shape.text[:40],
            "name": shape.name,
        })
    return elements


def _detect_overlaps(elements):
    """Detect text elements that overlap vertically by more than the threshold."""
    issues = []
    flagged_shapes = set()

    for i, a in enumerate(elements):
        for j, b in enumerate(elements[i + 1:], start=i + 1):
            h_overlap = a["left"] < b["right"] and b["left"] < a["right"]
            v_overlap = a["top"] < b["bottom"] and b["top"] < a["bottom"]
            if h_overlap and v_overlap:
                v_amount = min(a["bottom"], b["bottom"]) - max(a["top"], b["top"])
                if v_amount > OVERLAP_THRESHOLD_IN:
                    issues.append({
                        "type": "overlap",
                        "description": (
                            f'"{a["text"]}" overlaps "{b["text"]}" '
                            f"by {v_amount:.2f}\" vertically"
                        ),
                    })
                    flagged_shapes.add(id(a["shape"]))
                    flagged_shapes.add(id(b["shape"]))

    return issues, flagged_shapes


def _detect_margin_violations(elements, slide_w_in, slide_h_in):
    """Detect text elements within MARGIN_THRESHOLD_IN of the slide edge."""
    issues = []
    flagged_shapes = set()

    for el in elements:
        violations = []
        if el["left"] < MARGIN_THRESHOLD_IN:
            violations.append(f"left={el['left']:.2f}\"")
        if el["top"] < MARGIN_THRESHOLD_IN:
            violations.append(f"top={el['top']:.2f}\"")
        if el["right"] > slide_w_in - MARGIN_THRESHOLD_IN:
            violations.append(f"right={el['right']:.2f}\" (slide={slide_w_in:.2f}\")")
        if el["bottom"] > slide_h_in - MARGIN_THRESHOLD_IN:
            violations.append(f"bottom={el['bottom']:.2f}\" (slide={slide_h_in:.2f}\")")

        if violations:
            issues.append({
                "type": "margin",
                "description": (
                    f'"{el["text"]}" too close to edge: '
                    + ", ".join(violations)
                ),
            })
            flagged_shapes.add(id(el["shape"]))

    return issues, flagged_shapes


def _detect_containment_violations(slide, dpi):
    """Detect child shapes that extend beyond a parent group's bounds."""
    issues = []
    flagged_shapes = set()

    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP and shape.shape_type != 6:
            continue
        try:
            parent_top = emu_to_in(shape.top)
            parent_left = emu_to_in(shape.left)
            parent_bottom = parent_top + emu_to_in(shape.height)
            parent_right = parent_left + emu_to_in(shape.width)

            for child in shape.shapes:
                child_top = emu_to_in(child.top)
                child_left = emu_to_in(child.left)
                child_bottom = child_top + emu_to_in(child.height)
                child_right = child_left + emu_to_in(child.width)

                if (child_left < parent_left - 0.02 or
                        child_top < parent_top - 0.02 or
                        child_right > parent_right + 0.02 or
                        child_bottom > parent_bottom + 0.02):
                    issues.append({
                        "type": "containment",
                        "description": (
                            f'"{child.name}" extends beyond group '
                            f'"{shape.name}"'
                        ),
                    })
                    flagged_shapes.add(id(child))
        except (AttributeError, TypeError):
            pass

    return issues, flagged_shapes


# ---------------------------------------------------------------------------
# Issue overlay rendering
# ---------------------------------------------------------------------------

def _draw_issue_outlines(draw, flagged_shapes, slide, dpi):
    """Draw red outlines around all flagged shapes."""
    for shape in slide.shapes:
        if id(shape) in flagged_shapes:
            x = int(emu_to_px(shape.left, dpi))
            y = int(emu_to_px(shape.top, dpi))
            w = int(emu_to_px(shape.width, dpi))
            h = int(emu_to_px(shape.height, dpi))
            draw.rectangle(
                [x - 1, y - 1, x + w + 1, y + h + 1],
                outline=RED,
                width=ISSUE_OUTLINE_WIDTH,
            )

        # Also check group children
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP or shape.shape_type == 6:
            try:
                for child in shape.shapes:
                    if id(child) in flagged_shapes:
                        cx = int(emu_to_px(child.left, dpi))
                        cy = int(emu_to_px(child.top, dpi))
                        cw = int(emu_to_px(child.width, dpi))
                        ch = int(emu_to_px(child.height, dpi))
                        draw.rectangle(
                            [cx - 1, cy - 1, cx + cw + 1, cy + ch + 1],
                            outline=RED,
                            width=ISSUE_OUTLINE_WIDTH,
                        )
            except (AttributeError, TypeError):
                pass


def _draw_status_badge(draw, has_issues, img_w, img_h):
    """Draw a status badge in the top-right corner."""
    badge_w = 80
    badge_h = 28
    margin = 10
    x0 = img_w - badge_w - margin
    y0 = margin

    if has_issues:
        bg = RED
        label = "ISSUES"
    else:
        bg = GREEN
        label = "OK"

    # Badge background
    draw.rounded_rectangle(
        [x0, y0, x0 + badge_w, y0 + badge_h],
        radius=6,
        fill=bg,
    )

    # Badge text
    font = _load_font(16, bold=True)
    try:
        bbox = draw.textbbox((0, 0), label, font=font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
    except (AttributeError, TypeError):
        tw, th = len(label) * 9, 14

    tx = x0 + (badge_w - tw) // 2
    ty = y0 + (badge_h - th) // 2 - 1
    draw.text((tx, ty), label, fill=(255, 255, 255), font=font)


# ---------------------------------------------------------------------------
# Main rendering pipeline
# ---------------------------------------------------------------------------

def render_slide(slide, prs, dpi=DEFAULT_DPI):
    """Render a single slide to a PIL Image and return (image, issues_list).

    Returns:
        (PIL.Image, list[dict]): The rendered image and a list of issue dicts.
    """
    # Calculate image dimensions from actual slide size
    slide_w_in = emu_to_in(prs.slide_width)
    slide_h_in = emu_to_in(prs.slide_height)
    img_w = int(slide_w_in * dpi)
    img_h = int(slide_h_in * dpi)

    # Create canvas
    bg_color = _get_slide_background(slide, prs, dpi)
    img = Image.new("RGBA", (img_w, img_h), bg_color + (255,))
    draw = ImageDraw.Draw(img)

    # Render all shapes in z-order (shapes list is bottom-to-top)
    for shape in slide.shapes:
        _render_shape(draw, img, shape, prs, dpi, slide=slide)

    # --- Issue detection ---
    text_elements = _collect_text_elements(slide, dpi)

    all_issues = []
    all_flagged = set()

    overlap_issues, overlap_flagged = _detect_overlaps(text_elements)
    all_issues.extend(overlap_issues)
    all_flagged.update(overlap_flagged)

    margin_issues, margin_flagged = _detect_margin_violations(
        text_elements, slide_w_in, slide_h_in
    )
    all_issues.extend(margin_issues)
    all_flagged.update(margin_flagged)

    contain_issues, contain_flagged = _detect_containment_violations(slide, dpi)
    all_issues.extend(contain_issues)
    all_flagged.update(contain_flagged)

    # Draw red outlines on flagged elements
    if all_flagged:
        _draw_issue_outlines(draw, all_flagged, slide, dpi)

    # Draw status badge
    _draw_status_badge(draw, bool(all_issues), img_w, img_h)

    # Convert to RGB for PNG output (drop alpha)
    final = Image.new("RGB", (img_w, img_h), bg_color)
    final.paste(img, (0, 0), img)

    return final, all_issues


def render_proof_images(pptx_path, output_dir=None, dpi=DEFAULT_DPI):
    """Render all slides of a PPTX as proof PNG images.

    Args:
        pptx_path: Path to the PPTX file.
        output_dir: Directory for output PNGs. Defaults to <pptx_stem>-proof/.
        dpi: Output resolution (default 150).

    Returns:
        dict: Summary with keys "slides" (list of per-slide info) and
              "total_issues" (int).
    """
    pptx_path = os.path.abspath(pptx_path)
    if not os.path.isfile(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    if output_dir is None:
        stem = Path(pptx_path).stem
        output_dir = os.path.join(os.path.dirname(pptx_path), f"{stem}-proof")

    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    summary = {"file": pptx_path, "slides": [], "total_issues": 0}

    total_slides = len(prs.slides)
    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        print(f"  Rendering slide {slide_num}/{total_slides}...", end="", flush=True)

        img, issues = render_slide(slide, prs, dpi)

        # Save PNG
        png_name = f"slide-{slide_num:02d}.png"
        png_path = os.path.join(output_dir, png_name)
        img.save(png_path, "PNG")

        slide_info = {
            "slide": slide_num,
            "image": png_name,
            "layout": slide.slide_layout.name,
            "issues": issues,
        }
        summary["slides"].append(slide_info)
        summary["total_issues"] += len(issues)

        status = "OK" if not issues else f"{len(issues)} issue(s)"
        print(f" {status} -> {png_name}")

    # Write summary JSON
    json_path = os.path.join(output_dir, "proof-summary.json")
    with open(json_path, "w") as f:
        json.dump(summary, f, indent=2)
    print(f"  Summary: {json_path}")

    return summary


# ---------------------------------------------------------------------------
# Hi-Fi Renderer — LibreOffice headless for pixel-perfect PNGs
# ---------------------------------------------------------------------------

_SOFFICE_PATHS = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/usr/bin/soffice",
    "/usr/local/bin/soffice",
    "soffice",  # on PATH
]


def _find_soffice():
    """Find the LibreOffice soffice binary, or return None."""
    import shutil
    for path in _SOFFICE_PATHS:
        if os.path.isfile(path):
            return path
    # Try PATH
    found = shutil.which("soffice")
    return found


def is_hifi_available():
    """Check if LibreOffice is installed for hi-fi rendering."""
    return _find_soffice() is not None


def render_proof_images_hifi(pptx_path, output_dir=None, dpi=150):
    """Render all slides as pixel-perfect PNGs using LibreOffice headless.

    LibreOffice renders the PPTX exactly as PowerPoint would — correct fonts,
    shapes, transparency, gradients, and Unicode glyphs.

    Args:
        pptx_path: Path to the PPTX file.
        output_dir: Directory for output PNGs. Defaults to <pptx_stem>-proof/.
        dpi: Not directly supported by LO export, but we resize output to match.

    Returns:
        dict: Summary compatible with render_proof_images() output.

    Raises:
        RuntimeError: If LibreOffice is not installed.
    """
    import subprocess
    import glob
    import tempfile

    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError(
            "LibreOffice not found. Install with: brew install --cask libreoffice\n"
            "Falling back to PIL renderer."
        )

    pptx_path = os.path.abspath(pptx_path)
    if not os.path.isfile(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    if output_dir is None:
        stem = Path(pptx_path).stem
        output_dir = os.path.join(os.path.dirname(pptx_path), f"{stem}-proof")
    os.makedirs(output_dir, exist_ok=True)

    # LibreOffice exports to a temp dir then we rename
    with tempfile.TemporaryDirectory() as tmpdir:
        # Export PPTX → PNG using LibreOffice headless
        cmd = [
            soffice,
            "--headless",
            "--convert-to", "png",
            "--outdir", tmpdir,
            pptx_path,
        ]
        print(f"  LibreOffice rendering: {os.path.basename(pptx_path)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice export failed: {result.stderr}")

        # LO produces a single PNG of the first slide only with --convert-to png.
        # For per-slide export, we need to use the PDF route and split.
        # Strategy: export to PDF, then PDF → per-page PNGs via PIL

        # Try PDF export first (gives us all slides)
        pdf_cmd = [
            soffice,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", tmpdir,
            pptx_path,
        ]
        pdf_result = subprocess.run(pdf_cmd, capture_output=True, text=True, timeout=120)

        pdf_stem = Path(pptx_path).stem
        pdf_path = os.path.join(tmpdir, f"{pdf_stem}.pdf")

        summary = {"file": pptx_path, "slides": [], "total_issues": 0}

        if os.path.exists(pdf_path):
            # Convert PDF pages to PNGs
            try:
                # Try pdf2image (poppler-based) first
                from pdf2image import convert_from_path
                images = convert_from_path(pdf_path, dpi=dpi)
                for si, img in enumerate(images):
                    slide_num = si + 1
                    png_name = f"slide-{slide_num:02d}.png"
                    png_path = os.path.join(output_dir, png_name)
                    img.save(png_path, "PNG")
                    print(f"  Slide {slide_num}/{len(images)} -> {png_name}")
                    summary["slides"].append({
                        "slide": slide_num,
                        "image": png_name,
                        "layout": "",
                        "issues": [],
                    })
            except ImportError:
                # Fallback: use PIL to read PDF (limited but works for simple cases)
                # Or use Ghostscript if available
                try:
                    gs_cmd = ["gs", "-dNOPAUSE", "-dBATCH", "-sDEVICE=png16m",
                              f"-r{dpi}", f"-sOutputFile={output_dir}/slide-%02d.png",
                              pdf_path]
                    gs_result = subprocess.run(gs_cmd, capture_output=True, text=True, timeout=120)
                    if gs_result.returncode == 0:
                        pngs = sorted(glob.glob(os.path.join(output_dir, "slide-*.png")))
                        for si, png_path in enumerate(pngs):
                            slide_num = si + 1
                            png_name = f"slide-{slide_num:02d}.png"
                            print(f"  Slide {slide_num}/{len(pngs)} -> {png_name}")
                            summary["slides"].append({
                                "slide": slide_num,
                                "image": png_name,
                                "layout": "",
                                "issues": [],
                            })
                    else:
                        print(f"    Ghostscript failed: {gs_result.stderr[:200]}")
                except FileNotFoundError:
                    # No pdf2image, no ghostscript — use the single PNG from LO
                    single_png = os.path.join(tmpdir, f"{pdf_stem}.png")
                    if os.path.exists(single_png):
                        import shutil
                        dest = os.path.join(output_dir, "slide-01.png")
                        shutil.copy2(single_png, dest)
                        print(f"  Slide 1/1 -> slide-01.png (single-page fallback)")
                        summary["slides"].append({
                            "slide": 1, "image": "slide-01.png",
                            "layout": "", "issues": [],
                        })
                    print("    Install pdf2image or ghostscript for per-slide PNGs:")
                    print("    pip install pdf2image && brew install poppler")
        else:
            # PDF export failed — use the single PNG
            single_png = os.path.join(tmpdir, f"{pdf_stem}.png")
            if os.path.exists(single_png):
                import shutil
                dest = os.path.join(output_dir, "slide-01.png")
                shutil.copy2(single_png, dest)
                summary["slides"].append({
                    "slide": 1, "image": "slide-01.png",
                    "layout": "", "issues": [],
                })

    # Write summary JSON
    json_path = os.path.join(output_dir, "proof-summary.json")
    with open(json_path, "w") as f:
        json.dump(summary, f, indent=2)
    print(f"  Summary: {json_path}")

    return summary


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    import argparse

    parser = argparse.ArgumentParser(
        description="Visual Proof Renderer — render PPTX slides as PNG for QA"
    )
    parser.add_argument("pptx_file", help="Input PPTX file")
    parser.add_argument(
        "--output-dir", "-o", default=None,
        help="Output directory for PNGs (default: <name>-proof/)"
    )
    parser.add_argument(
        "--dpi", type=int, default=DEFAULT_DPI,
        help=f"Output DPI (default: {DEFAULT_DPI})"
    )
    parser.add_argument(
        "--hifi", action="store_true",
        help="Use LibreOffice for pixel-perfect rendering (requires LibreOffice)"
    )
    args = parser.parse_args()

    print(f"Rendering proof images for: {args.pptx_file}")
    if args.hifi:
        if not is_hifi_available():
            print("ERROR: LibreOffice not found. Install: brew install --cask libreoffice")
            sys.exit(1)
        print("  Mode: Hi-Fi (LibreOffice)")
        summary = render_proof_images_hifi(args.pptx_file, args.output_dir, args.dpi)
    else:
        print("  Mode: Fast (PIL)")
        summary = render_proof_images(args.pptx_file, args.output_dir, args.dpi)

    total = summary["total_issues"]
    slides = len(summary["slides"])
    if total == 0:
        print(f"\nDone! {slides} slide(s) rendered, no issues found.")
    else:
        print(f"\nDone! {slides} slide(s) rendered, {total} issue(s) found.")
        for si in summary["slides"]:
            if si["issues"]:
                print(f"  Slide {si['slide']}:")
                for iss in si["issues"]:
                    print(f"    [{iss['type']}] {iss['description']}")


if __name__ == "__main__":
    main()
