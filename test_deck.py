#!/usr/bin/env python3
"""
Deck Builder Test Harness — Build, verify, and preview decks.

Generates a proof report showing exactly what text landed in which placeholder
on each slide, flags any issues (empty required placeholders, overflow, etc.),
and opens the result for visual QA.

Usage:
    python3 test_deck.py <yaml_file> [--open] [--upload]

Outputs:
    <deck-name>-proof.md   — slide-by-slide verification report
    <deck-name>.pptx       — the built deck
"""

import argparse
import os
import subprocess
import sys

import yaml
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def emu_to_inches(emu):
    return round(emu / 914400, 2)


def generate_proof_report(pptx_path, output_path):
    """Generate a markdown proof report for visual QA."""
    prs = Presentation(pptx_path)
    lines = [
        f"# Deck Proof Report",
        f"",
        f"**File:** `{pptx_path}`",
        f"**Slides:** {len(prs.slides)}",
        f"**Slide size:** {emu_to_inches(prs.slide_width)}\" × {emu_to_inches(prs.slide_height)}\"",
        "",
        "---",
        "",
    ]

    issues = []

    def check_overlaps(slide, slide_num):
        """Check for overlapping text elements on a slide."""
        text_boxes = []
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue
            if shape.shape_type == 1:  # skip background rectangles
                continue
            top = emu_to_inches(shape.top)
            left = emu_to_inches(shape.left)
            bottom = top + emu_to_inches(shape.height)
            right = left + emu_to_inches(shape.width)
            text_boxes.append({
                "top": top, "left": left, "bottom": bottom, "right": right,
                "text": shape.text[:30],
            })

        overlaps = []
        for i, a in enumerate(text_boxes):
            for b in text_boxes[i+1:]:
                # Check if boxes overlap (both horizontally and vertically)
                h_overlap = a["left"] < b["right"] and b["left"] < a["right"]
                v_overlap = a["top"] < b["bottom"] and b["top"] < a["bottom"]
                if h_overlap and v_overlap:
                    # Calculate overlap amount
                    v_amount = min(a["bottom"], b["bottom"]) - max(a["top"], b["top"])
                    if v_amount > 0.02:  # ignore sub-pixel overlaps
                        overlaps.append(
                            f"Slide {slide_num}: \"{a['text']}\" overlaps \"{b['text']}\" "
                            f"by {v_amount:.2f}\" vertically"
                        )
        return overlaps

    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        layout_name = slide.slide_layout.name
        lines.append(f"## Slide {slide_num} — Layout: {layout_name}")
        lines.append("")

        placeholders = []
        images = []
        other_shapes = []

        for shape in slide.shapes:
            if shape.is_placeholder:
                idx = shape.placeholder_format.idx
                left = emu_to_inches(shape.left)
                top = emu_to_inches(shape.top)
                w = emu_to_inches(shape.width)
                h = emu_to_inches(shape.height)
                text = shape.text if hasattr(shape, "text") else ""
                placeholders.append({
                    "idx": idx, "left": left, "top": top,
                    "width": w, "height": h, "text": text,
                    "name": shape.name,
                })
            elif shape.shape_type == 13:  # PICTURE
                images.append({
                    "left": emu_to_inches(shape.left),
                    "top": emu_to_inches(shape.top),
                    "width": emu_to_inches(shape.width),
                    "height": emu_to_inches(shape.height),
                    "name": shape.name,
                })
            else:
                other_shapes.append(shape.name)

        # Sort placeholders by position (top then left)
        placeholders.sort(key=lambda p: (p["top"], p["left"]))

        # Categorize placeholders
        content_phs = [p for p in placeholders if p["top"] < 5.0]
        footer_phs = [p for p in placeholders if p["top"] >= 5.0]

        if content_phs:
            lines.append("**Content placeholders:**")
            lines.append("")
            for p in content_phs:
                text_preview = p["text"][:80].replace("\n", " ↵ ")
                if len(p["text"]) > 80:
                    text_preview += "…"
                status = "✓" if p["text"].strip() else "○"
                lines.append(
                    f"- {status} `idx={p['idx']}` at ({p['left']}\", {p['top']}\") "
                    f"{p['width']}\"×{p['height']}\" — "
                    f'"{text_preview}"'
                )

                # Check for potential overflow (text too long for small placeholders)
                text_len = len(p["text"])
                area = p["width"] * p["height"]
                if text_len > 0 and area < 0.5 and text_len > 5:
                    issue = f"Slide {slide_num}: idx={p['idx']} may overflow — {text_len} chars in {p['width']}\"×{p['height']}\" box"
                    issues.append(issue)
                    lines.append(f"  - ⚠️ **Possible overflow**: {text_len} chars in small area")

            lines.append("")

        if footer_phs:
            footer_texts = [f'idx={p["idx"]}: "{p["text"][:30]}"' for p in footer_phs]
            lines.append(f"**Footers:** {' | '.join(footer_texts)}")
            lines.append("")

        if images:
            for img in images:
                lines.append(
                    f"**Image:** ({img['left']}\", {img['top']}\") "
                    f"{img['width']}\"×{img['height']}\""
                )
            lines.append("")

        # Check for overlapping text elements
        overlaps = check_overlaps(slide, slide_num)
        if overlaps:
            lines.append("**⚠️ Overlapping elements:**")
            for ov in overlaps:
                lines.append(f"- {ov}")
                issues.append(ov)
            lines.append("")

        # Check for notes
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                lines.append(f"**Notes:** {notes[:100]}{'…' if len(notes) > 100 else ''}")
                lines.append("")
        except Exception:
            pass

        lines.append("---")
        lines.append("")

    # Summary
    if issues:
        lines.insert(7, "### ⚠️ Issues Found")
        lines.insert(8, "")
        for issue in issues:
            lines.insert(9, f"- {issue}")
        lines.insert(9 + len(issues), "")
        lines.insert(10 + len(issues), "---")
        lines.insert(11 + len(issues), "")

    with open(output_path, "w") as f:
        f.write("\n".join(lines))

    return issues


def main():
    parser = argparse.ArgumentParser(description="Deck Builder Test Harness")
    parser.add_argument("yaml_file", help="YAML deck definition file")
    parser.add_argument("--open", action="store_true",
                        help="Open the generated PPTX in default app")
    parser.add_argument("--upload", action="store_true",
                        help="Upload to Google Drive")
    parser.add_argument("--output-dir", "-o", default=None,
                        help="Output directory (default: same as YAML)")
    parser.add_argument("--proof-images", action="store_true",
                        help="Generate proof PNG images (LibreOffice headless, falls back to PIL)")
    parser.add_argument("--quick-proof", action="store_true",
                        help="Generate fast PIL-based proof images (structural checks only, no theme rendering)")
    parser.add_argument("--hifi", action="store_true",
                        help="(deprecated, same as --proof-images)")
    parser.add_argument("--ai-review", action="store_true",
                        help="Run Claude vision review on proof images")
    parser.add_argument("--strict", action="store_true",
                        help="Exit with error on critical QA issues")
    args = parser.parse_args()

    yaml_path = os.path.abspath(args.yaml_file)
    base_name = os.path.splitext(os.path.basename(yaml_path))[0]

    if args.output_dir:
        out_dir = os.path.abspath(args.output_dir)
    else:
        out_dir = os.path.dirname(yaml_path)

    os.makedirs(out_dir, exist_ok=True)

    pptx_path = os.path.join(out_dir, f"{base_name}.pptx")
    proof_path = os.path.join(out_dir, f"{base_name}-proof.md")

    # Step 1: Build the deck
    print("=" * 60)
    print("STEP 1: Building deck")
    print("=" * 60)

    build_script = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "build_deck.py")
    result = subprocess.run(
        [sys.executable, build_script, yaml_path, "--output", pptx_path],
        capture_output=False
    )
    if result.returncode != 0:
        print("BUILD FAILED")
        sys.exit(1)

    # Step 2: Generate proof report
    print()
    print("=" * 60)
    print("STEP 2: Generating proof report")
    print("=" * 60)

    issues = generate_proof_report(pptx_path, proof_path)
    print(f"  Proof report: {proof_path}")

    if issues:
        print(f"\n  ⚠️  {len(issues)} potential issue(s) found:")
        for issue in issues:
            print(f"    - {issue}")
    else:
        print("  ✓ No issues detected")

    # Step 2b: Generate proof images (optional)
    if args.proof_images or args.hifi or args.quick_proof:
        print()
        print("=" * 60)

        proof_img_dir = os.path.join(out_dir, f"{base_name}-proof")

        if args.quick_proof:
            # PIL-based fast structural proof (no theme rendering)
            print("STEP 2b: Generating quick proof images (PIL)")
            print("=" * 60)
            from proof_renderer import render_proof_images
            summary = render_proof_images(pptx_path, proof_img_dir)
        else:
            # LibreOffice headless — default for --proof-images and --hifi
            print("STEP 2b: Generating proof images (LibreOffice headless)")
            print("=" * 60)
            from proof_renderer import render_proof_images_hifi, is_hifi_available
            if is_hifi_available():
                summary = render_proof_images_hifi(pptx_path, proof_img_dir)
            else:
                print("  WARNING: LibreOffice not found, falling back to PIL")
                print("  Install with: brew install --cask libreoffice")
                from proof_renderer import render_proof_images
                summary = render_proof_images(pptx_path, proof_img_dir)

        total_visual = summary["total_issues"]
        if total_visual > 0:
            print(f"\n  ⚠️  {total_visual} visual issue(s) found in proof images")
        else:
            print("  ✓ No visual issues detected in proof images")

    # Step 3: Run unified QA pipeline
    print()
    print("=" * 60)
    print("STEP 3: Running QA pipeline")
    print("=" * 60)

    try:
        from qa_pipeline import QAPipeline
        from build_deck import BrandConfig

        proof_img_dir = os.path.join(out_dir, f"{base_name}-proof")
        # Load brand for QA pipeline (reads from YAML)
        with open(yaml_path) as _f:
            _deck_def = yaml.safe_load(_f)
        _brand = BrandConfig.load(_deck_def.get("brand")) if _deck_def else None
        qa = QAPipeline(pptx_path, proof_img_dir, yaml_path=yaml_path, brand=_brand)
        # Skip PIL visual rendering if LibreOffice was used (don't overwrite LO images)
        skip_visual = (args.proof_images or args.hifi) and not args.quick_proof
        qa_report = qa.run_all(ai_review=args.ai_review, skip_visual=skip_visual)
        json_path, md_path = qa.save_report(out_dir)
        s = qa_report.summary()

        print(f"  QA report: {md_path}")
        print(f"  Issues: {s['total_issues']} "
              f"({s['critical']} critical, {s['warnings']} warnings, {s['info']} info)")

        if s["blocked"]:
            print("  STATUS: BLOCKED — Critical issues found.")
            if args.strict:
                print("  Aborting due to --strict flag.")
                sys.exit(2)
        elif s["total_issues"] == 0:
            print("  STATUS: CLEAN")
        else:
            print(f"  STATUS: REVIEW — {s['warnings']} warning(s)")
    except ImportError as e:
        print(f"  QA pipeline not available: {e}")
    except Exception as e:
        print(f"  QA pipeline error: {e}")

    # Step 4: Open for visual QA
    if args.open:
        print()
        print("=" * 60)
        print("STEP 4: Opening for visual QA")
        print("=" * 60)
        subprocess.run(["open", pptx_path])
        print(f"  Opened: {pptx_path}")

    # Step 5: Upload if requested
    if args.upload:
        print()
        print("=" * 60)
        print("STEP 5: Uploading to Google Drive")
        print("=" * 60)

        title = base_name.replace("-", " ").replace("_", " ").title()
        # Read title from YAML if possible
        try:
            import yaml as _yaml
            with open(yaml_path) as f:
                deck_def = _yaml.safe_load(f)
            title = deck_def.get("title", title)
        except Exception:
            pass

        account = os.environ.get("GOG_ACCOUNT", "")
        if not account:
            print("Upload skipped: Set GOG_ACCOUNT environment variable to your Google Workspace email")
        else:
            subprocess.run([
                "gog", "-a", account,
                "drive", "upload", pptx_path,
                "--name", title,
                "--convert-to", "slides",
            ])

    print()
    print("Done!")
    print(f"  PPTX:  {pptx_path}")
    print(f"  Proof: {proof_path}")
    if args.proof_images:
        proof_img_dir = os.path.join(out_dir, f"{base_name}-proof")
        print(f"  Images: {proof_img_dir}/")


if __name__ == "__main__":
    main()
