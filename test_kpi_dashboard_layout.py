#!/usr/bin/env python3
"""
KPI Dashboard Layout Tests

Tests the structural correctness of the kpi_dashboard layout:
- Number textboxes span the full card width (so wide values don't wrap)
- Number textboxes are aligned with their parent card

Regression: claude-costs-deck slide 27 had a 4-card KPI row including
"+15-20%". The number textbox was sized to inner_w (card_w - 2*pad),
which fit the estimator's width prediction but overflowed in PowerPoint
where DM Sans + em-dash + percent sign render wider than the estimate.
The text wrapped to two lines, breaking the visual rhythm.

Fix: number textbox spans full card_w so centered text has maximum
horizontal room before any wrapping risk.

Run: python3 -m pytest test_kpi_dashboard_layout.py -v
  or: python3 test_kpi_dashboard_layout.py
"""

import os
import sys
import tempfile
from pathlib import Path

try:
    import pytest
except ImportError:
    pytest = None

try:
    from pptx import Presentation
    from pptx.util import Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


SCRIPT_DIR = Path(__file__).parent.resolve()
EMU_PER_INCH = 914400


# Fixture: 4-card KPI dashboard with values that historically wrapped:
# em-dashes (–), percent signs (%), plus signs (+), and dollar signs ($).
WIDE_VALUES_YAML = """\
title: "KPI Dashboard Wide Values Test"
brand: generic
slides:
- layout: kpi_dashboard
  headline: "Test KPIs"
  metrics:
  - number: "140"
    label: "Heads removed"
  - number: "$14M"
    label: "Annual payroll saved"
  - number: "~5–7%"
    label: "Output lost"
  - number: "+15–20%"
    label: "Net effective output"
"""


def _build_fixture(tmp_dir):
    """Build the fixture deck and return path to .pptx."""
    yaml_path = tmp_dir / "fixture.yaml"
    yaml_path.write_text(WIDE_VALUES_YAML)
    pptx_path = tmp_dir / "fixture.pptx"

    sys.path.insert(0, str(SCRIPT_DIR))
    try:
        import build_deck
        build_deck.build_deck(str(yaml_path), str(pptx_path))
    finally:
        sys.path.pop(0)
    return pptx_path


def _find_card_backgrounds(slide):
    """Find ROUNDED_RECTANGLE shapes (KPI card backgrounds), sorted by x."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    cards = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        try:
            ast = shape.auto_shape_type
        except (AttributeError, ValueError):
            continue
        # Rounded rectangle = MSO_SHAPE.ROUNDED_RECTANGLE (5)
        if ast == 5:
            cards.append(shape)
    return sorted(cards, key=lambda s: s.left)


def _find_textboxes_in_card(slide, card):
    """Find all textboxes whose center falls within the card's bounds."""
    cx_min = card.left
    cx_max = card.left + card.width
    cy_min = card.top
    cy_max = card.top + card.height

    textboxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Center of the textbox
        center_x = shape.left + shape.width // 2
        center_y = shape.top + shape.height // 2
        if cx_min <= center_x <= cx_max and cy_min <= center_y <= cy_max:
            textboxes.append(shape)
    return textboxes


def test_kpi_number_textbox_spans_full_card_width():
    """Number textbox must span the full card width so wide values
    (em-dashes, percent signs, dollar signs) don't risk wrapping.

    Regression: number textbox previously sized at card_w - 2*pad,
    causing "+15-20%" to wrap in PowerPoint despite fitting the estimator's
    prediction in LibreOffice.
    """
    if not HAS_PPTX:
        if pytest:
            pytest.skip("python-pptx not available")
        return

    with tempfile.TemporaryDirectory() as td:
        pptx = _build_fixture(Path(td))
        prs = Presentation(pptx)
        slide = prs.slides[0]
        cards = _find_card_backgrounds(slide)

    assert len(cards) == 4, f"Expected 4 card backgrounds, found {len(cards)}"

    violations = []
    for i, card in enumerate(cards):
        textboxes = _find_textboxes_in_card(slide, card)
        if not textboxes:
            continue

        # Find the number textbox: it's the larger-font one; we look for
        # the textbox whose width is closest to card_w (and not the small
        # label at the bottom).
        # Sort by font size descending — number is biggest.
        biggest = None
        biggest_size = 0
        for tb in textboxes:
            for para in tb.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        sz = run.font.size.pt
                        if sz > biggest_size:
                            biggest_size = sz
                            biggest = tb

        if biggest is None:
            continue

        card_w_in = card.width / EMU_PER_INCH
        num_w_in = biggest.width / EMU_PER_INCH

        # Number textbox should span (close to) full card width.
        # Allow 0.05" rounding tolerance.
        if num_w_in < card_w_in - 0.05:
            violations.append((
                i + 1,
                num_w_in,
                card_w_in,
                biggest.text_frame.text.strip()
            ))

    assert not violations, (
        f"\n{len(violations)} KPI card(s) have number textbox narrower than card:\n"
        + "\n".join(
            f"  Card {n}: number textbox {nw:.2f}\" "
            f"vs card {cw:.2f}\" — value: \"{v}\""
            for n, nw, cw, v in violations
        )
        + "\n\n  Fix: number textbox should use full card_w (not card_w - 2*pad)\n"
        "  so wide values like \"+15–20%\" don't wrap in PowerPoint."
    )


def test_kpi_number_textbox_word_wrap_disabled():
    """Number textbox must have word_wrap = False so wide values don't
    wrap regardless of textbox width.

    Even with full card width, very long numbers should auto-scale font
    rather than wrap to multiple lines.
    """
    if not HAS_PPTX:
        if pytest:
            pytest.skip("python-pptx not available")
        return

    with tempfile.TemporaryDirectory() as td:
        pptx = _build_fixture(Path(td))
        prs = Presentation(pptx)
        slide = prs.slides[0]
        cards = _find_card_backgrounds(slide)

    violations = []
    for i, card in enumerate(cards):
        textboxes = _find_textboxes_in_card(slide, card)
        if not textboxes:
            continue
        # The number textbox is the largest-font one
        biggest = None
        biggest_size = 0
        for tb in textboxes:
            for para in tb.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt > biggest_size:
                        biggest_size = run.font.size.pt
                        biggest = tb
        if biggest is None:
            continue

        if biggest.text_frame.word_wrap is True:
            violations.append((i + 1, biggest.text_frame.text.strip()))

    assert not violations, (
        f"\n{len(violations)} KPI number textbox(es) have word_wrap enabled:\n"
        + "\n".join(f"  Card {n}: \"{v}\"" for n, v in violations)
        + "\n\n  Fix: tf.word_wrap = False on the number textbox."
    )


# ─── CLI runner ─────────────────────────────────────────────


if __name__ == "__main__":
    tests = [
        test_kpi_number_textbox_spans_full_card_width,
        test_kpi_number_textbox_word_wrap_disabled,
    ]
    failures = 0
    for t in tests:
        name = t.__name__
        try:
            t()
            print(f"  PASS  {name}")
        except AssertionError as e:
            failures += 1
            print(f"  FAIL  {name}")
            print(str(e))
        except Exception as e:
            failures += 1
            print(f"  ERROR {name}: {e}")
    sys.exit(1 if failures else 0)
