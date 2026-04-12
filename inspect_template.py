#!/usr/bin/env python3
"""
Template Inspector — Dumps every placeholder's index, position, size, type,
and default text/formatting for all layouts in a PPTX template.

Outputs:
  1. template-inspect.md  — full markdown catalog of all layouts + placeholders
  2. template-visual-map.pptx — one slide per layout with labeled placeholder outlines
  3. template-inspect.json — machine-readable data for deck builder

Usage:
    python3 inspect_template.py [--template path/to/template.pptx]
"""

import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE


def emu_to_inches(emu):
    """Convert EMU to inches, rounded to 2 decimal places."""
    return round(emu / 914400, 2)


def emu_to_cm(emu):
    """Convert EMU to centimeters, rounded to 2 decimal places."""
    return round(emu / 360000, 2)


def get_font_info(paragraph):
    """Extract font info from the first run of a paragraph."""
    info = {}
    if paragraph.runs:
        run = paragraph.runs[0]
        font = run.font
        if font.name:
            info["font_name"] = font.name
        if font.size:
            info["font_size_pt"] = font.size.pt
        if font.bold is not None:
            info["bold"] = font.bold
        if font.italic is not None:
            info["italic"] = font.italic
        if font.color and font.color.rgb:
            info["color"] = str(font.color.rgb)
    if paragraph.alignment is not None:
        info["alignment"] = str(paragraph.alignment)
    return info


def inspect_layout(layout, layout_idx, slide_width, slide_height):
    """Inspect a single layout and return its data."""
    layout_data = {
        "index": layout_idx,
        "name": layout.name,
        "slide_width_in": emu_to_inches(slide_width),
        "slide_height_in": emu_to_inches(slide_height),
        "placeholders": [],
        "shapes": [],
    }

    # Inspect placeholders
    for ph in layout.placeholders:
        ph_data = {
            "idx": ph.placeholder_format.idx,
            "type": str(ph.placeholder_format.type),
            "name": ph.name,
            "left_in": emu_to_inches(ph.left),
            "top_in": emu_to_inches(ph.top),
            "width_in": emu_to_inches(ph.width),
            "height_in": emu_to_inches(ph.height),
            "left_emu": ph.left,
            "top_emu": ph.top,
            "width_emu": ph.width,
            "height_emu": ph.height,
        }

        # Get default text and formatting
        if hasattr(ph, "text_frame"):
            tf = ph.text_frame
            texts = []
            for para in tf.paragraphs:
                texts.append(para.text)
            ph_data["default_text"] = "\n".join(texts).strip()
            if tf.paragraphs:
                ph_data["font_info"] = get_font_info(tf.paragraphs[0])

        layout_data["placeholders"].append(ph_data)

    # Inspect non-placeholder shapes (background images, decorations, etc.)
    for shape in layout.shapes:
        if shape.is_placeholder:
            continue
        shape_data = {
            "name": shape.name,
            "shape_type": str(shape.shape_type) if shape.shape_type else "unknown",
            "left_in": emu_to_inches(shape.left),
            "top_in": emu_to_inches(shape.top),
            "width_in": emu_to_inches(shape.width),
            "height_in": emu_to_inches(shape.height),
        }
        layout_data["shapes"].append(shape_data)

    # Sort placeholders by index
    layout_data["placeholders"].sort(key=lambda p: p["idx"])

    return layout_data


def generate_markdown(all_layouts, output_path):
    """Generate a comprehensive markdown catalog."""
    lines = [
        "# Template — Layout & Placeholder Catalog",
        "",
        f"**Total layouts:** {len(all_layouts)}",
        f"**Slide size:** {all_layouts[0]['slide_width_in']}\" × {all_layouts[0]['slide_height_in']}\"",
        "",
        "---",
        "",
    ]

    # Table of contents
    lines.append("## Table of Contents")
    lines.append("")
    for layout in all_layouts:
        ph_count = len(layout["placeholders"])
        lines.append(f"- [{layout['index']:02d}. {layout['name']}](#{layout['index']:02d}-{layout['name'].lower().replace(' ', '-').replace('/', '-')}) — {ph_count} placeholders")
    lines.append("")
    lines.append("---")
    lines.append("")

    for layout in all_layouts:
        lines.append(f"## {layout['index']:02d}. {layout['name']}")
        lines.append("")

        if not layout["placeholders"]:
            lines.append("*No placeholders — fully decorated background only.*")
            lines.append("")
        else:
            lines.append(f"**Placeholders: {len(layout['placeholders'])}** | Non-placeholder shapes: {len(layout['shapes'])}")
            lines.append("")

            # Placeholder table
            lines.append("| idx | Name | Position (x,y) | Size (w×h) | Default Text | Font |")
            lines.append("|-----|------|-----------------|------------|--------------|------|")

            for ph in layout["placeholders"]:
                pos = f"({ph['left_in']}\", {ph['top_in']}\")"
                size = f"{ph['width_in']}\" × {ph['height_in']}\""
                text = ph.get("default_text", "").replace("\n", " ↵ ")[:60]
                if len(ph.get("default_text", "")) > 60:
                    text += "…"

                font_parts = []
                fi = ph.get("font_info", {})
                if fi.get("font_name"):
                    font_parts.append(fi["font_name"])
                if fi.get("font_size_pt"):
                    font_parts.append(f"{fi['font_size_pt']}pt")
                if fi.get("bold"):
                    font_parts.append("**B**")
                if fi.get("color"):
                    font_parts.append(f"#{fi['color']}")
                font_str = " ".join(font_parts)

                lines.append(f"| {ph['idx']} | {ph['name']} | {pos} | {size} | {text} | {font_str} |")

            lines.append("")

            # Spatial map (ASCII art showing rough positions)
            lines.append("**Spatial Map** (proportional positions on 10\" × 7.5\" slide):")
            lines.append("```")
            grid_w, grid_h = 60, 20
            grid = [[' ' for _ in range(grid_w)] for _ in range(grid_h)]

            # Draw border
            for x in range(grid_w):
                grid[0][x] = '─'
                grid[grid_h-1][x] = '─'
            for y in range(grid_h):
                grid[y][0] = '│'
                grid[y][grid_w-1] = '│'
            grid[0][0] = '┌'
            grid[0][grid_w-1] = '┐'
            grid[grid_h-1][0] = '└'
            grid[grid_h-1][grid_w-1] = '┘'

            sw = layout["slide_width_in"]
            sh = layout["slide_height_in"]

            for ph in layout["placeholders"]:
                # Map placeholder position to grid
                gx = max(1, min(grid_w-2, int(ph["left_in"] / sw * (grid_w-2)) + 1))
                gy = max(1, min(grid_h-2, int(ph["top_in"] / sh * (grid_h-2)) + 1))
                gw = max(1, int(ph["width_in"] / sw * (grid_w-2)))
                gh = max(1, int(ph["height_in"] / sh * (grid_h-2)))

                label = f"[{ph['idx']}]"

                # Draw corners of placeholder box
                for dx in range(gw):
                    x = gx + dx
                    if 1 <= x < grid_w-1:
                        if 1 <= gy < grid_h-1:
                            grid[gy][x] = '·'
                        ey = gy + gh
                        if 1 <= ey < grid_h-1:
                            grid[ey][x] = '·'

                for dy in range(gh):
                    y = gy + dy
                    if 1 <= y < grid_h-1:
                        if 1 <= gx < grid_w-1:
                            grid[y][gx] = '·'
                        ex = gx + gw
                        if 1 <= ex < grid_w-1:
                            grid[y][ex] = '·'

                # Place label inside the box
                ly = min(gy + 1, grid_h - 2)
                for ci, ch in enumerate(label):
                    lx = gx + 1 + ci
                    if 1 <= lx < grid_w-1 and 1 <= ly < grid_h-1:
                        grid[ly][lx] = ch

            for row in grid:
                lines.append("".join(row))
            lines.append("```")
            lines.append("")

        # Non-placeholder shapes summary
        if layout["shapes"]:
            lines.append(f"<details><summary>Non-placeholder shapes ({len(layout['shapes'])})</summary>")
            lines.append("")
            for s in layout["shapes"]:
                lines.append(f"- **{s['name']}** ({s['shape_type']}): ({s['left_in']}\", {s['top_in']}\") {s['width_in']}\" × {s['height_in']}\"")
            lines.append("")
            lines.append("</details>")
            lines.append("")

        lines.append("---")
        lines.append("")

    with open(output_path, "w") as f:
        f.write("\n".join(lines))
    print(f"  Markdown catalog: {output_path}")


def generate_visual_map(all_layouts, template_path, output_path):
    """Generate a PPTX with one slide per layout, showing labeled placeholder outlines."""
    from pptx.oxml.ns import qn

    prs = Presentation(template_path)

    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Colors for placeholder labels
    colors = [
        RGBColor(0xE7, 0x4C, 0x3C),  # red
        RGBColor(0x27, 0xAE, 0x60),  # green
        RGBColor(0x29, 0x80, 0xB9),  # blue
        RGBColor(0xF3, 0x9C, 0x12),  # orange
        RGBColor(0x8E, 0x44, 0xAD),  # purple
        RGBColor(0x16, 0xA0, 0x85),  # teal
        RGBColor(0xD3, 0x54, 0x00),  # burnt orange
        RGBColor(0x2C, 0x3E, 0x50),  # dark blue
        RGBColor(0xC0, 0x39, 0x2B),  # dark red
        RGBColor(0x1A, 0xBC, 0x9C),  # turquoise
        RGBColor(0xD4, 0xAC, 0x0D),  # gold
        RGBColor(0x7D, 0x3C, 0x98),  # dark purple
        RGBColor(0x2E, 0x86, 0xC1),  # medium blue
        RGBColor(0xA0, 0x42, 0x00),  # brown
        RGBColor(0x17, 0x80, 0x2E),  # dark green
        RGBColor(0xCA, 0x6F, 0x1E),  # dark orange
    ]

    for layout_data in all_layouts:
        layout_idx = layout_data["index"]
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # Add layout title banner at top
        banner = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.05), Inches(9.6), Inches(0.4)
        )
        tf = banner.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Layout #{layout_idx:02d}: {layout_data['name']}  |  {len(layout_data['placeholders'])} placeholders"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Semi-transparent dark background for banner
        from pptx.oxml.ns import qn as _qn
        fill = banner.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        # For each placeholder, draw a labeled outline
        for i, ph_data in enumerate(layout_data["placeholders"]):
            color = colors[i % len(colors)]

            left = ph_data["left_emu"]
            top = ph_data["top_emu"]
            width = ph_data["width_emu"]
            height = ph_data["height_emu"]

            # Draw outline rectangle
            from pptx.enum.shapes import MSO_SHAPE
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            rect.fill.background()  # transparent fill
            rect.line.color.rgb = color
            rect.line.width = Pt(2)

            # Add label near top-left of placeholder
            label_text = f"idx={ph_data['idx']}"
            if ph_data.get("default_text"):
                snippet = ph_data["default_text"][:25]
                label_text += f" \"{snippet}\""

            label_h = Inches(0.3)
            label_w = min(Inches(3), width)
            label = slide.shapes.add_textbox(left, max(0, top - Emu(int(label_h * 0.8))), label_w, label_h)
            ltf = label.text_frame
            ltf.word_wrap = False
            lp = ltf.paragraphs[0]
            lp.text = label_text
            lp.font.size = Pt(8)
            lp.font.bold = True
            lp.font.color.rgb = color

        # Add legend at bottom
        legend_top = Inches(6.8)
        legend_text = "  |  ".join(
            f"idx={ph['idx']}: {ph['name']}"
            for ph in layout_data["placeholders"]
        )
        if legend_text:
            legend = slide.shapes.add_textbox(
                Inches(0.2), legend_top, Inches(9.6), Inches(0.5)
            )
            ltf = legend.text_frame
            ltf.word_wrap = True
            lp = ltf.paragraphs[0]
            lp.text = legend_text
            lp.font.size = Pt(7)
            lp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(output_path)
    print(f"  Visual map PPTX: {output_path}")


def generate_json(all_layouts, output_path):
    """Write machine-readable JSON for the deck builder."""
    with open(output_path, "w") as f:
        json.dump(all_layouts, f, indent=2, default=str)
    print(f"  JSON data: {output_path}")


def main():
    import argparse
    parser = argparse.ArgumentParser(description="PPTX Template Inspector")
    parser.add_argument("--template", "-t",
                        default=os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                             "brands", "generic", "template.pptx"),
                        help="Path to template PPTX")
    parser.add_argument("--output-dir", "-o",
                        default=os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                             "inspection"),
                        help="Output directory")
    args = parser.parse_args()

    template_path = os.path.abspath(args.template)
    output_dir = os.path.abspath(args.output_dir)
    os.makedirs(output_dir, exist_ok=True)

    print(f"Inspecting template: {template_path}")
    print(f"Output directory: {output_dir}")
    print()

    prs = Presentation(template_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    print(f"Slide size: {emu_to_inches(slide_width)}\" × {emu_to_inches(slide_height)}\"")
    print(f"Layouts: {len(prs.slide_layouts)}")
    print()

    all_layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        ph_count = len(list(layout.placeholders))
        shape_count = len([s for s in layout.shapes if not s.is_placeholder])
        print(f"  [{i:02d}] {layout.name:30s}  {ph_count} placeholders, {shape_count} other shapes")
        layout_data = inspect_layout(layout, i, slide_width, slide_height)
        all_layouts.append(layout_data)

    print()

    # Generate outputs
    generate_markdown(all_layouts, os.path.join(output_dir, "template-inspect.md"))
    generate_json(all_layouts, os.path.join(output_dir, "template-inspect.json"))
    generate_visual_map(all_layouts, template_path, os.path.join(output_dir, "template-visual-map.pptx"))

    print()
    print("Done! Open template-visual-map.pptx to see labeled placeholder outlines on each layout.")


if __name__ == "__main__":
    main()
