#!/usr/bin/env python3
"""Extract and classify images from a PPTX template for brand onboarding.

Scans all slides in a PPTX file, extracts every image, and classifies them
by size into functional roles:
  - Full-bleed backgrounds (title, closing, section divider candidates)
  - Panel images (agenda side panel candidates)
  - Logos (small, repeated across many slides)
  - Icons (very small, catalog items)
  - Other (photos, diagrams, decorative)

Usage:
    python3 extract_brand_images.py --template path/to/template.pptx
    python3 extract_brand_images.py --template path/to/template.pptx --output-dir /tmp/extracted
"""

import argparse
import hashlib
import os
import sys
from collections import Counter, defaultdict
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx required")
    sys.exit(1)


# Size thresholds (inches) for classification
FULL_BLEED_MIN_W = 8.0
FULL_BLEED_MIN_H = 4.0
PANEL_MIN_W = 2.5
PANEL_MIN_H = 2.5
PANEL_MAX_W = 6.0
LOGO_MAX_W = 3.0
LOGO_MAX_H = 1.5
ICON_MAX_W = 0.6
ICON_MAX_H = 0.6


def extract_all_images(pptx_path):
    """Extract all images from a PPTX, deduplicate, and classify.

    Returns dict with:
        'images': list of image dicts (hash, blob, content_type, sizes, slides, category)
        'stats': summary counts
    """
    prs = Presentation(pptx_path)
    canvas_w = prs.slide_width / 914400
    canvas_h = prs.slide_height / 914400

    # Collect all image instances with deduplication by content hash
    image_map = {}  # hash -> image info
    image_slides = defaultdict(list)  # hash -> list of (slide_num, w_in, h_in)

    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        for shape in slide.shapes:
            if not hasattr(shape, 'image'):
                continue
            try:
                img = shape.image
                blob = img.blob
                content_type = img.content_type
                w_in = shape.width / 914400
                h_in = shape.height / 914400

                img_hash = hashlib.md5(blob).hexdigest()[:12]

                if img_hash not in image_map:
                    image_map[img_hash] = {
                        'hash': img_hash,
                        'blob': blob,
                        'content_type': content_type,
                        'size_bytes': len(blob),
                        'ext': 'png' if 'png' in content_type else 'jpg',
                    }

                image_slides[img_hash].append({
                    'slide': slide_num,
                    'width': round(w_in, 1),
                    'height': round(h_in, 1),
                })
            except Exception:
                continue

    # Classify each unique image
    results = []
    for img_hash, info in image_map.items():
        instances = image_slides[img_hash]
        # Use the largest instance dimensions for classification
        max_w = max(inst['width'] for inst in instances)
        max_h = max(inst['height'] for inst in instances)
        num_slides = len(set(inst['slide'] for inst in instances))
        slide_nums = sorted(set(inst['slide'] for inst in instances))

        # Classify
        if max_w <= ICON_MAX_W and max_h <= ICON_MAX_H:
            category = 'icon'
        elif (max_w <= LOGO_MAX_W and max_h <= LOGO_MAX_H
              and num_slides >= 3 and info['size_bytes'] < 200_000):
            category = 'logo'
        elif max_w >= FULL_BLEED_MIN_W and max_h >= FULL_BLEED_MIN_H:
            category = 'full-bleed'
        elif (PANEL_MIN_W <= max_w <= PANEL_MAX_W
              and max_h >= PANEL_MIN_H):
            category = 'panel'
        elif max_w > LOGO_MAX_W or max_h > LOGO_MAX_H:
            category = 'other'
        else:
            category = 'small'

        results.append({
            'hash': img_hash,
            'category': category,
            'content_type': info['content_type'],
            'ext': info['ext'],
            'size_kb': info['size_bytes'] // 1024,
            'width': max_w,
            'height': max_h,
            'num_slides': num_slides,
            'slides': slide_nums[:5],  # first 5 for display
            'blob': info['blob'],
        })

    # Sort by category then size
    cat_order = {'full-bleed': 0, 'panel': 1, 'logo': 2, 'other': 3, 'icon': 4, 'small': 5}
    results.sort(key=lambda x: (cat_order.get(x['category'], 9), -x['size_kb']))

    stats = Counter(r['category'] for r in results)
    return {'images': results, 'stats': dict(stats), 'canvas': (canvas_w, canvas_h)}


def save_extracted_images(results, output_dir):
    """Save extracted images to disk, organized by category."""
    os.makedirs(output_dir, exist_ok=True)

    for cat in ['full-bleed', 'panel', 'logo', 'other']:
        cat_dir = os.path.join(output_dir, cat)
        os.makedirs(cat_dir, exist_ok=True)

    saved = []
    for img in results['images']:
        if img['category'] in ('icon', 'small'):
            continue  # Skip icons and tiny images

        cat = img['category']
        filename = f"{cat}-{img['hash']}.{img['ext']}"
        filepath = os.path.join(output_dir, cat, filename)
        with open(filepath, 'wb') as f:
            f.write(img['blob'])
        saved.append({'path': filepath, **{k: v for k, v in img.items() if k != 'blob'}})

    return saved


def print_report(results):
    """Print a human-readable classification report."""
    stats = results['stats']
    canvas_w, canvas_h = results['canvas']

    print(f"\nCanvas: {canvas_w:.1f}\" x {canvas_h:.2f}\"")
    print(f"Total unique images: {len(results['images'])}")
    print(f"  Full-bleed backgrounds: {stats.get('full-bleed', 0)}")
    print(f"  Panel images: {stats.get('panel', 0)}")
    print(f"  Logos: {stats.get('logo', 0)}")
    print(f"  Icons: {stats.get('icon', 0)}")
    print(f"  Other: {stats.get('other', 0) + stats.get('small', 0)}")
    print()

    # Full-bleed backgrounds
    full_bleed = [img for img in results['images'] if img['category'] == 'full-bleed']
    if full_bleed:
        print("FULL-BLEED BACKGROUNDS (title/closing/section divider candidates):")
        for i, img in enumerate(full_bleed):
            slides_str = ', '.join(str(s) for s in img['slides'])
            print(f"  [{i+1}] {img['width']:.0f}\"x{img['height']:.0f}\" "
                  f"{img['content_type'].split('/')[-1]} {img['size_kb']}KB "
                  f"— slides: {slides_str}"
                  f"{' (+more)' if img['num_slides'] > 5 else ''}")
        print()

    # Panel images
    panels = [img for img in results['images'] if img['category'] == 'panel']
    if panels:
        print("PANEL IMAGES (agenda side panel candidates):")
        for i, img in enumerate(panels):
            slides_str = ', '.join(str(s) for s in img['slides'])
            print(f"  [{i+1}] {img['width']:.1f}\"x{img['height']:.1f}\" "
                  f"{img['content_type'].split('/')[-1]} {img['size_kb']}KB "
                  f"— slides: {slides_str}")
        print()

    # Logos
    logos = [img for img in results['images'] if img['category'] == 'logo']
    if logos:
        print("LOGOS (appears on multiple slides):")
        for i, img in enumerate(logos):
            print(f"  [{i+1}] {img['width']:.1f}\"x{img['height']:.1f}\" "
                  f"{img['size_kb']}KB — appears on {img['num_slides']} slides")
        print()

    # Other notable images
    other = [img for img in results['images']
             if img['category'] == 'other' and img['size_kb'] > 10]
    if other:
        print(f"OTHER IMAGES ({len(other)} notable):")
        for i, img in enumerate(other[:10]):
            slides_str = ', '.join(str(s) for s in img['slides'])
            print(f"  [{i+1}] {img['width']:.1f}\"x{img['height']:.1f}\" "
                  f"{img['content_type'].split('/')[-1]} {img['size_kb']}KB "
                  f"— slides: {slides_str}")
        if len(other) > 10:
            print(f"  ... and {len(other) - 10} more")
        print()


def main():
    parser = argparse.ArgumentParser(
        description="Extract and classify images from a PPTX template")
    parser.add_argument("--template", required=True, help="Path to PPTX template")
    parser.add_argument("--output-dir", "-o", default=None,
                        help="Directory to save extracted images (default: just report)")
    args = parser.parse_args()

    if not os.path.isfile(args.template):
        print(f"ERROR: File not found: {args.template}")
        sys.exit(1)

    print(f"Scanning: {args.template}")
    results = extract_all_images(args.template)
    print_report(results)

    if args.output_dir:
        saved = save_extracted_images(results, args.output_dir)
        print(f"Saved {len(saved)} images to {args.output_dir}/")
        for s in saved:
            print(f"  {s['category']}/{os.path.basename(s['path'])}")


if __name__ == "__main__":
    main()
