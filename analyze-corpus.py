#!/usr/bin/env python3
"""
Corpus Analyzer — Scans PPTX files, extracts metrics, generates analysis report.

Modes:
    --manifest <dir>       Scan directory for .pptx files, output inventory table
    --analyze <manifest>   Run full metric extraction on files listed in manifest JSON
    --include-vision       Add Claude vision spot-checks (requires ANTHROPIC_API_KEY)

Usage:
    python3 analyze-corpus.py --manifest /path/to/management/
    python3 analyze-corpus.py --analyze analysis/corpus-manifest.json
    python3 analyze-corpus.py --analyze analysis/corpus-manifest.json --include-vision
"""

import argparse
import json
import os
import sys
import statistics
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Default brand colors/fonts for corpus analysis — override via --brand-yaml
# to analyze against your own brand's palette.
BRAND_COLORS_HEX = {
    "1A365D", "3182CE", "63B3ED", "EBF4FF", "FFFFFF",
    "2D3748", "718096", "CBD5E0",
    "000000", "2F855A", "C05621", "C53030",
    "F7FAFC", "F5F5F5",
}

BRAND_FONTS = {"Arial", "Arial Bold", "Arial Black"}

# Canvas sizes for generation detection
GEN1_WIDTH_INCHES = 13.33   # Pre-template (13.33" x 7.5")
GEN2_WIDTH_INCHES = 10.0    # Template-based (10.0" x 5.62")

# Deck-builder output indicators (directory patterns)
BUILDER_OUTPUT_DIRS = {"showcase", "csuite", "mgmt", "deck-builder"}

SHAPE_TYPE_NAMES = {
    MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
    MSO_SHAPE_TYPE.PICTURE: "picture",
    MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
    MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
    MSO_SHAPE_TYPE.TABLE: "table",
    MSO_SHAPE_TYPE.GROUP: "group",
    MSO_SHAPE_TYPE.CHART: "chart",
    MSO_SHAPE_TYPE.FREEFORM: "freeform",
}


# ---------------------------------------------------------------------------
# Helpers (adapted from proof_renderer.py and qa_pipeline.py)
# ---------------------------------------------------------------------------

def _emu_to_inches(emu):
    return emu / 914400 if emu else 0


def _get_font_size_pt(run, paragraph=None, shape=None):
    """Resolve font size through inheritance chain: run -> paragraph -> placeholder."""
    if run.font.size is not None:
        return run.font.size.pt
    if paragraph and paragraph.font.size is not None:
        return paragraph.font.size.pt
    # Try shape placeholder default
    if shape and hasattr(shape, 'text_frame'):
        try:
            ph = shape.placeholder_format
            if ph is not None:
                # Check slide layout placeholder
                layout = shape._element.getparent().getparent()
                # Fallback: check XML for sz attribute
                for rpr in run._r.iterancestors():
                    sz = rpr.get('sz')
                    if sz:
                        return int(sz) / 12700
        except Exception:
            pass
    # Fallback: scan XML for sz attribute
    try:
        sz = run._r.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
        if sz is not None and sz.get('sz'):
            return int(sz.get('sz')) / 12700
    except Exception:
        pass
    return None


def _get_font_name(run):
    """Resolve font name, falling back to XML inspection."""
    if run.font.name:
        return run.font.name
    # XML fallback for <a:latin typeface="...">
    try:
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        rpr = run._r.find(f'{ns}rPr')
        if rpr is not None:
            latin = rpr.find(f'{ns}latin')
            if latin is not None and latin.get('typeface'):
                return latin.get('typeface')
    except Exception:
        pass
    return None


def _get_text_color_hex(run):
    """Get text color as hex string."""
    try:
        if run.font.color and run.font.color.rgb:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def _shape_fill_color_hex(shape):
    """Get shape fill color as hex string."""
    try:
        fill = shape.fill
        if fill.type is not None:
            if hasattr(fill, 'fore_color') and fill.fore_color and fill.fore_color.rgb:
                return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _is_footer_shape(shape, slide_height_inches):
    """Check if a shape is in the footer region (bottom 0.6 inches)."""
    if not hasattr(shape, 'top') or shape.top is None:
        return False
    top_in = _emu_to_inches(shape.top)
    return top_in > (slide_height_inches - 0.65)


def _get_shape_type_name(shape):
    """Get human-readable shape type name."""
    try:
        return SHAPE_TYPE_NAMES.get(shape.shape_type, f"other({shape.shape_type})")
    except Exception:
        return "unknown"


# ---------------------------------------------------------------------------
# Manifest mode — scan directory for PPTX files
# ---------------------------------------------------------------------------

def scan_for_pptx(root_dir):
    """Find all .pptx files under root_dir, return metadata list."""
    root = Path(root_dir)
    files = []

    for pptx_path in sorted(root.rglob("*.pptx")):
        # Skip temp files
        if pptx_path.name.startswith("~$"):
            continue
        # Skip python-pptx default template and venv files
        if ".venv/" in str(pptx_path) or "site-packages" in str(pptx_path):
            continue
        # Skip brand template files (they're references, not presentations)
        if pptx_path.name == "template.pptx" and "brands" in str(pptx_path):
            continue

        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            files.append({
                "path": str(pptx_path),
                "relative_path": str(pptx_path.relative_to(root)),
                "error": str(e),
            })
            continue

        # Canvas size
        width_in = _emu_to_inches(prs.slide_width)
        height_in = _emu_to_inches(prs.slide_height)

        # Generation detection
        if abs(width_in - GEN1_WIDTH_INCHES) < 0.5:
            generation = "gen1-pretemp"
        elif abs(width_in - GEN2_WIDTH_INCHES) < 0.5:
            generation = "gen2-template"
        else:
            generation = f"custom({width_in:.1f}x{height_in:.1f})"

        # Project detection (from path)
        rel = pptx_path.relative_to(root)
        parts = rel.parts
        project = parts[0] if len(parts) > 1 else "(root)"

        # Deck-builder output detection — only match known builder outputs,
        # NOT files downloaded into the corpus/ directory
        is_builder = False
        path_lower = str(pptx_path).lower()
        fname_lower = pptx_path.name.lower()
        # Files in corpus/ are Drive downloads, not builder outputs
        if "/corpus/" not in path_lower:
            if "showcase" in fname_lower:
                is_builder = True
            if any(d in path_lower for d in ["-proof/", "proof-"]):
                is_builder = True
            # Known builder output filenames
            if fname_lower in {"cloud-migration-csuite-deck.pptx",
                               "cloud-migration-mgmt-deck.pptx",
                               "template-visual-map.pptx"}:
                is_builder = True

        # Slide count
        slide_count = len(prs.slides)

        # Quick font scan for brand detection
        has_brand_fonts = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            fname = _get_font_name(run)
                            if fname and fname in BRAND_FONTS:
                                has_brand_fonts = True
                                break
                    if has_brand_fonts:
                        break
            if has_brand_fonts:
                break

        files.append({
            "path": str(pptx_path),
            "relative_path": str(pptx_path.relative_to(root)),
            "filename": pptx_path.name,
            "project": project,
            "generation": generation,
            "canvas": f"{width_in:.2f}x{height_in:.2f}",
            "slides": slide_count,
            "has_brand_fonts": has_brand_fonts,
            "is_builder_output": is_builder,
            "include": True,  # Default: include all; user curates
        })

    return files


def write_manifest(files, output_dir):
    """Write manifest as both JSON (machine) and markdown (human review)."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # JSON manifest
    json_path = output_dir / "corpus-manifest.json"
    with open(json_path, "w") as f:
        json.dump(files, f, indent=2)

    # Markdown manifest
    md_path = output_dir / "corpus-manifest.md"
    with open(md_path, "w") as f:
        f.write("# Corpus Manifest\n\n")
        f.write(f"**Total files:** {len(files)}\n\n")

        # Group by generation
        by_gen = defaultdict(list)
        for fi in files:
            gen = fi.get("generation", "unknown")
            by_gen[gen].append(fi)

        for gen in sorted(by_gen.keys()):
            gen_files = by_gen[gen]
            f.write(f"## {gen} ({len(gen_files)} files)\n\n")
            f.write("| # | File | Project | Slides | Canvas | Brand Fonts | Builder | Include |\n")
            f.write("|---|------|---------|--------|--------|-------------|---------|--------|\n")
            for i, fi in enumerate(gen_files, 1):
                if "error" in fi:
                    f.write(f"| {i} | {fi.get('relative_path', fi['path'])} | — | ERROR | — | — | — | No |\n")
                    continue
                brand = "Yes" if fi["has_brand_fonts"] else "No"
                builder = "Yes" if fi["is_builder_output"] else "No"
                include = "Yes" if fi["include"] else "No"
                f.write(f"| {i} | `{fi['relative_path']}` | {fi['project']} | {fi['slides']} | {fi['canvas']} | {brand} | {builder} | {include} |\n")
            f.write("\n")

        # Summary
        total_slides = sum(fi.get("slides", 0) for fi in files if "error" not in fi)
        builder_count = sum(1 for fi in files if fi.get("is_builder_output"))
        f.write(f"---\n\n**Summary:** {len(files)} files, {total_slides} total slides, {builder_count} deck-builder outputs\n")

    print(f"Manifest written: {json_path}")
    print(f"Manifest written: {md_path}")
    return json_path, md_path


# ---------------------------------------------------------------------------
# Extraction mode — per-slide and per-deck metrics
# ---------------------------------------------------------------------------

def extract_slide_metrics(slide, slide_idx, slide_width_in, slide_height_in):
    """Extract all metrics from a single slide."""
    metrics = {
        "slide_idx": slide_idx,
        "word_count": 0,
        "font_sizes": [],
        "font_families": Counter(),
        "text_colors": Counter(),
        "fill_colors": Counter(),
        "bold_runs": 0,
        "italic_runs": 0,
        "total_runs": 0,
        "shape_types": Counter(),
        "headline_text": "",
        "headline_font_size": 0,
        "layout_name": "",
        "image_count": 0,
        "table_count": 0,
        "shape_count": 0,
        "content_area_pct": 0.0,
    }

    # Layout name
    try:
        metrics["layout_name"] = slide.slide_layout.name
    except Exception:
        metrics["layout_name"] = "unknown"

    # Boolean grid for content density (100 x 56 resolution)
    grid_w, grid_h = 100, 56
    grid = [[False] * grid_w for _ in range(grid_h)]

    largest_font = 0
    headline_candidate = ""

    for shape in slide.shapes:
        metrics["shape_count"] += 1
        stype = _get_shape_type_name(shape)
        metrics["shape_types"][stype] += 1

        # Skip footer shapes
        if _is_footer_shape(shape, slide_height_in):
            continue

        # Content density grid
        if hasattr(shape, 'left') and shape.left is not None:
            left_in = _emu_to_inches(shape.left)
            top_in = _emu_to_inches(shape.top)
            w_in = _emu_to_inches(shape.width)
            h_in = _emu_to_inches(shape.height)

            x0 = max(0, int(left_in / slide_width_in * grid_w))
            y0 = max(0, int(top_in / slide_height_in * grid_h))
            x1 = min(grid_w, int((left_in + w_in) / slide_width_in * grid_w))
            y1 = min(grid_h, int((top_in + h_in) / slide_height_in * grid_h))

            for gy in range(y0, y1):
                for gx in range(x0, x1):
                    grid[gy][gx] = True

        # Image counting
        if stype == "picture":
            metrics["image_count"] += 1

        # Table counting
        if stype == "table":
            metrics["table_count"] += 1
            # Extract table text for word count
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        words = cell.text.strip().split()
                        metrics["word_count"] += len(words)

        # Text analysis
        if shape.has_text_frame:
            shape_text = shape.text_frame.text.strip()
            if shape_text:
                words = shape_text.split()
                metrics["word_count"] += len(words)

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue

                    metrics["total_runs"] += 1

                    # Font size
                    fs = _get_font_size_pt(run, para, shape)
                    if fs:
                        metrics["font_sizes"].append(fs)
                        # Headline detection: largest font near top
                        top_in = _emu_to_inches(shape.top) if shape.top else 999
                        if fs > largest_font and top_in < 1.5:
                            largest_font = fs
                            headline_candidate = run.text.strip()

                    # Font family
                    fname = _get_font_name(run)
                    if fname:
                        metrics["font_families"][fname] += 1

                    # Text color
                    color = _get_text_color_hex(run)
                    if color:
                        metrics["text_colors"][color] += 1

                    # Bold / italic
                    if run.font.bold:
                        metrics["bold_runs"] += 1
                    if run.font.italic:
                        metrics["italic_runs"] += 1

        # Fill color
        fill_color = _shape_fill_color_hex(shape)
        if fill_color:
            metrics["fill_colors"][fill_color] += 1

    # Headline
    metrics["headline_text"] = headline_candidate[:120]
    metrics["headline_font_size"] = largest_font

    # Content density percentage
    filled = sum(1 for row in grid for cell in row if cell)
    metrics["content_area_pct"] = round(filled / (grid_w * grid_h) * 100, 1)

    # Convert Counters to dicts for JSON
    metrics["font_families"] = dict(metrics["font_families"])
    metrics["text_colors"] = dict(metrics["text_colors"])
    metrics["fill_colors"] = dict(metrics["fill_colors"])
    metrics["shape_types"] = dict(metrics["shape_types"])

    return metrics


def extract_deck_metrics(pptx_path):
    """Extract metrics for an entire deck."""
    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        return {"path": str(pptx_path), "error": str(e)}

    width_in = _emu_to_inches(prs.slide_width)
    height_in = _emu_to_inches(prs.slide_height)

    slides = []
    for i, slide in enumerate(prs.slides):
        sm = extract_slide_metrics(slide, i + 1, width_in, height_in)
        slides.append(sm)

    # Deck-level aggregates
    word_counts = [s["word_count"] for s in slides]
    all_font_sizes = []
    all_fonts = Counter()
    all_text_colors = Counter()
    all_fill_colors = Counter()
    total_bold = 0
    total_italic = 0
    total_runs = 0
    total_images = 0
    total_tables = 0
    density_values = [s["content_area_pct"] for s in slides]
    headlines = [s["headline_text"] for s in slides if s["headline_text"]]

    for s in slides:
        all_font_sizes.extend(s["font_sizes"])
        for k, v in s["font_families"].items():
            all_fonts[k] += v
        for k, v in s["text_colors"].items():
            all_text_colors[k] += v
        for k, v in s["fill_colors"].items():
            all_fill_colors[k] += v
        total_bold += s["bold_runs"]
        total_italic += s["italic_runs"]
        total_runs += s["total_runs"]
        total_images += s["image_count"]
        total_tables += s["table_count"]

    # Brand adherence
    brand_font_runs = sum(v for k, v in all_fonts.items() if k in BRAND_FONTS)
    total_font_runs = sum(all_fonts.values())
    brand_color_runs = sum(v for k, v in all_text_colors.items() if k in BRAND_COLORS_HEX)
    total_color_runs = sum(all_text_colors.values())

    deck = {
        "path": str(pptx_path),
        "filename": Path(pptx_path).name,
        "canvas": f"{width_in:.2f}x{height_in:.2f}",
        "slide_count": len(slides),
        "word_counts": {
            "total": sum(word_counts),
            "per_slide_avg": round(statistics.mean(word_counts), 1) if word_counts else 0,
            "per_slide_median": round(statistics.median(word_counts), 1) if word_counts else 0,
            "per_slide_max": max(word_counts) if word_counts else 0,
            "per_slide_min": min(word_counts) if word_counts else 0,
            "per_slide_p90": round(sorted(word_counts)[int(len(word_counts) * 0.9)] if len(word_counts) > 1 else (word_counts[0] if word_counts else 0), 1),
        },
        "font_sizes": {
            "min": round(min(all_font_sizes), 1) if all_font_sizes else 0,
            "max": round(max(all_font_sizes), 1) if all_font_sizes else 0,
            "median": round(statistics.median(all_font_sizes), 1) if all_font_sizes else 0,
            "mean": round(statistics.mean(all_font_sizes), 1) if all_font_sizes else 0,
        },
        "font_families": dict(all_fonts.most_common(10)),
        "text_colors": dict(all_text_colors.most_common(15)),
        "fill_colors": dict(all_fill_colors.most_common(10)),
        "bold_pct": round(total_bold / total_runs * 100, 1) if total_runs else 0,
        "italic_pct": round(total_italic / total_runs * 100, 1) if total_runs else 0,
        "content_density": {
            "avg_pct": round(statistics.mean(density_values), 1) if density_values else 0,
            "median_pct": round(statistics.median(density_values), 1) if density_values else 0,
            "max_pct": round(max(density_values), 1) if density_values else 0,
            "min_pct": round(min(density_values), 1) if density_values else 0,
        },
        "brand_adherence": {
            "font_brand_pct": round(brand_font_runs / total_font_runs * 100, 1) if total_font_runs else 0,
            "color_brand_pct": round(brand_color_runs / total_color_runs * 100, 1) if total_color_runs else 0,
        },
        "images": total_images,
        "tables": total_tables,
        "headlines": headlines[:5],  # Sample for review
        "slides": slides,
    }

    return deck


def aggregate_corpus(deck_metrics_list):
    """Compute corpus-wide aggregates from per-deck metrics."""
    # Split by generation
    all_decks = [d for d in deck_metrics_list if "error" not in d]

    all_word_avgs = [d["word_counts"]["per_slide_avg"] for d in all_decks]
    all_font_mins = [d["font_sizes"]["min"] for d in all_decks if d["font_sizes"]["min"] > 0]
    all_font_medians = [d["font_sizes"]["median"] for d in all_decks if d["font_sizes"]["median"] > 0]
    all_density_avgs = [d["content_density"]["avg_pct"] for d in all_decks]
    all_bold_pcts = [d["bold_pct"] for d in all_decks]
    all_slide_counts = [d["slide_count"] for d in all_decks]

    # Font family frequency across all decks
    corpus_fonts = Counter()
    corpus_text_colors = Counter()
    corpus_fill_colors = Counter()
    total_images = 0
    total_tables = 0
    total_slides = 0

    for d in all_decks:
        for k, v in d["font_families"].items():
            corpus_fonts[k] += v
        for k, v in d["text_colors"].items():
            corpus_text_colors[k] += v
        for k, v in d["fill_colors"].items():
            corpus_fill_colors[k] += v
        total_images += d["images"]
        total_tables += d["tables"]
        total_slides += d["slide_count"]

    # All headlines
    all_headlines = []
    for d in all_decks:
        for s in d["slides"]:
            if s["headline_text"]:
                all_headlines.append(s["headline_text"])

    # Headline length stats
    headline_lengths = [len(h.split()) for h in all_headlines]

    agg = {
        "corpus_size": len(all_decks),
        "total_slides": total_slides,
        "slides_per_deck": {
            "avg": round(statistics.mean(all_slide_counts), 1) if all_slide_counts else 0,
            "median": round(statistics.median(all_slide_counts), 1) if all_slide_counts else 0,
            "min": min(all_slide_counts) if all_slide_counts else 0,
            "max": max(all_slide_counts) if all_slide_counts else 0,
        },
        "words_per_slide": {
            "avg_of_avgs": round(statistics.mean(all_word_avgs), 1) if all_word_avgs else 0,
            "median_of_avgs": round(statistics.median(all_word_avgs), 1) if all_word_avgs else 0,
            "range": f"{min(all_word_avgs):.0f}–{max(all_word_avgs):.0f}" if all_word_avgs else "—",
        },
        "font_sizes": {
            "corpus_min": round(min(all_font_mins), 1) if all_font_mins else 0,
            "median_of_medians": round(statistics.median(all_font_medians), 1) if all_font_medians else 0,
        },
        "font_families_top10": dict(corpus_fonts.most_common(10)),
        "text_colors_top10": dict(corpus_text_colors.most_common(10)),
        "fill_colors_top10": dict(corpus_fill_colors.most_common(10)),
        "bold_pct": {
            "avg": round(statistics.mean(all_bold_pcts), 1) if all_bold_pcts else 0,
            "range": f"{min(all_bold_pcts):.0f}–{max(all_bold_pcts):.0f}%" if all_bold_pcts else "—",
        },
        "content_density_pct": {
            "avg_of_avgs": round(statistics.mean(all_density_avgs), 1) if all_density_avgs else 0,
            "range": f"{min(all_density_avgs):.0f}–{max(all_density_avgs):.0f}%" if all_density_avgs else "—",
        },
        "total_images": total_images,
        "total_tables": total_tables,
        "images_per_slide": round(total_images / total_slides, 2) if total_slides else 0,
        "tables_per_slide": round(total_tables / total_slides, 3) if total_slides else 0,
        "headline_word_count": {
            "avg": round(statistics.mean(headline_lengths), 1) if headline_lengths else 0,
            "median": round(statistics.median(headline_lengths), 1) if headline_lengths else 0,
            "max": max(headline_lengths) if headline_lengths else 0,
        },
        "sample_headlines": all_headlines[:20],
    }

    return agg


def write_analysis_report(corpus_agg, deck_metrics, manifest_files, output_dir):
    """Write the human-readable analysis report."""
    output_dir = Path(output_dir)
    report_path = output_dir / "corpus-analysis-report.md"

    with open(report_path, "w") as f:
        f.write("# Corpus Analysis Report\n\n")

        # Executive summary
        f.write("## Executive Summary\n\n")
        f.write(f"Analyzed **{corpus_agg['corpus_size']} presentations** containing **{corpus_agg['total_slides']} slides**.\n\n")
        f.write(f"- **Words per slide:** avg {corpus_agg['words_per_slide']['avg_of_avgs']}, range {corpus_agg['words_per_slide']['range']}\n")
        f.write(f"- **Content density:** avg {corpus_agg['content_density_pct']['avg_of_avgs']}%, range {corpus_agg['content_density_pct']['range']}\n")
        f.write(f"- **Font size floor:** {corpus_agg['font_sizes']['corpus_min']}pt (median {corpus_agg['font_sizes']['median_of_medians']}pt)\n")
        f.write(f"- **Bold usage:** avg {corpus_agg['bold_pct']['avg']}% of text runs\n")
        f.write(f"- **Headline length:** avg {corpus_agg['headline_word_count']['avg']} words (median {corpus_agg['headline_word_count']['median']})\n")
        f.write(f"- **Visual elements:** {corpus_agg['images_per_slide']} images/slide, {corpus_agg['tables_per_slide']} tables/slide\n\n")

        # Corpus overview
        f.write("## Corpus Overview\n\n")
        f.write("| Deck | Slides | Words/Slide (avg) | Density (avg) | Min Font | Bold % | Images |\n")
        f.write("|------|--------|-------------------|---------------|----------|--------|--------|\n")
        valid = [d for d in deck_metrics if "error" not in d]
        for d in sorted(valid, key=lambda x: x["filename"]):
            f.write(f"| {d['filename']} | {d['slide_count']} | {d['word_counts']['per_slide_avg']} | {d['content_density']['avg_pct']}% | {d['font_sizes']['min']}pt | {d['bold_pct']}% | {d['images']} |\n")
        f.write("\n")

        # Typography
        f.write("## Typography Patterns\n\n")
        f.write("### Font Families\n\n")
        f.write("| Font | Usage (runs) |\n")
        f.write("|------|--------------|\n")
        for font, count in corpus_agg["font_families_top10"].items():
            brand_tag = " **(brand)**" if font in BRAND_FONTS else ""
            f.write(f"| {font}{brand_tag} | {count} |\n")
        f.write("\n")

        f.write("### Font Sizes\n\n")
        f.write(f"- Corpus minimum: **{corpus_agg['font_sizes']['corpus_min']}pt**\n")
        f.write(f"- Median across decks: **{corpus_agg['font_sizes']['median_of_medians']}pt**\n\n")

        f.write("### Bold/Italic Usage\n\n")
        f.write(f"- Bold: avg **{corpus_agg['bold_pct']['avg']}%** of runs, range {corpus_agg['bold_pct']['range']}\n\n")

        # Color
        f.write("## Color Patterns\n\n")
        f.write("### Text Colors\n\n")
        f.write("| Color | Usage | Brand? |\n")
        f.write("|-------|-------|--------|\n")
        for color, count in corpus_agg["text_colors_top10"].items():
            brand = "Yes" if color in BRAND_COLORS_HEX else "No"
            f.write(f"| #{color} | {count} | {brand} |\n")
        f.write("\n")

        f.write("### Fill Colors\n\n")
        f.write("| Color | Usage | Brand? |\n")
        f.write("|-------|-------|--------|\n")
        for color, count in corpus_agg["fill_colors_top10"].items():
            brand = "Yes" if color in BRAND_COLORS_HEX else "No"
            f.write(f"| #{color} | {count} | {brand} |\n")
        f.write("\n")

        # Content density
        f.write("## Content Density\n\n")
        f.write(f"- **Average:** {corpus_agg['content_density_pct']['avg_of_avgs']}%\n")
        f.write(f"- **Range:** {corpus_agg['content_density_pct']['range']}\n")
        f.write(f"- **65% rule comparison:** ", )
        avg_density = corpus_agg['content_density_pct']['avg_of_avgs']
        if avg_density <= 65:
            f.write(f"Corpus average ({avg_density}%) is within the 65% guideline.\n\n")
        else:
            f.write(f"Corpus average ({avg_density}%) **exceeds** the 65% guideline.\n\n")

        # Headlines
        f.write("## Headline Patterns\n\n")
        f.write(f"- **Average length:** {corpus_agg['headline_word_count']['avg']} words\n")
        f.write(f"- **Median length:** {corpus_agg['headline_word_count']['median']} words\n")
        f.write(f"- **Longest:** {corpus_agg['headline_word_count']['max']} words\n\n")
        f.write("### Sample Headlines\n\n")
        for h in corpus_agg["sample_headlines"][:20]:
            f.write(f"- {h}\n")
        f.write("\n")

        # Structure
        f.write("## Structure Patterns\n\n")
        f.write(f"- **Slides per deck:** avg {corpus_agg['slides_per_deck']['avg']}, range {corpus_agg['slides_per_deck']['min']}–{corpus_agg['slides_per_deck']['max']}\n")
        f.write(f"- **Images per slide:** {corpus_agg['images_per_slide']}\n")
        f.write(f"- **Tables per slide:** {corpus_agg['tables_per_slide']}\n\n")

        # Principles alignment placeholder
        f.write("## Principles Alignment\n\n")
        f.write("*To be completed after vision spot-checks and discussion.*\n\n")
        f.write("| Principle | Alignment | Notes |\n")
        f.write("|-----------|-----------|-------|\n")
        f.write("| One idea per slide | TBD | Check headlines |\n")
        f.write("| Audience density profiles | TBD | Compare word counts to profile limits |\n")
        f.write("| Storytelling frameworks | TBD | Requires slide sequence analysis |\n")
        f.write("| Icon rules | TBD | Requires vision review |\n")
        f.write("| Typography hierarchy | TBD | Check font size distribution |\n")
        f.write("| 65% content rule | TBD | Compare density measurements |\n")
        f.write("| Split not shrink | TBD | Check min font sizes |\n")
        f.write("| Section structure | TBD | Check slide counts and dividers |\n\n")

        f.write("## Recommendations\n\n")
        f.write("*To be populated after full analysis and discussion.*\n\n")
        f.write("- **(A) Rules to tighten:** TBD\n")
        f.write("- **(B) Rules to relax:** TBD\n")
        f.write("- **(C) New rules to add:** TBD\n")
        f.write("- **(D) Rules to keep as-is:** TBD\n")
        f.write("- **(E) Cannot determine:** TBD\n")

    print(f"Report written: {report_path}")
    return report_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Corpus Analyzer for PPTX presentations")
    parser.add_argument("--manifest", metavar="DIR",
                        help="Scan directory for .pptx files and output inventory")
    parser.add_argument("--analyze", metavar="MANIFEST_JSON",
                        help="Run full extraction on files in manifest JSON")
    parser.add_argument("--include-vision", action="store_true",
                        help="Include Claude vision spot-checks")
    parser.add_argument("--output-dir", default="analysis",
                        help="Output directory (default: analysis/)")
    parser.add_argument("--add-file", metavar="PATH", action="append", default=[],
                        help="Add individual PPTX files to the manifest (can be used multiple times)")
    args = parser.parse_args()

    output_dir = Path(args.output_dir)

    if args.manifest:
        print(f"Scanning {args.manifest} for PPTX files...")
        files = scan_for_pptx(args.manifest)

        # Add any individually specified files
        for extra in args.add_file:
            p = Path(extra)
            if p.exists() and p.suffix.lower() == '.pptx':
                try:
                    prs = Presentation(str(p))
                    width_in = _emu_to_inches(prs.slide_width)
                    height_in = _emu_to_inches(prs.slide_height)
                    if abs(width_in - GEN1_WIDTH_INCHES) < 0.5:
                        gen = "gen1-pretemp"
                    elif abs(width_in - GEN2_WIDTH_INCHES) < 0.5:
                        gen = "gen2-template"
                    else:
                        gen = f"custom({width_in:.1f}x{height_in:.1f})"

                    has_brand = False
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    for run in para.runs:
                                        if _get_font_name(run) in BRAND_FONTS:
                                            has_brand = True
                                            break
                                if has_brand:
                                    break
                        if has_brand:
                            break

                    files.append({
                        "path": str(p),
                        "relative_path": str(p.name),
                        "filename": p.name,
                        "project": "(external)",
                        "generation": gen,
                        "canvas": f"{width_in:.2f}x{height_in:.2f}",
                        "slides": len(prs.slides),
                        "has_brand_fonts": has_brand,
                        "is_builder_output": False,
                        "include": True,
                        "source": "external",
                    })
                except Exception as e:
                    print(f"Warning: Could not read {extra}: {e}")

        write_manifest(files, output_dir)
        print(f"\nFound {len(files)} PPTX files. Review {output_dir}/corpus-manifest.md to curate.")

    elif args.analyze:
        manifest_path = Path(args.analyze)
        if not manifest_path.exists():
            print(f"Error: {manifest_path} not found")
            sys.exit(1)

        with open(manifest_path) as f:
            manifest = json.load(f)

        # Filter to included files
        included = [fi for fi in manifest if fi.get("include", True) and "error" not in fi]
        print(f"Analyzing {len(included)} files...")

        deck_metrics = []
        for fi in included:
            pptx_path = fi["path"]
            print(f"  Extracting: {fi.get('filename', pptx_path)}")
            dm = extract_deck_metrics(pptx_path)
            dm["project"] = fi.get("project", "unknown")
            dm["generation"] = fi.get("generation", "unknown")
            dm["source"] = fi.get("source", "local")
            deck_metrics.append(dm)

        # Aggregate
        print("Computing aggregates...")
        corpus_agg = aggregate_corpus(deck_metrics)

        # Save metrics JSON
        output_dir.mkdir(parents=True, exist_ok=True)
        metrics_path = output_dir / "corpus-metrics.json"
        with open(metrics_path, "w") as f:
            json.dump({
                "aggregate": corpus_agg,
                "decks": [{k: v for k, v in d.items() if k != "slides"} for d in deck_metrics],  # Exclude per-slide for summary
                "decks_full": deck_metrics,
            }, f, indent=2)
        print(f"Metrics written: {metrics_path}")

        # Write report
        write_analysis_report(corpus_agg, deck_metrics, manifest, output_dir)

    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
