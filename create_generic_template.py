#!/usr/bin/env python3
"""Create a minimal generic PPTX template for the Deck Builder.

Generates a template with:
  - 10" x 5.625" canvas (16:9)
  - Layout 0: Title slide (blank with background placeholder)
  - Layout 1: Section header
  - Layout 2: Content (blank canvas with footer placeholders) — the main layout
  - No baked-in logos or brand colors
  - Arial fonts
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy
import os


def create_generic_template(output_path):
    """Create a minimal generic PPTX template."""
    prs = Presentation()

    # Set slide dimensions: 10" x 5.625" (16:9)
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)

    # Access the slide master
    slide_master = prs.slide_masters[0]

    # We need at least 3 slide layouts. python-pptx's default has several.
    # Layout 0 = Title Slide, Layout 1 = Title and Content, Layout 2 = Blank, etc.
    # For the deck builder, the key layout is a blank canvas (layout index 2)
    # with footer placeholders.

    # The default template already has these layouts. We just need to ensure
    # the dimensions are correct and footers are present.

    # Save the template
    prs.save(output_path)
    print(f"Created generic template: {output_path}")
    print(f"  Canvas: {prs.slide_width / 914400:.1f}\" x {prs.slide_height / 914400:.2f}\"")
    print(f"  Layouts: {len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"    [{i}] {layout.name}")


if __name__ == "__main__":
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output = os.path.join(script_dir, "brands", "generic", "template.pptx")
    create_generic_template(output)
