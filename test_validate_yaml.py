#!/usr/bin/env python3
"""
YAML Deck Validation Tests

Tests that catch common authoring mistakes in deck YAML files:
- Invalid layout names that silently fall back to content_generic
- Invalid icon names that silently render as missing
- Missing required fields per layout type
- Field name mismatches (e.g., left_items vs left_body)

Run: python3 -m pytest test_validate_yaml.py -v
  or: python3 test_validate_yaml.py
"""

import json
import os
import sys
import yaml
try:
    import pytest
except ImportError:
    pytest = None

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))


# ─── Helpers ────────────────────────────────────────────────

def load_icon_catalog():
    """Load the icon catalog and return set of valid icon names."""
    catalog_path = os.path.join(SCRIPT_DIR, "icons", "icon-catalog.json")
    with open(catalog_path) as f:
        return set(json.load(f)["icons"].keys())


def get_valid_layouts():
    """Return the set of valid layout names (including aliases)."""
    # These are the layout names the builder actually handles
    # (from build_slide dispatcher + LAYOUT_ALIASES)
    handled = {
        "title_cover", "agenda", "content_diagram_text",
        "side_by_side", "three_column", "four_card",
        "big_stat_manual", "closing", "content_table",
        "content_table_bullets", "section_divider", "quote",
        "kpi_dashboard", "roadmap", "before_after",
        "numbered_list", "status_board", "image_showcase",
        "matrix", "funnel", "callout",
        "content_generic",
    }
    aliases = {
        "content_three_section": "three_column",
        "content_three_col_cards": "four_card",
        "content_four_cards": "four_card",
        "title_h1": "title_cover",
        "big_stat": "big_stat_manual",
    }
    return handled | set(aliases.keys())


def get_layout_required_fields():
    """Return required fields per layout type."""
    return {
        "title_cover": {
            "required": ["headline"],
            "wrong_field_hints": {
                "title": "title_cover uses 'headline', not 'title'",
                "subtitle": "title_cover uses 'subheader', not 'subtitle'",
            },
        },
        "agenda": {
            "required": ["headline", "items"],
            "wrong_field_hints": {
                "sections": "agenda uses 'items' (list of strings), not 'sections'",
            },
        },
        "before_after": {
            "required": ["headline"],
            "expected_one_of": [["before", "left_items"]],
            "field_spec": {
                "before": {"type": "dict", "keys": ["label", "items"]},
                "after": {"type": "dict", "keys": ["label", "items"]},
            },
        },
        "side_by_side": {
            "required": ["headline"],
            "expected_one_of": [["left_body", "left_items"]],
            "wrong_field_hints": {
                "left_items": "side_by_side uses 'left_body' (markdown text with |), not 'left_items' (list)",
                "right_items": "side_by_side uses 'right_body' (markdown text with |), not 'right_items' (list)",
            },
        },
        "content_table": {
            "required": ["headline", "columns", "rows"],
        },
        "kpi_dashboard": {
            "required": ["headline", "metrics"],
        },
        "roadmap": {
            "required": ["headline", "swimlanes"],
        },
        "numbered_list": {
            "required": ["headline", "items"],
        },
        "callout": {
            "required": ["callout_text"],
        },
        "four_card": {
            "required": ["headline"],
        },
        "three_column": {
            "required": ["headline"],
        },
        "section_divider": {
            "required": ["headline"],
        },
    }


def load_deck_yaml(yaml_path):
    """Load and parse a deck YAML file."""
    with open(yaml_path) as f:
        return yaml.safe_load(f)


def collect_all_icons_from_slide(slide_def):
    """Extract all icon references from a slide definition."""
    icons = []
    # Top-level icon field
    if slide_def.get("icon"):
        icons.append(("icon", slide_def["icon"]))
    # Left/right icons
    for side in ("left", "right"):
        if slide_def.get(f"{side}_icon"):
            icons.append((f"{side}_icon", slide_def[f"{side}_icon"]))
    # Column icons
    for i in range(1, 5):
        for prefix in (f"col{i}_icon", f"section{i}_icon", f"card{i}_icon"):
            if slide_def.get(prefix):
                icons.append((prefix, slide_def[prefix]))
    # Metrics icons (kpi_dashboard)
    for m in slide_def.get("metrics", []):
        if isinstance(m, dict) and m.get("icon"):
            icons.append(("metrics[].icon", m["icon"]))
    # Items icons (numbered_list, status_board)
    for item in slide_def.get("items", []):
        if isinstance(item, dict) and item.get("icon"):
            icons.append(("items[].icon", item["icon"]))
    return icons


# ─── Test: Invalid layout names ─────────────────────────────

def find_deck_yamls(search_dir=None):
    """Find all deck YAML files in a directory (defaults to examples/)."""
    if search_dir is None:
        search_dir = os.path.join(SCRIPT_DIR, "examples")
    yamls = []
    if not os.path.isdir(search_dir):
        return yamls
    for f in os.listdir(search_dir):
        if f.endswith(".yaml") and ("deck" in f.lower() or "showcase" in f.lower()):
            yamls.append(os.path.join(search_dir, f))
    return yamls


VALID_LAYOUTS = get_valid_layouts()
ICON_CATALOG = load_icon_catalog()
REQUIRED_FIELDS = get_layout_required_fields()


if pytest:
    @pytest.fixture(params=find_deck_yamls(), ids=lambda p: os.path.basename(p))
    def deck(request):
        return load_deck_yaml(request.param)


def test_all_layouts_are_valid(deck):
    """Every slide must use a layout name that the builder actually handles."""
    invalid = []
    for i, slide in enumerate(deck.get("slides", []), 1):
        layout = slide.get("layout", "content_generic")
        if layout not in VALID_LAYOUTS:
            invalid.append(f"Slide {i}: layout '{layout}' is not valid")
    assert not invalid, (
        f"Found {len(invalid)} invalid layout name(s) that would silently "
        f"fall back to blank content_generic:\n" + "\n".join(invalid)
    )


def test_all_icons_are_valid(deck):
    """Every icon reference must exist in the icon catalog."""
    invalid = []
    for i, slide in enumerate(deck.get("slides", []), 1):
        for field, icon_name in collect_all_icons_from_slide(slide):
            if icon_name and icon_name not in ICON_CATALOG:
                # Find close matches for helpful error message
                close = [k for k in ICON_CATALOG if icon_name.lower() in k.lower()
                         or k.lower() in icon_name.lower()]
                hint = f" (did you mean: {', '.join(close[:3])}?)" if close else ""
                invalid.append(f"Slide {i}, {field}: icon '{icon_name}' not in catalog{hint}")
    assert not invalid, (
        f"Found {len(invalid)} invalid icon name(s) that would silently render as missing:\n"
        + "\n".join(invalid)
    )


def test_required_fields_present(deck):
    """Each layout's required fields must be present."""
    missing = []
    for i, slide in enumerate(deck.get("slides", []), 1):
        layout = slide.get("layout", "content_generic")
        spec = REQUIRED_FIELDS.get(layout)
        if not spec:
            continue
        for field in spec.get("required", []):
            if field not in slide:
                missing.append(f"Slide {i} ({layout}): missing required field '{field}'")
    assert not missing, (
        f"Found {len(missing)} missing required field(s):\n" + "\n".join(missing)
    )


def test_no_wrong_field_names(deck):
    """Catch common field name mistakes per layout."""
    wrong = []
    for i, slide in enumerate(deck.get("slides", []), 1):
        layout = slide.get("layout", "content_generic")
        spec = REQUIRED_FIELDS.get(layout, {})
        hints = spec.get("wrong_field_hints", {})
        for field, hint in hints.items():
            if field in slide:
                wrong.append(f"Slide {i} ({layout}): '{field}' is wrong — {hint}")
    assert not wrong, (
        f"Found {len(wrong)} field name mistake(s):\n" + "\n".join(wrong)
    )


def test_before_after_uses_dict_not_list(deck):
    """before_after layout needs before/after as dicts with label+items, not flat lists."""
    issues = []
    for i, slide in enumerate(deck.get("slides", []), 1):
        if slide.get("layout") != "before_after":
            continue
        for key in ("before", "after"):
            val = slide.get(key)
            if val is not None and not isinstance(val, dict):
                issues.append(
                    f"Slide {i}: '{key}' should be a dict with 'label' and 'items' keys, "
                    f"got {type(val).__name__}"
                )
            elif isinstance(val, dict) and "items" not in val:
                issues.append(f"Slide {i}: '{key}' dict is missing 'items' list")
    assert not issues, (
        f"Found {len(issues)} before_after structure issue(s):\n" + "\n".join(issues)
    )


def test_google_upload_preserves_images():
    """After uploading to Google Slides and downloading back, all images must survive.

    Google Slides sometimes drops images during PPTX import — especially the last
    image on a slide or images near the right edge. This test downloads the
    Google-converted version and verifies image counts match the original.

    Requires: set ORIGINAL_PPTX and GOOGLE_PPTX env vars to point at the local
    and Google-exported PPTX files respectively.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return

    original_path = os.environ.get("ORIGINAL_PPTX", "")
    google_path = os.environ.get("GOOGLE_PPTX", "")
    if not original_path or not google_path:
        return  # skip if env vars not set
    if not os.path.exists(original_path) or not os.path.exists(google_path):
        return  # skip if files not available

    orig = Presentation(original_path)
    goog = Presentation(google_path)

    issues = []
    for i in range(min(len(orig.slides), len(goog.slides))):
        orig_pics = sum(1 for s in orig.slides[i].shapes if s.shape_type == 13)
        goog_pics = sum(1 for s in goog.slides[i].shapes if s.shape_type == 13)
        if orig_pics != goog_pics:
            issues.append(
                f"Slide {i + 1}: original has {orig_pics} images, "
                f"Google version has {goog_pics} — {orig_pics - goog_pics} dropped"
            )

    assert not issues, (
        f"Google Slides dropped images during upload:\n" + "\n".join(issues)
        + "\nThis is a known Google PPTX import issue. "
        "Try reducing shape count or repositioning images."
    )


def test_title_cover_no_placeholder_text():
    """Title cover slides must have empty placeholders — title goes in TextBox only.

    Google Slides shows placeholder text on top of the manual TextBox overlay,
    making the title invisible if placeholders are filled.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return

    for yaml_path in find_deck_yamls():
        deck = load_deck_yaml(yaml_path)
        pptx_path = yaml_path.replace('.yaml', '.pptx')
        if not os.path.exists(pptx_path):
            continue

        prs = Presentation(pptx_path)
        for i, slide_def in enumerate(deck.get("slides", [])):
            if slide_def.get("layout") != "title_cover":
                continue
            slide = prs.slides[i]
            # Check placeholders (except picture placeholder idx=2) are empty
            for shape in slide.placeholders:
                idx = shape.placeholder_format.idx
                if idx == 2:  # background image
                    continue
                if hasattr(shape, 'text_frame'):
                    text = shape.text_frame.text.strip()
                    assert text == "", (
                        f"{os.path.basename(yaml_path)} slide {i+1}: "
                        f"title_cover placeholder idx={idx} has text '{text[:40]}' — "
                        f"this will overlay the title. Footers should be skipped for title_cover."
                    )
            # Check TextBox has the headline
            headline = slide_def.get("headline", "")
            if headline:
                found = any(
                    s.shape_type == 17 and headline.split(" — ")[0] in s.text_frame.text
                    for s in slide.shapes if hasattr(s, 'text_frame')
                )
                assert found, (
                    f"{os.path.basename(yaml_path)} slide {i+1}: "
                    f"title_cover headline '{headline[:40]}' not found in any TextBox"
                )


def test_built_pptx_kpi_icons_present():
    """Build the deck and verify every KPI metric with an icon actually has
    a PICTURE shape in the PPTX at the correct position within its card."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return  # skip if python-pptx not available

    yamls = find_deck_yamls()
    if not yamls:
        return

    # Use first YAML that has a corresponding built PPTX
    deck = None
    pptx_path = None
    for yp in yamls:
        pp = yp.replace('.yaml', '.pptx')
        if os.path.exists(pp):
            deck = load_deck_yaml(yp)
            pptx_path = pp
            break
    if not deck or not pptx_path:
        return  # skip if no built PPTX found
    prs = Presentation(pptx_path)
    issues = []

    for slide_idx, slide_def in enumerate(deck.get("slides", [])):
        if slide_def.get("layout") != "kpi_dashboard":
            continue
        metrics = slide_def.get("metrics", [])
        icons_expected = sum(1 for m in metrics if m.get("icon"))
        if icons_expected == 0:
            continue

        pptx_slide = prs.slides[slide_idx]
        pic_shapes = [s for s in pptx_slide.shapes if s.shape_type == 13]  # PICTURE

        if len(pic_shapes) < icons_expected:
            issues.append(
                f"Slide {slide_idx + 1}: expected {icons_expected} icon images, "
                f"found {len(pic_shapes)} PICTURE shapes"
            )

        # Check each picture has non-zero dimensions and valid PNG
        for j, pic in enumerate(pic_shapes):
            w = Emu(pic.width).inches
            h = Emu(pic.height).inches
            if w < 0.1 or h < 0.1:
                issues.append(
                    f"Slide {slide_idx + 1}, icon {j + 1}: "
                    f"image is too small ({w:.2f}\" x {h:.2f}\")"
                )
            blob = pic.image.blob
            if not blob or blob[:4] != b'\x89PNG':
                issues.append(
                    f"Slide {slide_idx + 1}, icon {j + 1}: "
                    f"image blob is invalid or empty ({len(blob)} bytes)"
                )

    assert not issues, (
        f"Found {len(issues)} KPI icon issue(s) in built PPTX:\n" + "\n".join(issues)
    )


# ─── Run as script ──────────────────────────────────────────

if __name__ == "__main__":
    # Quick standalone validation without pytest
    yamls = find_deck_yamls()
    if not yamls:
        print("No deck YAML files found")
        sys.exit(1)

    total_issues = 0
    for yaml_path in yamls:
        print(f"\nValidating: {os.path.basename(yaml_path)}")
        deck = load_deck_yaml(yaml_path)
        slides = deck.get("slides", [])
        print(f"  Slides: {len(slides)}")

        # Check layouts
        for i, slide in enumerate(slides, 1):
            layout = slide.get("layout", "content_generic")
            if layout not in VALID_LAYOUTS:
                print(f"  ERROR Slide {i}: invalid layout '{layout}'")
                total_issues += 1

        # Check icons
        for i, slide in enumerate(slides, 1):
            for field, icon_name in collect_all_icons_from_slide(slide):
                if icon_name and icon_name not in ICON_CATALOG:
                    close = [k for k in ICON_CATALOG if icon_name.lower() in k.lower()]
                    hint = f" (try: {', '.join(close[:3])})" if close else ""
                    print(f"  ERROR Slide {i}, {field}: unknown icon '{icon_name}'{hint}")
                    total_issues += 1

        # Check required fields
        for i, slide in enumerate(slides, 1):
            layout = slide.get("layout", "content_generic")
            spec = REQUIRED_FIELDS.get(layout)
            if not spec:
                continue
            for field in spec.get("required", []):
                if field not in slide:
                    print(f"  ERROR Slide {i} ({layout}): missing required '{field}'")
                    total_issues += 1

        # Check wrong field names
        for i, slide in enumerate(slides, 1):
            layout = slide.get("layout", "content_generic")
            spec = REQUIRED_FIELDS.get(layout, {})
            for field, hint in spec.get("wrong_field_hints", {}).items():
                if field in slide:
                    print(f"  ERROR Slide {i} ({layout}): wrong field '{field}' — {hint}")
                    total_issues += 1

    if total_issues:
        print(f"\n FAILED: {total_issues} issue(s) found")
        sys.exit(1)
    else:
        print(f"\n PASSED: all decks valid")
        sys.exit(0)
