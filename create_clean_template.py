#!/usr/bin/env python3
"""Create a logo-free version of a PPTX template.

Removes small logo images (< 1.5" x 0.5") from every slide layout
while preserving all layout structure, placeholders, and dimensions.

Usage:
    python3 create_clean_template.py --source input.pptx --output clean.pptx
"""

import argparse
from pptx import Presentation
import os
import sys


def remove_logos(prs):
    """Remove small logo images from all slide layouts."""
    removed = 0

    for i, layout in enumerate(prs.slide_layouts):
        shapes_to_remove = []
        for shape in layout.shapes:
            # Find small images (logos are typically < 1.5"x0.5")
            if hasattr(shape, 'image'):
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                if w_in < 1.5 and h_in < 0.5:
                    shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            removed += 1

    # Also check slide master
    master = prs.slide_masters[0]
    shapes_to_remove = []
    for shape in master.shapes:
        if hasattr(shape, 'image'):
            w_in = shape.width / 914400
            h_in = shape.height / 914400
            if w_in < 1.5 and h_in < 0.5:
                shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        removed += 1

    return removed


def main():
    parser = argparse.ArgumentParser(
        description="Create a logo-free version of a PPTX template")
    parser.add_argument("--source", required=True, help="Source PPTX template")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    args = parser.parse_args()

    if not os.path.isfile(args.source):
        print(f"ERROR: Source template not found: {args.source}")
        sys.exit(1)

    prs = Presentation(args.source)
    removed = remove_logos(prs)
    prs.save(args.output)

    print(f"Created logo-free template: {args.output}")
    print(f"  Removed {removed} logo images")
    print(f"  Layouts: {len(prs.slide_layouts)}")
    print(f"  Canvas: {prs.slide_width/914400:.1f}\" x {prs.slide_height/914400:.2f}\"")


if __name__ == "__main__":
    main()
