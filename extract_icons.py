#!/usr/bin/env python3
"""Icon Extraction Pipeline — Extract, deduplicate, and AI-label icons from PPTX files.

Pipeline:
  1. Extract all small images (<0.6") from template and/or corpus
  2. Deduplicate by visual similarity (perceptual hash)
  3. Label with AI vision (Claude) for semantic names + categories
  4. Generate icon-catalog.json

Usage:
    python3 extract_icons.py --template path/to/template.pptx
    python3 extract_icons.py --template t.pptx --corpus path/to/decks/
    python3 extract_icons.py --template t.pptx --label  # with AI labeling
"""

import argparse
import base64
import hashlib
import io
import json
import os
import sys
from collections import defaultdict
from pathlib import Path

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx required")
    sys.exit(1)


# Size thresholds
ICON_MAX_W_IN = 0.6
ICON_MAX_H_IN = 0.6
ICON_MIN_SIZE_BYTES = 200  # Skip tiny/empty images


def _perceptual_hash(blob, hash_size=8):
    """Compute a perceptual hash for deduplication.

    Resizes image to hash_size x hash_size grayscale, then compares
    each pixel to the mean. Returns a hex string.
    """
    if not Image:
        # Fallback: use content hash
        return hashlib.md5(blob).hexdigest()[:16]

    try:
        img = Image.open(io.BytesIO(blob)).convert('L').resize(
            (hash_size, hash_size), Image.LANCZOS)
        pixels = list(img.getdata())
        mean = sum(pixels) / len(pixels)
        bits = ''.join('1' if p > mean else '0' for p in pixels)
        return hex(int(bits, 2))[2:].zfill(hash_size * hash_size // 4)
    except Exception:
        return hashlib.md5(blob).hexdigest()[:16]


def _hamming_distance(h1, h2):
    """Count differing bits between two hex hash strings."""
    try:
        i1 = int(h1, 16)
        i2 = int(h2, 16)
        return bin(i1 ^ i2).count('1')
    except (ValueError, TypeError):
        return 999


def extract_icons_from_pptx(pptx_path):
    """Extract all icon-sized images from a PPTX file.

    Returns list of dicts: hash, blob, content_type, width, height, phash
    """
    prs = Presentation(pptx_path)
    icons = []
    seen_content = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, 'image'):
                continue
            try:
                img = shape.image
                w_in = shape.width / 914400
                h_in = shape.height / 914400

                if w_in > ICON_MAX_W_IN or h_in > ICON_MAX_H_IN:
                    continue

                blob = img.blob
                if len(blob) < ICON_MIN_SIZE_BYTES:
                    continue

                content_hash = hashlib.md5(blob).hexdigest()[:12]
                if content_hash in seen_content:
                    continue
                seen_content.add(content_hash)

                ext = 'png' if 'png' in img.content_type else 'jpg'
                phash = _perceptual_hash(blob)

                icons.append({
                    'hash': content_hash,
                    'blob': blob,
                    'content_type': img.content_type,
                    'ext': ext,
                    'size_bytes': len(blob),
                    'width': round(w_in, 2),
                    'height': round(h_in, 2),
                    'phash': phash,
                })
            except Exception:
                continue

    return icons


def extract_icons_from_corpus(corpus_dir):
    """Extract icons from all PPTX files in a directory."""
    all_icons = []
    seen = set()

    for pptx_file in sorted(Path(corpus_dir).glob("*.pptx")):
        try:
            icons = extract_icons_from_pptx(str(pptx_file))
            for icon in icons:
                if icon['hash'] not in seen:
                    seen.add(icon['hash'])
                    all_icons.append(icon)
        except Exception:
            continue

    return all_icons


def deduplicate_icons(icons, threshold=5):
    """Group visually similar icons using perceptual hash.

    Icons within `threshold` hamming distance are considered variants
    (e.g., pink and lightpink versions of the same icon).
    Returns list of groups, each group is a list of icon dicts.
    The first icon in each group is the "representative".
    """
    groups = []
    assigned = set()

    for i, icon in enumerate(icons):
        if i in assigned:
            continue

        group = [icon]
        assigned.add(i)

        for j, other in enumerate(icons):
            if j in assigned:
                continue
            dist = _hamming_distance(icon['phash'], other['phash'])
            if dist <= threshold:
                group.append(other)
                assigned.add(j)

        groups.append(group)

    return groups


def label_icons_with_ai(icon_groups, max_batch=20):
    """Use Claude vision to label icon groups with semantic names and categories.

    Sends icons in batches for efficiency.
    Returns list of dicts: {group_idx, name, category, confidence}
    """
    try:
        import anthropic
        client = anthropic.Anthropic()
    except (ImportError, Exception) as e:
        print(f"  AI labeling unavailable ({e}). Using generic labels.")
        return _generic_labels(icon_groups)

    labels = []
    # Process in batches
    for batch_start in range(0, len(icon_groups), max_batch):
        batch = icon_groups[batch_start:batch_start + max_batch]
        batch_end = min(batch_start + max_batch, len(icon_groups))
        print(f"  Labeling icons {batch_start + 1}-{batch_end} of {len(icon_groups)}...")

        # Build message with all icons in the batch
        content = []
        for i, group in enumerate(batch):
            representative = group[0]
            img_b64 = base64.b64encode(representative['blob']).decode()
            media = representative['content_type']
            content.append({
                "type": "image",
                "source": {"type": "base64", "media_type": media, "data": img_b64}
            })

        content.append({
            "type": "text",
            "text": f"""I'm showing you {len(batch)} icons extracted from a presentation template.
For each icon (in order, 1 to {len(batch)}), provide:
- A short semantic name (2-3 words, lowercase, hyphenated, e.g., "chart-bar", "team-meeting", "cloud-upload")
- A category (one of: business, technology, analytics, people, transport, finance, ui, security, infrastructure, communication, other)

Return ONLY valid JSON array, no explanation:
[{{"name": "icon-name", "category": "category"}}, ...]"""
        })

        try:
            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=2048,
                messages=[{"role": "user", "content": content}],
            )
            result_text = response.content[0].text.strip()
            # Parse JSON from response
            if result_text.startswith('['):
                batch_labels = json.loads(result_text)
            else:
                # Try to find JSON in the response
                import re
                match = re.search(r'\[.*\]', result_text, re.DOTALL)
                if match:
                    batch_labels = json.loads(match.group())
                else:
                    batch_labels = _generic_labels(batch)

            for idx, label in enumerate(batch_labels):
                if idx < len(batch):
                    labels.append({
                        'group_idx': batch_start + idx,
                        'name': label.get('name', f'icon-{batch_start + idx + 1:03d}'),
                        'category': label.get('category', 'other'),
                    })
        except Exception as e:
            print(f"    AI batch failed: {e}")
            for idx in range(len(batch)):
                labels.append({
                    'group_idx': batch_start + idx,
                    'name': f'icon-{batch_start + idx + 1:03d}',
                    'category': 'other',
                })

    return labels


def label_icons_with_claude_code(icon_groups, temp_dir=None, batch_size=25,
                                 max_parallel=5):
    """Label icons by shelling out to Claude Code CLI in parallel batches.

    Saves icon images to a temp directory, then runs parallel `claude -p`
    instances to label each batch. No API key needed — uses the user's
    Claude Code auth.

    Returns list of dicts: {group_idx, name, category}
    """
    import subprocess
    import tempfile
    import concurrent.futures

    if not temp_dir:
        temp_dir = tempfile.mkdtemp(prefix='icon-label-')

    # Save representative icons to temp directory
    icon_dir = os.path.join(temp_dir, 'icons')
    os.makedirs(icon_dir, exist_ok=True)
    for i, group in enumerate(icon_groups):
        rep = group[0]
        path = os.path.join(icon_dir, f'icon-{i+1:03d}.{rep["ext"]}')
        with open(path, 'wb') as f:
            f.write(rep['blob'])

    # Split into batches
    total = len(icon_groups)
    batches = []
    for start in range(0, total, batch_size):
        end = min(start + batch_size, total)
        batch_icons = [f'icon-{i+1:03d}' for i in range(start, end)]
        batches.append((start, end, batch_icons))

    print(f"  Labeling {total} icons in {len(batches)} batches (max {max_parallel} parallel)...")

    all_labels = [None] * total

    def _label_batch(batch_info):
        start, end, icon_names = batch_info
        count = end - start
        icon_list = ', '.join(icon_names[:5]) + (f'... ({count} total)' if count > 5 else '')

        # Build the prompt
        prompt = f"""Look at the {count} icon PNG files in {icon_dir}/ named {', '.join(icon_names)}.

For each icon file (in order), identify what it depicts and provide:
- name: short semantic name (2-3 words, lowercase, hyphenated, e.g. "chart-bar", "team-meeting")
- category: one of business, technology, analytics, people, transport, finance, ui, security, infrastructure, communication, other

Return ONLY a valid JSON array, nothing else:
[{{"name": "icon-name", "category": "category"}}, ...]"""

        try:
            result = subprocess.run(
                ['claude', '-p', prompt, '--output-format', 'text'],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode == 0:
                text = result.stdout.strip()
                # Find JSON array in response
                import re
                match = re.search(r'\[.*\]', text, re.DOTALL)
                if match:
                    batch_labels = json.loads(match.group())
                    return start, batch_labels
        except Exception as e:
            print(f"    Batch {start}-{end} failed: {e}")
        return start, None

    # Run batches in parallel
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_parallel) as executor:
        futures = {executor.submit(_label_batch, b): b for b in batches}
        for future in concurrent.futures.as_completed(futures):
            start, batch_labels = future.result()
            batch_info = futures[future]
            _, end, _ = batch_info
            if batch_labels:
                for idx, label in enumerate(batch_labels):
                    global_idx = start + idx
                    if global_idx < total:
                        all_labels[global_idx] = {
                            'group_idx': global_idx,
                            'name': label.get('name', f'icon-{global_idx+1:03d}'),
                            'category': label.get('category', 'other'),
                        }
                print(f"    Batch {start+1}-{end}: labeled")
            else:
                for idx in range(start, min(end, total)):
                    all_labels[idx] = {
                        'group_idx': idx,
                        'name': f'icon-{idx+1:03d}',
                        'category': 'other',
                    }
                print(f"    Batch {start+1}-{end}: fallback labels")

    # Fill any gaps
    for i in range(total):
        if all_labels[i] is None:
            all_labels[i] = {'group_idx': i, 'name': f'icon-{i+1:03d}', 'category': 'other'}

    labeled = sum(1 for l in all_labels if not l['name'].startswith('icon-'))
    print(f"  Labeled {labeled}/{total} icons with AI names")

    return all_labels


def _generic_labels(icon_groups):
    """Generate generic numbered labels as fallback."""
    return [{'group_idx': i, 'name': f'icon-{i + 1:03d}', 'category': 'other'}
            for i in range(len(icon_groups))]


def save_icon_catalog(icon_groups, labels, output_dir):
    """Save deduplicated icons and generate icon-catalog.json.

    Saves the representative icon from each group.
    Returns the catalog dict.
    """
    icons_dir = os.path.join(output_dir, 'icons')
    os.makedirs(icons_dir, exist_ok=True)

    catalog = {
        "_meta": {
            "total": len(icon_groups),
            "generated_by": "extract_icons.py"
        },
        "icons": {}
    }

    for label in labels:
        idx = label['group_idx']
        if idx >= len(icon_groups):
            continue

        group = icon_groups[idx]
        representative = group[0]
        name = label['name']
        category = label['category']

        # Save the icon file
        filename = f"{name}"
        filepath = os.path.join(icons_dir, f"{filename}.{representative['ext']}")
        with open(filepath, 'wb') as f:
            f.write(representative['blob'])

        # Save color variants too
        for vi, variant in enumerate(group[1:], 1):
            var_path = os.path.join(icons_dir, f"{filename}-v{vi}.{variant['ext']}")
            with open(var_path, 'wb') as f:
                f.write(variant['blob'])

        catalog['icons'][name] = {
            'file': filename,
            'category': category,
            'variants': len(group),
        }

    # Save catalog
    catalog_path = os.path.join(icons_dir, 'icon-catalog.json')
    with open(catalog_path, 'w') as f:
        json.dump(catalog, f, indent=2)

    return catalog


def extract_and_catalog(template_path=None, corpus_dir=None, output_dir=None,
                        use_ai=True):
    """Full pipeline: extract → dedup → label → save.

    Returns (icon_groups, labels, catalog).
    """
    all_icons = []

    if template_path:
        print(f"  Extracting icons from template...")
        icons = extract_icons_from_pptx(template_path)
        all_icons.extend(icons)
        print(f"  Found {len(icons)} unique icons in template")

    if corpus_dir:
        print(f"  Extracting icons from corpus...")
        icons = extract_icons_from_corpus(corpus_dir)
        # Dedup against template icons
        seen = {ic['hash'] for ic in all_icons}
        new = [ic for ic in icons if ic['hash'] not in seen]
        all_icons.extend(new)
        print(f"  Found {len(new)} additional icons from corpus")

    if not all_icons:
        print("  No icons found")
        return [], [], {}

    print(f"  Total unique icons: {len(all_icons)}")

    # Deduplicate by visual similarity
    print(f"  Deduplicating by visual similarity...")
    groups = deduplicate_icons(all_icons)
    print(f"  {len(groups)} unique icon designs ({len(all_icons) - len(groups)} variants merged)")

    # Label with AI
    if use_ai:
        print(f"  Labeling with AI vision...")
        labels = label_icons_with_ai(groups)
    else:
        labels = _generic_labels(groups)

    # Save
    if output_dir:
        catalog = save_icon_catalog(groups, labels, output_dir)
        print(f"  Saved {len(catalog['icons'])} icons to {output_dir}/icons/")
    else:
        catalog = {}

    return groups, labels, catalog


def main():
    parser = argparse.ArgumentParser(description="Extract and catalog icons from PPTX")
    parser.add_argument("--template", help="PPTX template file")
    parser.add_argument("--corpus", help="Directory of PPTX files")
    parser.add_argument("--output", "-o", help="Output directory for icons + catalog")
    parser.add_argument("--label", action="store_true", help="Use AI vision for labeling")
    parser.add_argument("--no-ai", action="store_true", help="Skip AI labeling, use generic names")
    args = parser.parse_args()

    if not args.template and not args.corpus:
        parser.print_help()
        sys.exit(1)

    groups, labels, catalog = extract_and_catalog(
        template_path=args.template,
        corpus_dir=args.corpus,
        output_dir=args.output or '/tmp/extracted-icons',
        use_ai=args.label and not args.no_ai,
    )

    print(f"\nResults: {len(groups)} unique icons, {len(labels)} labeled")


if __name__ == "__main__":
    main()
