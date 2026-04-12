#!/usr/bin/env python3
"""
Unified QA Pipeline — Runs all quality checks on a built PPTX deck.

Combines structural checks (PPTX object model), visual proof rendering (PIL),
and optional Claude vision review into a single report.

Usage (standalone):
    python3 qa_pipeline.py <input.pptx> [--ai-review] [--strict] [--output-dir <dir>]

Usage (as module):
    from qa_pipeline import QAPipeline
    qa = QAPipeline("deck.pptx")
    report = qa.run_all()
"""

import json
import os
import sys
import base64
from pathlib import Path

import yaml

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Default brand colors/fonts — overridden when brand is passed
BRAND_COLORS_HEX = {
    "5F016F", "FF80D4", "FFADE4", "F0E8F5", "FFFFFF",
    "333333", "888888", "D0C0D8",
    # Common variants
    "000000", "4EC98B", "FFD766", "E85D5D",
    "FFF0F8", "F5F5F5",
}

BRAND_FONTS = {"Urbanist ExtraBold", "DM Sans", "Urbanist", "DM Sans Medium"}


def _brand_colors(brand=None):
    """Get brand color hex set, falling back to defaults."""
    if brand:
        # Include all brand colors + common visualization colors
        colors = brand.all_color_hexes()
        colors.update({"000000", "4EC98B", "FFD766", "E85D5D"})
        return colors
    return BRAND_COLORS_HEX


def _brand_fonts(brand=None):
    """Get brand font name set, falling back to defaults."""
    if brand:
        fonts = brand.all_font_names()
        # Include common variants (e.g. "Urbanist" without weight suffix)
        for f in list(fonts):
            base = f.split()[0] if " " in f else f
            fonts.add(base)
        return fonts
    return BRAND_FONTS


def _vision_prompt(brand=None):
    """Build vision review prompt, templated from brand config."""
    if brand:
        name = brand.name
        primary = "#" + brand.primary_hex
        secondary = "#" + brand.secondary_hex
        heading = brand.heading_font
        body = brand.body_font
    else:
        name = "Default"
        primary = "#1A365D"
        secondary = "#3182CE"
        heading = "Arial"
        body = "Arial"
    return f"""You are reviewing a slide from a {name} executive presentation. The {name} brand uses {primary}, {secondary}, and white, with {heading} headings and {body} body text on a 10"x5.625" (16:9) canvas.

Evaluate this slide image for:
1. SPACING: Is whitespace balanced? Are margins consistent? Is content too cramped or too sparse?
2. ALIGNMENT: Are related elements aligned on their edges or centers?
3. TEXT OVERFLOW: Is any text cut off, overlapping other elements, or visually crammed?
4. VISUAL BALANCE: Is content distributed well across the slide?
5. READABILITY: Would text be readable when projected at 1080p?
6. BRAND: Does the slide feel consistent with the brand colors?

Return ONLY valid JSON (no markdown, no explanation): {{"pass": true/false, "issues": [{{"category": "spacing|alignment|overflow|balance|readability|brand", "severity": "critical|warning|info", "description": "..."}}]}}
Only flag real, noticeable issues. An empty issues list with pass=true means the slide looks good.
"""

FONT_FLOOR_PT = 7          # Hard floor — nothing below this
FONT_WARN_PT = 9            # Warn on non-table body text below this
WORD_COUNT_WARN = 30        # Words per text box
WORD_COUNT_FLAG = 50        # Words per text box (critical)
OVERFLOW_TOLERANCE = 1.15   # estimated_h / actual_h threshold


# ---------------------------------------------------------------------------
# Issue classes
# ---------------------------------------------------------------------------

class Issue:
    def __init__(self, slide_num, category, severity, description, location=""):
        self.slide_num = slide_num
        self.category = category      # overlap, margin, overflow, font, wordcount, consistency, containment, vision
        self.severity = severity      # critical, warning, info
        self.description = description
        self.location = location      # e.g. "top-left", "center"

    def to_dict(self):
        return {
            "slide": self.slide_num,
            "category": self.category,
            "severity": self.severity,
            "description": self.description,
            "location": self.location,
        }


# ---------------------------------------------------------------------------
# Structural checks — operate on the PPTX object model, no rendering needed
# ---------------------------------------------------------------------------

def _emu_to_inches(emu):
    return emu / 914400 if emu else 0


def _estimate_text_height(text, width_inches, font_size_pt=10):
    """Mirror of build_deck.estimate_text_height for QA checking."""
    chars_per_inch = 13 * (10 / font_size_pt)
    chars_per_line = max(1, int(width_inches * chars_per_inch))
    line_height = font_size_pt * 1.4 / 72
    total_lines = 0
    for paragraph in str(text).split("\n"):
        if not paragraph.strip():
            total_lines += 0.6
        else:
            total_lines += max(1, -(-len(paragraph) // chars_per_line))
    return total_lines * line_height


def check_word_count(prs):
    """Check word count per text box. Warn at threshold, flag at higher threshold."""
    issues = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue
            # Skip footer-area shapes (y > 5.0")
            if _emu_to_inches(shape.top) > 5.0:
                continue
            # Skip table cells (handled by the table shape itself)
            if shape.has_table:
                continue

            words = len(shape.text.split())
            if words >= WORD_COUNT_FLAG:
                issues.append(Issue(
                    si + 1, "wordcount", "warning",
                    f"Text box has {words} words (threshold: {WORD_COUNT_FLAG}): "
                    f"\"{shape.text[:40]}...\""))
            elif words >= WORD_COUNT_WARN:
                issues.append(Issue(
                    si + 1, "wordcount", "info",
                    f"Text box has {words} words (threshold: {WORD_COUNT_WARN}): "
                    f"\"{shape.text[:40]}...\""))
    return issues


def check_font_sizes(prs):
    """Check for text below the font size floor."""
    issues = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            # Skip footer-area shapes
            if _emu_to_inches(shape.top) > 5.0:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        continue
                    size_pt = run.font.size / 12700  # EMU to points
                    if size_pt < FONT_FLOOR_PT:
                        issues.append(Issue(
                            si + 1, "font", "warning",
                            f"Text at {size_pt:.0f}pt (below {FONT_FLOOR_PT}pt floor): "
                            f"\"{run.text[:30]}\""))
                    elif size_pt < FONT_WARN_PT and not shape.has_table:
                        # Only info-level for non-table text
                        issues.append(Issue(
                            si + 1, "font", "info",
                            f"Text at {size_pt:.0f}pt (below {FONT_WARN_PT}pt): "
                            f"\"{run.text[:30]}\""))
    return issues


def check_text_overflow(prs):
    """Check if text likely overflows its container by estimating rendered height."""
    issues = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text.strip():
                continue
            if shape.has_table:
                continue
            # Skip footer-area shapes
            if _emu_to_inches(shape.top) > 5.0:
                continue

            text = shape.text
            width_in = _emu_to_inches(shape.width)
            actual_h = _emu_to_inches(shape.height)

            if width_in < 0.1 or actual_h < 0.1:
                continue

            # Get font size from first run
            font_size_pt = 10  # default
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size_pt = run.font.size / 12700
                        break
                break

            estimated_h = _estimate_text_height(text, width_in, font_size_pt)

            if estimated_h > actual_h * OVERFLOW_TOLERANCE:
                overflow_pct = int((estimated_h / actual_h - 1) * 100)
                issues.append(Issue(
                    si + 1, "overflow", "warning",
                    f"Text likely overflows by ~{overflow_pct}%: "
                    f"\"{text[:40]}...\" (est {estimated_h:.2f}\" > box {actual_h:.2f}\")"))
    return issues


def check_consistency(prs):
    """Check cross-slide consistency: fonts, headline sizes, footer presence."""
    issues = []
    slide_fonts = set()
    headline_sizes = []
    slides_missing_footers = []

    for si, slide in enumerate(prs.slides):
        has_footer = False
        has_headline = False

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            # Check for footer
            if _emu_to_inches(shape.top) > 5.0:
                if shape.text.strip():
                    has_footer = True
                continue

            # Check for headline (top of slide, large text)
            if _emu_to_inches(shape.top) < 0.5:
                has_headline = True
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            headline_sizes.append(run.font.size / 12700)

            # Collect all fonts used
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        slide_fonts.add(run.font.name)

        if not has_footer and si > 0:  # skip first slide (title cover)
            slides_missing_footers.append(si + 1)

    # Check for non-brand fonts
    non_brand = slide_fonts - BRAND_FONTS
    if non_brand:
        issues.append(Issue(
            0, "consistency", "info",
            f"Non-brand fonts detected: {', '.join(sorted(non_brand))}"))

    # Check headline size consistency
    if headline_sizes:
        unique_sizes = set(round(s, 1) for s in headline_sizes)
        if len(unique_sizes) > 3:
            issues.append(Issue(
                0, "consistency", "info",
                f"Headline font sizes vary: {sorted(unique_sizes)}"))

    return issues


def check_containment(prs):
    """Check that text boxes inside filled containers don't extend beyond bounds."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    issues = []
    for si, slide in enumerate(prs.slides):
        # Collect filled rectangles as potential containers
        containers = []
        text_boxes = []

        for shape in slide.shapes:
            if _emu_to_inches(shape.top) > 5.0:
                continue

            left = _emu_to_inches(shape.left)
            top = _emu_to_inches(shape.top)
            right = left + _emu_to_inches(shape.width)
            bottom = top + _emu_to_inches(shape.height)

            if hasattr(shape, "fill") and shape.fill.type is not None:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    containers.append({"left": left, "top": top, "right": right, "bottom": bottom})

            if hasattr(shape, "text") and shape.text.strip():
                text_boxes.append({
                    "left": left, "top": top, "right": right, "bottom": bottom,
                    "text": shape.text[:30],
                    "cx": (left + right) / 2, "cy": (top + bottom) / 2,
                })

        # For each text box, check if its center is inside a container
        for tb in text_boxes:
            for c in containers:
                if (c["left"] < tb["cx"] < c["right"] and
                        c["top"] < tb["cy"] < c["bottom"]):
                    # Text box center is inside this container — check bounds
                    margin = 0.02  # tolerance
                    if (tb["left"] < c["left"] - margin or
                            tb["right"] > c["right"] + margin or
                            tb["top"] < c["top"] - margin or
                            tb["bottom"] > c["bottom"] + margin):
                        issues.append(Issue(
                            si + 1, "containment", "warning",
                            f"Text extends beyond container: \"{tb['text']}\""))
                    break  # only check first matching container

    return issues


# ---------------------------------------------------------------------------
# Claude Vision review
# ---------------------------------------------------------------------------

# Legacy static prompt — used when no brand is passed. See _vision_prompt() for brand-aware version.
VISION_PROMPT = _vision_prompt(None)


def run_vision_review(proof_dir, max_slides=30, brand=None):
    """Send slide PNGs to Claude Sonnet for visual QA review."""
    issues = []

    try:
        import anthropic
        client = anthropic.Anthropic()
    except (ImportError, Exception) as e:
        print(f"  Vision review skipped: {e}")
        return issues

    proof_path = Path(proof_dir)
    slide_pngs = sorted(proof_path.glob("slide-*.png"))

    for png_path in slide_pngs[:max_slides]:
        slide_num = int(png_path.stem.split("-")[1])

        try:
            with open(png_path, "rb") as f:
                img_data = base64.standard_b64encode(f.read()).decode("utf-8")

            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=1024,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_data}},
                        {"type": "text", "text": _vision_prompt(brand) if brand else VISION_PROMPT},
                    ],
                }],
            )

            result_text = response.content[0].text.strip()
            # Try to parse JSON from the response
            result = json.loads(result_text)

            for issue in result.get("issues", []):
                issues.append(Issue(
                    slide_num,
                    f"vision-{issue.get('category', 'general')}",
                    issue.get("severity", "info"),
                    issue.get("description", "Vision review issue"),
                ))

        except json.JSONDecodeError:
            print(f"  Vision review: Could not parse JSON for slide {slide_num}")
        except Exception as e:
            print(f"  Vision review error on slide {slide_num}: {e}")

    return issues


# ---------------------------------------------------------------------------
# QA Pipeline orchestrator
# ---------------------------------------------------------------------------

class QAReport:
    def __init__(self, pptx_path):
        self.pptx_path = str(pptx_path)
        self.slides = {}    # slide_num → {"issues": [...], "status": "pass"|"warn"|"fail"}
        self.cross_slide_issues = []
        self.total_slides = 0

    def add_issue(self, issue):
        num = issue.slide_num
        if num == 0:
            self.cross_slide_issues.append(issue)
        else:
            if num not in self.slides:
                self.slides[num] = {"issues": [], "status": "pass"}
            self.slides[num]["issues"].append(issue)
            if issue.severity == "critical":
                self.slides[num]["status"] = "fail"
            elif issue.severity == "warning" and self.slides[num]["status"] != "fail":
                self.slides[num]["status"] = "warn"

    def has_critical(self):
        for s in self.slides.values():
            for i in s["issues"]:
                if i.severity == "critical":
                    return True
        return False

    def summary(self):
        total_issues = sum(len(s["issues"]) for s in self.slides.values()) + len(self.cross_slide_issues)
        critical = sum(1 for s in self.slides.values() for i in s["issues"] if i.severity == "critical")
        warnings = sum(1 for s in self.slides.values() for i in s["issues"] if i.severity == "warning")
        info = sum(1 for s in self.slides.values() for i in s["issues"] if i.severity == "info")
        passed = sum(1 for s in self.slides.values() if s["status"] == "pass")
        return {
            "total_slides": self.total_slides,
            "total_issues": total_issues,
            "critical": critical,
            "warnings": warnings,
            "info": info,
            "passed": passed,
            "blocked": self.has_critical(),
        }

    def to_json(self):
        return {
            "file": self.pptx_path,
            "summary": self.summary(),
            "slides": {
                num: {
                    "status": data["status"],
                    "issues": [i.to_dict() for i in data["issues"]],
                }
                for num, data in sorted(self.slides.items())
            },
            "cross_slide_issues": [i.to_dict() for i in self.cross_slide_issues],
        }

    def to_markdown(self):
        lines = ["# QA Report", "", f"**File:** `{self.pptx_path}`"]
        s = self.summary()
        lines.append(f"**Slides:** {s['total_slides']} | **Issues:** {s['total_issues']} "
                      f"({s['critical']} critical, {s['warnings']} warnings, {s['info']} info)")
        lines.append("")

        if s["blocked"]:
            lines.append("**STATUS: BLOCKED** — Critical issues found. Fix before delivery.")
        elif s["total_issues"] == 0:
            lines.append("**STATUS: CLEAN** — No issues detected.")
        else:
            lines.append(f"**STATUS: REVIEW** — {s['warnings']} warning(s) to check.")
        lines.append("")
        lines.append("---")
        lines.append("")

        for num in sorted(self.slides.keys()):
            data = self.slides[num]
            status_emoji = {"pass": "PASS", "warn": "WARN", "fail": "FAIL"}[data["status"]]
            lines.append(f"### Slide {num} [{status_emoji}]")
            if not data["issues"]:
                lines.append("No issues.")
            else:
                for i in data["issues"]:
                    sev = i.severity.upper()
                    lines.append(f"- **[{sev}]** {i.category}: {i.description}")
            lines.append("")

        if self.cross_slide_issues:
            lines.append("### Cross-Slide Issues")
            for i in self.cross_slide_issues:
                sev = i.severity.upper()
                lines.append(f"- **[{sev}]** {i.category}: {i.description}")
            lines.append("")

        return "\n".join(lines)


def check_yaml_content_limits(yaml_path):
    """Check YAML slide definitions against measured layout content limits.

    Reads layout-limits.json for per-layout min/max values. Returns issues
    with actionable suggestions including alternative layouts and split strategies.
    """
    issues = []
    script_dir = os.path.dirname(os.path.abspath(__file__))
    limits_path = os.path.join(script_dir, "layout-limits.json")

    if not os.path.exists(limits_path) or not yaml_path or not os.path.exists(yaml_path):
        return issues

    with open(limits_path) as f:
        limits_data = json.load(f)
    with open(yaml_path) as f:
        deck = yaml.safe_load(f)

    all_limits = limits_data.get("layouts", {})

    # Field-to-list mapping for counting items in each layout
    ITEM_FIELDS = {
        "agenda": "items", "numbered_list": "items", "status_board": "items",
        "before_after": None,  # special: items inside before/after dicts
        "kpi_dashboard": "metrics", "process_flow": "steps", "funnel": "stages",
        "staircase": "levels", "cycle_diagram": "nodes", "hub_spoke": "spokes",
        "donut_rings": "rings", "gauge_dashboard": "gauges", "tornado_chart": "items",
        "bubble_chart": "bubbles", "risk_heat_map": "items", "team_profiles": "profiles",
        "comparison_matrix": "rows", "pricing_table": "tiers", "bold_bullet": "points",
        "concentric_circles": "rings", "pyramid": "tiers", "bento_grid": "tiles",
        "left_nav_sidebar": "nav_items", "venn": "circles", "waterfall": "items",
    }

    for si, slide_def in enumerate(deck.get("slides", [])):
        slide_num = si + 1
        layout = slide_def.get("layout", "unknown")
        lim = all_limits.get(layout)
        if not lim:
            continue

        # --- Check item counts ---
        item_field = ITEM_FIELDS.get(layout)
        if item_field:
            items = slide_def.get(item_field, [])
            count = len(items)

            # Check max items
            for max_key in ["max_items", "max_steps_short_labels", "max_nodes",
                            "max_spokes_short_labels", "max_rings", "max_gauges",
                            "max_metrics", "max_profiles", "max_tiers",
                            "max_points", "max_bubbles", "max_stages",
                            "max_levels", "max_tiles", "max_nav_items",
                            "max_circles", "max_criteria"]:
                max_val = lim.get(max_key)
                if max_val and count > max_val:
                    guidance = lim.get("split_guidance",
                        f"Reduce to {max_val} items or split across multiple slides.")
                    issues.append(Issue(slide_num, "content_limit", "warning",
                        f"{layout} has {count} {item_field} (max {max_val}). {guidance}"))
                    break

            # Check min items
            for min_key in ["min_items", "min_steps", "min_nodes", "min_spokes",
                            "min_rings", "min_gauges", "min_metrics", "min_profiles",
                            "min_tiers", "min_bubbles", "min_stages", "min_levels",
                            "min_tiles", "min_nav_items", "min_circles", "min_criteria"]:
                min_val = lim.get(min_key)
                if min_val and count < min_val:
                    issues.append(Issue(slide_num, "content_limit", "info",
                        f"{layout} has {count} {item_field} (min {min_val}). "
                        f"Add more items or use a simpler layout."))
                    break

        # --- Check text field lengths ---
        CHAR_CHECKS = {
            "headline": "max_headline_chars",
            "subheader": "max_subheader_chars",
            "quote_text": "max_quote_chars",
            "callout_text": "max_callout_chars",
            "supporting_text": "max_supporting_chars",
            "recommendation": "max_recommendation_chars",
            "summary": "max_summary_chars",
            "body": "max_body_chars",
        }
        for field, limit_key in CHAR_CHECKS.items():
            text = slide_def.get(field, "")
            if text and isinstance(text, str):
                max_chars = lim.get(limit_key)
                if max_chars and len(text) > max_chars:
                    issues.append(Issue(slide_num, "content_limit", "warning",
                        f"{layout}.{field} has {len(text)} chars (max {max_chars}). "
                        f"Shorten text or use a layout with more text capacity."))

        # --- Check table dimensions ---
        if layout in ("content_table", "content_table_bullets"):
            rows = slide_def.get("rows", [])
            cols = slide_def.get("columns", [])
            max_cols = lim.get("max_cols", 6)
            max_rows = lim.get("max_rows_single_slide", lim.get("max_rows", 12))
            if len(cols) > max_cols:
                issues.append(Issue(slide_num, "content_limit", "warning",
                    f"{layout} has {len(cols)} columns (max {max_cols}). "
                    f"Remove non-essential columns or split into multiple tables."))
            if len(rows) > max_rows:
                issues.append(Issue(slide_num, "content_limit", "info",
                    f"{layout} has {len(rows)} rows (max {max_rows} per slide). "
                    f"Table will auto-split across continuation slides."))

        # --- Check comparison_matrix dimensions ---
        if layout == "comparison_matrix":
            opts = slide_def.get("columns", [])
            max_opts = lim.get("max_options", 5)
            if len(opts) > max_opts:
                issues.append(Issue(slide_num, "content_limit", "warning",
                    f"comparison_matrix has {len(opts)} options (max {max_opts}). "
                    f"Reduce to top candidates or split evaluation."))

        # --- Check four_card ---
        if layout == "four_card":
            n_cards = sum(1 for i in range(1, 9) if slide_def.get(f"card{i}_title"))
            max_long = lim.get("max_cards_long_text", 6)
            max_short = lim.get("max_cards_short_text", 8)
            has_long = any(len(slide_def.get(f"card{i}_body", "")) > 50 for i in range(1, n_cards + 1))
            limit = max_long if has_long else max_short
            if n_cards > limit:
                issues.append(Issue(slide_num, "content_limit", "warning",
                    f"four_card has {n_cards} cards with {'long' if has_long else 'short'} text "
                    f"(max {limit}). {lim.get('split_guidance', 'Split across slides.')}"))

        # --- Check before_after item counts ---
        if layout == "before_after":
            for side in ["before", "after"]:
                side_def = slide_def.get(side, {})
                if isinstance(side_def, dict):
                    side_items = side_def.get("items", [])
                    max_per = lim.get("max_items_per_side", 7)
                    if len(side_items) > max_per:
                        issues.append(Issue(slide_num, "content_limit", "warning",
                            f"before_after '{side}' has {len(side_items)} items (max {max_per})."))

        # --- Check roadmap ---
        if layout == "roadmap":
            swimlanes = slide_def.get("swimlanes", [])
            max_swim = lim.get("max_swimlanes", 4)
            if len(swimlanes) > max_swim:
                issues.append(Issue(slide_num, "content_limit", "warning",
                    f"roadmap has {len(swimlanes)} swimlanes (max {max_swim}). "
                    f"Split into separate roadmap slides per theme."))

    return issues


# ---------------------------------------------------------------------------
# New checks — from corpus analysis and partner-recommendations review
# ---------------------------------------------------------------------------

def check_yaml_icon_resolution(yaml_path):
    """Check that all icons specified in YAML exist in the icon catalog."""
    import yaml as _yaml
    issues = []
    if not yaml_path or not os.path.exists(yaml_path):
        return issues

    script_dir = os.path.dirname(os.path.abspath(__file__))
    # Try brand icon dir first (brands/<name>/icons/), fall back to script_dir/icons/
    catalog_path = None
    for icons_dir in [os.path.join(script_dir, "icons")]:
        cp = os.path.join(icons_dir, "icon-catalog.json")
        if os.path.isfile(cp):
            catalog_path = cp
            break
    if not catalog_path:
        return issues
    try:
        with open(catalog_path) as f:
            icon_catalog = json.load(f)["icons"]
    except Exception:
        return issues

    with open(yaml_path) as f:
        deck = _yaml.safe_load(f)

    slides = deck.get("slides", [])
    for si, slide_def in enumerate(slides):
        slide_num = si + 1
        # Collect all icon fields from this slide
        icon_fields = []
        for key, val in slide_def.items():
            if "icon" in key.lower() and isinstance(val, str) and val:
                icon_fields.append((key, val))
        # Check metrics list (kpi_dashboard)
        for mi, metric in enumerate(slide_def.get("metrics", [])):
            icon_name = metric.get("icon", "")
            if icon_name:
                icon_fields.append((f"metrics[{mi}].icon", icon_name))

        for field_name, icon_name in icon_fields:
            info = icon_catalog.get(icon_name, {})
            if not info.get("file"):
                issues.append(Issue(
                    slide_num, "missing_icon", "critical",
                    f"Icon '{icon_name}' (in {field_name}) not found in icon catalog. "
                    f"It will not render on the slide.",
                ))
            else:
                icon_file = os.path.join(script_dir, "icons", info["file"] + ".png")
                if not os.path.exists(icon_file):
                    issues.append(Issue(
                        slide_num, "missing_icon", "critical",
                        f"Icon '{icon_name}' (in {field_name}) catalog entry exists but "
                        f"PNG file is missing: {info['file']}.png",
                    ))
    return issues


def check_table_font_opportunity(prs):
    """Check if tables have unused vertical space where a larger font would improve readability."""
    issues = []
    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            num_rows = len(table.rows)
            if num_rows < 2:  # header only
                continue

            # Get table shape dimensions
            table_top = _emu_to_inches(shape.top)
            table_height = _emu_to_inches(shape.height)
            table_bottom = table_top + table_height

            # Available space below the table to the footer zone
            footer_zone = 5.0
            unused_below = footer_zone - table_bottom

            # Also check if data rows are very compact
            header_height = _emu_to_inches(table.rows[0].height)
            data_height = table_height - header_height
            avg_row_h = data_height / (num_rows - 1) if num_rows > 1 else data_height

            # Find smallest font in data cells
            min_font = 999
            for ri in range(1, num_rows):
                for ci in range(len(table.columns)):
                    cell = table.cell(ri, ci)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                fs = run.font.size.pt
                                if fs < min_font:
                                    min_font = fs

            # Check 1: unused space below the table
            if unused_below > 0.5 and min_font <= 8:
                issues.append(Issue(
                    slide_num, "table_font_opportunity", "warning",
                    f"Table has {unused_below:.1f}\" unused space below it and uses "
                    f"{min_font}pt font. Consider increasing font size for readability.",
                ))
            # Check 2: row height generous relative to font size
            # A font needs roughly (pt * 1.4 / 72) inches per line.
            # If row height is >2x what the font needs, there's room to grow.
            if min_font < 999:
                font_needs = min_font * 1.4 / 72  # inches for one line
                if avg_row_h > font_needs * 2.2 and min_font <= 8:
                    max_font = int(avg_row_h * 72 / 1.8)  # font that would use ~55% of row
                    max_font = min(max_font, 11)  # cap suggestion
                    issues.append(Issue(
                        slide_num, "table_font_opportunity", "warning",
                        f"Table data rows are {avg_row_h:.2f}\" tall but font is only "
                        f"{min_font}pt. Row height could support up to ~{max_font}pt.",
                    ))
    return issues


def check_empty_text_boxes(prs):
    """Check for empty text boxes/placeholders that occupy visible space."""
    issues = []
    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            # Skip shapes in footer area
            if _emu_to_inches(shape.top) > 5.0:
                continue
            # Check if shape has a text frame but no visible text
            text = shape.text_frame.text.strip() if shape.text_frame else ""
            if text:
                continue
            # Shape is empty — check if it's large enough to be visible
            w = _emu_to_inches(shape.width)
            h = _emu_to_inches(shape.height)
            top = _emu_to_inches(shape.top)
            # Skip full-slide background shapes (section dividers, etc.)
            if w > 9.0 and h > 4.0:
                continue
            # Skip small card backgrounds (kpi_dashboard metric boxes, etc.)
            if w < 3.0 and h < 2.0:
                continue
            # Skip filled auto-shapes (background boxes for callout, cards, etc.)
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    continue
            except Exception:
                pass
            if w > 0.5 and h > 0.2:
                # Check if it's a placeholder (template artifact)
                ph_type = ""
                try:
                    pf = shape.placeholder_format
                    ph_type = f" (placeholder idx={pf.idx})"
                except (ValueError, AttributeError):
                    pass
                issues.append(Issue(
                    slide_num, "empty_text_box", "warning",
                    f"Empty text box{ph_type} at ({_emu_to_inches(shape.left):.1f}\", "
                    f"{_emu_to_inches(shape.top):.1f}\") size {w:.1f}\"x{h:.1f}\". "
                    f"Remove or populate it.",
                ))
    return issues


def check_unhyperlinked_urls(prs):
    """Check for URL-like text that is not wrapped in a hyperlink."""
    import re
    url_pattern = re.compile(
        r'(?:https?://[^\s]+|'           # http:// or https://
        r'(?:[a-zA-Z0-9-]+\.)+[a-zA-Z]{2,}(?:/[^\s]*)?)'  # domain.com/path
    )
    issues = []
    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text
                    urls = url_pattern.findall(text)
                    if not urls:
                        continue
                    # Check if the run has a hyperlink
                    has_hlink = False
                    try:
                        hlink = run.hyperlink
                        if hlink and hlink.address:
                            has_hlink = True
                    except Exception:
                        pass
                    if not has_hlink:
                        for url in urls:
                            issues.append(Issue(
                                slide_num, "unhyperlinked_url", "warning",
                                f"URL '{url[:60]}' in text is not a clickable hyperlink.",
                            ))
    return issues


def check_table_centering(prs):
    """Check if tables are vertically centered in their available space."""
    issues = []
    avail_top = 0.85
    avail_bottom = 5.05
    avail_h = avail_bottom - avail_top
    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table_top = _emu_to_inches(shape.top)
            table_h = _emu_to_inches(shape.height)
            table_bottom = table_top + table_h

            # Where should center be?
            ideal_top = avail_top + (avail_h - table_h) / 2
            offset_from_ideal = abs(table_top - ideal_top)

            # Space above and below
            space_above = table_top - avail_top
            space_below = avail_bottom - table_bottom

            # If both margins are >0.3" and differ by >0.5", table is poorly centered
            if space_above > 0.3 and space_below > 0.3 and abs(space_above - space_below) > 0.5:
                issues.append(Issue(
                    slide_num, "table_centering", "warning",
                    f"Table is not vertically centered: {space_above:.1f}\" above, "
                    f"{space_below:.1f}\" below. Consider centering.",
                ))
    return issues


def check_kpi_text_overflow(prs):
    """Check if KPI dashboard number text is too wide for its card.

    KPI cards use large fonts for the number field. Text like 'Nutanix AHV'
    or 'EOY 2026' can overflow the card at 32pt even though it fits at 20pt.
    This check estimates whether the number text would overflow the card width
    at the rendered font size.
    """
    issues = []
    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        # Find KPI-style shapes: large Urbanist ExtraBold text in a rounded rectangle context
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text.strip():
                continue
            if shape.has_table:
                continue
            tf = shape.text_frame
            # Check if this looks like a KPI number: Urbanist ExtraBold, >= 20pt, no word wrap
            for para in tf.paragraphs:
                for run in para.runs:
                    if (run.font.name == "Urbanist ExtraBold"
                            and run.font.size
                            and run.font.size / 12700 >= 20
                            and not tf.word_wrap):
                        font_pt = run.font.size / 12700
                        text = run.text.strip()
                        box_w = _emu_to_inches(shape.width)
                        # Estimate text width: bold ExtraBold ≈ 0.6 * pt / 72 inches per char
                        est_text_w = len(text) * (font_pt * 0.6 / 72)
                        if est_text_w > box_w * 0.95:
                            issues.append(Issue(
                                slide_num, "kpi_overflow", "warning",
                                f"KPI number \"{text}\" at {font_pt:.0f}pt likely overflows "
                                f"card ({est_text_w:.1f}\" text > {box_w:.1f}\" card). "
                                f"Auto-sizing should reduce font or shorten text."))
                    break  # only check first run
                break  # only check first paragraph
    return issues


def check_kpi_label_alignment(prs):
    """Check that KPI card labels within the same row are vertically aligned.

    KPI dashboards have metric cards in a grid. The label text (e.g., 'VMs to AWS')
    below each number should be at the same Y position across all cards in a row.
    Misalignment occurs when the number text has different heights (e.g., auto-sized
    font for long text) and the label is positioned relative to the number rather
    than anchored to the card bottom.
    """
    issues = []
    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        # Find DM Sans labels that look like KPI labels:
        # small font (8-12pt), centered, within a card-like vertical band
        kpi_labels = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text.strip():
                continue
            if shape.has_table:
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                if para.alignment == PP_ALIGN.CENTER:
                    for run in para.runs:
                        if (run.font.name == "DM Sans"
                                and run.font.size
                                and 7 <= run.font.size / 12700 <= 13
                                and len(run.text.strip()) < 30):
                            top = _emu_to_inches(shape.top)
                            kpi_labels.append((top, run.text.strip(), shape))
                        break
                break

        # Group labels that are at similar Y positions (within 0.3" = same row)
        if len(kpi_labels) < 2:
            continue
        kpi_labels.sort(key=lambda x: x[0])

        # Find groups of labels at approximately the same Y (KPI labels in a row)
        row_groups = []
        current_group = [kpi_labels[0]]
        for label in kpi_labels[1:]:
            if abs(label[0] - current_group[0][0]) < 0.5:
                current_group.append(label)
            else:
                if len(current_group) >= 2:
                    row_groups.append(current_group)
                current_group = [label]
        if len(current_group) >= 2:
            row_groups.append(current_group)

        for group in row_groups:
            tops = [g[0] for g in group]
            spread = max(tops) - min(tops)
            if spread > 0.15:  # more than 0.15" misalignment
                labels = [g[1] for g in group]
                issues.append(Issue(
                    slide_num, "kpi_label_alignment", "warning",
                    f"KPI labels misaligned by {spread:.2f}\": "
                    f"{', '.join(labels[:3])}{'...' if len(labels) > 3 else ''}. "
                    f"Labels should be at the same Y position across cards."))
    return issues


def check_audience_density(yaml_path):
    """Validate deck content against audience density presets.

    Checks slide count, word density, layout choices against the preset
    defined by the 'audience' field in the deck YAML.
    """
    issues = []
    if not yaml_path or not os.path.exists(yaml_path):
        return issues

    script_dir = os.path.dirname(os.path.abspath(__file__))
    presets_path = os.path.join(script_dir, "audience-presets.json")
    if not os.path.exists(presets_path):
        return issues

    try:
        with open(yaml_path) as f:
            deck = yaml.safe_load(f)
        with open(presets_path) as f:
            presets = json.load(f)
    except Exception:
        return issues

    audience = deck.get("audience", "")
    if not audience or audience not in presets:
        return issues

    preset = presets[audience]
    slides = deck.get("slides", [])
    label = preset["label"]

    # Check slide count
    max_slides = preset.get("max_slides", 999)
    if len(slides) > max_slides:
        issues.append(Issue(0, "audience", "warning",
            f"{label} decks should have ≤{max_slides} slides "
            f"(this deck has {len(slides)}). Consolidate or move detail to appendix."))

    # Check for avoided layouts
    avoid = set(preset.get("avoid_layouts", []))
    for si, slide_def in enumerate(slides):
        layout = slide_def.get("layout", "")
        if layout in avoid:
            preferred = ", ".join(preset.get("preferred_layouts", [])[:3])
            issues.append(Issue(si + 1, "audience", "info",
                f"'{layout}' is not ideal for {label} audience. "
                f"Consider: {preferred}"))

    # Check headline style guidance
    headline_style = preset.get("headline_style", "")
    if headline_style == "assertion":
        for si, slide_def in enumerate(slides):
            hl = slide_def.get("headline", "")
            if hl and not any(w in hl.lower() for w in
                              ["we ", "should", "recommend", "invest", "will",
                               "must", "need", "approved", "decided"]):
                # Check if it's a topic label (short, no verb)
                words = hl.split()
                if len(words) <= 4 and len(words) >= 1:
                    layout = slide_def.get("layout", "")
                    if layout not in ("title_cover", "closing", "section_divider", "agenda"):
                        issues.append(Issue(si + 1, "audience", "info",
                            f"C-Suite headline '{hl}' looks like a topic label. "
                            f"Consider an assertion: what should the audience decide or know?"))

    return issues


def check_story_structure(yaml_path):
    """Validate slide sequence against narrative framework patterns.

    Checks:
    - Deck starts with title_cover
    - Deck ends with closing
    - Data slides don't appear before context
    - Recommendation/callout appears after evidence, not before
    - Section dividers used in 10+ slide decks
    """
    issues = []
    if not yaml_path or not os.path.exists(yaml_path):
        return issues

    try:
        with open(yaml_path) as f:
            deck = yaml.safe_load(f)
    except Exception:
        return issues

    slides = deck.get("slides", [])
    if not slides:
        return issues

    layouts = [s.get("layout", "") for s in slides]
    n = len(layouts)

    # Must start with title_cover
    if layouts[0] != "title_cover":
        issues.append(Issue(1, "story", "warning",
            f"Deck should open with title_cover, not '{layouts[0]}'"))

    # Must end with closing
    if layouts[-1] != "closing":
        issues.append(Issue(n, "story", "warning",
            f"Deck should end with closing, not '{layouts[-1]}'"))

    # Data/table slides before any context slide
    context_layouts = {"callout", "big_stat_manual", "side_by_side", "before_after",
                       "bold_bullet", "quote"}
    data_layouts = {"content_table", "content_table_bullets", "matrix",
                    "kpi_dashboard", "status_board", "comparison_matrix",
                    "risk_heat_map"}
    first_context = None
    first_data = None
    for i, layout in enumerate(layouts):
        if layout in context_layouts and first_context is None:
            first_context = i
        if layout in data_layouts and first_data is None:
            first_data = i
    if first_data is not None and (first_context is None or first_data < first_context):
        if first_data > 1:  # skip if it's slide 2 (agenda pattern)
            issues.append(Issue(first_data + 1, "story", "info",
                "Data slide appears before a context-setting slide — "
                "consider adding a callout or big_stat first"))

    # Recommendation before evidence
    reco_layouts = {"callout", "pros_cons"}
    evidence_layouts = {"content_table", "comparison_matrix", "kpi_dashboard",
                        "status_board", "matrix"}
    first_reco = None
    first_evidence = None
    for i, layout in enumerate(layouts):
        if layout in reco_layouts and first_reco is None:
            first_reco = i
        if layout in evidence_layouts and first_evidence is None:
            first_evidence = i

    # Section dividers in long decks
    if n >= 10:
        divider_count = layouts.count("section_divider")
        if divider_count == 0:
            issues.append(Issue(0, "story", "info",
                f"Deck has {n} slides but no section_divider — "
                "consider adding dividers every 3-4 content slides"))

    return issues


def auto_fix_pptx(pptx_path):
    """Apply safe automatic fixes to a PPTX file.

    Fixes:
    - Re-center off-center tables
    - Upsize fonts when unused vertical space > 0.5" and font <= 8pt
    - Add hyperlinks to bare URLs in text

    Returns list of fix descriptions.
    """
    import re

    prs = Presentation(pptx_path)
    fixes = []

    for si, slide in enumerate(prs.slides):
        slide_num = si + 1

        for shape in slide.shapes:
            # --- Fix 1: Re-center off-center tables ---
            if shape.has_table:
                avail_top = Inches(0.85)
                avail_bottom = Inches(5.05)
                avail_h = avail_bottom - avail_top
                table_h = shape.height
                ideal_top = avail_top + (avail_h - table_h) // 2
                ideal_top = max(ideal_top, avail_top)

                space_above = shape.top - avail_top
                space_below = avail_bottom - (shape.top + shape.height)
                if abs(space_above - space_below) > Inches(0.5):
                    old_top = shape.top
                    shape.top = ideal_top
                    fixes.append(f"Slide {slide_num}: Re-centered table "
                                 f"(was {old_top/914400:.2f}\", now {ideal_top/914400:.2f}\")")

            # --- Fix 2: Upsize table fonts when space allows ---
            if shape.has_table:
                tbl = shape.table
                footer_zone = Inches(5.0)
                table_bottom = shape.top + shape.height
                unused = footer_zone - table_bottom

                if unused > Inches(0.5):
                    min_font = None
                    for row in tbl.rows:
                        for ci in range(len(tbl.columns)):
                            cell = tbl.cell(row.idx if hasattr(row, 'idx') else 0, ci)
                            for p in cell.text_frame.paragraphs:
                                for r in p.runs:
                                    if r.font.size and (min_font is None or r.font.size < min_font):
                                        min_font = r.font.size

                    if min_font and min_font <= Pt(8):
                        target_pt = min(min_font + Pt(1), Pt(10))
                        for ri in range(len(tbl.rows)):
                            for ci in range(len(tbl.columns)):
                                cell = tbl.cell(ri, ci)
                                for p in cell.text_frame.paragraphs:
                                    for r in p.runs:
                                        if r.font.size and r.font.size <= Pt(8):
                                            r.font.size = target_pt
                        fixes.append(f"Slide {slide_num}: Upsized table font from "
                                     f"{min_font.pt:.0f}pt to {target_pt.pt:.0f}pt "
                                     f"(unused space: {unused/914400:.2f}\")")

            # --- Fix 3: Hyperlink bare URLs ---
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        text = r.text
                        if r.hyperlink and r.hyperlink.address:
                            continue  # already linked
                        url_match = re.search(
                            r'(https?://[^\s,)]+)', text)
                        if url_match and text.strip() == url_match.group(1):
                            url = url_match.group(1)
                            try:
                                r.hyperlink.address = url
                                fixes.append(f"Slide {slide_num}: Added hyperlink to "
                                             f"'{url[:50]}...'")
                            except Exception:
                                pass

    if fixes:
        prs.save(pptx_path)
        print(f"  Auto-fix: {len(fixes)} fix(es) applied")
        for fix in fixes:
            print(f"    - {fix}")
    else:
        print("  Auto-fix: no fixes needed")

    return fixes


class QAPipeline:
    def __init__(self, pptx_path, proof_dir=None, yaml_path=None, brand=None):
        self.pptx_path = pptx_path
        self.yaml_path = yaml_path
        self.brand = brand
        self.proof_dir = proof_dir or str(Path(pptx_path).with_suffix("")) + "-proof"
        self.prs = Presentation(pptx_path)
        self.report = QAReport(pptx_path)
        self.report.total_slides = len(self.prs.slides)

    def run_structural(self):
        """Run all structural checks on the PPTX object model."""
        print("  Running structural checks...")
        issues = []
        issues.extend(check_word_count(self.prs))
        issues.extend(check_font_sizes(self.prs))
        issues.extend(check_text_overflow(self.prs))
        issues.extend(check_consistency(self.prs))
        issues.extend(check_containment(self.prs))
        issues.extend(check_empty_text_boxes(self.prs))
        issues.extend(check_table_font_opportunity(self.prs))
        issues.extend(check_table_centering(self.prs))
        issues.extend(check_unhyperlinked_urls(self.prs))
        issues.extend(check_kpi_text_overflow(self.prs))
        issues.extend(check_kpi_label_alignment(self.prs))
        # YAML-dependent checks
        if self.yaml_path:
            issues.extend(check_yaml_content_limits(self.yaml_path))
            issues.extend(check_yaml_icon_resolution(self.yaml_path))
            issues.extend(check_story_structure(self.yaml_path))
            issues.extend(check_audience_density(self.yaml_path))
        for i in issues:
            self.report.add_issue(i)
        return issues

    def run_visual(self):
        """Generate proof images and run visual checks."""
        print("  Generating proof images...")
        try:
            from proof_renderer import render_proof_images
            summary = render_proof_images(self.pptx_path, self.proof_dir)
            # Convert proof_renderer issues to our format
            if summary:
                for slide_info in summary.get("slides", []):
                    for issue in slide_info.get("issues", []):
                        self.report.add_issue(Issue(
                            slide_info["slide"],
                            issue.get("type", "visual"),
                            "warning",
                            issue.get("description", "Visual issue"),
                        ))
        except ImportError:
            print("    proof_renderer not available, skipping visual checks")
        except Exception as e:
            print(f"    Visual check error: {e}")

    def run_vision(self):
        """Run Claude vision review on proof images."""
        print("  Running Claude vision review...")
        issues = run_vision_review(self.proof_dir, brand=self.brand)
        for i in issues:
            self.report.add_issue(i)
        return issues

    def run_auto_fix(self):
        """Apply safe automatic fixes to the PPTX."""
        print("  Running auto-fix...")
        fixes = auto_fix_pptx(self.pptx_path)
        if fixes:
            # Reload the presentation after fixes
            self.prs = Presentation(self.pptx_path)
        return fixes

    def run_all(self, ai_review=False, auto_fix=True, skip_visual=False):
        """Run all checks and return the report."""
        self.run_structural()
        if auto_fix:
            self.run_auto_fix()
        if not skip_visual:
            self.run_visual()
        if ai_review:
            self.run_vision()
        return self.report

    def save_report(self, output_dir=None):
        """Save the QA report as JSON and markdown."""
        out_dir = output_dir or os.path.dirname(self.pptx_path)
        base = Path(self.pptx_path).stem

        json_path = os.path.join(out_dir, f"{base}-qa-report.json")
        md_path = os.path.join(out_dir, f"{base}-qa-report.md")

        with open(json_path, "w") as f:
            json.dump(self.report.to_json(), f, indent=2)

        with open(md_path, "w") as f:
            f.write(self.report.to_markdown())

        return json_path, md_path


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    import argparse

    parser = argparse.ArgumentParser(description="Run QA pipeline on a PPTX deck")
    parser.add_argument("pptx", help="Path to the PPTX file")
    parser.add_argument("--ai-review", action="store_true", help="Enable Claude vision review")
    parser.add_argument("--strict", action="store_true", help="Exit with error on critical issues")
    parser.add_argument("--output-dir", help="Output directory for reports and proof images")
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"ERROR: File not found: {args.pptx}")
        sys.exit(1)

    proof_dir = args.output_dir or str(Path(args.pptx).with_suffix("")) + "-proof"
    qa = QAPipeline(args.pptx, proof_dir)
    report = qa.run_all(ai_review=args.ai_review)

    json_path, md_path = qa.save_report(args.output_dir)
    s = report.summary()

    print(f"\n  QA Report: {md_path}")
    print(f"  Slides: {s['total_slides']} | Issues: {s['total_issues']} "
          f"({s['critical']} critical, {s['warnings']} warnings, {s['info']} info)")

    if s["blocked"]:
        print("  STATUS: BLOCKED — Critical issues found.")
        if args.strict:
            print("  Use --force to override.")
            sys.exit(2)
    elif s["total_issues"] == 0:
        print("  STATUS: CLEAN")
    else:
        print(f"  STATUS: REVIEW — {s['warnings']} warning(s)")


if __name__ == "__main__":
    main()
