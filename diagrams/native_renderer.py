"""
Approach A: PPTX Native Shape Renderer

Generates diagrams directly as python-pptx shapes (rectangles, connectors,
text boxes). These are added directly to the slide — no external files needed.

Returns a callable that takes a slide object and adds shapes to it.
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# Style palettes — each style maps to a set of role-based colors
STYLE_PALETTES = {
    # === PRIMARY STYLES (brand-aware — uses brand colors) ===
    "corporate": {
        "root_fill": "#1A365D",
        "root_text": "#FFFFFF",
        "root_sub": "#A3C4EC",
        "child_fill": "#EBF4FF",
        "child_border": "#1A365D",
        "child_text": "#1A365D",
        "child_sub": "#2D3748",
        "connector": "#1A365D",
        "accent": "#3182CE",
        "bg": "#FFFFFF",
    },
    "tech-gradient": {
        "root_fill": "#1A0533",
        "root_text": "#E0D0FF",
        "root_sub": "#A855F7",
        "child_fill": "#2D1B4E",
        "child_border": "#7C3AED",
        "child_text": "#E0D0FF",
        "child_sub": "#A78BDA",
        "connector": "#7C3AED",
        "accent": "#06B6D4",
        "bg": "#0F0720",
    },
    "blueprint": {
        "root_fill": "#1E3A5F",
        "root_text": "#FFFFFF",
        "root_sub": "#7CB3D8",
        "child_fill": "#0D2137",
        "child_border": "#4A90C4",
        "child_text": "#C8DDF0",
        "child_sub": "#7CB3D8",
        "connector": "#4A90C4",
        "accent": "#FF6B35",
        "bg": "#0A1628",
    },
    # === ADDITIONAL STYLES ===
    "isometric": {
        "root_fill": "#2C3E50",
        "root_text": "#ECF0F1",
        "root_sub": "#3498DB",
        "child_fill": "#EBF5FB",
        "child_border": "#2980B9",
        "child_text": "#2C3E50",
        "child_sub": "#566573",
        "connector": "#2980B9",
        "accent": "#1ABC9C",
        "bg": "#FFFFFF",
    },
    "glassmorphism": {
        "root_fill": "#1A365D",
        "root_text": "#FFFFFF",
        "root_sub": "#A3C4EC",
        "child_fill": "#F0F7FF",
        "child_border": "#B3D4FC",
        "child_text": "#1A365D",
        "child_sub": "#4A6A8A",
        "connector": "#90B8E0",
        "accent": "#3182CE",
        "bg": "#FAFCFF",
    },
    "neon-wireframe": {
        "root_fill": "#1A1A2E",
        "root_text": "#00FFAA",
        "root_sub": "#00CC88",
        "child_fill": "#16213E",
        "child_border": "#00FFAA",
        "child_text": "#00FFAA",
        "child_sub": "#00CC88",
        "connector": "#00FFAA",
        "accent": "#FF00FF",
        "bg": "#0F0F23",
    },
    "paper-cut": {
        "root_fill": "#1A365D",
        "root_text": "#FFFFFF",
        "root_sub": "#A3C4EC",
        "child_fill": "#F5F9FF",
        "child_border": "#C5D8E8",
        "child_text": "#1A365D",
        "child_sub": "#5A6A7A",
        "connector": "#A0B8D4",
        "accent": "#3182CE",
        "bg": "#FBFDFF",
    },
    "minimal-line": {
        "root_fill": "#FFFFFF",
        "root_text": "#1A1A1A",
        "root_sub": "#666666",
        "child_fill": "#FFFFFF",
        "child_border": "#1A1A1A",
        "child_text": "#1A1A1A",
        "child_sub": "#666666",
        "connector": "#1A1A1A",
        "accent": "#3182CE",
        "bg": "#FFFFFF",
    },
    "hand-drawn": {
        "root_fill": "#34495E",
        "root_text": "#FFFFFF",
        "root_sub": "#95A5A6",
        "child_fill": "#FDFEFE",
        "child_border": "#7F8C8D",
        "child_text": "#2C3E50",
        "child_sub": "#7F8C8D",
        "connector": "#95A5A6",
        "accent": "#E74C3C",
        "bg": "#FFFFFF",
    },
}


def _build_brand_palettes(tokens):
    """Build brand-aware palettes for corporate, glassmorphism, and paper-cut styles."""
    c = tokens.get("colors", {})
    primary = c.get("purple", "#1A365D")
    secondary = c.get("pink", "#3182CE")
    accent = c.get("light_pink", "#63B3ED")
    light_bg = c.get("light_bg", "#EBF4FF")
    dark = c.get("dark", "#2D3748")
    return {
        "corporate": {
            "root_fill": primary, "root_text": "#FFFFFF", "root_sub": accent,
            "child_fill": light_bg, "child_border": primary,
            "child_text": primary, "child_sub": dark,
            "connector": primary, "accent": secondary, "bg": "#FFFFFF",
        },
        "glassmorphism": {
            "root_fill": primary, "root_text": "#FFFFFF", "root_sub": accent,
            "child_fill": "#F5EEFA", "child_border": "#D4B5E8",
            "child_text": primary, "child_sub": "#7A5A8A",
            "connector": "#C9A0DC", "accent": secondary, "bg": "#FAFAFE",
        },
        "paper-cut": {
            "root_fill": primary, "root_text": "#FFFFFF", "root_sub": accent,
            "child_fill": "#FFF5FA", "child_border": "#E8C5D8",
            "child_text": primary, "child_sub": "#8A5A7A",
            "connector": "#D4A0C0", "accent": secondary, "bg": "#FFFBFD",
        },
    }


def _get_style_colors(style, brand_palettes=None):
    """Get RGBColor objects for a given style palette."""
    # Brand-aware palettes override the static ones for corporate/glassmorphism/paper-cut
    if brand_palettes and style in brand_palettes:
        palette = brand_palettes[style]
    else:
        palette = STYLE_PALETTES.get(style, STYLE_PALETTES["corporate"])
    return {k: RGBColor.from_string(v.lstrip("#")) for k, v in palette.items()}


class NativeRenderer:

    def __init__(self, tokens):
        self.tokens = tokens
        self.colors = {k: RGBColor.from_string(v.lstrip("#"))
                       for k, v in tokens["colors"].items()}
        self._brand_palettes = _build_brand_palettes(tokens)

    def render(self, diagram_type, data, target_width_in, target_height_in,
               style="corporate", output_dir=None):
        """Render a diagram as PPTX shapes.

        Returns dict with type="shapes" and a shapes_fn callable.
        """
        handler = {
            "org-hierarchy": self._render_org_hierarchy,
            "flow": self._render_flow,
            "comparison": self._render_comparison,
            "timeline": self._render_timeline,
            "key-stats": self._render_key_stats,
            "labeled-boxes": self._render_labeled_boxes,
            "process-steps": self._render_process_steps,
        }.get(diagram_type)

        if not handler:
            return None

        shapes_fn = handler(data, target_width_in, target_height_in, style)
        return {
            "type": "shapes",
            "shapes_fn": shapes_fn,
            "label": f"Native PPTX shapes ({diagram_type})",
        }

    # ------------------------------------------------------------------
    # org-hierarchy: tree/folder structure with boxes and connectors
    # ------------------------------------------------------------------
    def _render_org_hierarchy(self, data, w, h, style):
        sc = _get_style_colors(style, self._brand_palettes)  # style-specific colors
        colors = self.colors  # brand colors (fallback)
        tokens = self.tokens

        root_name = data.get("root", "")
        root_owner = data.get("owner", "")
        children = data.get("children", [])

        def shapes_fn(slide, x_offset, y_offset):
            """Add org hierarchy shapes to slide at the given offset."""
            n = len(children)

            # Adaptive sizing based on available space and child count
            root_box_h = min(0.4, h * 0.25)
            connector_h = min(0.2, h * 0.1)
            child_box_h = min(0.35, h * 0.2)
            child_gap = 0.08

            # Determine child layout: single row or two rows
            # Use two rows only if single-row boxes would be narrower than 1.0"
            usable_w = w - 0.3  # margins
            single_row_child_w = (usable_w - (n - 1) * child_gap) / n if n > 0 else 2.0
            min_child_w = 1.0  # minimum readable box width

            if single_row_child_w >= min_child_w or n <= 2:
                rows = 1
                children_per_row = [children]
            else:
                rows = 2
                mid = (n + 1) // 2
                children_per_row = [children[:mid], children[mid:]]

            # Calculate child box width based on the widest row
            max_per_row = max(len(r) for r in children_per_row) if children_per_row else 0
            if max_per_row > 0:
                child_w = (usable_w - (max_per_row - 1) * child_gap) / max_per_row
                child_w = min(child_w, 2.0)
            else:
                child_w = 1.5

            # Adaptive font size for children based on box width
            max_child_name = max((len(c) if isinstance(c, str) else len(c.get("name", "")))
                                 for c in children) if children else 0
            chars_per_inch = 10  # approximate for typical body font bold
            fits_at_8pt = child_w * chars_per_inch >= max_child_name
            child_font_size = 8 if fits_at_8pt else max(6, min(8, int(child_w * chars_per_inch * 8 / max_child_name)))

            # Root font size — scale with available width
            root_w = min(w * 0.5, 3.0)
            root_font = 10 if root_w > 1.5 else 8

            # --- Connector spacing rule ---
            # Lines centered between boxes. Whitespace = 4× line visual width.
            line_pt = 1.5
            line_visual_w = line_pt / 72
            min_connector_gap = line_visual_w * 4 + line_visual_w + line_visual_w * 4

            # --- Calculate height, expand gaps, but cap at box height ---
            num_gaps = rows
            min_total_h = root_box_h + num_gaps * min_connector_gap + rows * child_box_h

            # Distribute extra space across gaps, but never expand a gap
            # beyond the height of the boxes in a row
            max_gap = child_box_h  # gap should not exceed box height
            extra_h = max(0, h - min_total_h)
            expanded_gap = min(min_connector_gap + extra_h / num_gaps, max_gap)

            # Actual total height used (may be less than available h)
            actual_total_h = root_box_h + num_gaps * expanded_gap + rows * child_box_h

            # Center vertically within available space
            v_offset = y_offset + max(0, (h - actual_total_h) / 2)

            # Position root box (rendered in Pass 2 after lines)
            root_x = x_offset + (w - root_w) / 2
            root_y = v_offset

            if n == 0:
                # No children — just render root box and return
                root_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(root_x), Inches(root_y),
                    Inches(root_w), Inches(root_box_h))
                root_box.fill.solid()
                root_box.fill.fore_color.rgb = sc["root_fill"]
                root_box.line.fill.background()
                tf = root_box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = root_name
                r.font.name = tokens["fonts"]["heading"]
                r.font.size = Pt(root_font); r.font.color.rgb = sc["root_text"]
                r.font.bold = True
                if root_owner:
                    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
                    r2 = p2.add_run(); r2.text = root_owner
                    r2.font.name = tokens["fonts"]["body"]
                    r2.font.size = Pt(max(6, root_font - 3))
                    r2.font.color.rgb = sc["root_sub"]
                return

            root_center_x = root_x + root_w / 2
            root_bottom = root_y + root_box_h

            # --- Two-pass rendering: lines first (behind), then boxes (in front) ---
            # Pre-calculate all positions, then render lines, then render boxes.

            # Calculate positions for all rows
            row_layouts = []
            prev_element_bottom = root_bottom

            for row_idx, row_children in enumerate(children_per_row):
                row_n = len(row_children)
                child_y = prev_element_bottom + expanded_gap
                bar_y = prev_element_bottom + expanded_gap / 2

                use_rotated = (row_idx > 0 and row_n > 4)

                if use_rotated:
                    actual_w = 0.5
                    actual_h = min(child_box_h * 2.5, 0.9)
                else:
                    actual_w = child_w
                    actual_h = child_box_h

                total_row_w = row_n * actual_w + (row_n - 1) * child_gap
                row_start_x = x_offset + (w - total_row_w) / 2

                row_layouts.append({
                    "children": row_children,
                    "child_y": child_y,
                    "bar_y": bar_y,
                    "row_start_x": row_start_x,
                    "actual_w": actual_w,
                    "actual_h": actual_h,
                    "use_rotated": use_rotated,
                    "row_n": row_n,
                    "prev_bottom": prev_element_bottom,
                })

                prev_element_bottom = child_y + actual_h

            # PASS 1: Draw all connector lines (behind)
            for rl in row_layouts:
                # Vertical line from previous element to bar
                vline = slide.shapes.add_connector(
                    1, Inches(root_center_x), Inches(rl["prev_bottom"]),
                    Inches(root_center_x), Inches(rl["bar_y"]))
                vline.line.color.rgb = sc["connector"]
                vline.line.width = Pt(line_pt)

                # Horizontal bar
                if rl["row_n"] > 1:
                    first_c = rl["row_start_x"] + rl["actual_w"] / 2
                    last_c = rl["row_start_x"] + (rl["row_n"] - 1) * (rl["actual_w"] + child_gap) + rl["actual_w"] / 2

                    hline = slide.shapes.add_connector(
                        1, Inches(first_c), Inches(rl["bar_y"]),
                        Inches(last_c), Inches(rl["bar_y"]))
                    hline.line.color.rgb = sc["connector"]
                    hline.line.width = Pt(line_pt)

                # Vertical drops to each child
                for i in range(rl["row_n"]):
                    cx = rl["row_start_x"] + i * (rl["actual_w"] + child_gap)
                    cc = cx + rl["actual_w"] / 2

                    vdrop = slide.shapes.add_connector(
                        1, Inches(cc), Inches(rl["bar_y"]),
                        Inches(cc), Inches(rl["child_y"]))
                    vdrop.line.color.rgb = sc["connector"]
                    vdrop.line.width = Pt(line_pt)

            # PASS 2: Draw root box (in front of lines from root)
            root_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(root_x), Inches(root_y),
                Inches(root_w), Inches(root_box_h))
            root_box.fill.solid()
            root_box.fill.fore_color.rgb = sc["root_fill"]
            root_box.line.fill.background()

            tf = root_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = root_name
            r.font.name = tokens["fonts"]["heading"]
            r.font.size = Pt(root_font)
            r.font.color.rgb = colors["white"]
            r.font.bold = True

            if root_owner:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run()
                r2.text = root_owner
                r2.font.name = tokens["fonts"]["body"]
                r2.font.size = Pt(max(6, root_font - 3))
                r2.font.color.rgb = colors["light_pink"]

            # PASS 3: Draw all child boxes (in front of lines)
            for rl in row_layouts:
                for i, child in enumerate(rl["children"]):
                    if isinstance(child, str):
                        child_name = child
                        child_owner = ""
                    else:
                        child_name = child.get("name", "")
                        child_owner = child.get("owner", "")

                    cx = rl["row_start_x"] + i * (rl["actual_w"] + child_gap)

                    cbox = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(cx), Inches(rl["child_y"]),
                        Inches(rl["actual_w"]), Inches(rl["actual_h"]))
                    cbox.fill.solid()
                    cbox.fill.fore_color.rgb = sc["child_fill"]
                    cbox.line.color.rgb = sc["child_border"]
                    cbox.line.width = Pt(0.75)

                    ctf = cbox.text_frame
                    ctf.word_wrap = True
                    cp = ctf.paragraphs[0]
                    cp.alignment = PP_ALIGN.CENTER
                    cr = cp.add_run()
                    cr.text = child_name
                    cr.font.name = tokens["fonts"]["body"]
                    cr.font.color.rgb = sc["child_text"]
                    cr.font.bold = True
                    cr.font.size = Pt(max(5, child_font_size - 1) if rl["use_rotated"] else child_font_size)

                    if child_owner:
                        cp2 = ctf.add_paragraph()
                        cp2.alignment = PP_ALIGN.CENTER
                        cr2 = cp2.add_run()
                        cr2.text = child_owner
                        cr2.font.name = tokens["fonts"]["body"]
                        cr2.font.size = Pt(max(5, child_font_size - 2))
                        cr2.font.color.rgb = sc["child_sub"]

        return shapes_fn

    # ------------------------------------------------------------------
    # flow: process steps A → B → C
    # ------------------------------------------------------------------
    def _render_flow(self, data, w, h, style):
        sc = _get_style_colors(style, self._brand_palettes)
        colors = self.colors
        tokens = self.tokens
        steps = data.get("steps", [])

        def shapes_fn(slide, x_offset, y_offset):
            n = len(steps)
            if n == 0:
                return

            arrow_w = 0.2
            total_arrow_w = (n - 1) * arrow_w
            total_gap = (n - 1) * 0.15
            step_w = (w - total_arrow_w - total_gap - 0.2) / n
            step_w = min(step_w, 2.0)
            step_h = min(h * 0.6, 0.6)

            total_used = n * step_w + (n - 1) * (arrow_w + 0.15)
            start_x = x_offset + (w - total_used) / 2
            step_y = y_offset + (h - step_h) / 2

            for i, step in enumerate(steps):
                if isinstance(step, str):
                    step_text = step
                    step_sub = ""
                else:
                    step_text = step.get("name", "")
                    step_sub = step.get("detail", "")

                sx = start_x + i * (step_w + arrow_w + 0.15)

                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(sx), Inches(step_y),
                    Inches(step_w), Inches(step_h))
                box.fill.solid()
                box.fill.fore_color.rgb = sc["root_fill"] if i == 0 else sc["child_fill"]
                box.line.color.rgb = sc["child_border"]
                box.line.width = Pt(0.75)

                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = step_text
                r.font.name = tokens["fonts"]["heading"]
                r.font.size = Pt(9)
                r.font.color.rgb = sc["root_text"] if i == 0 else sc["child_text"]
                r.font.bold = True

                if step_sub:
                    p2 = tf.add_paragraph()
                    p2.alignment = PP_ALIGN.CENTER
                    r2 = p2.add_run()
                    r2.text = step_sub
                    r2.font.name = tokens["fonts"]["body"]
                    r2.font.size = Pt(7)
                    r2.font.color.rgb = sc["root_sub"] if i == 0 else sc["child_sub"]

                # Arrow to next step
                if i < n - 1:
                    ax = sx + step_w + 0.03
                    ay = step_y + step_h / 2
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Inches(ax), Inches(ay - 0.08),
                        Inches(arrow_w + 0.09), Inches(0.16))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = sc["accent"]
                    arrow.line.fill.background()

        return shapes_fn

    # ------------------------------------------------------------------
    # comparison: two options side by side with indicators
    # ------------------------------------------------------------------
    def _render_comparison(self, data, w, h, style):
        sc = _get_style_colors(style, self._brand_palettes)
        colors = self.colors
        tokens = self.tokens

        options = data.get("options", [])

        def shapes_fn(slide, x_offset, y_offset):
            n = len(options)
            if n == 0:
                return

            col_gap = 0.15
            col_w = (w - (n - 1) * col_gap) / n
            header_h = 0.35
            item_h = 0.22

            for oi, opt in enumerate(options):
                cx = x_offset + oi * (col_w + col_gap)
                opt_name = opt.get("name", f"Option {oi+1}")
                pros = opt.get("pros", [])
                cons = opt.get("cons", [])

                # Header
                hbox = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(y_offset),
                    Inches(col_w), Inches(header_h))
                hbox.fill.solid()
                hbox.fill.fore_color.rgb = sc["root_fill"]
                hbox.line.fill.background()
                tf = hbox.text_frame
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = opt_name
                r.font.name = tokens["fonts"]["heading"]
                r.font.size = Pt(10); r.font.color.rgb = sc["root_text"]
                r.font.bold = True

                # Pros
                cy = y_offset + header_h + 0.1
                for item in pros:
                    tb = slide.shapes.add_textbox(
                        Inches(cx + 0.05), Inches(cy),
                        Inches(col_w - 0.1), Inches(item_h))
                    tf = tb.text_frame; tf.word_wrap = True
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = f"✓ {item}"
                    r.font.name = tokens["fonts"]["body"]
                    r.font.size = Pt(7)
                    r.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)
                    cy += item_h

                # Cons
                for item in cons:
                    tb = slide.shapes.add_textbox(
                        Inches(cx + 0.05), Inches(cy),
                        Inches(col_w - 0.1), Inches(item_h))
                    tf = tb.text_frame; tf.word_wrap = True
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = f"✗ {item}"
                    r.font.name = tokens["fonts"]["body"]
                    r.font.size = Pt(7)
                    r.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
                    cy += item_h

        return shapes_fn

    # ------------------------------------------------------------------
    # timeline: horizontal milestones
    # ------------------------------------------------------------------
    def _render_timeline(self, data, w, h, style):
        sc = _get_style_colors(style, self._brand_palettes)
        colors = self.colors
        tokens = self.tokens
        milestones = data.get("milestones", [])

        def shapes_fn(slide, x_offset, y_offset):
            n = len(milestones)
            if n == 0:
                return

            line_y = y_offset + h * 0.4
            margin = 0.3
            line_w = w - 2 * margin

            # Main timeline line
            tline = slide.shapes.add_connector(
                1, Inches(x_offset + margin), Inches(line_y),
                Inches(x_offset + w - margin), Inches(line_y))
            tline.line.color.rgb = sc["connector"]
            tline.line.width = Pt(2)

            spacing = line_w / (n - 1) if n > 1 else 0

            for i, ms in enumerate(milestones):
                if isinstance(ms, str):
                    ms_label = ms
                    ms_date = ""
                else:
                    ms_label = ms.get("label", "")
                    ms_date = ms.get("date", "")

                mx = x_offset + margin + i * spacing

                # Dot on timeline
                dot_r = 0.08
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(mx - dot_r), Inches(line_y - dot_r),
                    Inches(dot_r * 2), Inches(dot_r * 2))
                dot.fill.solid()
                dot.fill.fore_color.rgb = sc["accent"]
                dot.line.fill.background()

                # Date above
                if ms_date:
                    tb = slide.shapes.add_textbox(
                        Inches(mx - 0.5), Inches(line_y - 0.35),
                        Inches(1.0), Inches(0.22))
                    tf = tb.text_frame
                    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                    r = p.add_run(); r.text = ms_date
                    r.font.name = tokens["fonts"]["heading"]
                    r.font.size = Pt(7); r.font.color.rgb = sc["child_text"]
                    r.font.bold = True

                # Label below
                tb = slide.shapes.add_textbox(
                    Inches(mx - 0.6), Inches(line_y + 0.15),
                    Inches(1.2), Inches(0.4))
                tf = tb.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = ms_label
                r.font.name = tokens["fonts"]["body"]
                r.font.size = Pt(7); r.font.color.rgb = sc["child_sub"]

        return shapes_fn

    # ------------------------------------------------------------------
    # key-stats: numbers with icons (reusable from stats graphic)
    # ------------------------------------------------------------------
    def _render_key_stats(self, data, w, h, style):
        sc = _get_style_colors(style, self._brand_palettes)
        tokens = self.tokens
        stats = data.get("stats", [])

        def shapes_fn(slide, x_offset, y_offset):
            n = len(stats)
            if n == 0:
                return

            stat_h = h / n
            for i, stat in enumerate(stats):
                sy = y_offset + i * stat_h
                num = stat.get("number", "")
                label = stat.get("label", "")

                tb = slide.shapes.add_textbox(
                    Inches(x_offset), Inches(sy),
                    Inches(w), Inches(stat_h * 0.55))
                tf = tb.text_frame
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = num
                r.font.name = tokens["fonts"]["heading"]
                r.font.size = Pt(28); r.font.color.rgb = sc["child_text"]
                r.font.bold = True

                tb2 = slide.shapes.add_textbox(
                    Inches(x_offset), Inches(sy + stat_h * 0.55),
                    Inches(w), Inches(stat_h * 0.35))
                tf2 = tb2.text_frame
                p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run(); r2.text = label
                r2.font.name = tokens["fonts"]["body"]
                r2.font.size = Pt(9); r2.font.color.rgb = sc["child_sub"]

        return shapes_fn

    # ------------------------------------------------------------------
    # labeled-boxes: generic labeled rectangles in a grid
    # ------------------------------------------------------------------
    def _render_labeled_boxes(self, data, w, h, style):
        sc = _get_style_colors(style, self._brand_palettes)
        colors = self.colors
        tokens = self.tokens
        boxes = data.get("boxes", [])
        cols = data.get("columns", 3)

        def shapes_fn(slide, x_offset, y_offset):
            n = len(boxes)
            if n == 0:
                return

            rows = -(-n // cols)  # ceil division
            gap = 0.1
            box_w = (w - (cols - 1) * gap) / cols
            box_h = (h - (rows - 1) * gap) / rows
            box_h = min(box_h, 0.5)

            for i, box_data in enumerate(boxes):
                if isinstance(box_data, str):
                    label = box_data
                else:
                    label = box_data.get("label", "")

                col = i % cols
                row = i // cols
                bx = x_offset + col * (box_w + gap)
                by = y_offset + row * (box_h + gap)

                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(bx), Inches(by),
                    Inches(box_w), Inches(box_h))
                box.fill.solid()
                box.fill.fore_color.rgb = sc["child_fill"]
                box.line.color.rgb = sc["child_border"]
                box.line.width = Pt(0.75)

                tf = box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = label
                r.font.name = tokens["fonts"]["body"]
                r.font.size = Pt(8); r.font.color.rgb = sc["child_text"]
                r.font.bold = True

        return shapes_fn

    # ------------------------------------------------------------------
    # process-steps: numbered vertical or horizontal steps
    # ------------------------------------------------------------------
    def _render_process_steps(self, data, w, h, style):
        # Reuse flow renderer with numbered steps
        steps = data.get("steps", [])
        numbered = [{"name": f"{i+1}. {s}" if isinstance(s, str) else f"{i+1}. {s.get('name','')}",
                      "detail": "" if isinstance(s, str) else s.get("detail", "")}
                     for i, s in enumerate(steps)]
        return self._render_flow({"steps": numbered}, w, h, style)
